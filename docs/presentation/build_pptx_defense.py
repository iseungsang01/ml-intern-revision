# -*- coding: utf-8 -*-
"""Build the research-audit / defense deck (Korean, abstract register).

Output: docs/presentation/KSTAR_CES_종합방어.pptx  (25 slides, 16:9)

연구 종합 정리 + novelty 검증 + 예상 질문 방어표. 2026-08-27 전면 재작성:
(1) 승상님 지시에 따라 모든 슬라이드 본문·표·노트를 논문 초록 문체(서술형 종결, 객관·비인칭)로
    통일하였다. 발표 대사 인용문과 명령형 메모는 서술문으로 바꾸었다.
(2) 2026-08-16 §8ab 이후의 기록 — B.9 도달 범위 사다리·계열 비교·연산자 비용(§8ac–§8aj),
    밀집 사다리·통합 재채점·승패 방전 분석(§8ak–§8an), μs shot 집합 동결(§8ao), 양자 가지
    종결(§8ap), 논문 프레이밍 §9 — 를 판정표·B.9 슬라이드 2장·예상 질문 ④·결정 기록에 반영하였다.

확정 프로토콜(W = 2 · held-free · 두 모집단 공동 1차 · 백본 `seq_v2`) 기준이며, 이전 W = 4
초고의 수치·서사(held 포함/제외 이중 보고, +0.18~+0.28, MNAR 1/4·캠페인 0/4, anchor+Δ 31.5%,
seq +0.045, CPU p99 6.4 ms, iter2→iter9 progression)는 전부 폐기되었다. 배치 슬라이드의 윈도
대조군 p99 18.9 ms는 §8ac에서 오염된 측정으로 판정되어 같은 세션 값(4.46 ms)으로 교체하였다.

수치 원천: docs/paper/paper_numbers.json(= main_ko.tex; B.1–B.5 동결 산출물에서
collect_paper_numbers.py가 자동 수집)과 THESIS_RESULTS.md §8ac–§8ap의 표. 프로토콜은
ces_prediction/experiments/PREREGISTRATION_W2.md · PREREGISTRATION_B9.md · PREREGISTRATION_B6.md.

Standalone on purpose: importing build_pptx.py would rebuild the 1-hour deck as a side
effect, so the palette/layout helpers are copied here instead. Layout QC uses
preview_pptx.py's own font metrics, so a card that overflows is reported at build time
(`FIT WARNING`) rather than discovered in PowerPoint.

Usage (from repo root):
    py docs/presentation/build_pptx_defense.py
    py docs/presentation/preview_pptx.py docs/presentation/KSTAR_CES_종합방어.pptx
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
if HERE not in sys.path:
    sys.path.insert(0, HERE)

from preview_pptx import load_font, _TOKEN  # noqa: E402  (same metrics as the QC renderer)
from PIL import Image, ImageDraw  # noqa: E402

FIG = os.path.join(HERE, "figures")

NAVY = RGBColor(0x13, 0x33, 0x5F)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
TEAL = RGBColor(0x1B, 0x9E, 0x8A)
ORANGE = RGBColor(0xE8, 0x74, 0x3B)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x5B, 0x66, 0x70)
LGRAY = RGBColor(0xE8, 0xEC, 0xF1)
MGRAY = RGBColor(0x9A, 0xA5, 0xB1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2B, 0x35)
CARDBG = RGBColor(0xF4, 0xF7, 0xFA)

FONT = "Malgun Gothic"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_pageno = {"n": 0}

_MEASURE = ImageDraw.Draw(Image.new("RGB", (8, 8)))
_MARGIN_IN = 2 * (2 / 72.0)   # text() sets 2 pt left + 2 pt right margins
_WARNED = []


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _set_fill(bg, WHITE)
    bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
        round_=False):
    sp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, dict):
            opts = para
            segs = opts["segs"]
            if "align" in opts:
                p.alignment = opts["align"]
            if "space_after" in opts:
                p.space_after = Pt(opts["space_after"])
            if "space_before" in opts:
                p.space_before = Pt(opts["space_before"])
            if "level" in opts:
                p.level = opts["level"]
        else:
            segs = para
        for seg in segs:
            txt, size, color, bold, italic, font = (seg + (None,) * 6)[:6]
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color if color is not None else DARK
            r.font.bold = bool(bold)
            r.font.italic = bool(italic)
            r.font.name = font if font else FONT
    return tb


def header(s, kicker, title, accent=NAVY):
    box(s, Inches(0), Inches(0), EMU_W, Inches(1.28), fill=WHITE)
    box(s, Inches(0.0), Inches(0.0), Inches(0.22), Inches(1.28), fill=accent)
    text(s, Inches(0.55), Inches(0.20), Inches(12.3), Inches(0.34),
         [[(kicker, 13, accent, True, False, None)]])
    size = 26.0
    while size > 18 and _n_lines(title, 12.4 - _MARGIN_IN, size, bold=True) > 1:
        size -= 1.0
    if _n_lines(title, 12.4 - _MARGIN_IN, size, bold=True) > 1:
        _WARNED.append(f"header title wraps at {size} pt: {title!r}")
    text(s, Inches(0.52), Inches(0.50), Inches(12.4), Inches(0.72),
         [[(title, size, NAVY, True, False, None)]])
    box(s, Inches(0.55), Inches(1.18), Inches(12.25), Pt(2), fill=LGRAY)
    footer(s)


def footer(s):
    _pageno["n"] += 1
    text(s, Inches(0.55), Inches(7.06), Inches(9.6), Inches(0.3),
         [[("KSTAR CES Nowcasting — 연구 종합 정리·방어 · 확정 프로토콜(W=2 · held-free · 두 모집단 · 백본 seq_v2) · B.9 반영 · 2026-08-27",
            9, MGRAY, False, False, None)]])
    text(s, Inches(11.6), Inches(7.06), Inches(1.2), Inches(0.3),
         [[(str(_pageno["n"]), 10, MGRAY, False, False, None)]], align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=15, gap=7, color=DARK):
    paras = []
    for it in items:
        txt, lvl, col, bold = (list(it) + [None, None])[:4]
        lvl = lvl or 0
        mark = "●  " if lvl == 0 else ("–  " if lvl == 1 else "·  ")
        sz = size if lvl == 0 else (size - 2 if lvl == 1 else size - 3)
        paras.append({"segs": [(mark + txt, sz,
                                col if col else (NAVY if lvl == 0 else color),
                                bold if bold is not None else (lvl == 0), False, None)],
                      "level": lvl, "space_after": gap})
    return text(s, x, y, w, h, paras, line_spacing=1.08)


def card(s, x, y, w, h, title, lines, accent=BLUE, title_size=14, body_size=12,
         body_color=DARK):
    c = box(s, x, y, w, h, fill=CARDBG, round_=True)
    box(s, x, y, Inches(0.10), h, fill=accent)
    text(s, x + Inches(0.26), y + Inches(0.12), w - Inches(0.4), Inches(0.4),
         [[(title, title_size, accent, True, False, None)]])
    paras = [[(ln, body_size, body_color, False, False, None)] for ln in lines]
    text(s, x + Inches(0.26), y + Inches(0.56), w - Inches(0.42), h - Inches(0.7),
         paras, line_spacing=1.1, space_after=3)
    return c


def add_image_fit(s, path, x, y, w, h):
    from PIL import Image as _Image
    try:
        iw, ih = _Image.open(path).size
    except Exception:
        return s.shapes.add_picture(path, x, y, width=w)
    ar = iw / ih
    bw, bh = w, h
    if bw / bh > ar:
        nh = bh
        nw = int(bh * ar)
    else:
        nw = bw
        nh = int(bw / ar)
    nx = x + (bw - nw) // 2
    ny = y + (bh - nh) // 2
    return s.shapes.add_picture(path, nx, ny, width=nw, height=nh)


def table(s, x, y, col_w, head, rows, row_h=Inches(0.44), head_h=Inches(0.44),
          head_fill=NAVY, head_color=WHITE, size=13, head_size=12.5,
          zebra=LGRAY, emphasis=None, emphasis_fill=None, label_align_left=True,
          left_cols=None):
    # left_cols: column indices whose *body* cells are left-aligned (prose columns
    # look ragged when centred); column 0 follows label_align_left as before.
    left_cols = set(left_cols or ())

    def _align(j):
        return (PP_ALIGN.LEFT if ((j == 0 and label_align_left) or j in left_cols)
                else PP_ALIGN.CENTER)

    total_w = sum(col_w)
    box(s, x, y, total_w, head_h, fill=head_fill)
    cx = x
    for j, h in enumerate(head):
        text(s, cx + Inches(0.10), y + Inches(0.07), col_w[j] - Inches(0.16),
             head_h - Inches(0.1),
             [[(h, head_size, head_color, True, False, None)]],
             align=PP_ALIGN.LEFT if (j == 0 and label_align_left) else PP_ALIGN.CENTER)
        cx += col_w[j]
    emphasis = emphasis or set()
    yy = y + head_h
    for i, row in enumerate(rows):
        if i in emphasis and emphasis_fill is not None:
            box(s, x, yy, total_w, row_h, fill=emphasis_fill)
        elif zebra is not None and i % 2 == 1:
            box(s, x, yy, total_w, row_h, fill=CARDBG)
        cx = x
        for j, cell in enumerate(row):
            if isinstance(cell, str):
                txt, col, bold, font = cell, DARK, (i in emphasis), None
            else:
                txt, col, bold, font = (tuple(cell) + (None,) * 4)[:4]
                col = col if col is not None else DARK
                bold = bold if bold is not None else (i in emphasis)
            text(s, cx + Inches(0.10), yy + Inches(0.08), col_w[j] - Inches(0.16),
                 row_h - Inches(0.1),
                 [[(txt, size, col, bold, False, font)]], align=_align(j))
            cx += col_w[j]
        box(s, x, yy + row_h - Pt(0.75), total_w, Pt(0.75), fill=LGRAY)
        yy += row_h
    return yy


# ====================== LAYOUT QC (preview_pptx metrics) ===================

_QC_DPI = 110.0   # preview_pptx.py's default render scale; font px rounding must match


def _n_lines(txt, avail_in, size_pt, bold=False):
    """Line count after wrapping, using preview_pptx's exact tokenizer and metrics.

    The renderer loads fonts at `int(pt * dpi / 72)` pixels, so measuring at
    `px == pt` under-reports width by a few percent and lets boxes slip through.
    Measure at the same pixel scale the QC renderer uses.
    """
    font = load_font(FONT, bold, size_pt * _QC_DPI / 72.0)
    avail_px = max(avail_in, 0.1) * _QC_DPI
    n, cur = 1, 0.0
    for tok in _TOKEN.findall(txt):
        tw = _MEASURE.textlength(tok, font=font)
        if tok.isspace():
            if cur + tw > avail_px:
                n += 1
                cur = 0.0
                continue
            cur += tw
            continue
        if cur + tw > avail_px and cur > 0:
            n += 1
            cur = 0.0
        cur += tw
    return n


def _block_h(lines, avail_in, size_pt, line_spacing, space_after_pt, bold=False):
    lh = size_pt * 1.24 * line_spacing / 72.0
    total = 0.0
    for ln in lines:
        total += _n_lines(ln, avail_in, size_pt, bold) * lh + space_after_pt / 72.0
    return total


def _in(v):
    """Accept either inches (float) or Emu (what Inches() returns)."""
    return v / 914400.0 if isinstance(v, int) and v > 1000 else float(v)


def fcard(s, x, y, w, h, title, lines, accent=BLUE, body_size=12.0, tag=""):
    """card() that shrinks the type until the wrapped text actually fits the box."""
    x, y, w, h = _in(x), _in(y), _in(w), _in(h)
    body_avail = w - 0.42 - _MARGIN_IN
    title_avail = w - 0.40 - _MARGIN_IN
    body_room = h - 0.70

    title_size = 14.0
    while title_size > 10.5 and _n_lines(title, title_avail, title_size, bold=True) > 1:
        title_size -= 0.5
    if _n_lines(title, title_avail, title_size, bold=True) > 1:
        _WARNED.append(f"card title too long: {title!r} ({tag})")

    size = body_size
    while size > 9.5 and _block_h(lines, body_avail, size, 1.1, 2) > body_room:
        size -= 0.5
    need = _block_h(lines, body_avail, size, 1.1, 2)
    if need > body_room:
        lh = size * 1.24 * 1.1 / 72.0
        cut = int((need - body_room) / lh) + 1
        _WARNED.append(f"card {title!r}: cut ~{cut} line(s) at {size} pt ({tag})")

    c = box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=CARDBG, round_=True)
    box(s, Inches(x), Inches(y), Inches(0.10), Inches(h), fill=accent)
    text(s, Inches(x + 0.26), Inches(y + 0.12), Inches(w - 0.40), Inches(0.4),
         [[(title, title_size, accent, True, False, None)]])
    text(s, Inches(x + 0.26), Inches(y + 0.56), Inches(w - 0.42), Inches(h - 0.70),
         [[(ln, size, DARK, False, False, None)] for ln in lines],
         line_spacing=1.1, space_after=2)
    return c


def fbullets(s, x, y, w, h, items, size=14.0, gap=6, tag=""):
    """bullets() with an auto-shrink pass measured against the real box height."""
    x, y, w, h = _in(x), _in(y), _in(w), _in(h)
    avail = w - _MARGIN_IN

    def _height(sz):
        total = 0.0
        for it in items:
            txt, lvl = (list(it) + [None])[:2]
            lvl = lvl or 0
            mark = "●  " if lvl == 0 else ("–  " if lvl == 1 else "·  ")
            s_ = sz if lvl == 0 else (sz - 2 if lvl == 1 else sz - 3)
            total += (_n_lines(mark + txt, avail, s_, bold=(lvl == 0))
                      * s_ * 1.24 * 1.08 / 72.0) + gap / 72.0
        return total

    while size > 9.5 and _height(size) > h:
        size -= 0.5
    if _height(size) > h:
        _WARNED.append(f"bullets overflow at {size} pt ({tag})")
    return bullets(s, Inches(x), Inches(y), Inches(w), Inches(h), items,
                   size=size, gap=gap)


def ftable(s, x, y, col_w_in, head, rows, row_h=0.44, head_h=0.44, size=13.0,
           head_size=12.5, tag="", **kw):
    """table() that shrinks the type until every cell fits its row height."""
    x, y = _in(x), _in(y)

    # table() leaves text()'s default 4 pt space_after on every cell paragraph
    def _cell_h(txt, avail, sz):
        return _n_lines(txt, avail, sz) * sz * 1.24 * 1.06 / 72.0 + 4 / 72.0

    def _fits(sz):
        room = row_h - 0.10
        for row in rows:
            for j, cell in enumerate(row):
                txt = cell if isinstance(cell, str) else cell[0]
                if _cell_h(txt, col_w_in[j] - 0.16 - _MARGIN_IN, sz) > room:
                    return False
        return True

    while size > 9.0 and not _fits(size):
        size -= 0.5
        head_size = min(head_size, size)
    if not _fits(size):
        bad = []
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                txt = cell if isinstance(cell, str) else cell[0]
                if _cell_h(txt, col_w_in[j] - 0.16 - _MARGIN_IN, size) > row_h - 0.10:
                    bad.append((i, j, txt[:28]))
        _WARNED.append(f"table cells overflow at {size} pt ({tag}): {bad[:3]}")
    return table(s, Inches(x), Inches(y), [Inches(c) for c in col_w_in], head, rows,
                 row_h=Inches(row_h), head_h=Inches(head_h), size=size,
                 head_size=head_size, **kw)


def fband(s, x, y, w, h, paras, tag="", align=PP_ALIGN.LEFT, space_after=3,
          line_spacing=1.06):
    """Free-standing text block (caption / footnote) that shrinks to fit its box.

    `paras` = [(text, size_pt, color, bold), ...]. Sizes scale together.
    """
    x, y, w, h = _in(x), _in(y), _in(w), _in(h)
    avail = w - _MARGIN_IN

    def _height(ratio):
        total = 0.0
        for txt, sz, _c, bold in paras:
            s_ = sz * ratio
            total += (_n_lines(txt, avail, s_, bold=bool(bold))
                      * s_ * 1.24 * line_spacing / 72.0) + space_after / 72.0
        return total

    ratio = 1.0
    while ratio > 0.72 and _height(ratio) > h:
        ratio -= 0.04
    if _height(ratio) > h:
        _WARNED.append(f"band overflows at ratio {ratio:.2f} ({tag}): {paras[0][0][:36]!r}")
    return text(s, Inches(x), Inches(y), Inches(w), Inches(h),
                [[(txt, round(sz * ratio, 1), col, bool(bold), False, None)]
                 for txt, sz, col, bold in paras],
                align=align, space_after=space_after, line_spacing=line_spacing)


def note(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip("\n")
    return s


# ============================ SLIDES ======================================
PAPERFIG = os.path.join(HERE, "..", "paper", "figures")

# --- 1. Title -------------------------------------------------------------
s = slide()
box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
text(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(1.7),
     [[("KSTAR 희소 CES 신호의 멀티모달 나우캐스팅", 34, WHITE, True, False, None)],
      [("연구 종합 정리 · novelty 검증 · 예상 질문 방어", 24, RGBColor(0xBF, 0xD3, 0xEE), True, False, None)]],
     space_after=10)
text(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(2.2),
     [[("확정 프로토콜: W = 2 · held-free(학습·평가) · 파일당 500 · 두 모집단 공동 1차(컷 / 포함) · 인과 GP 기준선 · 백본 seq_v2",
        15, RGBColor(0xBF, 0xD3, 0xEE), False, False, None)],
      [("재실험 B.1–B.5(2026-08-16)와 B.9 문맥·구조·비용 분석(2026-08-17 ~ 08-24)을 반영한다. 이전 W = 4 초고의 수치는 전부 대체되었다.",
        14, RGBColor(0x8F, 0xA8, 0xC8), False, False, None)],
      [("", 8, WHITE, False, False, None)],
      [("이승상 · 서울대학교 원자핵공학과 · 2026년 8월 27일", 14, WHITE, False, False, None)],
      [("근거: THESIS_RESULTS.md §8v–§8ap · §9 · docs/paper/main_ko.tex · paper_numbers.json(동결 산출물 자동 수집) · experiments/PREREGISTRATION_{W2,B9,B6}.md",
        11, RGBColor(0x8F, 0xA8, 0xC8), False, False, None)]],
     space_after=7)
footer(s)
note(s, """
이 덱은 확정 프로토콜(W = 2 · held-free · 두 모집단 공동 1차 · 백본 seq_v2)과 B.9 이후의 기록을
종합한 방어 자료이다. 2026-08-12 방향 재정립(§8v) → B.7 감사(§8w) → B.1 관문(§8x) → B.2(§8y) →
B.3(§8z) → B.4(§8aa) → B.5 전면 재채점(§8ab) → B.9 4축(§8ac–§8aj) → 밀집 사다리·통합·승패 방전
(§8ak–§8an) → μs shot 동결(§8ao) → 양자 가지 종결(§8ap)의 결과만 인용한다.

2분 요약 경로는 2번(두 주장·두 모집단) → 3번(판정표) → 11번(스트레스 2종) → 18번(B.9)이다.
심사 대비의 핵심은 21~24번(예상 질문 16문항 방어표)이다.

폐기 목록(인용 금지): +0.18~+0.28 headline, MNAR 1/4·캠페인 0/4, anchor+Δ 31.5%, seq +0.045,
CPU p99 6.4 ms, 윈도 W=2 p99 18.9 ms(§8ac 오염 판정), held 포함/제외 이중 보고, iter2→iter9
progression, "T_i는 500 ms 문맥이 필요하다"(§8ac 절단 사다리, §8ae·§8af로 대체), "70 ms 문턱"
(§8al로 50 ms), "V_rot는 검정력 부족"(§8al §4로 대체).
""")

# --- 2. One-page summary: two claims + two populations ---------------------
s = slide()
header(s, "한 장 요약", "두 개의 주장, 두 개의 모집단, 그리고 문맥이 추가한 한 문장")
fcard(s, 0.55, 1.45, 6.03, 2.2, "주장 1 (오프라인) — Tᵢ는 미래를 읽는 보간을 이긴다", [
    "PCHIP 대비 컷 +0.174/+0.248/+0.257/+0.264(4/4, 평균 +0.236),",
    "포함 +0.225/+0.238/+0.292/+0.316(4/4, 평균 +0.268)으로 무조건부이다.",
    "이 주장의 상한도 측정되었다. 최강 오프라인 팔인 오프라인 GP와는",
    "동률(-0.05~+0.11, 1/8 유의)이다. headline 비교 대상은 사전등록(PR1)대로",
    "PCHIP를 유지하고 GP는 사다리에 병기한다.",
], accent=BLUE, body_size=12, tag="s2-claim1")
fcard(s, 6.75, 1.45, 6.03, 2.2, "주장 2 (인과·배치) — 배치 가능한 모든 인과 방법을 이긴다", [
    "인과(past-only) GP 대비 4/4 + 4/4(+0.08~+0.17),",
    "persistence 대비 4/4 + 4/4(+0.36~+0.46)이다.",
    "두 스트레스를 모두 견딘다. 실제 결측점 재가중에서 인과 대비 8/8,",
    "캠페인 시간 분할에서 PCHIP·인과 GP 대비 4/4 + 4/4이다.",
    "온라인 가상 센서가 실제로 경쟁하는 상대는 이 사다리이다.",
], accent=TEAL, body_size=12, tag="s2-claim2")
fcard(s, 0.55, 3.8, 6.03, 1.95, "두 모집단 규칙 (2026-08-14 개정 · 08-16 유지 결정)", [
    "컷은 관측 Tᵢ > 3 keV(피팅 실패, 1,197행 = 0.53%)를 로드 시점에 결측",
    "처리한 모집단이며 지도 타깃·이력·정규화 통계·모든 보간 앵커에 동일 적용된다.",
    "포함은 컷 없음(p100)이다. 각 모집단 안에서 전 arm의 처리가 일관된다.",
    "무조건부 주장은 두 모집단 모두에서 성립할 때만 하며, 한쪽만이면 모집단을",
    "명시한 조건부로만 보고한다. 단일 headline보다 엄격한 기준이다.",
], accent=ORANGE, body_size=12, tag="s2-pop")
fcard(s, 6.75, 3.8, 6.03, 1.95, "V_rot — 동률이며, 우위는 회전이 변하는 방전에 집중된다", [
    "PCHIP 대비 PR4 1/4(컷) · 2/4(포함)이므로 회전 승리는 주장하지 않는다.",
    "persistence 대비는 3/4 양쪽(+0.30~+0.50), 인과 GP 대비 2/4이다.",
    "Δt > 15 ms에서는 두 모집단 모두 PCHIP를 이긴다(+0.418* / +0.432*, 130 방전).",
    "방전 단위 승률은 0.48이며 상위 5개 방전을 빼면 4 분할 전부 0 이하이다(§8al).",
    "승패를 예측하는 유일한 공변량은 방전 내 타깃 산포이다(§8an).",
], accent=GRAY, body_size=12, tag="s2-vrot")
fband(s, 0.55, 5.9, 12.25, 0.95,
      [("B.9가 추가한 한 문장: 약 50 ms의 연속 인과 문맥이 최강 배치 기준선에 대한 우위를 전형적으로 만들며(승리 방전 비율 0.52 → 0.66), "
        "그 문맥을 어떤 계열로 넘느냐는 skill이 아니라 비용을 정한다(순환·확장 합성곱·attention은 0.023 이내 동률, 지연은 연산자당 2–3 µs).",
        13, NAVY, True)], tag="s2-caption")
note(s, """
[정확한 수치]
Tᵢ vs PCHIP (seq_v2, TEST 4 분할): 컷 +0.174 [+0.097,+0.236] / +0.248 [+0.188,+0.295] /
+0.257 [+0.199,+0.302] / +0.264 [+0.188,+0.320]; 포함 +0.225 / +0.238 / +0.292 / +0.316.
vs 인과 GP(컷/포함): +0.078/+0.154 · +0.133/+0.169 · +0.138/+0.123 · +0.105/+0.149.
V_rot vs PCHIP: 컷 +0.390* / +0.183 / +0.135 / +0.305 (1/4), 포함 +0.384* / +0.195* /
+0.132 / +0.304 (2/4).

[제시 순서]
두 주장을 먼저 분리한다(오프라인은 관측 모집단의 성질, 인과는 배치의 문장). 두 모집단 규칙을
말한다("무조건부 = 둘 다"가 이 연구의 판정 언어이다). 오프라인 GP 동률을 먼저 밝힌다. 상한을
스스로 측정하였다는 사실이 신뢰의 근거이다. B.9의 한 문장은 §9의 lead claim이다.

[반박 대비]
"컷은 어려운 행을 버린 것이 아닌가"에는 포함 모집단이 민감도가 아니라 공동 1차라는 사실로 답한다.
"p100만 쓰면 되지 않는가"에는 21번 슬라이드 Q1(절제 근거)으로 답한다.
출처: main_ko.tex §6.2, THESIS_RESULTS.md §8ab·§8al·§8am·§8an·§9, PREREGISTRATION_W2.md §1.1.
""")

# --- 3. Verdict table ------------------------------------------------------
s = slide()
header(s, "판정표", "무엇이 무조건부이고, 무엇이 모집단 조건부이며, B.9가 무엇을 더했는가")
ftable(s, 0.55, 1.42, [5.9, 1.85, 1.85, 2.65],
       ["판정 항목", "컷", "포함", "지위"],
       [
           ["백본 Tᵢ가 오프라인 PCHIP를 이김 (PR4)", "4/4", "4/4", ("무조건부", GREEN, True)],
           ["백본 Tᵢ가 인과 GP를 이김 (주장 2)", "4/4", "4/4", ("무조건부", GREEN, True)],
           ["Tᵢ 과도(peak) 층 포착 · Δt > 15 ms 통합", "4/4 · PASS", "4/4 · PASS", ("무조건부", GREEN, True)],
           ["캠페인 이동 — 백본 Tᵢ vs PCHIP · 인과 GP", "4/4", "4/4", ("무조건부 †", GREEN, True)],
           ["conformal 구간 우위 (Winkler 32/32 셀) · V_rot 라우팅 구조적", "전부", "전부", ("무조건부", GREEN, True)],
           ["Tᵢ 문맥 포화 약 50 ms (§3.4 규칙) · 문맥 10배당 +0.050 상승", ("통합 301 방전", NAVY, True), "—", ("무조건부 (B.9)", GREEN, True)],
           ["세 계열 같은 문맥 동률 (|Δ| ≤ 0.023) · 비용 = 연산자 수", ("paired 4 분할", NAVY, True), "—", ("무조건부 (B.9)", GREEN, True)],
           ["윈도 대조군 Tᵢ가 PCHIP를 이김", "3/4", "4/4", ("조건부", ORANGE, True)],
           ["b3k8 (21k) = 백본 Tᵢ", "+0.002", ("-0.194*", RED, True), ("컷 조건부", ORANGE, True)],
           ["타깃 자신의 10 ms 재현성 상한 (Tᵢ 46~130 eV) 대비 백본 157.8 eV",
            ("같은 자릿수", NAVY, True), "—", ("보고 (§8aq)", BLUE, True)],
           ["본류(peak 외) Tᵢ / Δt > 45 ms", "4/4 · PASS", "2/4 · 동률", ("조건부", ORANGE, True)],
           ["MNAR 재가중 Tᵢ vs PCHIP", "2/4", "4/4", ("조건부 ‡", ORANGE, True)],
           ["1 ms 지연 판정", ("보류", GRAY, True), "—", ("미결 (p99 산포 21.84배)", GRAY, True)],
           ["V_rot vs PCHIP (PR4) · 방전 단위 승률 0.48", "1/4", "2/4", ("미확립 §", GRAY, True)],
       ],
       row_h=0.35, head_h=0.40, size=11.5, head_size=12.0, tag="s3-verdict")
fband(s, 0.55, 6.60, 12.25, 0.40,
      [("†  캠페인 = 한 시간 블록 × 초기화 4개.   ‡  persistence 대비는 두 모집단 4/4.   §  V_rot는 persistence 3/4 · 캠페인 4/4 · Δt > 15 ms PASS(양쪽); 우위는 소수 방전에 집중.   B.9 행은 컷 모집단·통합 301 방전.   §8aq 행은 641파일 전수 차분 상한이다.",
        10.5, GRAY, False)], tag="s3-footnote", space_after=1)
note(s, """
이 표가 이 덱의 중심이다. THESIS_RESULTS.md §8ab "Verdict — what is unconditional"의 국문판에
B.9(§8ac–§8an)의 세 행(문맥 포화, 계열 동률·비용, 1 ms 보류)을 더한 것이다.

[읽는 법]
"무조건부"는 컷·포함 두 모집단에서 모두 성립함을 뜻한다. "조건부"는 한쪽에서만 성립하므로
반드시 모집단을 붙여 말한다. "미확립"은 점추정은 양수이나 PR4를 넘지 못함을 뜻한다. B.9 행은
컷 모집단의 동결 B.1 매니페스트 위에서 측정되었고, 통합 재채점은 301 방전의 union이다.

[자주 나오는 추궁]
"conformal 32/32는 무엇의 32인가"에는 두 기준선 × 두 타깃 × 4 분할 × 두 모집단의 Winkler 비교 셀
전부에서 모델 구간이 이겼다고 답한다. "캠페인이 무조건부인데 왜 †가 붙는가"에는 시간 분할은 하나뿐이고
재현은 초기화 4개에 대한 것이라고 먼저 밝힌다. "b3k8 행의 -0.194는 실패가 아닌가"에는 실패가 아니라
컷 조건부이며 22번 Q7이 답한다. "1 ms는 왜 미결인가"에는 5세션 프로토콜의 p99 산포가 21.84배에 이르러
사전등록 규칙이 판정을 보류하였다고 답한다(24번 Q15).
출처: THESIS_RESULTS.md §8ab·§8al·§8am·§8aj, main_ko.tex §11.
""")

# --- 4. Why the full re-experiment ----------------------------------------
s = slide()
header(s, "재실험의 이유", "2026-08-12 방향 재정립: W = 4 프로토콜은 자체 장부의 세 결과로 폐기되었다", accent=ORANGE)
fcard(s, 0.55, 1.45, 3.95, 2.15, "폐기 근거 ① §8f 윈도 스윕", [
    "24-run 스윕의 plateau-최소 규칙이 반환하는",
    "값은 W = 2이다.",
    "관측 하나가 두 타깃을 plateau로 올리고",
    "이후 곡선은 평평하다(Tᵢ 0.190–0.246).",
    "W = 4는 skill로 정당화되지 않는다는 것이",
    "먼저 보고되었다.",
], accent=ORANGE, body_size=12, tag="s4-f")
fcard(s, 4.70, 1.45, 3.95, 2.15, "폐기 근거 ② §8c held-free", [
    "관측 V_rot의 54%가 계측 유지값이며,",
    "이는 평가뿐 아니라 학습도 오염한다.",
    "held 제거 학습이 V_rot을 4/4 개선하였다.",
    "확정 프로토콜은 held를 지도 타깃·이력·",
    "정규화 통계·모든 기준선의 보간 앵커에서",
    "제거한다.",
], accent=ORANGE, body_size=12, tag="s4-c")
fcard(s, 8.85, 1.45, 3.95, 2.15, "폐기 근거 ③ §8q 피팅 실패", [
    "관측 Tᵢ > 3 keV 스파이크가 측정 skill을",
    "대략 절반으로 만든다.",
    "재실행 전에 규칙을 고정하면 이 선택은",
    "사후 선별이 아니다. 따라서 컷을 사전등록",
    "하되, 배치된 시스템은 그런 행을 실제로",
    "만나므로 포함도 1차로 둔다.",
], accent=ORANGE, body_size=12, tag="s4-q")
fbullets(s, 0.55, 3.85, 12.25, 2.9, [
    ("결정: 결과를 하나씩 기워 붙이지 않고 단일 확정 프로토콜에서 전부 한 번에 재실행하였다. 교체되기 전의 수치는 잠정이며 인용하지 않는다.", 0),
    ("개정(2026-08-14 승상님): 포함(p100) 모집단은 민감도 행이 아니라 공동 1차이다. 무조건부 주장은 두 모집단 모두에서 성립해야 한다.", 1),
    ("기록: 이 개정은 컷 체제 윈도 계열 3/4 seed의 TEST 수치가 부분 언블라인딩된 뒤 결정되었다. 개정 방향이 기준 강화이므로 선택 편향은 반대로 작동한다.", 2),
    ("실행 순서: ① B.7 프로토콜 감사 → ② B.1 백본 관문(+ 인과 GP 병렬) → ③ B.2 탐색 → ④ B.3 해석가능 단 → ⑤ B.4 크기 → ⑥ B.5 전면 재채점 → ⑦ B.9 문맥·구조·비용. B.6(kHz Mirnov)은 비동기이다.", 0),
    ("B.7 감사(21개 제약, 무학습): 결측 원장 · MC lag-1 자기상관 · Tₑ/NBI 프로브가 641파일에서 정확히 재현되었고, 0.5 s 세그먼트 문턱은 쌍봉 분포의 골짜기(247k 간격 중 82개만 (0.1, 0.5) s)이다.", 0),
    ("교정 4건은 전부 실행 전 사전등록되었다. 컷을 로드 시점 결측 처리로 구현하였고(커밋 3598760, 캐시 서명 v5, 35/35 테스트 통과), 관문 대조군을 W = 4 계열에서 W = 2 계열로 교체하였다.", 1),
    ("§8h 감사 서사: 수기 전사가 논문·산출물 드리프트를 만든 적이 있으므로, 이후 모든 수치는 동결 산출물에서 자동 수집된다.", 0, GRAY),
], size=15, gap=8, tag="s4-bullets")
note(s, """
[요지]
이전 초고의 숫자는 방어하지 않는다. 자체 장부의 세 결과가 그 프로토콜을 독립적으로 기소하였고,
그래서 한 번에 전부 다시 실행하였다.

[B.7 감사에서 재확인된 사실]
관측 Tᵢ p99 = 2,089 eV, > 3 keV = 1,197행 = 관측 Tᵢ의 0.53%이다. CES_VT는 소수점 5자리·최소 간격
약 0.00004이므로 "연속 동일값 = forward-fill" 규칙의 오탐 채널이 없다. Tᵢ의 held는 226,991행 중
1행(경험적 상한)이다. 0.5 s 세그먼트 문턱은 쌍봉 분포의 골짜기에 있어 문턱 자체가 둔감하다.

[사전등록 문서와의 관계]
PREREGISTRATION_W2.md는 커밋으로만 변경되고, 해당 실험의 TEST 수치를 본 뒤의 변경은 금지된다.
각 배치 규칙 아래에 "집행 결과" 블록이 §8x/§8y/§8z/§8aa/§8ab로 이어 붙는다. B.9는 별도 문서
PREREGISTRATION_B9.md(축 A–D의 가설 H1–H6과 §3.2·§3.4 판정 규칙)로 사전등록되었다.
출처: THESIS_RESULTS.md §8v·§8w, PREREGISTRATION_W2.md §1·§1.1, PREREGISTRATION_B9.md.
""")

# --- 5. Protocol -----------------------------------------------------------
s = slide()
header(s, "프로토콜", "사전등록 · TEST 동결 · shot 군집 · 두 모집단 · 인과 GP")
fcard(s, 0.55, 1.45, 6.03, 2.5, "공통 고정 상수 (실행 전 커밋, 배치 전체 불변)", [
    "W = 2(§8f plateau-최소 규칙), held-free 학습·평가,",
    "Tᵢ > 3 keV 로드 시점 결측 처리(전 arm 자동 동일), 파일당 500,",
    "파일(shot) 단위 분할 상속 + TEST 격리 assert,",
    "10 epoch · batch 512 · lr 1e-3 AdamW · Tᵢ 음수 소프트 페널티,",
    "train-파일 전용 정규화(타깃 NaN-aware), seq_v2는 정의상 shot별 표준화 ON,",
    "V_rot 라우팅(빠른 진단 차단) 유지 — B.5의 W = 2 held-free 절제로 재확인되었다.",
], accent=BLUE, body_size=12, tag="s5-const")
fcard(s, 6.75, 1.45, 6.03, 2.5, "PR1–PR4 (기존 사전등록 계승)", [
    "PR1: headline 비교 대상은 PCHIP이며 사다리 전체의 병기가 의무이다.",
    "PR2: 보간은 모든 채점 지점에서 예측하며(미래 이웃 없으면 persistence 폴백),",
    "모집단 솎아내기는 금지되고 폴백률 보고가 의무이다.",
    "폴백률은 Tᵢ 0.3–0.4%, V_rot 40–44%이므로 V_rot의 'vs PCHIP'는",
    "5분의 2가 사실상 'vs persistence'이다.",
    "PR3: TEST 하한은 15 방전 · 3,000 관측 Tᵢ 이상이다(실제 96 방전).",
    "PR4: 유의는 shot 군집 paired bootstrap 10,000회의 95% CI가 0을 배제함이다.",
], accent=NAVY, body_size=12, tag="s5-pr")
fcard(s, 0.55, 4.13, 3.95, 2.05, "TEST 동결", [
    "B.2 탐색의 모든 모델 선택은 val에서만 하였다.",
    "TEST는 후보당 1회, 판정 규칙을 채점 전에",
    "커밋한 확증 실험에서만 채점하였다.",
    "B.2의 v3가 그 규율의 시험대였다. val 2/2",
    "유의였으나 TEST 1/4로 미승격되었다(22번 Q6).",
], accent=TEAL, body_size=12, tag="s5-freeze")
fcard(s, 4.70, 4.13, 3.95, 2.05, "shot 군집이 독립 단위이다", [
    "방전 내 인접 행은 강하게 자기상관되므로",
    "표본 3만 개를 독립 취급하면 CI가 가짜로",
    "좁아진다. SE 차이를 shot으로 집계해 shot을",
    "복원추출한다. 유효 표본은 방전 수(Tᵢ 약 96,",
    "V_rot 60–66)이며 검정력의 상한이다.",
], accent=TEAL, body_size=12, tag="s5-cluster")
fcard(s, 8.85, 4.13, 3.95, 2.05, "사다리에 인과 GP를 신설하였다", [
    "기존 GP에서 미래 이웃을 제거하고 과거 16",
    "이웃만 사용한다(NaN 조건 동일 → 채점 모집단",
    "불변, bit-identical 검증).",
    "이것이 최강 배치 기준선이다. seed 42 Tᵢ RMSE",
    "164.3 vs persistence 197.2(컷)이다.",
], accent=TEAL, body_size=12, tag="s5-gp")
note(s, """
[한 문장 정의]
TEST를 보기 전에 규칙을 문서로 고정하고(PR1–4 + 확정 프로토콜 + 배치별 판정 규칙), 방전을 재표본
단위로 하는 paired bootstrap 95% CI가 0을 제외하는지를, 모델 선택이 접근하지 않은 4개 독립 분할에서
반복 확인하였다.

[PR2 폴백률의 중요성 — V_rot 질문의 핵심]
W = 2 held-free에서 V_rot의 보간 폴백률이 40–44%이다. 즉 V_rot의 "vs PCHIP" 비교는 5분의 2가
사실상 "vs persistence"이다. V_rot 동률을 논할 때 반드시 함께 진술한다.

[B.9의 추가 규칙]
PREREGISTRATION_B9.md §3.2: |paired| < 0.02이고 유의 ≤ 1/4이면 동률, ≥ 3/4 유의이고 같은 방향이면
차이이다. §3.4: 전체 대비 결손이 0.02 미만이고 유의 결손이 ≤ 1/4인 최소 도달 범위가 포화점이다.
§4: 5개 독립 세션의 p99가 전부 예산 안일 때만 통과이며, [0.8 B, 1.25 B]는 경계이다.

[분할 규모]
기준 seed 42 TEST = 관측 Tᵢ 32,589행/96 방전, genuine V_rot 10,463행/60 방전(컷; 포함 32,721/10,461).
4 분할에서 Tᵢ 32.6–35.9k행/96 방전, V_rot 10.5–14.5k행/60–66 방전이다.
출처: PREREGISTRATION_W2.md §1·§2·§3, PREREGISTRATION_B9.md §3·§4, main_ko.tex §5.
""")

# --- 6. Novelty ------------------------------------------------------------
s = slide()
header(s, "novelty 검증", "계보 위에 서서 세 축으로 확장한다: 부재가 아니라 확장이다")
fbullets(s, 0.55, 1.42, 12.25, 1.2, [
    ("표현 규칙: \"관련 연구가 없다\"를 앞세우지 않는다. 부재 프레임은 반례 하나에 무너지고, 확장 프레임은 계보를 인정할수록 강해진다.", 0, ORANGE),
    ("교차 진단 재구성·시간 조밀화는 활발한 계열이다(Diag2Diag·COMPASS 시간 초해상·EAST 결측 Tₑ·FusionMAE·PanoMHD·RTCAKENN·EAST XCS·KSTAR EPED). 전부 인용하고 그 프로그램을 이어받는다고 쓴다.", 1),
], size=13.5, gap=6, tag="s6-frame")
fcard(s, 0.55, 2.72, 3.95, 2.15, "확장 ① 타깃", [
    "전자 채널(Thomson·ECE의 Tₑ·nₑ) 또는",
    "임의의 조밀 채널에서 희소 이온 채널",
    "(CES Tᵢ·V_rot)로 확장한다.",
    "선행에서 CER/CES가 등장할 때는 어디까지나",
    "입력이었고 희소 타깃이 아니었다.",
], accent=BLUE, body_size=12, tag="s6-t")
fcard(s, 4.70, 2.72, 3.95, 2.15, "확장 ② 매핑", [
    "동시각·기억 없는 매핑에서 타깃 자신의",
    "불규칙한 과거에 조건화된 인과 추정으로",
    "확장한다. 전체격자 인과 시퀀스를 쓰므로",
    "도달 범위가 고정 윈도가 아니라 세그먼트",
    "전체이며, B.9는 그중 약 50 ms가 실제로",
    "쓰임을 측정하였다.",
], accent=TEAL, body_size=12, tag="s6-m")
fcard(s, 8.85, 2.72, 3.95, 2.15, "확장 ③ 평가", [
    "가정된 재구성 가능성에서 사전등록 기준선에",
    "대한 타깃별 검정으로 확장한다. 현직 방법",
    "(미래를 읽는 보간)을 통계 bar로 세운다.",
    "이 프로토콜이 음성을 낼 수 있다는 것 자체가",
    "기여이며, V_rot 승리를 부정하였다.",
], accent=ORANGE, body_size=12, tag="s6-e")
fcard(s, 0.55, 5.05, 12.25, 1.6, "한정은 한 문장뿐이다 (논문 관련연구 절 마지막 문장)", [
    "\"우리가 아는 한 이 세 확장의 결합 — 희소 CXRS 이온 측정의 인과적 시간 결측 채움 — 은 아직 다뤄지지 않았으며, 이를 계열로부터의 이탈이 아니라 그 자연스러운 다음 단계로 제시한다.\"",
    "가장 가까운 선행과의 대비: Diag2Diag는 CER이 입력이고 타깃은 Thomson뿐이며 명시적으로 기억이 없다. COMPASS·EAST·FusionMAE는 전자 채널 또는 윈도 내 마스킹 복원이다. RTCAKENN은 dropout 강건성이지 타깃 희소성이 아니며 보간 기준선과의 통계 비교가 없다. EAST XCS는 같은 물리 채널의 순간 매핑이다. KSTAR EPED 연구는 CES를 입력으로 쓰고 회전을 향후 과제로 남긴다. 2차 문헌 검토(2026-08-20): Bishop 1993이 계보의 기원이고 RTCAKENN이 가장 가까운 친척이며, UF-CES는 실재하는 하드웨어이다.",
], accent=NAVY, body_size=12, tag="s6-caveat")
note(s, """
[제시 순서 — 승상님 2026-08-05 규칙]
① 계보 인정 → ② 3축 확장 → ③ 관례적 한정 한 문장. 이 순서를 뒤집지 않는다.

[정식 대비 — 원문에서 확인 가능한 범위]
Diag2Diag(Jalalvand et al., Nat. Commun. 2025, DIII-D)는 입력에 CER 포함(동시각), 타깃은 Thomson
Tₑ/nₑ이며 모델은 명시적으로 memory-less이다. COMPASS 시간 초해상(Imríšek et al., PPCF 2026)은 고속
복사 진단 → Thomson 프로파일이다. EAST 결측 Tₑ(Wang et al., NF 2025)는 타 진단 시계열 외생 회귀로
ECE형 Tₑ를 복구한다. FusionMAE(Yang et al., 2026, HL-3)는 88채널 마스킹 복원이며 윈도 내 비인과이다.
RTCAKENN(Shousha et al., NF 2024)은 실시간 운동학 프로파일이며 결측은 입력 강건성 조건이다. EAST XCS
(Lin et al., NF 2024)는 같은 Tᵢ·V_rot이지만 입력이 또 다른 도플러 분광기이다. KSTAR EPED 프로파일
(Jung, Kim & Kang, JKPS 2026)은 CES가 입력이고 회전은 향후 과제로 명시한다.

[비대칭이 novelty의 산출물이다]
Tᵢ↔V_rot 정보 비대칭은 확장 ③(검정 가능한 평가)이 만들어낸 발견이다. 가정하였다면 나올 수 없었던
결과이다. B.9의 계열 동률 결과도 같은 계보에 서 있다: 다중 센서 논문의 합성곱 관행에 대해 "계열은 skill을
정하지 않고 비용을 정한다"는 측정된 답을 준다.
출처: main_ko.tex §2 관련 연구, docs/paper/NOVELTY.md.
""")

# --- 7. B.1 backbone gate --------------------------------------------------
s = slide()
header(s, "모델 선택 ①", "B.1 백본 관문: 사전등록 4조건을 전부 충족하였다 (§8x)")
add_image_fit(s, os.path.join(FIG, "fig_gate_b1.png"),
              Inches(0.55), Inches(1.45), Inches(6.5), Inches(5.15))
fcard(s, 7.25, 1.45, 5.55, 2.15, "판정 규칙 (실행 전 고정, 전부 충족 시에만 채택)", [
    "① split별 init-평균 paired Tᵢ 부호가 4/4 양수여야 한다.",
    "② 16-run paired 평균의 run-클러스터 bootstrap 95% CI가 0을 배제해야 한다.",
    "③ 예산 균등화 arm(고정 10 epoch)에서도 부호가 유지되어야 한다.",
    "④ V_rot 유의 열세 run이 0개여야 한다.",
    "설계: seq_v2 16 run = 분할 seed 4 × 초기화 seed 4이며, 각 run은",
    "자기 분할의 W = 2 held-free 윈도 대조군과 paired 비교된다.",
], accent=NAVY, body_size=11.5, tag="s7-rule")
fbullets(s, 7.25, 3.75, 5.55, 2.85, [
    ("paired Tᵢ는 16/16 양수, 13/16 개별 유의였다.", 0, GREEN),
    ("분할별 초기화 평균은 +0.129 / +0.059 / +0.078 / +0.058이며 초기화 산포가 분할 산포보다 훨씬 작다.", 1),
    ("pooled +0.081, run-클러스터 CI [+0.067, +0.096]이다.", 1, GREEN),
    ("예산 균등화 +0.063 / +0.033 / +0.045 / +0.030으로 4/4 양수이며, 이득은 예산이 아니라 아키텍처 효과이다.", 1),
    ("V_rot 유의 열세는 0/16이다(유의 우세 8/16).", 1),
    ("확증 4 분할의 seq − 윈도는 컷 +0.130/+0.058/+0.062/+0.044, 포함 +0.053/+0.024/+0.047/+0.029로 8/8 양수이다.", 0),
    ("윈도 대조군은 인과 GP와 동률(1/4)이고 시퀀스 백본은 4/4+4/4이다. B.9는 이 +0.081을 도달 범위 −0.065와 구조 ≈ −0.016으로 분해하였다(§8af).", 0, TEAL),
], size=12.5, gap=6, tag="s7-res")
note(s, """
[관문의 설계 의도]
"왜 시퀀스 모델인가"를 사후 정당화로 만들지 않기 위해, 채택 조건 4개를 실행 전에 문서로 고정하고
그 문서를 커밋하였다(PREREGISTRATION_W2.md §4). 대조군도 감사 (B)-4에 따라 W = 4 계열이 아니라
확정 프로토콜과 정합한 W = 2 계열로 교체하였고, 컷 도입이 모집단을 바꾸므로 양쪽 모두 재학습하였다.

[Stage A — 윈도 대조군(w2cut) 4 seed]
Tᵢ vs PCHIP +0.051 n.s. / +0.203 / +0.208 / +0.231로 3/4 PASS이다(컷 모집단은 앵커가 깨끗해져
보간을 강화한다). vs persistence 4/4, V_rot vs persistence 3/4이다.

[claim-2 게이트]
윈도 모델 vs 인과 GP는 1/4이므로 사전 규칙대로 주장 2를 일단 한정 재서술하고 "인과 GP 초과"를 B.2
목표로 올렸다. §8y에서 백본 자신이 4/4로 이겨 복권되었다(8번 슬라이드).

[B.9가 더한 분해]
§8ae·§8af: 백본−윈도 +0.081은 도달 범위(−0.065, 20 ms에서 학습한 seq_v2)와 구조(≈ −0.016)로
분해된다. 즉 백본이 사는 것의 4/5는 얼마나 멀리 보는가이고 1/5는 무엇인가이다.

[수치 확인]
2.2 h wall, 산출물 data/.b1_*, 판정 data/.b1_gate_summary.json이다. 분할은 동결 W = 2 스윕 manifest에
seed별로 고정하였고 test 격리 assert가 4 seed 전부 통과하였다.
출처: THESIS_RESULTS.md §8x·§8af, PREREGISTRATION_W2.md §4.
""")

# --- 8. B.2 + backbone architecture ---------------------------------------
s = slide()
header(s, "모델 선택 ②", "B.2 후보는 미승격이며 주 모델은 seq_v2로 남는다 (§8y)")
add_image_fit(s, os.path.join(FIG, "fig_architecture_seq.png"),
              Inches(0.55), Inches(1.45), Inches(6.6), Inches(5.15))
fcard(s, 7.35, 1.45, 5.45, 2.3, "백본 seq_v2 (357,570 파라미터)", [
    "22채널 전체격자 시퀀스 위의 독립 인과 LSTM 2개이다.",
    "Tᵢ 분기(2층 160)는 전체 상태(빠른 진단 15 + 두 타깃 이월값·신선도·flag +",
    "log Δt)를, V_rot 분기(1층 64)는 비-빠른 7채널만 읽는다.",
    "라벨 없는 행도 맥락으로 유지하고 희소성은 masked loss로 처리한다.",
    "라우팅은 head가 아니라 인코더에서 성립한다. 공유 상태였다면 배선을",
    "어떻게 하든 빠른 진단이 V_rot로 누출된다.",
], accent=TEAL, body_size=11.5, tag="s8-arch")
fbullets(s, 7.35, 3.9, 5.45, 2.7, [
    ("후보 seq v3는 v2 + 관측마스킹 인과 attention readout(396,930, 출력 사영 0 초기화 → 학습이 정확히 v2 함수에서 출발)이다.", 0),
    ("탐색은 val 전용(분할 42/7)이었다. paired Tᵢ +0.024* / +0.037*, vs 인과 GP +0.170* / +0.132*였다.", 1),
    ("TEST 확증(규칙 사전 커밋): +0.009 / +0.013 / +0.033* / +0.020으로 4/4 양수 · 유의 1/4이며 기준 3/4 이상에 미달하여 미승격되었다.", 0, RED),
    ("val에서 2/2 유의였다는 사실이 승격 bar를 TEST에 두는 이유이다. 탐색 분할 쌍이 유리하였던 것이다.", 1),
    ("부수 소득: v3가 인과 GP를 4/4로 넘고, 교차 검증에서 백본 자신이 인과 GP를 Tᵢ 4/4 유의로 이겼다(+0.078/+0.133/+0.138/+0.105). claim 2가 복권되었다.", 0, GREEN),
], size=12.5, gap=6, tag="s8-b2")
note(s, """
[이 슬라이드의 방어 가치]
"모델을 고르고 나서 이유를 붙인 것이 아닌가"에 대한 실물 반증이다. val에서 유의하였던 후보를 사전
규칙 때문에 떨어뜨렸고 그 기록을 §8y에 남겼다.

[v3의 성격]
iter009의 관측마스킹 attention을 시퀀스 프레임으로 옮긴 것이다(각 타깃 경로가 자기 과거의 신선한 관측
스텝들에만 attend, shifted mask·strict lower-triangular). 작지만 일관되게 양수인 실제 메커니즘이므로
폐기가 아니라 "B.4 인접 스케일링 질문"으로 남겼다. B.9의 계열 비교(§8ak)는 attention 계열이 70 ms에서
같은 문맥의 LSTM보다 −0.023(3/4 유의) 뒤지고 150 ms부터 동률임을 보였으며, 이 문제에서 attention을
추천할 근거는 없다.

[백본 구조의 검증 가능한 주장]
빠른 채널 15개를 어떻게 섭동해도 V_rot 출력이 bit-identical이다(seq_v2·b3k8·윈도 대조군 모두 같은
시험 통과). 라우팅은 학습된 것이 아니라 구조이다.

[학습]
AdamW 1e-3, batch 16 세그먼트, val masked MSE 조기 종료(patience 6, 상한 30; 확증 run은 14–25 epoch에서
종료)이다.
출처: THESIS_RESULTS.md §8y·§8ak, main_ko.tex §4·§7.
""")

# --- 9. RMSE ladder --------------------------------------------------------
s = slide()
header(s, "결과 ①", "RMSE 사다리: 정보 접근이 다른 8개 팔을 같은 행·같은 마스크로 채점하였다")
add_image_fit(s, os.path.join(FIG, "fig_rmse_ladder.png"),
              Inches(0.55), Inches(1.45), Inches(6.3), Inches(5.15))
ftable(s, 7.05, 1.45, [2.5, 1.65, 1.65],
       ["arm (seed 42, 컷)", "Tᵢ eV", "V_rot km/s"],
       [
           [("seq_v2 (나우캐스터)", TEAL, True), ("157.8", TEAL, True), ("23.6", TEAL, True)],
           ["윈도 대조군 (W = 2)", "169.2", "26.1"],
           [("인과 GP", NAVY, True), "164.3", "28.8"],
           ["persistence", "197.2", "33.4"],
           ["AR (국소, 인과)", "472.2", "51.0"],
           [("GP (오프라인)", GRAY, False), "153.8", "24.7"],
           [("선형 보간 (오프라인)", GRAY, False), "169.8", "29.0"],
           [("PCHIP (오프라인)", GRAY, False), "173.6", "30.2"],
       ],
       row_h=0.40, head_h=0.42, size=12, head_size=11.5, tag="s9-ladder")
fbullets(s, 7.05, 5.2, 5.75, 1.45, [
    ("백본은 인과 GP보다 RMSE가 Tᵢ 4% · V_rot 18% 낮고, 오프라인 GP와는 동률(153.8 vs 157.8)이다.", 0),
    ("포함 모집단은 seq_v2 363.0 / 23.7, PCHIP 412.4 / 30.2, 인과 GP 394.6 / 28.8, persistence 478.0 / 33.4이다. 스파이크가 Tᵢ RMSE를 두 배 이상 키우지만 순서는 불변이다.", 1),
    ("점선 위는 미래를 읽는 팔(배치 불가), 아래는 배치 가능한 팔이다.", 1, GRAY),
], size=12.5, gap=5, tag="s9-notes")
note(s, """
[사다리의 요점]
같은 (file, row) 집합·같은 per-target mask에서 8개 팔을 채점한다. 보간은 타깃 자기 값을 읽지 않고
세그먼트 경계를 넘지 않으며, 경계 밖 이웃이 필요하면 persistence로 폴백한다(PR2).

[인과 GP를 새로 만든 이유]
§8p에서 오프라인 GP가 최강 오프라인 팔이고 모델과 동률임이 드러났다. 배치 주장을 persistence 기준으로
말하는 것은 너무 관대하므로, 미래 이웃을 제거한 인과 GP를 사다리의 최상단 배치 가능 팔로 세우고 그
기준에서 판정한다. NaN 조건을 ar_local과 동일하게 유지해 채점 모집단이 바뀌지 않도록 구조적으로
보장하였다(기존 키 bit-identical 검증).

[AR이 나쁜 이유]
국소 AR은 불규칙 격자에서 발산에 가까운 외삽을 한다. 사다리에 남기는 이유는 "인과 시계열 모형 대표"라는
자리를 비워두지 않기 위해서이다.

[B.9와의 연결]
§8ac의 실시간 사다리(같은 세션)에서 인과 GP의 p99는 2.34 ms, seq_v2 1-step은 1.49 ms, 윈도 W=2는
4.46 ms이며 순서는 세션과 무관하게 유지되었다.
출처: main_ko.tex §6.1, THESIS_RESULTS.md §8ac §2.
""")

# --- 10. Headline forest + gap --------------------------------------------
s = slide()
header(s, "결과 ②", "Headline: Tᵢ는 두 모집단·네 분할 전부에서 미래를 읽는 보간을 이겼다")
add_image_fit(s, os.path.join(FIG, "fig_forest.png"),
              Inches(0.55), Inches(1.45), Inches(6.6), Inches(5.15))
fbullets(s, 7.25, 1.45, 5.55, 2.9, [
    ("Tᵢ vs PCHIP는 컷 +0.174 / +0.248 / +0.257 / +0.264(4/4, 평균 +0.236)이다.", 0, GREEN),
    ("포함은 +0.225 / +0.238 / +0.292 / +0.316(4/4, 평균 +0.268)이다.", 0, GREEN),
    ("8개 셀 전부 인과 GP(+0.08~+0.17)와 persistence(+0.36~+0.46)도 이기므로 무조건부이다.", 1),
    ("오프라인 GP와는 동률(-0.05~+0.11, 1/8 유의)이며 오프라인 주장의 상한을 그대로 보고한다.", 1, GRAY),
    ("포함 수치가 더 높은 이유는 컷이 없으면 스파이크가 보간 앵커를 오염시켜 모든 arm이 PCHIP 대비 좋아 보이기 때문이며, 어느 한 모집단만 인용하지 않는 이유이다.", 1, ORANGE),
    ("V_rot는 점추정 8/8 양수이나 PR4는 1/4(컷) · 2/4(포함)이므로 동률로 보고한다.", 0),
], size=12.5, gap=6, tag="s10-head")
fcard(s, 7.25, 4.5, 5.55, 2.1, "간극 층화 (4 분할 통합, 방전 군집) — 컷 / 포함", [
    "Tᵢ ≤ 15 ms: +0.239* / +0.299* (n = 134,546 / 135,317)",
    "Tᵢ > 15 ms: +0.268* / +0.206* (n = 3,422 / 3,334, 265 / 263 방전)",
    "Tᵢ > 45 ms: +0.267* / -0.004 n.s. (429–460행, 101–104 방전)",
    "V_rot > 15 ms: +0.418* / +0.432* (130 방전)이며 유일한 V_rot 무조건부 양성이다.",
    "가장 넓은 층(V_rot > 45 ms)은 14행뿐이므로 채점하지 않는다.",
], accent=BLUE, body_size=11.5, tag="s10-gap")
note(s, """
[정확한 CI]
컷: +0.174 [+0.097,+0.236] / +0.248 [+0.188,+0.295] / +0.257 [+0.199,+0.302] / +0.264 [+0.188,+0.320].
포함: +0.225 [+0.109,+0.293] / +0.238 [+0.153,+0.302] / +0.292 [+0.232,+0.344] / +0.316 [+0.186,+0.392].
V_rot 컷: +0.390 [+0.077,+0.591]* / +0.183 [-0.028,+0.280] / +0.135 [-0.358,+0.269] / +0.305 [-0.049,+0.437].
포함: +0.384* / +0.195* / +0.132 / +0.304.

[간극 슬라이드의 핵심 문장]
비인접 영역(Δt > 15 ms)에서도 미래를 읽는 PCHIP을 두 모집단 모두에서 이긴다. > 45 ms가 포함에서
동률인 이유는 429행/101 방전 규모에서 스파이크 앵커 몇 행이 층 전체를 지배할 수 있기 때문이며, 그것을
그대로 조건부로 보고한다.

[V_rot > 15 ms를 반드시 진술한다]
전역 동률만 말하면 "회전은 아무것도 안 된다"로 오해된다. 보간이 가장 어려운 영역에서는 두 모집단 모두
PCHIP을 이긴다. 이것이 논문의 유일한 V_rot 무조건부 양성이다. B.9의 승패 방전 분석(§8an)은 이 양성이
회전이 실제로 변하는 방전에 집중됨을 보였다.
출처: main_ko.tex §6.2·§6.4, THESIS_RESULTS.md §8ab.
""")

# --- 11. Two stress tests --------------------------------------------------
s = slide()
header(s, "결과 ③", "스트레스 2종: 두 주장을 가르던 시험을 백본은 둘 다 통과하였다", accent=ORANGE)
add_image_fit(s, os.path.join(FIG, "fig_campaign.png"),
              Inches(0.55), Inches(1.42), Inches(6.15), Inches(3.05))
fcard(s, 0.55, 4.6, 6.15, 2.05, "캠페인(시간) 분할 — 어떤 test 방전도 train보다 앞서지 않는다", [
    "train 416 [30801–31991] / val 128 [32002–32310] / test 97 [32312–32751], 초기화 4개이다.",
    "윈도 OFF는 컷 2/4(+0.027 / +0.091* / -0.001 / +0.061*), 포함 0/4, 인과 GP 0/4이다.",
    "윈도 ON(shot별 표준화, 컷)은 4/4(+0.103 / +0.107 / +0.094 / +0.107)로 수리가 작동하였다.",
    "seq_v2는 컷 +0.187 / +0.174 / +0.181 / +0.177, 포함 +0.173 / +0.202 / +0.198 / +0.184이다.",
    "PCHIP 4/4+4/4, 인과 GP 4/4+4/4, seq − 윈도 8/8 유의이다.",
    "V_rot vs persistence도 seq는 4/4 양쪽이다(윈도 대조군은 0/4).",
], accent=TEAL, body_size=11, tag="s11-camp")
fcard(s, 6.75, 1.42, 6.05, 2.3, "MNAR — 실제 결측점 분포로 재가중 (W = 2 in-domain)", [
    "층은 Δt(15/25/45 ms) × 입력만의 활동 flag이며, 30 미만 층 기각, 가중 격자에도 컷 적용이다.",
    "도달 범위(사전등록이 요구한 보고)는 결측 Tᵢ의 54–68%, V_rot의 4–6%만 도메인 내이므로",
    "재가중 V_rot는 결측 질량의 1/20에 대한 답이라 결론을 내지 않는다.",
    "Tᵢ vs PCHIP는 컷 +0.140 / +0.164* / +0.203 / +0.283*(2/4),",
    "포함 +0.140* / +0.217* / +0.167* / +0.221*(4/4)이다.",
    "Tᵢ vs persistence는 4/4 양쪽(+0.28~+0.44, 보정 비용 최대 0.12)이다.",
], accent=BLUE, body_size=11.5, tag="s11-mnar")
fcard(s, 6.75, 3.82, 6.05, 1.4, "원인은 추측이 아니라 측정되었다", [
    "train → test 드리프트는 BES 1.22σ · ECEI 0.53σ(스케일비 0.75 / 0.62)인 반면",
    "CES 타깃은 0.115σ(1.06)이다. 빠른 진단 경로가 5–11배 더 이동한다.",
    "train-파일 전용 정규화가 캠페인 이동에서 깨지는 지점이다.",
], accent=ORANGE, body_size=11.5, tag="s11-drift")
fcard(s, 6.75, 5.32, 6.05, 1.33, "어느 주장이 어디서 성립하는가", [
    "무작위 분할·관측점: PCHIP 4/4·4/4, 인과 GP 4/4·4/4",
    "결측점 재가중: PCHIP 2/4·4/4, persistence 4/4·4/4",
    "캠페인 시간 분할: PCHIP 4/4·4/4, 인과 GP 4/4·4/4",
], accent=NAVY, body_size=11.5, tag="s11-sum")
note(s, """
[이 슬라이드가 이전 덱과 가장 크게 달라진 곳]
W = 4 초고에서는 두 스트레스가 오프라인 우위를 죽였다(MNAR 1/4, 캠페인 0/4). 확정 프로토콜의 백본은
둘 다 통과한다. 단, 윈도 대조군은 여전히 캠페인에서 붕괴하며 그것을 같은 슬라이드에 함께 싣는다.
달라진 것은 프로토콜이 아니라 모델이라는 것이 요점이다.

[MNAR 방법]
관측·결측 양쪽에서 계산 가능한 공변량으로 사후 층화하고 결측 행의 층 분포로 재가중한다. 층 커버리지는
Tᵢ 0.99–1.00, V_rot 0.73–0.76이다. 가정(층 내 교환 가능성)은 명시한다. 재가중 CI가 넓은 이유는 > 45 ms
층이 얇기 때문이다.

[캠페인의 남는 주의 — 먼저 진술한다]
분할은 하나(97 파일의 시간 블록)이고 반복은 초기화 4개이다. 컷 seq_v2 run 2/4가 30-epoch 상한에서
종료되었다. shot별 표준화는 현재 오프라인 형태이며 배치 가능한 인과 running 형태(EWMA)는 미측정이다.
§8am의 통합 재채점도 "방법의 기대 skill"을 추정하며 단일 체크포인트의 배치 주장은 여전히 이 시간 분할이
담당한다.
출처: main_ko.tex §6.5·§6.6, THESIS_RESULTS.md §8ab·§8am.
""")

# --- 12. Ablation / asymmetry ---------------------------------------------
s = slide()
header(s, "결과 ④", "정보 비대칭: 절제가 두 모집단의 필요성까지 설명한다")
add_image_fit(s, os.path.join(FIG, "fig_ablation.png"),
              Inches(0.55), Inches(1.45), Inches(6.1), Inches(4.0))
ftable(s, 6.85, 1.45, [1.95, 2.0, 1.95],
       ["절제 (평가 시)", "컷", "포함"],
       [
           [("Tᵢ 전체", NAVY, True), "+0.173", "+0.238"],
           ["Tᵢ  no_fast", ("-0.125", RED, True), ("+0.201", ORANGE, True)],
           ["Tᵢ  no_history", ("-2.11", RED, True), ("-1.16", RED, True)],
           [("V_rot 전체", NAVY, True), "+0.213", "+0.206"],
           ["V_rot  no_fast", ("+0.000 ×4", TEAL, True), ("+0.000 ×4", TEAL, True)],
           ["V_rot  no_history", ("-2.89", RED, True), ("-3.51", RED, True)],
       ],
       row_h=0.42, head_h=0.42, size=12, head_size=11.5, tag="s12-abl")
fbullets(s, 6.85, 4.5, 5.95, 2.1, [
    ("no_fast Tᵢ paired는 컷 -0.25*/-0.38*/-0.42*/-0.43*, 포함 -0.03*/-0.04/-0.03/-0.09*이다.", 1, GRAY),
    ("컷에서 Tᵢ 마진은 빠른 진단 정보로 만들어진다. 빠른 채널을 0으로 두면 보간 아래로 떨어진다(-0.10~-0.18). 물리 경로는 전자–이온 충돌 결합(ECEI Tₑ, BES nₑ)이다.", 0),
    ("포함에서는 이력-전용 모델도 PCHIP를 +0.15~+0.23 이기고 빠른 채널은 0.03–0.09만 더한다. p100 마진에는 스파이크-강건성 성분이 섞여 있다.", 0, ORANGE),
    ("이것이 두 모집단을 함께 보고하는 측정된 이유이자, 빠른 진단 기여를 분리하는 것이 컷 모집단인 이유이다.", 1, ORANGE),
    ("V_rot 정보는 전부 CES 이력이다. 빠른 채널 섭동에 출력이 bit-identical(8/8)이며, NBI 토크 미관측 + Mirnov 100 Hz 앨리어싱이 원인이다.", 0, TEAL),
], size=12, gap=5, tag="s12-read")
note(s, """
[세 번째 행이 이 슬라이드의 발견이다]
"p100만 쓰면 되지 않는가"에 대한 측정된 답이다. 포함 모집단에서 이력-전용(빠른 진단 제거) 모델조차
PCHIP를 +0.15~+0.23 이긴다. 보간의 앵커가 스파이크라서 학습된 모델이 그것을 할인할 수 있기 때문이다.
즉 p100 마진의 일부는 "빠른 진단이 나르는 정보"가 아니라 "스파이크에 강건함"이다. 빠른 진단의 기여를
분리하는 모집단은 컷이다.

[재확인된 두 라우팅 사실 — held-free, W = 2, 두 모집단]
① 빠른 채널을 0으로 만들어도 V_rot 출력이 bit-identical이므로 라우팅은 학습이 아니라 구조이다.
② 이력을 없애면 두 타깃이 모두 붕괴(skill -1~-8)하므로 100 Hz 진단만으로는 CES 앵커 없이 어느 타깃도
   예측되지 않는다.

[물리 근거와 B.9의 세 번째 증거]
Tᵢ 경로는 충돌성 e–i 결합(τ_ei ∝ Tₑ^1.5/nₑ)이며 ECEI(Tₑ)·BES(nₑ)가 정보를 나른다. V_rot 경로 부재는
NBI 토크 미관측(Tₑ~V_rot r = +0.024, p = 0.58) + Mirnov 앨리어싱이다. §8an은 세 번째 독립 증거를
더한다: 방전 단위 Tₑ 수준은 V_rot 승패를 예측하지 못하며(ρ = -0.031), 승리 방전은 뜨거운 방전이 아니라
움직이는 방전이다.
출처: main_ko.tex §6.7, THESIS_RESULTS.md §8ab-8·§8an.
""")

# --- 13. Complexity ladder + width sweep -----------------------------------
s = slide()
header(s, "결과 ⑤", "복잡도 사다리와 크기 축: 상한은 추정기가 아니라 정보에 있다")
add_image_fit(s, os.path.join(FIG, "fig_ladder_scaling.png"),
              Inches(0.55), Inches(1.45), Inches(6.6), Inches(4.1))
ftable(s, 7.35, 1.45, [2.45, 1.55, 1.45],
       ["arm", "Tᵢ 컷", "Tᵢ 포함"],
       [
           ["persistence (0)", "-0.264", "-0.288"],
           ["anchor+Δ (1,258)", "-0.261", "-0.287"],
           [("b3k8 (21,498)", TEAL, True), ("+0.237", TEAL, True), ("+0.126", ORANGE, True)],
           ["윈도 대조군 (201,258)", "+0.173", "+0.238"],
           [("seq_v2 백본 (357,570)", NAVY, True), "+0.236", "+0.268"],
       ],
       row_h=0.42, head_h=0.42, size=12, head_size=11.5, tag="s13-ladder")
fbullets(s, 7.35, 4.1, 5.45, 2.5, [
    ("b3k8은 persistence + 8개 유계 latent의 선형 보정(정확 분해, readout 0 초기화 → 학습이 정확히 persistence에서 출발)이다.", 0),
    ("b3 − anchor는 컷 +0.35~+0.42 · 포함 +0.29~+0.34로 4/4* 양쪽이다.", 1),
    ("b3 − seq_v2는 컷 평균 +0.002(CI 전부 0 포함)로 백본 skill 전부가 유계 수 8개로 압축되고, 포함 평균 -0.194(4/4*)이다.", 1, ORANGE),
    ("컷 조건부인 이유: 포함 test의 0.6–1.3% 행(persistence 오차 > 2 keV)이 b3 Tᵢ SSE의 73–83%(모든 arm 70–83%)이며, 유계 보정은 스파이크 이월값을 살리지 못한다.", 1, GRAY),
    ("B.4 폭 스윕 24/40/80/160/260(34k/49k/114k/358k/879k)은 +0.230/+0.236/+0.235/+0.236/+0.230으로 160 대비 ±0.008, V_rot 불변이다. B.9는 여기에 1,808 파라미터 tcn2k가 인과 GP를 4/4로 이긴다는 사실을 더하였다(§8ai).", 0, RED),
], size=12, gap=5, tag="s13-notes")
note(s, """
[두 문장 요약]
① 컷 모집단에서 백본의 Tᵢ skill 전부가 "persistence + 8개의 유계 수"로 압축된다.
② 26배 폭 스윕이 평평하다. 두 사실을 합치면 상한은 추정기 용량이 아니라 정보이다.

[정확한 paired 수치]
b3 − seq_v2(컷): -0.009 / -0.005 / +0.026 / -0.004로 평균 +0.002, PR4 4/4, 인과 GP 4/4이다.
포함: -0.16~-0.21, 4/4 유의이며 윈도 대조군보다도 3/4 아래이다.
anchor+Δ는 W = 2에서 persistence로 붕괴한다(기울기 항이 관측 2행을 요구하기 때문).

[B.9가 더한 하한]
§8ai: 크기를 맞춘 순환 arm 대비 합성곱 arm은 +0.027~+0.040(3–4/4 유의)이며, tcn2k(1,808)는 인과 GP를
4/4로 이기고 백본과 +0.001 차이이다. 계열은 상한이 아니라 상한에 이르는 데 필요한 파라미터 수를 정한다.
"작은 모델이 필요하다"는 배치 논거는 §8ah가 무효화하였고(357k 백본도 비용이 아님), 남는 것은 측정 자체이다.

[해석 가능성 주장의 정확한 형태]
"1,258 파라미터가 마진의 31.5%를 회수한다"(W = 4 시대 문장)는 폐기되었다. 현재 문장은 "21,498 파라미터의
정확 분해 모델이 컷 모집단에서 백본과 동급(+0.002)이고 포함 모집단에서는 -0.194로 벌어진다"이며 조건을
반드시 붙인다.
출처: main_ko.tex §6.9, THESIS_RESULTS.md §8z·§8aa·§8ai.
""")

# --- 14. Peak stratification ----------------------------------------------
s = slide()
header(s, "결과 ⑥", "우위는 고변동 구간에 집중된다 (peak 층화, TEST, seq_v2)")
add_image_fit(s, os.path.join(FIG, "fig_peak.png"),
              Inches(0.55), Inches(1.45), Inches(6.5), Inches(5.15))
fcard(s, 7.25, 1.45, 5.55, 2.15, "Tᵢ — 무조건부", [
    "peak 층은 컷 +0.45~+0.61, 포함 +0.62~+0.72로 8/8 PASS이다.",
    "본류는 +0.09~+0.20(컷 4/4) · +0.06~+0.19(포함 2/4)이다.",
    "보간은 매끄러운 본류에서 사실상 최적이고 모델의 값은 활동 구간에 있다는",
    "문장은 두 모집단 모두에서 성립한다.",
], accent=GREEN, body_size=12, tag="s14-ti")
fcard(s, 7.25, 3.75, 5.55, 2.15, "V_rot — 비대칭은 지역적이다", [
    "peak 층은 +0.54~+0.79(8/8 양수, PASS 각 2/4; persistence 대비",
    "+0.75~+0.86으로 8/8 PASS)이다.",
    "본류는 약 0(-0.07~+0.15, 0/8)이다.",
    "전역으로는 동률이어도 고활동 이웃에서는 이력 기반 V_rot 예측기도",
    "값을 더한다. §8an은 같은 사실을 방전 단위에서 확인하였다: 승률은",
    "조용한 방전 34%, 변동 큰 방전 55%이다.",
], accent=TEAL, body_size=12, tag="s14-vrot")
fband(s, 7.25, 6.08, 5.55, 0.56,
      [("peak 층은 입력만으로 계산되는 활동도 프록시로 정의되며 타깃 값을 보고 고른 층이 아니다. 방전 단위 승패의 예측 변수는 peak 비율이 아니라 타깃 산포이다(§8an).",
        11.5, GRAY, False)], tag="s14-caption")
note(s, """
[peak 층화가 ELM/천이 주장을 운반하는 이유]
사전등록은 "과도 포착" 주장을 모집단 선택(컷/포함)이 아니라 peak 층화 분석이 나르도록 정하였다. 그래서
8/8 PASS라는 결과가 모집단 논쟁과 독립적으로 성립한다.

[층의 정의]
입력 전용 활동도 프록시(빠른 진단의 국소 변동)로 층을 나눈다. 타깃 값을 사용하지 않으므로 "결과를 보고
유리한 구간을 골랐다"는 반박이 구조적으로 차단된다.

[V_rot 해석과 §8an의 교정]
전역 동률 + peak 양수 + 본류 약 0의 조합은 "회전 정보가 전혀 없다"가 아니라 "이력이 나르는 회전 정보는
활동 구간에서만 보간을 넘는다"는 뜻이다. 단, 행 단위 peak 비율은 방전 단위 승패를 예측하지 못하였다
(고-peak 절반 0.435 vs 저-peak 절반 0.528, 방향이 반대). 승패를 예측하는 변수는 방전 내 타깃의 산포이다.
두 양은 다른 양이며, 이 슬라이드의 결과는 행 단위 PCHIP 대비 결과로 한정하여 인용한다.
출처: main_ko.tex §6.10, THESIS_RESULTS.md §8ab·§8al §4·§8an.
""")

# --- 15. Transient case study ---------------------------------------------
s = slide()
header(s, "결과 ⑦", "held-out TEST 방전 #31815: 사례로 보는 과도 추적 (seq_v2)")
add_image_fit(s, os.path.join(FIG, "fig_transient_seq_31815.png"),
              Inches(0.55), Inches(1.42), Inches(8.9), Inches(5.3))
fcard(s, 9.65, 1.45, 3.15, 2.4, "이 방전의 수치", [
    "Tᵢ skill +0.42",
    "V_rot skill +0.21",
    "",
    "모델은 과도 구간을 추적하고,",
    "미래를 읽는 보간은 관측 앵커",
    "사이를 매끄럽게 지나간다.",
], accent=TEAL, body_size=12.5, tag="s15-num")
fcard(s, 9.65, 4.0, 3.15, 2.6, "사례의 지위", [
    "사례는 예시일 뿐이며,",
    "통계 판정은 peak 층화(8/8 PASS)와",
    "간극 층화가 운반한다.",
    "",
    "이 방전은 TEST 분할에 속하며",
    "모델 선택 과정에서 한 번도",
    "열람되지 않았다.",
], accent=GRAY, body_size=12.5, tag="s15-status")
note(s, """
[사례를 쓰는 방식]
"이 그림이 증거이다"라고 말하지 않는다. 통계 판정은 peak 층화(8/8 PASS)와 간극 층화가 운반하고, 사례는
그 판정이 시간 축에서 어떻게 보이는지를 보여줄 뿐이다. "체리피킹이 아닌가"에는 이 순서로 답한다. §8an의
3분위 표는 이런 방전이 전체의 1/3(변동 큰 3분위)에 해당하며 거기서 Tᵢ 승률이 85%임을 보인다.

[이 그림이 W = 4 시대 그림과 다른 점]
fig_transient_31815.png(구 윈도 모델, W = 4)은 폐기되었다. 여기 쓰는 것은 fig_transient_seq_31815.png이며
확정 프로토콜의 백본 seq_v2가 낸 예측이다. 같은 계열의 두 번째 방전 그림(#30842)도 있으나 이 덱에는 넣지
않았다. #31815는 문헌 검토에서 vetted된 과도 시연 방전이다.
출처: docs/presentation/figures/fig_transient_seq_31815.png (§8ab 재생성), THESIS_RESULTS.md §8an.
""")

# --- 16. Data audits -------------------------------------------------------
s = slide()
header(s, "데이터 감사", "두 개의 감사가 프로토콜을 강제하였다: 유지값과 피팅 실패", accent=RED)
add_image_fit(s, os.path.join(FIG, "fig_missing.png"),
              Inches(0.55), Inches(1.42), Inches(5.45), Inches(2.6))
fcard(s, 0.55, 4.15, 5.45, 2.5, "감사 1 — 유지(forward-fill) 값", [
    "관측 V_rot의 54%가 계측 유지값이다(run 최대 1,214행, 499/641 파일).",
    "NaN 23.9% + held 41.1%로 격자의 65.0%에 독립 V_rot 정보가 없다.",
    "V_rot는 5자리·최소 간격 0.00004이므로 규칙의 오탐 채널이 없다.",
    "Tᵢ의 held는 226,991행 중 1행으로 경험적 상한이다.",
    "확정 프로토콜은 held를 지도 타깃·이력·정규화 통계·모든 기준선의",
    "보간 앵커에서 제거한다.",
], accent=RED, body_size=11.5, tag="s16-held")
fcard(s, 6.15, 1.42, 6.65, 2.35, "감사 2 — Tᵢ 피팅 실패의 구조 (B.7 후속 감사)", [
    "관측 Tᵢ p99 = 2,089 eV, p99.9 = 9,601, 최대 14,984이며 꼬리는 피팅 실패이다.",
    "> 3 keV 1,197행(0.53%)은 951 run / 274 방전이고 85%가 단일행(813)이다.",
    "2–4행 run 121개, ≥ 5행은 17개(2%, 최대 15행), 70%가 고립 run이다.",
    "run 피크 중앙값은 이웃 평균의 13배(IQR 6–26배)이다.",
    "따라서 컷이 제거하는 것은 거의 전부 단일-표본 사건이며 피팅 실패의 모습이지",
    "ELM·천이의 모습이 아니다. 그래서 컷을 사전등록할 수 있었다.",
], accent=ORANGE, body_size=11.5, tag="s16-fit")
fcard(s, 6.15, 3.9, 6.65, 1.5, "그러나 값 컷은 일방향 프록시이다", [
    "관측된 두 이웃 기준 상향 ≥ 2배 이상치 3,845행, 하향 ≤ ½배 dip 4,965행이다.",
    "3 keV 컷은 상향의 19%(731행)만 제거하고 dip은 하나도 건드리지 못한다.",
    "CES 피팅 품질 메타(fit χ²·신호 세기)가 오면 품질 컷이 값 컷을 대체·동반한다.",
], accent=RED, body_size=11.5, tag="s16-proxy")
fcard(s, 6.15, 5.5, 6.65, 1.15, "V_rot는 컷하지 않는다 — 값·점프 규칙 없음, 재학습 없음 (2026-08-16 결정)", [
    "> 1,000 km/s 119행 / 16 방전이며 그중 101행이 s31181의 한 블록이다(p99 = 223 km/s).",
    "대신 persistence 앵커 비교마다 스파이크 행의 SSE 비중을 병기한다.",
], accent=GRAY, body_size=11, tag="s16-vt")
note(s, """
[감사는 신뢰를 얻는 지점이다]
두 감사 모두 자체적으로 발견하였고, 전수 정량화하였고, 프로토콜로 강제하였다.

[유지값의 결과]
평가에서 제외하면 V_rot RMSE가 올라간다. 학습에서도 제거해야 한다는 것은 4-seed paired가 단일 seed A/B의
"무해" 판정을 뒤집어 확인하였다(§8c). forward-fill이 "이력 복사가 최적"이라고 가르치고 있었다.

[스파이크 구조 감사(§8ab 말미)의 정확한 숫자]
1,197행 = 951 run / 274 방전, 단일행 813(85%), 2–4행 121 run, ≥5행 17 run(최대 15), 고립 run 70%,
run 피크 중앙값 = 이웃 평균의 13배(IQR 6–26배)이다. 일방향성: 225,580행 기준 상향 ≥2배 3,845 ·
하향 ≤½배 4,965, 컷 제거분 731 = 19%이다.

[V_rot 스파이크와 b3k8 논쟁의 연결]
TEST 분할마다 0–4행의 스파이크 앵커가 b3의 V_rot SSE의 28–72%를 나른다. 그래서 anchored 비교에는
항상 SSE 비중을 병기하라는 결정이 붙었다.
출처: THESIS_RESULTS.md §8ab, main_ko.tex §3.4·§3.5.
""")

# --- 17. Deployment --------------------------------------------------------
s = slide()
header(s, "배치", "상태 유지 추론과 전 셀 우위의 conformal 구간, 그리고 지연 절댓값의 한계")
fcard(s, 0.55, 1.42, 6.03, 2.5, "지연 — 상태 유지 1-step (은닉 상태 이월, 배치 1)", [
    "seq_v2 스텝은 CPU 중앙값 1.05 ms / p95 1.35 / p99 1.61 ms로 예산의 16%이다.",
    "GPU 1.21 / 2.31 ms이며 이 크기에서는 배치 1에 GPU가 사주는 것이 없다.",
    "세그먼트 재실행은 100행 2.9 / 5.6 ms, 300행 6.4 / 8.9 ms(35–47k 행/s)이다.",
    "같은 세션에서 인과 GP p99 2.34 ms, 윈도 W = 2 p99 4.46 ms(44.6%)였고, 순서",
    "seq_v2 < 인과 GP < 윈도 W = 2 < W = 4는 세션과 무관하게 유지되었다(§8ac).",
    "권고: 제어 계산기의 CPU에서 상태 유지형으로 실행하면 예산의 80%가 남는다.",
], accent=TEAL, body_size=11.5, tag="s17-lat")
fcard(s, 6.75, 1.42, 6.03, 2.5, "불확실성 — split conformal (α = 0.10, val 교정, TEST)", [
    "모델 구간이 32/32 셀에서 두 기준선을 Winkler 점수로 이겼다.",
    "Tᵢ 컷은 1,272 vs 1,554(PCHIP) vs 1,727(persistence)이다.",
    "Tᵢ 포함은 2,290 vs 2,851 vs 3,120, V_rot는 150 vs 164 vs 179이다.",
    "포함에서는 모델 Tᵢ 구간이 PCHIP보다 넓으면서도(반폭 224–255 vs",
    "211–241 eV) 점수가 더 좋다. 스파이크가 miss 페널티를 키우기 때문이다.",
    "Mondrian 층화는 Tᵢ arm을 4–5% 조이고 V_rot는 불변이며 판정은 불변이다.",
], accent=BLUE, body_size=11.5, tag="s17-conf")
fcard(s, 0.55, 4.05, 6.03, 1.75, "인정하는 한계 ① — 커버리지는 주변적이다", [
    "coverage는 Tᵢ 0.87–0.92, V_rot 0.91–0.94로 명목 0.90 주변이지만",
    "방전 조건부 보장은 아니다.",
    "방전별 보정은 현재 방전 수(96 / 60–66)로는 통계가 서지 않으며",
    "다중 캠페인이 확보되어야 열리는 수리이다.",
], accent=GRAY, body_size=11.5, tag="s17-cov")
fcard(s, 6.75, 4.05, 6.03, 1.75, "인정하는 한계 ② — 절댓값은 기계 종속이다 (§8aj)", [
    "5세션 프로토콜의 p99 산포가 21.84배까지 관측되어 1 ms 판정은 보류되었다.",
    "따라서 밀리초가 아니라 순서와 디스패치 연산자 수를 주장한다.",
    "비용 모델은 t ≈ N_ops × 2–3 µs이며 융합 백본은 111 ops, 최소 중앙값 0.352 ms이다.",
    "지연은 네트워크 순전파만 측정하였고 특징 조립·수집 지연은 모델 밖이다.",
], accent=ORANGE, body_size=11.5, tag="s17-rt")
fband(s, 0.55, 5.95, 12.25, 0.7,
      [("배치 형태 3종: ① 온라인 가상 CES 센서(결측 Tᵢ를 인과적으로 채워 프로파일 피팅·제어 입력으로) ② 아카이브 재분석의 간극 채움(과도 구간에서 이점 최대) ③ 다른 희소 진단의 ML 채움 연구가 그대로 쓸 수 있는 평가 bar이다. 10 ms 예산은 어느 arm에도 구속 조건이 아니다(§8ah).",
        12.5, NAVY, True)], tag="s17-uses")
note(s, """
[지연 측정의 조건 — 함께 진술한다]
유휴 노트북급 머신, 워밍업 후 1,000회, 배치 1, 은닉 상태 이월이다. 전원 상태에 따라 절댓값이 최대 2배
흔들리며(§8ac에서는 세션 간 4.2배, §8aj에서는 21.84배), arm 사이의 순서만 불변이다. 그래서 "1.05 ms"를
절대 성능이 아니라 "예산의 16%"로 말하고, B.9 이후로는 연산자 수를 주장한다.

[윈도 대조군 p99 18.9 ms는 인용하지 않는다]
§8ac가 같은 세션에서 재측정하여 4.455 ms(예산의 44.6%)를 얻었고, 동결값은 W = 2가 W = 4보다 2.3배
느리다는 물리적으로 불가능한 순서를 담고 있어 오염으로 판정되었다. 살아남는 주장은 순서이다: 백본의
꼬리가 윈도 대조군의 1/3이다.

[conformal의 공정성]
동일 절차(split conformal, 같은 val 교정 집합)를 PCHIP·persistence에도 적용한 뒤 Winkler 점수로
비교한다. 넓은 구간이 자동으로 이기지 않도록 하는 지표이다. 포함 모집단에서 모델 구간이 더 넓은데도
점수가 좋은 이유를 그 자리에서 설명한다.
출처: main_ko.tex §8, THESIS_RESULTS.md §8ab·§8ac·§8ah·§8aj.
""")

# --- 18. B.9 I: reach ladder, pooled re-scoring -----------------------------
s = slide()
header(s, "B.9 ①", "도달 범위 사다리: 문맥 약 50 ms에서 포화하며, 문맥이 사는 것은 전형성이다 (§8ac–§8am)", accent=TEAL)
fband(s, 0.55, 1.40, 12.25, 0.78,
      [("설계: §8ac는 동결 백본의 순환 상태를 ctx 스텝 전에 리셋하여 재채점하였고(재학습 없음), 그 결손의 87%가 cold start였음을 §8ae·§8af가 밝혔다(20 ms: -0.510 → -0.065). "
        "따라서 seq_v2를 2·3·4·5·6·7·10·15·31·63 스텝의 도달 범위에서 각각 학습·채점하였다(GATE_ENV, 동결 B.1 매니페스트, 4 분할, 행 수 고정 배치).",
        11.5, DARK, False)], tag="s18-design")
add_image_fit(s, os.path.join(PAPERFIG, "fig_context_family_ladder.png"),
              Inches(0.45), Inches(2.2), Inches(6.15), Inches(4.5))
ftable(s, 6.75, 2.2, [0.95, 2.25, 0.85, 0.9, 0.9],
       ["문맥", "skill vs 인과 GP [95% CI]", "승률", "-top10", "결손"],
       [
           ["20 ms", "+0.057 [+0.027, +0.085]", ("0.52", ORANGE, True), "+0.028", ("-0.066", RED, False)],
           ["30 ms", "+0.087 [+0.061, +0.111]", "0.60", "+0.060", "-0.033"],
           ["50 ms", "+0.104 [+0.079, +0.128]", "0.64", "+0.077", ("-0.017", GREEN, True)],
           ["70 ms", "+0.119 [+0.095, +0.142]", ("0.66", GREEN, True), "+0.092", "+0.002"],
           ["150 ms", "+0.123 [+0.096, +0.148]", "0.66", "+0.096", "+0.005"],
           ["630 ms", "+0.143 [+0.118, +0.168]", "0.67", "+0.116", "+0.029"],
       ],
       row_h=0.34, head_h=0.36, size=10.5, head_size=10, tag="s18-pooled")
fcard(s, 6.75, 4.7, 5.85, 1.97, "판정", [
    "사전등록 §3.4 규칙(전체 대비 결손 < 0.02, 유의 ≤ 1/4)은 4·5·6·10스텝을 채운 뒤",
    "50 ms를 반환한다. 70 ms는 그 칸들이 학습되지 않아 생긴 값이었다.",
    "모델은 20 ms에서도 인과 GP를 이기며(통합 CI [+0.027, +0.085]), 문맥이 사는 것은",
    "평균이 아니라 전형성이다(승률 0.52 → 0.66). 문맥 10배당 skill은 +0.050 [+0.036, +0.064]로",
    "상승하며, 요동하는 4/4 계수는 문턱 추정기에서 퇴출되었다.",
], accent=NAVY, body_size=10.5, tag="s18-verdict")
note(s, """
[세 번의 교정을 순서대로 진술한다]
§8ac(절단 사다리)는 "Tᵢ는 500 ms 문맥이 필요하다"고 읽었다. §8ae는 그 결손의 84% 이상이 warm-up임을
동결 산출물로 경계지었고, §8af는 각 도달 범위에서 학습한 seq_v2로 70 ms를 얻었다. §8al은 4·5·6·10스텝을
채워 규칙이 50 ms를 반환함을 보였다(70은 4·5·6이 학습되지 않아 생긴 값). §8am은 301 방전 통합
재채점으로 "N ms가 있어야 인과 GP를 이긴다"를 철회하였다. 모델은 20 ms에서도 이긴다.

[통합 재채점의 정당성]
각 분할은 301 방전 union 중 96개를 test하며 224개는 한 분할에만 나타난다. 두 분할에 나타나는 방전은
하나의 군집으로 묶여 보수적으로 처리된다. 추세는 각 bootstrap 재표본 안에서 재적합되어 군집 구조를
그대로 실어 나른다.

[표 열의 의미]
"승률"은 모델이 인과 GP를 이기는 방전의 비율이고, "-top10"은 기여 상위 10개 방전을 제거한 뒤의 통합
skill이며, "결손"은 전체 블록 백본 대비 paired 차이(4 분할 평균)이다. 좁은 CI가 곧 전형적 효과는 아니라는
사실을 보이기 위해 세 열을 나란히 둔다(§8am §3).

[V_rot]
V_rot는 20 ms에서만 -0.013(2/4)의 작은 결손을 보이고 70 ms부터 0이며, 네 계열 전부에서 문맥이 길수록
오히려 나빠진다(통합 추세 음수). 이월값을 타는 분기에 긴 문맥은 분산만 더한다는 라우팅 결과와 일관된다.
출처: THESIS_RESULTS.md §8ac·§8ae·§8af·§8al·§8am, PREREGISTRATION_B9.md §3.4.
""")

# --- 19. B.9 II: family, cost, which discharges ---------------------------
s = slide()
header(s, "B.9 ②", "계열은 skill을 정하지 않고 비용을 정하며, 모델은 타깃이 움직이는 방전에서 이긴다 (§8ag–§8an)", accent=TEAL)
ftable(s, 0.55, 1.42, [1.15, 0.85, 1.05, 1.15, 0.8, 0.9],
       ["arm", "문맥", "파라미터", "Tᵢ paired", "승/패", "판정"],
       [
           ["tcn3", "30 ms", "71,442", "-0.004", "0/0", ("동률", GREEN, True)],
           ["tcn7", "70 ms", "128,034", "-0.004", "0/1", ("동률", GREEN, True)],
           ["tcn63", "630 ms", "297,810", "-0.016", "0/2", "미결"],
           [("xfmr7", RED, True), "70 ms", "295,746", ("-0.023", RED, True), ("0/3", RED, True), ("차이", RED, True)],
           ["xfmr15", "150 ms", "295,746", "+0.002", "0/0", ("동률", GREEN, True)],
           ["xfmr63", "630 ms", "295,746", "-0.019", "0/1", ("동률", GREEN, True)],
       ],
       row_h=0.33, head_h=0.36, size=10.5, head_size=10, tag="s19-family")
fband(s, 0.55, 3.95, 5.9, 0.5,
      [("같은 도달 범위의 LSTM 칸 대비 paired Tᵢ skill(4 분할 평균, §3.2 규칙). 계열 효과의 최댓값 0.023은 문맥 효과 +0.060(20→70 ms)의 1/2.6이다.",
        10.5, GRAY, False)], tag="s19-cap")
fcard(s, 0.55, 4.5, 5.9, 2.15, "10k 파라미터 아래에서는 동률이 깨진다 (§8ai)", [
    "크기를 맞춘 순환 arm 대비 합성곱 arm은 +0.027~+0.040(3–4/4 유의)이다.",
    "tcn2k(1,808 파라미터)는 인과 GP를 4/4로 이기며 백본과 +0.001 차이이다.",
    "계열은 상한이 아니라 상한에 이르는 데 필요한 파라미터 수를 정한다.",
    "diagonal SSM은 문맥을 70 ms까지 변환한 뒤 낮은 천장(+0.105)에서 멈춘다(§8am 부록).",
], accent=NAVY, body_size=11, tag="s19-small")
fcard(s, 6.65, 1.42, 6.15, 2.5, "비용 = 디스패치 연산자 수 (§8ah·§8aj)", [
    "지연은 t ≈ N_ops × 2–3 µs이며 파라미터 151배는 지연 1.6배만 산다.",
    "순환은 도달 범위에 O(1)이다(융합 111 ops가 모든 칸에서 같다).",
    "확장 합성곱은 O(log R)이다(층당 +48 ops, RF = 2^(L+1) − 1).",
    "attention은 O(1)이나 상수가 4.3배이다(473 ops). 이 문제에서는 엄격히 열세이다.",
    "백본의 구현 사다리는 0.793 → 0.688 → 0.378 → 0.352 ms(2.25배)이고 융합은 거의 소진되었다.",
    "10 ms 예산은 전 arm에 2배 이상 여유이며, 1 ms는 세션 산포 21.84배로 판정이 보류되었다.",
], accent=BLUE, body_size=11, tag="s19-cost")
fcard(s, 6.65, 4.05, 6.15, 2.6, "승패 방전 분석 (§8al §4 · §8an)", [
    "방전 단위 승률은 Tᵢ 0.695(96 방전), V_rot 0.481(약 62 방전)이다.",
    "V_rot는 상위 5개 방전을 빼면 4 분할 전부 0 이하이다(split 42: +0.331 → -0.006).",
    "11개 공변량 중 방전 내 타깃 산포만 승패를 예측한다(순열·Bonferroni: Tᵢ ρ +0.401, V_rot +0.281).",
    "3분위 승률은 Tᵢ 42/83/85%, V_rot 34/48/55%이다. 조용한 방전은 인과 GP가 이미 최적이다.",
    "변동 큰 3분위에서도 V_rot가 55%에 그치는 잔차는 구동 변수(NBI 토크)의 부재를 가리키며,",
    "Tₑ~V_rot 방전 단위 null(ρ = -0.031)이 세 번째 독립 증거이다.",
], accent=ORANGE, body_size=11, tag="s19-wins")
note(s, """
[한 문장]
문맥은 skill을 정하고 계열은 비용을 정한다(§9 lead claim). 두 한정이 있다: 10k 파라미터 아래에서는
합성곱이 낫고(§8ai), 70 ms 칸에서는 attention이 -0.023 뒤진다(§8ak). 둘 다 문맥 효과 +0.060보다 작다.

[계열 비교의 설계]
세 계열 모두 seq_v2의 라우팅(V_rot 분기는 빠른 진단을 읽지 않음)을 유지하고, 각자의 수용 범위가 선언하는
도달 범위에서 학습·채점하였다. attention의 밴드는 두 층이 63폭으로 겹치면 125가 되므로 목표 도달 범위에서
역산하였고, 인과성과 수용 범위는 테스트가 수치로 단언한다. 각 arm은 백본이 아니라 같은 도달 범위의 축 A
LSTM 칸과 짝지었다. 백본과만 비교하면 계열과 도달 범위가 §8ae처럼 뒤섞인다.

[비용의 측정 단위]
op_count.py는 온라인 1스텝을 torch.profiler 아래 실행하여 디스패치된 aten:: 연산자를 센다. 이 수는 기계와
무관하게 재현되며 기계가 바쁠 때도 움직이지 않았다. 변환율 2.1–3.2 µs/op는 파라미터 151배 범위와 세 계열
전부에서 거의 상수였다. stock nn.LSTM만 6.72 µs/op로 예외이며 nn.Module 호출 프로토콜이 aten 연산이
아니어서 수에 잡히지 않기 때문이다. tcn2k의 tight 스텝은 106 ops로 융합 백본(111) 아래이며, torch.jit
경로는 202 ops로 기각되었다.

[승패 분석의 지위]
탐색적이다. 11개 공변량, 한 arm, 사전등록 규칙 없음이며 두 생존 공변량(타깃 산포와 인과 GP RMSE)은
ρ = 0.85로 공선적이다. 비율의 분모 함정은 이진 승/패 검정으로 제거하였다(Tᵢ ρ +0.393, V_rot +0.198).
이를 뒤집을 측정은 B.6이다: μs 재획득이 모드 회전 주파수를 제공하면 V_rot 승률이 변동 3분위에서 먼저
올라야 한다.
출처: THESIS_RESULTS.md §8ag·§8ah·§8ai·§8aj·§8ak·§8al §4·§8am·§8an, PREREGISTRATION_B9.md.
""")

# --- 20. Headroom + limitations -------------------------------------------
s = slide()
header(s, "개선 여지 · 한계", "남은 레버는 전부 데이터이며, 인정하는 한계의 목록이다")
fcard(s, 0.55, 1.45, 3.95, 2.5, "① CES 피팅 품질 메타데이터", [
    "값 컷은 일방향 프록시이다(상향 이상치의",
    "19%만 제거, dip 4,965행은 그대로).",
    "fit χ²·신호 세기가 오면 품질 컷이 값 컷을",
    "대체하고 두 모집단이 하나로 합쳐진다.",
    "이 연구의 판정 언어를 가장 단순하게",
    "만드는 레버이다.",
    "V_rot 스파이크도 같은 규칙으로 처리된다.",
], accent=BLUE, body_size=11.5, tag="s18-h1")
fcard(s, 4.70, 1.45, 3.95, 2.5, "② 원본 kHz Mirnov (V_rot 최상위)", [
    "블록 내 lag-1 자기상관은 BES +0.568 ·",
    "ECEI +0.572 vs Mirnov -0.009(블록 82%가",
    "|r| < 0.1)이며, anti-aliasing 없이 kHz dB/dt를",
    "100 Hz로 데시메이트한 서명이다.",
    "원본 스트림에서 윈도 RMS·대역 파워·모드수·",
    "모드 회전 주파수를 뽑아 V_rot 분기에 투입한다.",
    "B.6 shot 집합은 동결되었고 예측이 기록되었다.",
], accent=TEAL, body_size=11.5, tag="s18-h2")
fcard(s, 8.85, 1.45, 3.95, 2.5, "③ NBI 토크 채널 부재", [
    "Tₑ~Tᵢ r = +0.353(p = 3e-17) vs Tₑ~V_rot",
    "r = +0.024(p = 0.58)이며 파워는 토크가 아니다.",
    "회전의 원인 변수가 데이터셋에 없다.",
    "§8an의 승패 분석이 세 번째 독립 증거이다.",
    "아키텍처가 아니라 데이터 획득 과제이며,",
    "DIII-D의 액추에이터 조건부 예보가 양성",
    "대조군이다.",
], accent=ORANGE, body_size=11.5, tag="s18-h3")
fbullets(s, 0.55, 4.2, 12.25, 2.45, [
    ("검정력: 유효 표본은 방전 수(Tᵢ 약 96, V_rot 60–66)이며, 포함 모집단에서는 약 1% 행이 모든 arm SSE의 70–83%를 나른다.", 0),
    ("MNAR 낙관: 재가중은 인과 비교엔 두 모집단, 오프라인 비교엔 모집단 조건부이며 도달 범위는 Tᵢ 54–68% · V_rot 4–6%이다.", 0),
    ("오프라인 주장의 상한은 오프라인 GP와 동률(1/8 유의)이고, 값 컷은 일방향 프록시이며, V_rot는 컷하지 않는다.", 0),
    ("캠페인은 한 시간 블록 위의 초기화 4개(분할 4개 아님)이고 컷 run 2/4가 30-epoch 상한 종료였으며, shot별 표준화는 오프라인 형태이다(인과 EWMA 미측정).", 0),
    ("B.9: 통합 재채점은 방법의 기대 skill이지 단일 체크포인트의 배치 주장이 아니고, 승패 공변량 분석은 탐색적이며, 1 ms 판정은 조용한 기계의 5세션 재실행을 요구한다.", 0),
    ("conformal은 marginal 커버리지, 지연은 네트워크만, 단일 장치, 페데스탈 상단 프레이밍이며 이벤트-위상 분석은 후속이다. 규칙: 음성 결과는 그것을 뒤집을 측정을 지목할 때만 보고한다(§8j).", 0, ORANGE),
], size=12.5, gap=5, tag="s18-lim")
note(s, """
[세 레버의 공통점]
전부 데이터 레버이다. B.4가 크기 축을 닫았고, B.3이 21k로도 충분함을 보였으며, B.9가 문맥(50 ms 포화)과
계열(동률) 축을 닫았으므로, 남은 개선은 모델이 아니라 입력에서 온다. 이 문장이 결론의 마지막 축이다.

[Mirnov를 진술하는 법]
"Mirnov가 잡음이라 뺐다"가 아니라 "전처리가 정보를 파괴하였음을 측정으로 밝혔다"이다. 파생 특징 4종
(적분·PCHIP 적분·|MC|·이동 RMS)은 4-seed paired에서 전부 무효였고, 수리는 상류(원 kHz 스트림)에 있다.
B.6의 shot 집합(test 31921·31873·31114·31902, pool 31097·31359·31747·32027·32092·32097, companion
31923·31357)은 2026-08-21 동결되었고 #32092(EHO n = 1, ~4/~8 kHz)가 다섯 번째 양성 대조이다.

[한계를 먼저 진술하는 이유]
심사에서 한계를 상대가 먼저 지적하면 방어가 되고, 스스로 먼저 말하면 프로토콜의 일부가 된다. 특히
캠페인 †(한 시간 블록), conformal marginal, 1 ms 보류는 먼저 진술한다.
출처: main_ko.tex §9·§10, THESIS_RESULTS.md §8ab·§8aj·§8am·§8an·§8ao·§9.
""")

_QA_COLS = [3.3, 7.75, 1.2]
_QA_HEAD = ["예상 질문", "답변 요지 (수치는 전부 §8ab·§8ac–§8an 기준)", "방어"]

# --- 21. Q&A I -------------------------------------------------------------
s = slide()
header(s, "예상 질문 ①", "프로토콜 · 모집단 · 사전등록", accent=RED)
ftable(s, 0.55, 1.42, _QA_COLS, _QA_HEAD,
       [
           ["왜 두 모집단인가? p100 하나로 쓰면 안 되는가",
            "절제가 측정된 답을 준다. 포함 모집단에서는 이력-전용 모델도 PCHIP를 +0.15~+0.23 이기고 빠른 채널이 더하는 것은 0.03–0.09뿐이므로, p100 마진에는 빠른 진단 정보가 아닌 스파이크-강건성 성분이 섞여 있다. 반대로 컷만 쓰면 \"어려운 행을 버렸다\"는 비판에 노출된다. 그래서 두 모집단을 공동 1차로 두고 무조건부 주장의 기준을 \"둘 다 성립\"으로 올렸다. 단일 headline보다 엄격한 규칙이다.",
            ("●●●", GREEN, True)],
           ["값 컷은 결국 임의적 프록시가 아닌가",
            "그렇게 보고한다. 두 이웃이 관측된 225,580행 중 상향 ≥ 2배 이상치 3,845 · 하향 dip 4,965인데 3 keV 컷이 제거하는 것은 상향의 19%(731행)뿐이고 dip은 하나도 건드리지 못한다. 근거는 물리(> 3 keV는 KSTAR 이온온도가 아니다)와 구조 감사(85%가 단일행, run 피크가 이웃 평균의 13배)이다. 문턱 2.5/3/4 keV에서 Tᵢ 평균 +0.230/+0.236/+0.232, PR4 판정 전부 동일이므로 문턱이 아니라 두 모집단이 본질이다.",
            ("●●○", ORANGE, True)],
           ["왜 W = 2인가",
            "24-run 스윕이 먼저 답하였다. history-0은 붕괴하고(Tᵢ -0.026, V_rot -0.783), 관측 하나가 곧바로 plateau를 만들며(W = 2: Tᵢ +0.238 4/4), 이후 곡선은 평평하다(Tᵢ 0.190–0.246; 한 점 안의 seed 산포 0.07–0.16이 곡선 전체 폭보다 크다). plateau-최소 규칙이 W = 2를 반환한다. B.9는 이 '관측 수'와 별개의 자원인 '연속 빠른 진단 문맥'이 약 50 ms 필요함을 보였다(§9.2). 윈도 계열은 구성상 두 번째 자원에 닿지 못한다.",
            ("●●●", GREEN, True)],
           ["MNAR 재가중은 어디까지 말해 주는가",
            "도달 범위를 사전등록이 요구한 대로 보고한다. 결측 Tᵢ의 54–68%, V_rot의 4–6%만 W = 2 도메인 안에 있다(과거 관측이 2행 이내). 그래서 재가중 V_rot는 결측 질량의 1/20에 대한 답이라 결론을 내지 않는다. Tᵢ는 persistence 대비 두 모집단 4/4(+0.28~+0.44, 보정 비용 최대 0.12), PCHIP 대비는 컷 2/4·포함 4/4이다. 인과 대비는 무조건부, 오프라인 대비는 모집단 조건부이다.",
            ("●●○", ORANGE, True)],
       ],
       row_h=1.24, head_h=0.42, size=12, head_size=12, left_cols={1}, tag="s19-qa")
note(s, """
[Q1 답변의 구조]
p100만 쓰면 결과가 더 좋아 보인다(+0.268 vs +0.236). 그러나 절제 결과 포함 모집단에서는 빠른 진단을 전부
끈 이력-전용 모델도 PCHIP을 +0.15~+0.23 이긴다. 그 마진의 상당 부분은 빠른 진단 정보가 아니라 보간 앵커가
오염되어 있다는 사실이다. 빠른 진단의 기여를 분리하는 모집단은 컷이고, 배치된 시스템이 실제로 만나는
모집단은 포함이다. 그래서 둘 다 1차로 두고 무조건부는 둘 다 성립할 때만 쓴다.

[Q2에서 피할 것]
컷을 정당화하려고 방어적으로 길게 말하는 것이다. 일방향성을 먼저 인정하고, 품질 메타데이터가 오면
대체된다는 후속을 지목하는 순서가 강하다.

[Q3 보강]
W = 4는 이전 초고의 기본값이었고 그것이 정당화되지 않는다는 것을 먼저 보고하였다(§8f). §9.2의 "두 자원"
프레이밍: 과거 CES 관측은 하나면 충분하고(§8f), 연속 빠른 진단 문맥은 약 50 ms 필요하다(§8af·§8al).
윈도 계열은 AdaptiveAvgPool이 순서를 버리고 시간 부분집합 증강이 비연속 부분집합을 만들며 라벨 없는 행을
버리므로 어떤 W에서도 두 번째 자원에 닿지 못한다. 윈도 iter009 W=2의 +0.041(1/4)과 seq_v2 2스텝 절단의
+0.055(2/4)가 같은 굶주림의 같은 결과이다.

[Q4 보강]
"in-domain"은 정확도 한계가 아니라 커버리지 한계이다. 이 구분을 흐리지 않는다.
출처: THESIS_RESULTS.md §8ab-8·§8ab-9·§9.2, main_ko.tex §6.1·§8·§11.
""")

# --- 22. Q&A II ------------------------------------------------------------
s = slide()
header(s, "예상 질문 ②", "모델 선택 · 복잡도 · 기준선", accent=RED)
ftable(s, 0.55, 1.42, _QA_COLS, _QA_HEAD,
       [
           ["왜 seq_v2가 주 모델인가",
            "관문 4조건을 실행 전에 커밋하고 전부 충족하였다(PREREGISTRATION_W2.md §4): ① 분할별 초기화 평균 paired Tᵢ 부호 4/4 양수 ② 16-run pooled +0.081, run-클러스터 CI [+0.067, +0.096]가 0 배제 ③ 예산 균등화(고정 10 epoch)에서도 4/4 양수(+0.063/+0.033/+0.045/+0.030) ④ V_rot 유의 열세 0/16. 16 run 전부 양수·13/16 유의였고 학습 비용은 1/10이다. B.9는 이 +0.081의 4/5가 도달 범위이고 1/5가 구조임을 분해하였다.",
            ("●●●", GREEN, True)],
           ["val에서 좋았던 attention 후보는 왜 올리지 않았는가",
            "사전 고정한 승격 기준(유의 ≥ 3/4)에 미달하였기 때문이다. seq v3는 val(분할 42/7)에서 paired +0.024*/+0.037*로 2/2 유의였지만 TEST 확증에서 +0.009/+0.013/+0.033*/+0.020 = 4/4 양수·1/4 유의였다. 규칙대로 미승격하고 §8y에 음성 기록으로 남겼다. \"모델을 고르고 이유를 붙였다\"는 의심에 대한 실물 반증이다. v3가 인과 GP를 4/4로 넘은 덕에 백본 자신의 claim 2가 복권되었다.",
            ("●●●", GREEN, True)],
           ["b3k8 = 백본이라면서 왜 컷 조건부인가",
            "포함 모집단에서 persistence 오차 > 2 keV인 행이 TEST의 0.6–1.3%인데 그 행들이 b3 Tᵢ SSE의 73–83%(모든 arm 70–83%)를 나른다. 유계 보정(latent ∈ [-1,1] 8개)은 스파이크 이월값을 되살릴 수 없으므로 포함에서 -0.194(4/4 유의)로 벌어진다. 컷에서는 평균 +0.002이고 CI가 전부 0을 포함한다. 단(rung) 자체는 양쪽에서 성립하고 조건부인 것은 \"백본과 동급\"이라는 허용치뿐이다.",
            ("●●●", GREEN, True)],
           ["크기를 키우면 되지 않는가",
            "B.4가 그 축을 닫았다. Tᵢ 인코더 폭 24 → 260(34k → 879k, 26배)에서 평균 skill이 +0.230/+0.236/+0.235/+0.236/+0.230으로 160 대비 ±0.008이고, 최대 폭의 유의 우세는 1/4이며 V_rot는 불변(+0.250~+0.254)이다. 반대편에서는 21,498 파라미터 b3k8이 컷 모집단에서 백본과 동급이고, B.9의 1,808 파라미터 tcn2k도 인과 GP를 4/4로 이긴다. 상한은 추정기 용량이 아니라 정보이다. 다만 이 문장은 총합 MSE에 한정된다(§8aq): 타깃 자신의 10 ms 재현성은 본류에서 46 eV, 꼬리를 포함하면 130 eV이고 4차 차분 제곱질량의 46.6%가 상위 1%에 있으므로, 폭·계열이 평평한 것은 지표가 꼬리 통계이기 때문이며 본류에는 아직 2.3~3.4배의 여지가 있다.",
            ("●●●", GREEN, True)],
       ],
       row_h=1.24, head_h=0.42, size=12, head_size=12, left_cols={1}, tag="s20-qa")
note(s, """
[Q5 보강 — 관문의 대조군]
감사 (B)-4에 따라 대조군을 W = 4 계열에서 W = 2 계열로 바꾸고, 컷 도입이 모집단을 바꾸므로 양쪽을
재학습하였다. "유리한 대조군을 골랐다"는 반박을 사전에 차단한 부분이다.

[Q6 보강 — 답변의 태도]
후보를 떨어뜨린 것을 강조하지 않고, 규칙이 있었기 때문에 가능하였다는 순서로 진술한다. val에서 2/2
유의였으므로 TEST에 bar를 둔 것이다. §8ak가 attention 계열이 70 ms에서 LSTM보다 -0.023 뒤진다는 사실을
더하였으므로, attention 후보의 미승격은 사후적으로도 손실이 아니다.

[Q7 보강]
b3 − seq_v2 컷 -0.009/-0.005/+0.026/-0.004(평균 +0.002), 포함 -0.16~-0.21이다. b3 − anchor는 두 모집단
4/4 유의(컷 +0.35~+0.42, 포함 +0.29~+0.34)이므로 "단이 무너진 것"이 아니다.

[Q8 보강 — probe와 B.9]
b3의 latent은 직전 관측 Tᵢ(R² 0.47–0.75)와 ECEI Tₑ 프록시(0.31–0.48)를 분산 부호화하고 보정 항이
예측 분산의 25–39%를 설명한다. §8ai: 10k 파라미터 아래에서는 합성곱이 크기를 맞춘 순환 arm보다
+0.027~+0.040 낫다. 즉 최소 모델의 기본 형태는 합성곱이며, 이는 배치 논거가 아니라 측정이다(357k 백본도
비용이 아니다, §8ah).
출처: THESIS_RESULTS.md §8x·§8y·§8z·§8aa·§8ai·§8ak.
""")

# --- 23. Q&A III -----------------------------------------------------------
s = slide()
header(s, "예상 질문 ③", "물리 · 일반화 · 증거의 소재", accent=RED)
ftable(s, 0.55, 1.42, _QA_COLS, _QA_HEAD,
       [
           ["V_rot는 왜 동률인가 — 실패가 아닌가",
            "네 가지가 겹친다. ① 라우팅이 인코더에서 성립하여 빠른 채널을 어떻게 섭동해도 V_rot 출력이 bit-identical이다(no_fast paired +0.000 ×4, 두 모집단). ② 회전의 원인 변수인 NBI 토크가 데이터셋에 없다(Tₑ~V_rot r = +0.024, p = 0.58). ③ 유일한 회전 대리인 Mirnov가 100 Hz 앨리어싱으로 파괴되었다(lag-1 -0.009). ④ PR2 폴백률이 40–44%라 \"vs PCHIP\"의 2/5가 사실상 vs persistence이다. 그래도 Δt > 15 ms에서는 두 모집단 모두 PCHIP을 이기고 peak 층은 +0.54~+0.79이다.",
            ("●●●", GREEN, True)],
           ["캠페인 결과가 왜 뒤집혔는가",
            "프로토콜이 아니라 모델이 바뀌었다. 윈도 대조군은 여전히 붕괴한다(컷 2/4·포함 0/4, 인과 GP 0/4). 이전 초고의 0/4가 그대로 재현된다. 백본은 PCHIP·인과 GP 모두 4/4+4/4이고 대조군 대비 8/8 유의이다. 원인은 측정되었다: train→test 드리프트가 BES 1.22σ·ECEI 0.53σ인 데 반해 타깃은 0.115σ이다. 지목하였던 수리(shot별 표준화)는 윈도에서 2/4→4/4로 실제 작동하였고, seq_v2는 정의상 그 표준화 + 세그먼트 도달 범위를 갖는다.",
            ("●●●", GREEN, True)],
           ["오프라인 GP와 동률인데 무슨 의미가 있는가",
            "오프라인 주장의 상한을 먼저 측정하여 보고한 것이다. 오프라인 GP는 미래 이웃 16개를 보고 잡음을 평균해 내는 평활기이므로 배치할 수 없으며 실시간에는 존재하지 않는 팔이다. 배치에 관한 문장은 인과 사다리 위에서 판정하며 거기서는 백본이 인과 GP를 4/4+4/4(+0.08~+0.17)로 이긴다. headline 비교 대상은 사전등록(PR1)대로 PCHIP를 유지하고 GP는 사다리에 병기하여 결과를 그대로 싣는다.",
            ("●●●", GREEN, True)],
           ["TEST 동결·사전등록의 증거는 어디에 있는가",
            "ces_prediction/experiments/PREREGISTRATION_W2.md(B.1–B.5)와 PREREGISTRATION_B9.md(B.9 축 A–D, H1–H6), PREREGISTRATION_B6.md(μs 재획득)이다. 배치별 판정 규칙이 실행 전에 커밋되어 있고 그 아래 \"집행 결과\" 블록이 §8x·§8y·§8z·§8aa·§8ab·§8af–§8aj로 이어 붙는다. 스파이크 컷 구현은 커밋 3598760(로드 시점 NaN 처리, 캐시 서명 v5, 35/35 테스트 통과)이다. B.2 탐색은 val 전용이었고 TEST는 후보당 1회만 채점하였으며, 모든 수치는 동결 run 디렉터리에서 단일 collector가 읽는다.",
            ("●●●", GREEN, True)],
       ],
       row_h=1.24, head_h=0.42, size=12, head_size=12, left_cols={1}, tag="s21-qa")
note(s, """
[Q9 — V_rot 답변의 태도]
"실패한 타깃이 아니라 의미 있는 비교가 제한된 타깃"으로 진술한다. 물리로 사전에 예측된 비대칭이고, 절제가
기전을 확인하였으며, 데이터 레버(NBI 토크·kHz Mirnov)를 지목한다. §8an의 승패 분석이 추가 답이다:
모델은 회전이 실제로 변하는 방전에서만 이기며 거기서도 55%에 그치므로, 회전을 구동하는 입력이 없다는
결론과 세 방향(행 단위 Tₑ 상관·절제·방전 단위 승패)에서 일치한다. 반증 가능한 형태로 진술한다: 토크가
들어오면 학습 가능해야 하며 DIII-D의 액추에이터 조건부 예보가 양성 대조군이다.

[Q10 — 하지 말아야 할 진술]
"프로토콜을 바꿔서 좋아졌다"는 사실이 아니다. 같은 캠페인 분할에서 윈도 대조군의 붕괴는 그대로 재현되며,
달라진 것은 모델(도달 범위 + shot별 표준화)이다. 캠페인 †(한 시간 블록 × 초기화 4개, 컷 run 2/4 상한
종료)를 먼저 붙인다.

[Q11 보강]
GP 팔은 §8p에서 사후에 실측하였고 그 결과가 불리하였음에도 그대로 보고하였다. 기준선이 약하다는 비판을
측정으로 선제 차단하였고 결과를 그대로 싣는다.

[Q12 보강]
질문이 나오면 PREREGISTRATION_W2.md의 §4(관문 4조건)와 §6(v3 승격 규칙), PREREGISTRATION_B9.md의
§3.2·§3.4·§4를 실제로 열어 "이 문장이 TEST 채점 전에 커밋되었다"를 커밋 시각과 함께 보이는 것이 가장 강한
답이다. B.9의 H1 규칙이 잘못 명세되었다는 사실(효과 크기와 유의 계수를 disjunction으로 묶어 작은 강건
효과가 '약한 의존' 가설을 기각함)도 §8af §3에 그대로 기록되어 있으며 사후 수정하지 않았다.
출처: PREREGISTRATION_W2.md §3·§4·§6, PREREGISTRATION_B9.md, THESIS_RESULTS.md §8w·§8ab·§8af·§8an.
""")

# --- 24. Q&A IV (B.9) ------------------------------------------------------
s = slide()
header(s, "예상 질문 ④", "문맥 · 계열 · 비용 · V_rot의 승패 (B.9)", accent=RED)
ftable(s, 0.55, 1.42, _QA_COLS, _QA_HEAD,
       [
           ["다중 센서 논문은 합성곱을 쓰는데 왜 LSTM인가",
            "같은 문맥에서 계열은 skill을 정하지 않는다. 순환·확장 합성곱·attention을 seq_v2의 라우팅을 유지한 채 같은 도달 범위의 LSTM 칸과 paired 비교하면 모든 arm이 ±0.023 안에 있고, 사전등록 §3.2 규칙이 '차이'를 반환한 칸은 70 ms의 attention 하나뿐이다(-0.023, 3/4). 계열은 비용을 정한다: 순환은 도달 범위에 O(1)(111 ops), 합성곱은 O(log R)(층당 +48), attention은 상수 4.3배(473)이다. 다만 10k 파라미터 아래에서는 합성곱이 +0.027~+0.040 낫다(§8ai).",
            ("●●●", GREEN, True)],
           ["얼마나 긴 문맥이 필요한가 — 500 ms인가 70 ms인가",
            "약 50 ms이다. §8ac의 절단 사다리(500 ms)는 학습되지 않은 상태 재구축(cold start)을 정보 부족으로 읽은 것이며 §8ae·§8af가 그 결손의 87%를 warm-up으로 분해하였다. 각 도달 범위에서 학습한 seq_v2에 §3.4 규칙을 적용하면 4·5·6·10스텝을 채운 뒤 50 ms를 반환한다(§8al). 통합 301 방전에서 모델은 20 ms에서도 인과 GP를 이기므로(+0.057 [+0.027, +0.085]) 문맥은 승리 조건이 아니라 포화·전형성 조건이다(승률 0.52 → 0.66). 무한 문맥은 필요하지 않으며 더 낫지도 않다.",
            ("●●●", GREEN, True)],
           ["1 ms 제어 주기에도 쓸 수 있는가",
            "판정은 보류되었다. §8ah의 첫 패스는 lean 구현으로 백본이 max p99 0.591 ms를 기록하여 통과하였으나, 같은 5세션 프로토콜의 재실행(§8aj)에서 p99 세션 산포가 21.84배에 이르러 사전등록 §4 규칙이 어떤 arm에도 통과를 주지 않았다. 살아남는 진술은 기계와 무관한 것이다: 비용은 디스패치 연산자 수이며(2–3 µs/op), 융합 백본 111 ops·tcn2k 106 ops의 투영 시간은 0.2–0.35 ms이다. 10 ms는 어느 arm에도 구속 조건이 아니다. 지목된 측정은 조용한 기계의 5세션 재실행과 컴파일 런타임(ONNX Runtime)이다.",
            ("●●○", ORANGE, True)],
           ["V_rot는 방전이 더 많으면 유의해지지 않는가",
            "아니다. 방전 단위 승률이 0.48(동전 던지기)이고 중앙 방전의 skill이 약 0이며, 약 62개 중 상위 5개 방전을 제거하면 4 분할 전부 통합 우위가 0 이하이다(§8al §4). 승률 0.48 위의 shot 군집 bootstrap은 어떤 n에서도 0을 걸친다. 승패를 예측하는 유일한 공변량은 방전 내 타깃 산포이며(§8an), 조용한 방전은 인과 GP가 이미 최적이고 변동 큰 방전에서도 V_rot는 55%만 이긴다. 그 잔차는 구동 변수의 부재이며, 이를 뒤집을 측정은 B.6의 μs Mirnov(변동 3분위에서 먼저 개선되어야 함)이다.",
            ("●●●", GREEN, True)],
       ],
       row_h=1.24, head_h=0.42, size=12, head_size=12, left_cols={1}, tag="s24-qa")
note(s, """
[Q13 보강 — 계열 비교의 정직한 범위]
세 계열은 "모든 계열"이 아니다. diagonal SSM은 §8am 부록에서 여섯 칸이 채워졌고, 문맥을 70 ms까지 변환한
뒤 낮은 천장(+0.105 vs 순환 +0.143)에서 멈춘다. 같은 도달 범위의 LSTM에 7스텝에서 -0.022, 63스텝에서
-0.044로 뒤지므로 채택되지 않는다. 20 ms에서는 SSM이 통합 최고점(+0.065)이나 같은 칸 paired 검정은
수행되지 않았으므로 보고만 한다.

[Q14 보강 — 두 자원]
§9.2: 과거 CES 관측은 하나면 충분하고(§8f), 연속 빠른 진단 문맥은 약 50 ms 필요하다. 윈도 계열은 구성상
두 번째 자원에 닿지 못하며, 이것이 §8f가 평평하였던 이유를 소급하여 설명한다. "N ms가 있어야 인과 GP를
이긴다"는 문장은 §8am에서 철회되었다.

[Q15 보강 — 절댓값을 인용하지 않는 규칙]
§8ac(4.2배)·§8ah(1.32배, 2.56배)·§8aj(21.84배)의 세션 산포 기록이 그 자체로 답이다. 최솟값 통계는 arm
비교에는 옳고 통과 판정에는 틀리므로, 사다리의 순서(백본 < 인과 GP < 윈도)와 연산자 수만 주장한다.
torch.jit 경로는 1.8k 파라미터에서 202 ops로 lean(209)과 다르지 않아 기각되었고, 남은 레버는 커널 융합
런타임이다.

[Q16 보강 — 세 방향의 일치]
행 단위 Tₑ~V_rot r = +0.024(§8b.3), 절제의 bit-identical(§8ab), 방전 단위 Tₑ 수준 null(§8an 부록)이 같은
결론을 가리킨다. "V_rot는 검정력 부족"이라는 이전 문장은 §8al §4에서 철회되었다.
출처: THESIS_RESULTS.md §8ac–§8an, §9.1–§9.4, PREREGISTRATION_B9.md.
""")

# --- 25. Reproducibility + what remains ------------------------------------
s = slide()
header(s, "재현성 · 남은 일", "수치는 손으로 옮기지 않으며, 결정 기록은 세 날짜에 걸쳐 있다")
fcard(s, 0.55, 1.45, 6.03, 2.5, "수치 파이프라인 (§8h 재발 방지 장치)", [
    "동결 run 디렉터리 → collect_paper_numbers.py(v2) → paper_numbers.json →",
    "main.tex · main_ko.tex · make_figures_en.py이며 수기 전사는 금지된다.",
    "B.9의 수치는 data/.b9_*.json(reach_ladder · family · minimal_family · pooled_ladder ·",
    "latency · op_counts)과 shot_covariates.py 산출물에서 읽는다.",
    "동결 실행 재채점은 가산 키 + 기존 키 bit-identical 확인이 규칙이며 ctx = full은",
    "동결 se_model을 비트 단위로 재현한다(§8ac). 채점 모집단 키는 매번 확인된다.",
], accent=NAVY, body_size=11.5, tag="s22-pipe")
fcard(s, 6.75, 1.45, 6.03, 2.5, "코드와 산출물", [
    "백본 experiments/seq/model_seq_v2.py, 해석가능 단 model_seq_b3.py, 윈도 대조군",
    "model_iter009.py(SHA-256 고정, 테스트가 검사)이다.",
    "B.9는 experiments/reach·b9_reach·b9_family·b9_minimal·b9_latency이며 계열 모델의",
    "인과성·수용 범위와 스트리밍 캐시 동치(3e-7)를 테스트가 단언한다(46/46).",
    "프로토콜은 PREREGISTRATION_{W2,B9,B6}.md, 감사는 PROTOCOL_AUDIT.md이다.",
    "μs shot 집합은 hires_shots/SELECTION.md · folds.py, 양자 가지는 experiments/quantum/이다.",
], accent=BLUE, body_size=11.5, tag="s22-code")
fcard(s, 0.55, 4.15, 6.03, 2.5, "남은 일", [
    "□ 논문 본문에 B.9 절(문맥·구조·비용)과 §9 프레이밍을 반영한다. main_ko.tex는",
    "   2026-08-16 이후 개정되지 않았으며 그림 fig_context_family_ladder만 추가되었다.",
    "□ B.6: μs 재획득 데이터가 도착하면 사전등록 §1.4 양성 대조 5종(#32092 EHO 포함)으로",
    "   파일럿 → 확대한다. 예측: V_rot 승률은 변동 3분위에서 먼저 올라야 한다.",
    "□ 조용한 기계에서 5세션 지연 프로토콜을 재실행하여 1 ms 판정을 해소한다.",
    "□ 논문 검토: 초록·기여 문구, 두 모집단 서술 톤, 제목 유지 여부이다.",
], accent=ORANGE, body_size=11.5, tag="s22-todo")
fcard(s, 6.75, 4.15, 6.03, 2.5, "결정 기록 (승상님)", [
    "08-16 ① 두 모집단 공동 1차를 유지한다(p100 단일 headline 금지). 근거는 §8ab의 절제이다.",
    "08-16 ② V_rot 프로토콜은 불변이다. 값 컷·점프 규칙·재학습 없이 SSE 비중만 병기한다.",
    "08-21 ③ μs shot 집합을 동결한다: test 4(31921·31873·31114·31902) / pool 6 / companion 2.",
    "   #32092는 같은 날 confirmed 되었고 B.6 양성 대조 5번이 되었다.",
    "08-24 ④ 양자 가지는 하드웨어 검증 후 음성으로 종결한다(λ = 0.661 ± 0.040, 22.9 s/예측).",
    "이 네 결정이 현재 판정표(3번 슬라이드)의 전제이다.",
], accent=TEAL, body_size=11.5, tag="s22-dec")
note(s, """
[마무리 진술]
관측 지점에서 이온온도는 미래를 읽는 보간을 두 모집단·네 분할 전부에서 유의하게 이기고, 최강 오프라인
평활기와는 동률이다. 배치에 관한 주장은 두 스트레스를 모두 견딘다. 약 50 ms의 연속 인과 문맥이 최강 배치
기준선에 대한 우위를 전형적으로 만들며, 계열은 skill이 아니라 비용을 정한다. 회전의 동률은 실패가 아니라
물리로 예측되고 절제와 승패 분석으로 확인된 발견이며, 그 발견이 다음 데이터 과제를 지목한다. 작동하지
않는 지점도 전부 판정표에 있고, 각각 그것을 뒤집을 측정과 함께 있다.

[숫자 최종 점검 목록]
Tᵢ vs PCHIP 컷 +0.174/+0.248/+0.257/+0.264 · 포함 +0.225/+0.238/+0.292/+0.316 · vs 인과 GP 4/4+4/4 ·
V_rot 1/4·2/4 · B.1 pooled +0.081 [+0.067,+0.096] · b3k8 +0.002 / -0.194 · 폭 스윕 ±0.008 · MNAR 도달
54–68% / 4–6% · 캠페인 seq 4/4+4/4, 윈도 2/4·0/4 · peak 8/8 · conformal 32/32 · 지연 1.05 / 1.61 ms ·
held 54% · 스파이크 0.53% · 값 컷 19% · 문맥 포화 50 ms · 통합 20 ms +0.057, 630 ms +0.143 · 승률
0.52→0.66 · 추세 +0.050/decade · 계열 최대 차 0.023 · 순환 111 / 합성곱 +48/층 / attention 473 ops ·
2–3 µs/op · 세션 산포 21.84배 · 방전 승률 Tᵢ 0.695 / V_rot 0.481 · 3분위 42/83/85 · 34/48/55 ·
tcn2k 1,808 4/4 · 양자 λ 0.661 · 22.9 s.
전부 THESIS_RESULTS.md §8ab·§8ac–§8ap와 paper_numbers.json·data/.b9_*.json에서 재확인 가능하다.
""")

# ============================ SAVE =========================================

OUT = os.path.join(HERE, "KSTAR_CES_종합방어.pptx")

if __name__ == "__main__":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    n_slides = len(prs.slides._sldIdLst)
    n_notes = sum(1 for sl in prs.slides if sl.has_notes_slide
                  and sl.notes_slide.notes_text_frame.text.strip())
    assert n_notes == n_slides, f"notes({n_notes}) != slides({n_slides})"
    prs.save(OUT)
    print("saved:", OUT, "| slides:", n_slides, "| notes:", n_notes)
    for w in _WARNED:
        print("  FIT WARNING:", w)
    if not _WARNED:
        print("  layout: every card, bullet block and table cell fits at its chosen size")
