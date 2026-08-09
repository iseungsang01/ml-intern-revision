# -*- coding: utf-8 -*-
"""Approximate PPTX -> PNG renderer for layout QC (no PowerPoint/LibreOffice needed).

This is a *proofing* tool, not a faithful renderer. It draws rectangles, pictures and
text runs with the real fonts at the real coordinates so that overflow, collision and
clipping problems are visible before opening PowerPoint.

Usage (from repo root):
    python docs/presentation/preview_pptx.py docs/presentation/KSTAR_CES_발표자료_20분.pptx
    python docs/presentation/preview_pptx.py <deck.pptx> --out <dir> --slides 3,5,9 --dpi 110
"""
import argparse
import io
import os
import re
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU_PER_IN = 914400
FONT_DIR = os.path.join(os.environ.get("WINDIR", "C:/Windows"), "Fonts")
FONT_FILES = {
    ("Malgun Gothic", False): "malgun.ttf",
    ("Malgun Gothic", True): "malgunbd.ttf",
    ("Consolas", False): "consola.ttf",
    ("Consolas", True): "consolab.ttf",
}
_FONT_CACHE = {}

# words stay together; CJK breaks per character
_TOKEN = re.compile(r"[A-Za-z0-9_.,;:!?()\[\]{}<>/@#$%^&*+=~`'\"\\|·—–-]+|\s+|.")


def load_font(name, bold, px):
    key = (name or "Malgun Gothic", bool(bold), int(px))
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    fname = FONT_FILES.get((key[0], key[1])) or FONT_FILES.get((key[0], False)) \
        or FONT_FILES[("Malgun Gothic", key[1])]
    try:
        f = ImageFont.truetype(os.path.join(FONT_DIR, fname), max(int(px), 1))
    except OSError:
        f = ImageFont.load_default()
    _FONT_CACHE[key] = f
    return f


def rgb(color_fmt):
    try:
        c = color_fmt.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


def shape_fill(shape):
    try:
        if shape.fill.type is None or int(shape.fill.type) != 1:  # 1 = MSO_FILL.SOLID
            return None
        return rgb(shape.fill.fore_color)
    except Exception:
        return None


def shape_line(shape):
    try:
        if shape.line.fill.type is None or int(shape.line.fill.type) != 1:
            return None
        return rgb(shape.line.color)
    except Exception:
        return None


def draw_text_frame(draw, shape, scale):
    tf = shape.text_frame
    x0 = shape.left / EMU_PER_IN * scale
    y0 = shape.top / EMU_PER_IN * scale
    w = shape.width / EMU_PER_IN * scale
    h = shape.height / EMU_PER_IN * scale
    ml = (tf.margin_left or 0) / EMU_PER_IN * scale
    mr = (tf.margin_right or 0) / EMU_PER_IN * scale
    mt = (tf.margin_top or 0) / EMU_PER_IN * scale
    avail = max(w - ml - mr, 1)

    # ---- layout pass: build lines of (text, font, color) segments
    lines = []          # list of (segments, line_height, space_after, align)
    for para in tf.paragraphs:
        align = para.alignment
        ls = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        sa = (para.space_after.pt if para.space_after is not None else 0) * scale / 72.0
        segs_in_line, line_w, max_px = [], 0.0, 0.0
        pending = []

        def flush():
            nonlocal segs_in_line, line_w, max_px, pending
            if pending:
                segs_in_line.extend(pending)
                pending = []
            lines.append((segs_in_line, max_px * 1.24 * ls, 0.0, align))
            segs_in_line, line_w, max_px = [], 0.0, 0.0

        for run in para.runs:
            pt = run.font.size.pt if run.font.size is not None else 18.0
            px = pt * scale / 72.0
            font = load_font(run.font.name, run.font.bold, px)
            col = rgb(run.font.color) or (34, 43, 53)
            max_px = max(max_px, px)
            for tok in _TOKEN.findall(run.text):
                tw = draw.textlength(tok, font=font)
                if tok.isspace():
                    if line_w + tw > avail:
                        flush()
                        max_px = max(max_px, px)
                        continue
                    segs_in_line.append((tok, font, col, line_w))
                    line_w += tw
                    continue
                if line_w + tw > avail and line_w > 0:
                    flush()
                    max_px = max(max_px, px)
                segs_in_line.append((tok, font, col, line_w))
                line_w += tw
        # end of paragraph
        lines.append((segs_in_line, max_px * 1.24 * ls if max_px else 0, sa, align))

    total = sum(lh + sa for _, lh, sa, _ in lines)
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y = y0 + (h - total) / 2
    elif tf.vertical_anchor == MSO_ANCHOR.BOTTOM:
        y = y0 + h - total
    else:
        y = y0 + mt

    overflow = total > h + 2  # rendered text taller than its box
    for segs, lh, sa, align in lines:
        if segs:
            line_w = segs[-1][3] + draw.textlength(segs[-1][0], font=segs[-1][1])
            if align == PP_ALIGN.CENTER:
                off = (avail - line_w) / 2
            elif align == PP_ALIGN.RIGHT:
                off = avail - line_w
            else:
                off = 0
            for tok, font, col, dx in segs:
                draw.text((x0 + ml + off + dx, y), tok, font=font, fill=col)
        y += lh + sa
    return overflow


def render(path, out_dir, only=None, dpi=110):
    prs = Presentation(path)
    scale = dpi
    W = int(prs.slide_width / EMU_PER_IN * scale)
    H = int(prs.slide_height / EMU_PER_IN * scale)
    os.makedirs(out_dir, exist_ok=True)
    made, warnings = [], []

    for idx, slide in enumerate(prs.slides, start=1):
        if only and idx not in only:
            continue
        img = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            if shape.left is None:
                continue
            x = shape.left / EMU_PER_IN * scale
            y = shape.top / EMU_PER_IN * scale
            w = shape.width / EMU_PER_IN * scale
            h = shape.height / EMU_PER_IN * scale
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    pic = pic.resize((max(int(w), 1), max(int(h), 1)))
                    img.paste(pic, (int(x), int(y)))
                except Exception as e:  # pragma: no cover
                    warnings.append(f"slide {idx}: picture failed ({e})")
                continue
            fill = shape_fill(shape)
            line = shape_line(shape)
            if fill or line:
                draw.rectangle([x, y, x + w, y + h], fill=fill, outline=line)
            if shape.has_text_frame and shape.text_frame.text.strip():
                if draw_text_frame(draw, shape, scale):
                    warnings.append(
                        f"slide {idx}: text overflows its box @ "
                        f"({shape.left/EMU_PER_IN:.2f}, {shape.top/EMU_PER_IN:.2f}) "
                        f"-> {shape.text_frame.text[:42]!r}")
            if x < -1 or y < -1 or x + w > W + 1 or y + h > H + 1:
                warnings.append(
                    f"slide {idx}: shape outside slide @ "
                    f"({shape.left/EMU_PER_IN:.2f}, {shape.top/EMU_PER_IN:.2f})")
        out = os.path.join(out_dir, f"slide_{idx:02d}.png")
        img.save(out)
        made.append(out)

    print(f"rendered {len(made)} slide(s) -> {out_dir}")
    for w in warnings:
        print("  WARN", w)
    if not warnings:
        print("  no overflow/out-of-bounds warnings")
    return made


def main():
    try:  # warnings quote slide text; console may be cp949
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--out", default=None)
    ap.add_argument("--slides", default=None, help="e.g. 3,5,9")
    ap.add_argument("--dpi", type=int, default=110)
    a = ap.parse_args()
    only = {int(v) for v in a.slides.split(",")} if a.slides else None
    out = a.out or os.path.join(os.path.dirname(os.path.abspath(a.pptx)), ".preview")
    render(a.pptx, out, only=only, dpi=a.dpi)


if __name__ == "__main__":
    sys.exit(main())
