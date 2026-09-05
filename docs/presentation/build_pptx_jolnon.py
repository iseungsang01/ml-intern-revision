# -*- coding: utf-8 -*-
"""Build the **paper-flow** deck on the *졸논 정리* template.

Output: docs/presentation/KSTAR_CES_졸논정리_논문흐름.pptx

이 덱은 앞의 네 덱과 다른 자리를 차지한다. `build_pptx*.py`가 만드는 덱들은 카드와
밴드를 직접 그리는 자체 레이아웃이지만, 이 덱은 **승상님이 5월에 쓰신 제안 단계 덱
`docs/졸논 정리.pptx`의 형식**을 그대로 물려받는다: 같은 "분할" 테마, 같은 슬라이드
마스터와 레이아웃, 같은 표지 배치, 그리고 제목 개체 틀을 (4.87, 1.65)로 옮겨 48 pt
accent1으로 한 단어만 적는 같은 **구역 표지** 관례. 달라진 것은 내용이며, 그 순서는
`docs/paper/main_ko.tex`의 절 순서를 그대로 따른다:

    서론 -> 관련 연구 -> 데이터와 문제 설정 -> 모델 -> 평가 방법론 -> 결과
    -> 모델 선택 프로토콜 -> 배치 가능성 -> 남은 개선 여지 -> 한계 -> 결론

빌드 방식은 템플릿에서 슬라이드만 비우고 새로 채우는 것이다. 테마·마스터·레이아웃·
표 스타일이 전부 원본 파일에서 오므로 형식 일치는 재현이 아니라 상속이다. 한 가지만
바꾸었다: 이 테마의 한글 테마 글꼴은 휴먼매직체(H2MKPB)인데, 본문 밀도가 제안 덱보다
훨씬 높아 가독성이 떨어지므로 모든 런의 동아시아 글꼴을 **맑은 고딕**으로 지정한다.
원본 덱도 이미 맑은 고딕을 섞어 쓰고 있었다(docProps/app.xml). 라틴 글꼴(Gill Sans MT)과
테마 색은 그대로 둔다.

문체는 논문 초록체(서술형 종결, 객관·비인칭)이며, 모든 수치는 `docs/paper/main_ko.tex`
(= 동결 산출물 `paper_numbers.json`에서 수집된 값)에서 그대로 옮겼다. 각 슬라이드의
발표자 노트에 대응하는 절의 `\\label`을 적어 두어, 인용할 때 논문의 어느 자리인지
바로 찾을 수 있다.

Usage (from repo root):
    python docs/presentation/build_pptx_jolnon.py
"""
import copy
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
if HERE not in sys.path:
    sys.path.insert(0, HERE)

from preview_pptx import load_font, _TOKEN  # noqa: E402  (same metrics as the QC renderer)
import appendix_content as AP  # noqa: E402  (shared with build_pptx.py's appendix)
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.dml import MSO_THEME_COLOR  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

REPO = os.path.dirname(os.path.dirname(HERE))
TEMPLATE = os.path.join(REPO, "docs", "졸논 정리.pptx")
PAPER_FIG = os.path.join(REPO, "docs", "paper", "figures")
OUT = os.path.join(HERE, "KSTAR_CES_졸논정리_논문흐름.pptx")

EA_FONT = "맑은 고딕"
QC_FONT = "Malgun Gothic"          # preview_pptx's metric name for the same face

# Geometry inherited from the template (docs/졸논 정리.pptx).
TITLE_BOX = (0.64, 0.77, 12.06, 1.11)      # layout "제목 및 내용" title placeholder
BODY_BOX = (0.64, 2.38, 12.06, 4.02)       # layout "제목 및 내용" content placeholder
DIVIDER_BOX = (4.87, 1.65, 7.98, 4.05)     # where the original deck parks its section word
COVER_TITLE = (0.68, 0.95, 11.98, 1.94)
COVER_SUB = (0.98, 5.43, 11.43, 1.12)

# Body sizes per outline level, before auto-shrink (master defaults: 18 / 16 / 14).
LEVEL_SIZE = {0: 18.0, 1: 16.0, 2: 14.0}
LEVEL_INDENT_IN = {0: 0.335, 1: 0.689, 2: 0.984}
SPACE_AFTER_PT = {0: 6.0, 1: 3.0, 2: 3.0}

_MEASURE = ImageDraw.Draw(Image.new("RGB", (8, 8)))
_TEXT_MARGIN_IN = 2 * (0.1)   # placeholder lIns/rIns = 91440 EMU = 0.1 in each side
WARNINGS = []


# ---------------------------------------------------------------- text metrics
def _n_lines(txt, avail_in, size_pt, bold=False):
    """Line count after wrapping, using preview_pptx's tokenizer and font metrics."""
    font = load_font(QC_FONT, bold, size_pt)      # px == pt -> lengths are in points
    avail_px = max(avail_in, 0.1) * 72.0
    n, cur = 1, 0.0
    for tok in _TOKEN.findall(txt):
        tw = _MEASURE.textlength(tok, font=font)
        if tok.isspace():
            if cur + tw > avail_px:
                n += 1
                cur = 0.0
                continue
            cur += tw
            continue
        if cur + tw > avail_px and cur > 0:
            n += 1
            cur = 0.0
        cur += tw
    return n


def _block_h(items, width_in, scale):
    """Wrapped height in inches of a (level, text) list at the given size scale."""
    total = 0.0
    for level, txt in items:
        size = LEVEL_SIZE[level] * scale
        avail = width_in - _TEXT_MARGIN_IN - LEVEL_INDENT_IN[level]
        total += _n_lines(txt, avail, size) * size * 1.24 / 72.0
        total += SPACE_AFTER_PT[level] * scale / 72.0
    return total


# ------------------------------------------------------------------- xml bits
def _ea(run, name=EA_FONT):
    """Pin the East Asian face on a run (python-pptx only exposes the Latin one)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _drop(shape):
    shape._element.getparent().remove(shape._element)


def _move(shape, box):
    x, y, w, h = box
    shape.left, shape.top, shape.width, shape.height = (
        Inches(x), Inches(y), Inches(w), Inches(h))


# --------------------------------------------------------------- slide makers
prs = Presentation(TEMPLATE)


def _clear_template_slides():
    """Keep the theme, masters, layouts and table styles; drop the 18 May slides."""
    id_list = prs.slides._sldIdLst
    for sld_id in list(id_list):
        prs.part.drop_rel(sld_id.rId)
        id_list.remove(sld_id)


def _new(layout_idx=1):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def note(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()
    for para in s.notes_slide.notes_text_frame.paragraphs:
        for run in para.runs:
            _ea(run)
    return s


def _title(s, txt, box=TITLE_BOX, size=None, bold=False):
    ph = s.shapes.title
    _move(ph, box)
    tf = ph.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = txt
    run.font.bold = bold
    _ea(run)
    if size is None:                       # shrink to one line inside the title band
        size = 28.0
        avail = box[2] - _TEXT_MARGIN_IN
        while size > 17.0 and _n_lines(txt, avail, size, bold) > 1:
            size -= 0.5
        if _n_lines(txt, avail, size, bold) > 1:
            WARNINGS.append("title wraps: %r" % txt)
    run.font.size = Pt(size)
    return ph


def divider(word, note_txt=""):
    """Section header exactly as the template does it: title box moved, 48 pt accent1."""
    s = _new(1)
    _drop(s.placeholders[1])
    ph = _title(s, word, box=DIVIDER_BOX, size=48.0)
    ph.text_frame.paragraphs[0].runs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    ph.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if note_txt:
        note(s, note_txt)
    return s


def bullets(s, items, box=BODY_BOX, min_scale=0.62):
    """Fill the content placeholder, shrinking uniformly until the text fits the box."""
    ph = s.placeholders[1]
    _move(ph, box)
    tf = ph.text_frame
    tf.word_wrap = True

    scale = 1.0
    while scale > min_scale and _block_h(items, box[2], scale) > box[3]:
        scale -= 0.02
    if _block_h(items, box[2], scale) > box[3]:
        WARNINGS.append("body overflows at scale %.2f: %r" % (scale, items[0][1][:40]))

    first = True
    for level, txt in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.level = level
        para.space_after = Pt(SPACE_AFTER_PT[level] * scale)
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(round(LEVEL_SIZE[level] * scale, 1))
        _ea(run)
    return s


def content(title, items, note_txt=""):
    s = _new(1)
    _title(s, title)
    bullets(s, items)
    if note_txt:
        note(s, note_txt)
    return s


def _para_h(lines, width_in, size, space_after=3.0):
    """Wrapped height in inches of a plain text block."""
    total = 0.0
    for txt in lines:
        total += _n_lines(txt, width_in - _TEXT_MARGIN_IN, size) * size * 1.24 / 72.0
        total += space_after / 72.0
    return total + 0.06


def _textbox(s, x, y, w, lines, size, space_after=3.0):
    """Textbox sized to its own wrapped content. Returns the height it occupies."""
    h = _para_h(lines, w, size, space_after)
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(size)
        _ea(run)
    return h


def _row_heights(headers, rows, col_w_in, size, min_row_h=0.30):
    """Per-row height a table actually needs (PowerPoint grows rows, never shrinks)."""
    out = []
    for r, row in enumerate([headers] + rows):
        bold = (r == 0)
        n = 1
        for c, txt in enumerate(row):
            n = max(n, _n_lines(txt, col_w_in[c] - 0.14, size, bold))
        out.append(max(min_row_h, n * size * 1.24 / 72.0 + 0.06))
    return out


def table_slide(title, headers, rows, widths, lead=None, tail=None,
                size=11.0, note_txt=""):
    """Title + optional lead + table + optional tail, all measured so nothing collides."""
    s = _new(1)
    _title(s, title)
    _drop(s.placeholders[1])

    x, top, w, h = BODY_BOX
    y = top - 0.16
    if lead:
        y += _textbox(s, x, y, w, [lead], 13.0) + 0.08

    tail_h = (_para_h([tail], w, 12.5) + 0.10) if tail else 0.0
    avail = (top + h) - y - tail_h

    total_w = sum(widths)
    col_w_in = [w * cw / total_w for cw in widths]
    while size > 8.5:
        heights = _row_heights(headers, rows, col_w_in, size)
        if sum(heights) <= avail:
            break
        size -= 0.5
    heights = _row_heights(headers, rows, col_w_in, size)
    if sum(heights) > avail:
        WARNINGS.append("table too tall at %.1f pt: %r" % (size, title))

    n_rows = len(rows) + 1
    gtbl = s.shapes.add_table(n_rows, len(headers), Inches(x), Inches(y),
                              Inches(w), Inches(sum(heights)))
    tbl = gtbl.table
    for i, cw in enumerate(col_w_in):
        tbl.columns[i].width = Inches(cw)
    for r, rh in enumerate(heights):
        tbl.rows[r].height = Inches(rh)

    def _cell(cell, txt, bold):
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        cell.text_frame.word_wrap = True
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        _ea(run)

    for c, htxt in enumerate(headers):
        _cell(tbl.cell(0, c), htxt, True)
    for r, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            _cell(tbl.cell(r, c), txt, False)

    if tail:
        _textbox(s, x, y + sum(heights) + 0.10, w, [tail], 12.5)
    if note_txt:
        note(s, note_txt)
    return s


def figure_slide(title, lines, image, note_txt=""):
    """Title + a few lines + the paper's own figure, both inside the content box."""
    s = _new(1)
    _title(s, title)
    _drop(s.placeholders[1])

    x, top, w, h = BODY_BOX
    y = top - 0.16
    if lines:
        y += _textbox(s, x, y, w, lines, 14.0, space_after=4.0) + 0.12

    room_h = (top + h) - y
    im = Image.open(image)
    ratio = im.size[0] / float(im.size[1])
    pw = min(w, room_h * ratio)
    ph_in = pw / ratio
    s.shapes.add_picture(image, Inches(x + (w - pw) / 2.0), Inches(y),
                         Inches(pw), Inches(ph_in))
    if note_txt:
        note(s, note_txt)
    return s


def cover(title, subtitle_lines):
    s = _new(0)
    ph = s.shapes.title
    _move(ph, COVER_TITLE)
    ph.text_frame.word_wrap = True
    ph.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = ph.text_frame.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True
    _ea(run)

    sub = s.placeholders[1]
    _move(sub, COVER_SUB)
    sub.text_frame.word_wrap = True
    first = True
    for line, size in subtitle_lines:
        para = sub.text_frame.paragraphs[0] if first else sub.text_frame.add_paragraph()
        first = False
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        _ea(run)
    return s


# =============================================================================
# 1. 표지 · 목차
# =============================================================================
_clear_template_slides()

cover(
    "KSTAR 희소 전하교환분광(CES) 신호의 멀티모달 나우캐스팅",
    [("고속 진단은 미래를 보는 보간을 넘어 이온온도를 복원하지만 "
      "토로이달 회전은 복원하지 못한다", 17),
     ("Department of Nuclear Engineering, Seoul National University, Seoul, Korea", 14),
     ("이승상 (Seungsang Lee)", 14)],
)

content(
    "발표 순서는 논문의 절 순서를 그대로 따른다",
    [(0, "1. 서론 — CES 결측이라는 문제, 연구 질문, 그리고 의도적으로 불리하게 세운 기준선"),
     (0, "2. 관련 연구 — 세 갈래의 선행 연구와 그 안에서 본 연구가 놓이는 자리"),
     (0, "3. 데이터와 문제 설정 — 641개 방전, 두 가지 프레이밍, 두 번의 데이터 품질 감사"),
     (0, "4. 모델 — 손실과 환원 불가능한 하한, 보존식이 정하는 입력 라우팅, 시퀀스 백본"),
     (0, "5. 평가 방법론 — 사전등록 규칙, 기준선 사다리, shot 군집 paired bootstrap"),
     (0, "6. 결과 — 헤드라인, 두 개의 스트레스 테스트, 정보 비대칭, 복잡도 사다리, 문맥과 비용"),
     (0, "7. 모델 선택 프로토콜 — 모든 결정 규칙은 그것이 결정할 수치보다 먼저 적혔다"),
     (0, "8. 배치 가능한가 — 상태 유지형 1-스텝 지연시간과 분포 무가정 예측 구간"),
     (0, "9. 남은 개선 여지 · 10. 한계 — 부정 결과마다 그것을 뒤집을 측정을 지목한다"),
     (0, "11. 결론")],
    note_txt="논문: docs/paper/main_ko.tex. 이 덱의 모든 수치는 그 파일에서 그대로 옮겼다.",
)

# =============================================================================
# 2. 서론
# =============================================================================
divider("Introduction", "sec:intro")

content(
    "CES는 이온온도와 회전을 재는 표준 진단이지만 느리고 간헐적으로 결측된다",
    [(0, "전하교환분광(CES)은 전하교환 선방출의 도플러 확장에서 이온온도 T_i를, "
         "도플러 이동에서 토로이달 회전 V_rot를 추정한다."),
     (0, "충분한 신호대잡음비를 위해 광자를 적분해야 하므로 본질적으로 느리고, "
         "개별 시점 측정은 노출·신호품질 문제로 자주 소실된다."),
     (0, "본 연구의 균일한 10 ms 격자에서 T_i는 8.2 %, V_rot는 23.9 %가 서로 "
         "독립적으로 결측된다."),
     (0, "V_rot는 여기에 더해 41.1 %가 직전 관측을 반복하는 계측기 유지값이어서, "
         "격자의 65.0 %가 독립적인 V_rot 정보를 갖지 않는다."),
     (0, "결측은 저신호 구간·ELM·천이에 집중되므로 무작위적이지 않다(MNAR)."),
     (0, "한편 BES·ECEI·Mirnov 코일은 같은 플라즈마를 결측 없이 관측한다. "
         "다만 이들 중 어느 것도 T_i나 V_rot를 직접 측정하지 않는다.")],
    note_txt="sec:intro, 그림 fig:missing. 결측률은 확정 프로토콜(유지값 제거) 기준이다.",
)

content(
    "연구 질문과, 의도적으로 모델에 불리하게 세운 기준선",
    [(0, "질문: CES가 결측된 10 ms 시점에서, 동시각 고속 진단과 불규칙한 과거 CES "
         "이력만으로 CES 단독의 시간 보간이 복원할 수 없는 정보를 회복할 수 있는가."),
     (0, "비교는 모델에 불리하게 설계하였다. 보간 기준선(선형, PCHIP, 국소 AR, "
         "가우시안 과정)은 오프라인으로 타겟 시점 양쪽의 CES 이웃을 쓴다."),
     (1, "모델은 엄격히 인과적이며 타겟 시점까지의 고속 진단과 과거 CES 이력만 본다."),
     (0, "인과 모델이 미래를 쓰는 보간을 이긴다면, 고속 진단이 CES의 시간 구조만으로는 "
         "담을 수 없는 정보를 운반한다는 뜻이다."),
     (0, "배치 가능한 가장 강한 경쟁자는 persistence가 아니라 과거만 쓰는 평활기이므로, "
         "인과 GP 팔을 추가하고 그에 대한 성능을 따로 보고한다."),
     (0, "긍정적 답은 강한 역산 가정 없이 CES 결측을 온라인으로 채우는 데이터 기반 "
         "가상 센서를 정당화한다.")],
    note_txt="sec:intro, '의도적으로 어려운 기준선' 문단.",
)

content(
    "기여는 다섯 가지로 요약된다",
    [(0, "① 희소 타겟 나우캐스팅을 위한 인과적 전체격자 프레이밍. 라벨 없는 행을 버리는 "
         "대신 맥락으로 유지하고 희소성을 손실로 옮긴다."),
     (1, "짝지은 윈도 대조군 대비 통합 T_i skill +0.081(16회 실행, CI [+0.067, +0.096], "
         "16/16 양수)이며 학습 비용은 오히려 더 싸다."),
     (0, "② 사전등록된 두-모집단 shot 군집 평가 프로토콜. W=2, 유지값 전 구간 제거, "
         "피팅 실패 0.53 %에 대한 두 공동 1차 모집단, 선택이 읽지 않는 test 셋."),
     (0, "③ T_i에 대한 강건한 긍정 결과. 미래를 쓰는 PCHIP 대비 skill +0.17~+0.32로 "
         "두 모집단의 4개 분할 전부에서 PASS이고, 8개 셀 전부에서 인과 GP를 이긴다."),
     (0, "④ V_rot에 대한 물리 기반 부정 결과. 동률이며(1/4·2/4 유의) 고속 진단을 0으로 "
         "만들어도 출력이 비트 단위로 동일하다."),
     (0, "⑤ 두 개의 스트레스 테스트(MNAR 재가중·캠페인 이동), 복잡도와 크기 축의 종결, "
         "문맥과 구조의 분리, 그리고 측정된 배치 가능성.")],
    note_txt="sec:intro '기여' 열거. 논문은 아홉 항목이며 여기서는 다섯으로 묶었다.",
)

# =============================================================================
# 3. 관련 연구
# =============================================================================
divider("Related Works", "sec:related")

content(
    "선행 연구는 세 갈래이며, 세 갈래 모두 타겟이 측정되어 있다고 가정한다",
    [(0, "프로파일 피팅과 결측 보간. 가우시안 과정 회귀는 핵융합 프로파일의 커뮤니티 "
         "표준이고, PCHIP 같은 단조 보간은 과도 현상 근방에서 선호된다."),
     (1, "본 연구는 persistence·선형·PCHIP·국소 AR 사다리를 사전등록 기준선으로 쓰고 "
         "PCHIP을 헤드라인 비교 대상으로 삼는다."),
     (0, "신경망 진단 추론. 스펙트럼 피팅 가속, 한 진단의 추론 파이프라인 대체, 느린 "
         "재구성 코드의 종단간 대리모형의 세 갈래가 있다."),
     (1, "이들은 분석 지연시간을 다루지 측정 가용성을 다루지 않는다. CES 노출이 "
         "결측되면 입력 자체가 없다."),
     (0, "교차 진단 재구성과 시간 조밀화. Diag2Diag, COMPASS 시간 초해상, FusionMAE, "
         "RTCAKENN이 활발한 계열을 이룬다."),
     (0, "본 연구는 이 계열을 이어받되, 자기 자신의 측정이 없는 시점을 대상으로 한다.")],
    note_txt="sec:related. 표현 규칙은 docs/paper/NOVELTY.md — 부재가 아니라 확장으로 쓴다.",
)

table_slide(
    "본 연구는 가장 가까운 계열을 세 축으로 확장한다",
    ["선행 연구", "타겟 채널", "인과성", "미래 사용 보간을 기준선으로"],
    [["Diag2Diag (DIII-D)", "Thomson T_e·n_e (조밀 전자 채널)", "동시각, 기억 없음", "없음"],
     ["COMPASS 시간 초해상", "Thomson 프로파일", "동시각", "없음"],
     ["RTCAKENN (DIII-D)", "운동학 프로파일 (입력 dropout 강건성)", "실시간, 인과", "없음"],
     ["EAST XCS 추론", "T_i·회전 (같은 도플러 분광 입력)", "순간 매핑", "없음"],
     ["FusionMAE (HL-3)", "임의 마스킹 채널", "윈도 내 재구성", "없음"],
     ["KSTAR EPED 재구성", "프로파일 (CES를 입력으로 사용)", "재구성", "없음"],
     ["본 연구", "희소 이온 채널 T_i·V_rot", "타겟 자신의 불규칙 과거에 조건화", "PCHIP·오프라인 GP"]],
    widths=[2.4, 4.0, 3.2, 2.4],
    lead="계보를 인정한 뒤 세 축으로 확장한다: 전자 채널에서 희소 이온 채널로, "
         "동시각 매핑에서 인과 추정으로, 가정된 재구성 가능성에서 타겟별 검정으로.",
    tail="이 세 확장의 결합, 즉 희소 CXRS 이온 측정의 인과적 시간 결측 채움은 아직 "
         "다뤄지지 않았으며, 계열로부터의 이탈이 아니라 그 자연스러운 다음 단계다.",
    size=11.0,
    note_txt="sec:related, '교차 진단 재구성과 시간 조밀화' 문단의 (i)~(iv) 축.",
)

# =============================================================================
# 4. 데이터와 문제 설정
# =============================================================================
divider("Dataset", "sec:data")

content(
    "641개 KSTAR 방전을 공통 10 ms 격자에 정렬해 사용한다",
    [(0, "데이터 제공 측이 하드웨어 일관성과 H-mode ELM 억제(RMP) 구간을 기준으로 선정한 "
         "641개 방전(shot 30801~32751), 총 247,207행(파일당 중앙값 339행)이다."),
     (0, "각 행은 BES 9채널, ECEI 4채널, Mirnov 코일 2채널, 시간, 그리고 두 CES 타겟 "
         "T_i·V_rot를 제공한다."),
     (0, "행은 0.5 s 이상의 시간 간극에서 연속 세그먼트로 묶인다. 이 임계값은 이봉 delta "
         "분포의 골에 놓이며, 세그먼트 안에서 스텝의 99.4 %가 10 ms다."),
     (0, "전형적인 파일은 주 세그먼트 하나(중앙값 301행, 약 3.0 s)와 그로부터 멀리 떨어진 "
         "고립 단일행 몇 개로 이루어진다."),
     (0, "어떤 팔도, 보간이든 모델 입력이든 세그먼트 간극을 넘지 않는다.")],
    note_txt="sec:data 데이터셋 문단.",
)

figure_slide(
    "T_i와 V_rot는 독립적으로 결측되므로 행 필터링이 아니라 타겟별 마스킹을 요구한다",
    ["두 타겟이 함께 빠지지 않으므로, 두 타겟 모두를 요구하는 행 필터는 라벨 행의 "
         "약 28 %를 조용히 버린다.",
     "계측기 유지값까지 세면 격자의 65.0 %가 독립적인 V_rot 정보를 갖지 않는다."],
    os.path.join(PAPER_FIG, "fig_missing.png"),
    note_txt="그림 fig:missing (docs/paper/figures/fig_missing.png).",
)

content(
    "학습 예제에는 두 가지 프레이밍이 있고, 본 논문은 둘을 짝지어 비교한다",
    [(0, "윈도 프레이밍(대조군)은 타겟 인덱스 t의 샘플을 직전 W개 격자 행으로 구성한다: "
         "bes (W,9), ecei (W,4), mc (W,2), 시간 특징 (W,4), CES 이력 (W,4)."),
     (1, "시간 특징 넷은 lookback 초, 행간 delta 초와 각각의 log(1+x)이며, 과거 CES 값의 "
         "신뢰도가 10 ms 전인지 200 ms 전인지에 강하게 의존함을 노출한다."),
     (1, "CES 이력 넷은 직전 정규화 T_i, 직전 정규화 V_rot, 그리고 타겟마다 하나씩의 "
         "관측 플래그다. 본문 전체에서 W=2이다."),
     (0, "전체격자 시퀀스 프레이밍(주 모델)은 고속 진단에 결측이 없다는 점을 이용해, "
         "입력이 완전한 모든 행을 라벨 유무와 무관하게 맥락으로 유지한다."),
     (1, "각 스텝은 엄격히 인과적인 22개 채널이다: z-score된 고속 15채널, "
         "log(1+Δt_prev), 그리고 타겟별 이월값·경과시간·과거 관측 플래그."),
     (0, "두 프레이밍 어느 쪽도 모델이 타겟 행 자신의 값을 읽게 하지 않는다.")],
    note_txt="sec:framing.",
)

content(
    "설계 원칙은 셋이며, 셋 다 누수와 가짜 라벨을 막기 위한 것이다",
    [(0, "가짜 라벨 금지. 학습 행을 만들기 위해 타겟을 대체(impute)하지 않는다. 윈도 "
         "프레이밍의 행은 입력이 완전하고 CES 타겟 중 적어도 하나가 실제로 관측된 "
         "경우에만 유지되며, 시퀀스 프레이밍에서 라벨 없는 행은 맥락으로만 기여한다."),
     (0, "타겟별 마스킹 손실. 손실은 관측된 타겟에 대해서만 계산되므로 한 타겟만 관측된 "
         "행도 그 타겟은 학습에 기여한다. 두-타겟-필수 필터를 제거한 것은 순수한 데이터 "
         "이득이다."),
     (0, "누수 방지 셋. (i) 파일 단위 분할을 디스크에 고정하고, (ii) 모든 z-score 통계를 "
         "학습 파일에서만 추정하며(시퀀스 모델은 여기에 각 방전 자신의 통계로 하는 입력 "
         "표준화를 더한다), (iii) 타겟 행 자신의 값과 관측 플래그는 입력에 결코 넣지 않는다.")],
    note_txt="sec:data 설계 원칙. (ii)의 shot별 표준화가 sec:campaign에서 결정적으로 작동한다.",
)

content(
    "데이터 품질 감사 1: 관측된 V_rot의 54 %는 계측기 유지값이다",
    [(0, "같은 연속 시간 블록 안에서 직전 관측과 비트 단위로 동일한 반복이 최대 1,214행 "
         "이어지며, 641개 중 499개 파일이 영향을 받는다."),
     (0, "V_rot의 고유 측정 주기가 행 주기보다 느려 생기는 현상으로, 독립적인 측정이 "
         "아니다."),
     (0, "V_rot는 소수점 다섯 자리까지 기록되고 서로 다른 값 사이의 최소 간격이 4e-5이므로 "
         "'연속 동일'에는 측정 가능한 위양성 통로가 없다. T_i는 226,991행 중 1행뿐이다."),
     (0, "확정 프로토콜은 유지값을 모든 곳에서 제거한다: 지도 타겟, 이력·이월 입력과 그 "
         "관측 플래그, 정규화 통계, 그리고 모든 기준선의 보간 앵커."),
     (0, "짝지은 재학습은 유지값이 학습도 오염시킴을 보였다. forward-fill 계단은 이력을 "
         "복사하는 것이 거의 최적이라고 모델에 가르친다. 그래서 민감도 한 줄이 아니라 "
         "프로토콜이다.")],
    note_txt="sec:stuck.",
)

content(
    "데이터 품질 감사 2: T_i 스펙트럼 피팅 실패가 두 개의 공동 1차 모집단을 만든다",
    [(0, "관측된 T_i의 p99는 2,089 eV, p99.9는 9,601 eV, 최댓값은 14,984 eV다. 이 먼 꼬리는 "
         "플라즈마 물리가 아니라 실패한 스펙트럼 피팅이다."),
     (0, "1,197행(0.53 %)이 3 keV를 초과하며 274개 방전에서 951개의 런을 이룬다. 런의 85 %는 "
         "단일 행이고, 런 정점의 중앙값은 관측된 이웃 평균의 13배다."),
     (0, "두 대응이 모두 방어 가능하고 각각 비판의 여지가 있다. 컷은 어떤 방법도 예측할 수 "
         "없는 행을 제거하지만 '어려운 행을 없앴다'는 비판을 부르고, 포함은 스파이크 앵커를 "
         "먹인 채 오프라인 기준선에 핸디캡을 준다."),
     (0, "그래서 컷(모든 팔에서 결측 처리)과 포함(컷 없음)을 두 개의 공동 1차 모집단으로 "
         "사전등록하였다. 무조건적 주장은 두 모집단 모두에서 성립할 때만 조건 없이 진술한다."),
     (0, "임계값은 2.5~4 keV 범위에서 무의미하며, 값 기준 컷은 피팅 품질 메타데이터가 "
         "확보되기 전까지 쓰는 한쪽 방향 대리 지표다.")],
    note_txt="sec:spikes, 민감도는 sec:cutsens.",
)

# =============================================================================
# 5. 모델
# =============================================================================
divider("Methods", "sec:model")

content(
    "무엇을 추정하는가, 그리고 추정기가 제거할 수 없는 하한은 무엇인가",
    [(0, "모델은 각 격자 시각 t에서 정규화 타겟을 그 시각까지의 관측 이력으로 추정한다. "
         "손실은 타겟별 관측 마스크 위의 마스킹 MSE에, 정규화 영점 아래의 T_i 예측을 "
         "억제하는 물리 페널티(λ = 0.1)를 더한 것이다."),
     (0, "이 손실을 최소화하는 추정량은 조건부 기댓값이므로, 본 연구가 보고하는 것은 "
         "조건부 분포의 평균이지 그 분포 전체가 아니다. 불확실성은 모델을 건드리지 않는 "
         "분포 무가정 구간으로 따로 다룬다."),
     (0, "타겟 자신이 측정 잡음을 지니므로 RMSE² = σ_meas² + b²라는 하한이 있다. 차분 기반 "
         "추정기로 잰 σ_meas는 컷 모집단 T_i에서 46.4~129.9 eV이고, 백본의 RMSE 157.8 eV는 "
         "이 범위와 같은 자릿수다."),
     (0, "4차 차분 제곱질량의 46.6 %(V_rot는 65.6 %)가 상위 1 %의 한 스텝 변화에 몰려 있다. "
         "따라서 총합 손실은 꼬리 통계이며, 폭과 계열이 평평하다는 뒤의 결과는 정보의 "
         "고갈이 아니라 이 사실의 귀결로 읽어야 한다."),
     (0, "본류에서는 모델이 여전히 타겟 자체 산포의 2.3~3.4배 위에 있다.")],
    note_txt="sec:model:target, 식 (1)·(2). 재현성 상한 근거는 THESIS_RESULTS.md §8aq.",
)

table_slide(
    "V_rot 분기가 고속 진단을 보지 않는 것은 취향이 아니라 항별 감사의 결과다",
    ["운동량식의 항", "물리", "우리 데이터에 있는가", "판정"],
    [["∂L/∂t", "회전 변화율", "연속 관측쌍의 차분", "측정됨"],
     ["T_NBI", "빔 토크", "0-D 채널이 존재하지 않음", "부재"],
     ["T_NTV ~ δB²", "비축대칭 제동", "Mirnov, 100 Hz 데시메이트", "음성 종결"],
     ["∇·Π_φ", "난류 운동량 유속", "BES가 밀도요동을 봄", "측정 -> 널"],
     ["T_int", "잔여응력", "∇T_i 필요, 스칼라만 보유", "도달 불가"]],
    widths=[2.0, 2.6, 4.4, 3.0],
    lead="두 전송식은 값이 아니라 변화율을 구속한다. 그러므로 물어야 할 것은 고속 진단이 "
         "V_rot와 상관되는가가 아니라 dV_rot/dt와 상관되는가이다.",
    tail="측정값: r(BES, dV_rot/dt) = -0.006, r(ECEI, ·) = -0.003 대 에너지식의 +0.070, "
         "+0.078. 값 수준도 +0.341 / +0.311 대 +0.027 / +0.005. 같은 측정이 에너지식에서는 "
         "신호를 잡으므로 회전의 널은 방법의 한계가 아니다.",
    size=11.5,
    note_txt="sec:model:physics, 표 tab:torque. 산출물은 data/.physics_scales.json (§8ar).",
)

content(
    "세 개의 시간 척도가 이 문제를 규정한다",
    [(0, "격자 간격이 10 ms이므로 Nyquist 주파수는 50 Hz다. Mirnov가 담아야 할 모드 회전은 "
         "kHz 대역이고(이 데이터셋의 방전들에서 보고된 EHO 조화파는 약 4·8 kHz), 반앨리어싱 "
         "필터 없이 100 Hz로 데시메이트된 dB/dt는 그것을 표현할 수 없다."),
     (0, "전자와 이온을 결합시키는 온도 평형화 시간은 관측 T_i 중앙값 593 eV에 대응하는 "
         "구석에서 8~59 ms다. 본 연구가 측정한 연속 인과 문맥의 포화점은 약 50 ms이며, "
         "두 수는 같은 자릿수다."),
     (1, "이는 자릿수 정합성 진술이고 그 이상이 아니다. 플라즈마 파라미터는 인용값이며, "
         "BES와 ECEI는 계측기 단위이므로 n_e와 T_e를 물리 단위로 갖지 않는다."),
     (0, "신호가 스스로를 닮아 있는 시간은 T_i 159 ms, BES 161 ms, ECEI 147 ms로 서로 10 % "
         "이내에서 일치한다."),
     (0, "V_rot는 그렇지 않다. 유지값을 제거하면 16 ms, 유지하면 300 ms를 넘어 19배 이상 "
         "벌어지므로, V_rot의 완화 시간은 이 데이터셋에서 측정 불가능하다.")],
    note_txt="sec:model:scales, 식 (3)·(4).",
)

content(
    "주 모델은 전체격자 인과 시퀀스 백본이고, 대조군과 해석 가능한 칸이 함께 놓인다",
    [(0, "seq_v2 (357,570 파라미터). 세그먼트의 22채널 격자 위에서 서로 독립인 두 인과 LSTM을 "
         "돌린다. T_i 분기(2층·160 유닛)는 전체 상태를 읽고, V_rot 분기(1층·64 유닛)는 고속이 "
         "아닌 7채널만 읽는다."),
     (1, "도달 범위가 고정 윈도가 아니라 세그먼트 전체이며, 라우팅이 인코더 수준에서 "
         "성립하므로 고속 15채널을 전부 교란해도 V_rot 출력이 비트 단위로 동일하다."),
     (0, "윈도 대조군 (201,258 파라미터). 모달리티별 시간 인지 1-D CNN과 양방향 GRU를 거쳐, "
         "타겟마다 하나씩의 멀티헤드 가산 어텐션으로 풀링하되 그 타겟이 실제로 측정된 행에만 "
         "어텐션 질량이 놓이도록 하드 마스킹한다."),
     (1, "보간 자신의 귀납 편향(관측된 표본만 사용)을 파라미터 비용 0으로 어텐션에 써 넣은 "
         "것이다."),
     (0, "b3k8 (21,498 파라미터). 예측이 이월값 + Σ w_k z_k + b로 정확히 분해되고 readout이 "
         "0으로 초기화되어 학습이 정확히 persistence에서 시작한다."),
     (0, "두 모델이 공유하는 데이터 계약·유지값 처리·마스킹·분할·채점 모집단은 동일하게 "
         "고정되므로, 비교는 행 단위로 짝지어진다.")],
    note_txt="sec:model:backbone / sec:model:window / sec:model:rung.",
)

content(
    "최적화 설정은 물려받은 것이며, 본 논문은 그 사실을 감추지 않는다",
    [(0, "AdamW(η = 1e-3, weight decay 1e-4), 배치 16 세그먼트, gradient norm clipping 1.0, "
         "ReduceLROnPlateau와 조기 종료(최대 30 epoch, 확정 실행은 14~25에서 정지)."),
     (0, "데이터와 구조의 축은 전부 통제 실험으로 닫혔지만 최적화 축은 그런 측정을 거치지 "
         "않았고, 이 설정은 윈도 파이프라인에서 물려받았다."),
     (0, "현재 이를 방어하는 증거는 간접적이다. 예산 균등화 조건(고정 10 epoch, 검증 기반 "
         "체크포인트 선택 금지)에서도 백본의 짝지은 T_i 이득은 4개 분할 전부 양수였다"
         "(+0.063, +0.033, +0.045, +0.030). 결론은 학습 일정의 인공물이 아니다."),
     (0, "'왜 배치 16인가'는 이제 측정된 답을 갖는다. 경사 잡음 척도를 학습 블록 452개 전부에 "
         "대해 계산하면 중앙값 33.6 블록이므로, 배치 16은 측정된 잡음 척도의 절반쯤에 있다."),
     (1, "첫 epoch 이후 블록 하나의 그래디언트 에너지 가운데 86~98 %가 잡음이다. 이는 계산 "
         "비용과 스텝 수의 교환에 관한 진술이지 최종 skill에 관한 것이 아니다."),
     (0, "μP는 아직 수행하지 않았다. 그것을 쓰면 폭 스윕이 평평한 이유가 용량인지 학습률인지의 "
         "구분이 구성상 불필요해진다.")],
    note_txt="sec:model:optim, 식 (5).",
)

# =============================================================================
# 6. 평가 방법론
# =============================================================================
divider("Evaluation", "sec:eval")

content(
    "지표는 PCHIP 대비 skill이고, 판정은 인과 GP에 대해 내린다",
    [(0, "모든 오차는 물리 CES 단위로 역정규화되어 타겟별로 계산되며, 1차 점수는 "
         "skill = 1 - MSE_model / MSE_PCHIP이다(양수 = 모델 우위, 0 = 동률)."),
     (0, "모든 팔은 동일한 (파일, 행) 집합에서 동일한 타겟별 유지 마스크로 채점되고, 어떤 "
         "짝지은 비교 이전에도 모집단 키가 비트 단위로 동일함을 검증한다."),
     (0, "인과 기준선은 persistence, 국소 AR, 그리고 같은 GP를 과거 이웃 16개로 제한한 "
         "인과 GP다. 미래를 읽는 팔은 선형·PCHIP·오프라인 GP다."),
     (0, "인과 GP가 배치 가능한 가장 강한 기준선이므로(T_i RMSE 164.3 대 persistence 197.2), "
         "'배치 가능한 모든 인과 방법을 이긴다'는 persistence가 아니라 이 기준선에 대해 "
         "판정한다."),
     (0, "보간도 모델 입력도 세그먼트 경계를 넘지 않으며, 경계 너머의 이웃이 필요한 지점에서 "
         "보간은 persistence 값을 대신 예측한다.")],
    note_txt="sec:eval 지표·기준선 사다리 문단.",
)

content(
    "사전등록 규칙과, 유의성을 방전 단위로 세는 이유",
    [(0, "PR1 헤드라인 비교 대상은 PCHIP. PR2 보간은 모델이 채점되는 모든 곳에서 예측하고 "
         "미래 이웃이 없으면 persistence로 후퇴하며 그 비율을 보고한다. PR3 test 셋 최소 규모. "
         "PR4 유의성은 shot 군집 bootstrap 95 % CI가 0을 제외하는 것으로 정의한다."),
     (0, "확정 프로토콜은 여기에 유지값 제거 학습·채점, W=2, 윈도 계열의 파일당 상한 500, "
         "두 공동 1차 모집단, 그리고 TEST 채점 이전에 확정한 모델 결정 규칙을 더한다."),
     (0, "표준 test 셋은 컷 모집단에서 96개 shot의 관측 T_i 32,589행과 60개 shot의 진짜 "
         "V_rot 10,463행을 담으며, 4개 분할(시드 42·1·7·123)이 각각 자기 test 셋을 갖는다."),
     (0, "한 방전 내 인접 CES 행은 강하게 상관되므로, 짝지은 오차를 shot 단위로 집계하고 shot "
         "전체를 복원추출로 10,000회 재표본한다. 대가는 유효 표본 크기가 shot 수(T_i 약 96, "
         "V_rot 60~66)가 된다는 것이며, 이것이 전체 검정력을 제한한다."),
     (0, "채점된 V_rot 행의 40~44 %에서 PCHIP은 persistence로 후퇴하므로, V_rot의 'PCHIP 대비'는 "
         "5분의 2가 'persistence 대비'다."),
     (0, "결측이 MNAR이므로 관측 시점 skill은 낙관적 추정이며, 이를 재가중으로 따로 검정한다.")],
    note_txt="sec:eval 사전등록·bootstrap·MNAR 문단. 프로토콜 문서는 experiments/PREREGISTRATION_W2.md.",
)

# =============================================================================
# 7. 결과
# =============================================================================
divider("Results", "sec:results")

table_slide(
    "결과 ①: 나우캐스터는 미래를 읽지 않는 모든 방법 가운데 최저 RMSE를 갖는다",
    ["팔", "정보 접근", "T_i RMSE (eV)", "V_rot RMSE (km/s)"],
    [["seq_v2 (나우캐스터)", "고속 진단 + 과거 CES, 세그먼트 전체", "157.8", "23.6"],
     ["윈도 대조군 (W=2)", "고속 진단 + 과거 CES 2행", "169.2", "26.1"],
     ["인과 GP", "과거 CES 이웃 16개", "164.3", "28.8"],
     ["Persistence", "마지막 관측 CES", "197.2", "33.4"],
     ["AR (국소, 인과)", "과거 CES만", "472.2", "51.0"],
     ["GP (오프라인)", "과거 + 미래 CES", "153.8", "24.7"],
     ["선형 보간", "과거 + 미래 CES", "169.8", "29.0"],
     ["PCHIP 보간", "과거 + 미래 CES", "173.6", "30.2"]],
    widths=[2.7, 5.0, 2.2, 2.2],
    lead="표준 test 분할(시드 42), 컷 모집단, 진짜 측정만(T_i n = 32,589, V_rot n = 10,463). "
         "아래 세 행은 미래를 읽는다.",
    tail="유일한 동률 팔은 잡음 앵커를 평균해 내는 오프라인 GP다(153.8 대 157.8). 미래 CES가 "
         "정의상 없는 모든 온라인 환경에서 시퀀스 나우캐스터는 명백한 승자다.",
    size=11.5,
    note_txt="표 tab:ladder, 그림 fig:ladder. 포함 모집단은 순서를 바꾸지 않는다"
             "(seq_v2 363.0 / 23.7, PCHIP 412.4 / 30.2, 인과 GP 394.6 / 28.8).",
)

table_slide(
    "결과 ②: T_i는 두 모집단 모두에서 4개 독립 분할 전부 미래를 쓰는 보간을 이긴다",
    ["타겟", "분할", "컷: skill (95 % CI)", "포함: skill (95 % CI)", "인과 GP 대비 (컷/포함)"],
    [["T_i", "42", "+0.174 [+0.097, +0.236]", "+0.225 [+0.109, +0.293]", "+0.078 / +0.154"],
     ["T_i", "1", "+0.248 [+0.188, +0.295]", "+0.238 [+0.153, +0.302]", "+0.133 / +0.169"],
     ["T_i", "7", "+0.257 [+0.199, +0.302]", "+0.292 [+0.232, +0.344]", "+0.138 / +0.123"],
     ["T_i", "123", "+0.264 [+0.188, +0.320]", "+0.316 [+0.186, +0.392]", "+0.105 / +0.149"],
     ["V_rot", "42", "+0.390 [+0.077, +0.591]", "+0.384 [+0.066, +0.586]", "+0.331 / +0.324"],
     ["V_rot", "1", "+0.183 [-0.028, +0.280]", "+0.195 [+0.000, +0.286]", "+0.130 / +0.143"],
     ["V_rot", "7", "+0.135 [-0.358, +0.269]", "+0.132 [-0.362, +0.266]", "+0.020 / +0.018"],
     ["V_rot", "123", "+0.305 [-0.049, +0.437]", "+0.304 [-0.065, +0.433]", "+0.134 / +0.132"]],
    widths=[1.0, 0.9, 3.4, 3.4, 3.4],
    lead="held-out test skill과 shot 군집 95 % CI(10,000회 재표본). 시드 1·7·123은 어떤 선택 "
         "단계에도 쓰이지 않았다.",
    tail="T_i는 8개 셀 전부 PASS이며 인과 GP와 persistence도 이긴다. V_rot는 점추정이 모두 "
         "양수지만 PR4 통과가 1/4(컷)·2/4(포함)이므로 승리가 아니라 동률로 보고한다.",
    size=10.5,
    note_txt="표 tab:headline, sec:headline. 오프라인 GP에 대해서는 8개 중 1개 셀만 유의하므로 동률이다.",
)

figure_slide(
    "forest plot: T_i는 8개 셀 전부에서, V_rot는 3개 셀에서만 CI가 0을 제외한다",
    ["포함 모집단의 수치가 더 높은 이유는 스파이크가 보간 앵커를 오염시키기 때문이다. "
         "학습된 모델은 스파이크가 낀 이월값을 할인할 수 있지만 보간은 그럴 수 없다.",
     "그래서 본 논문은 더 좋아 보이는 한쪽이 아니라 두 모집단을 함께 보고한다."],
    os.path.join(PAPER_FIG, "fig_forest.png"),
    note_txt="그림 fig:forest.",
)

content(
    "결과 ③: 전체격자 프레이밍의 채택은 네 조건 관문을 먼저 고정하고 한 번만 채점하였다",
    [(0, "분할 시드 × 초기화 시드의 4×4 격자를 돌리고 각 실행을 해당 분할의 W=2 윈도 대조군과 "
         "행 단위로 짝지었다. 짝지은 T_i skill은 16/16 실행에서 양수였고 13/16에서 개별적으로 "
         "유의했다."),
     (0, "통합 평균은 +0.081이며 실행 군집 95 % CI는 [+0.067, +0.096]이다. 초기화 산포가 분할 "
         "산포보다 훨씬 작다."),
     (0, "예산을 균등화한 팔도 4/4 분할에서 부호를 유지하므로 이 이득은 학습 예산 효과가 아니라 "
         "아키텍처에서 온다. 어떤 실행도 유의한 V_rot 손실을 보이지 않았다."),
     (0, "관문은 프레이밍이 무엇을 사는지도 확정한다. 윈도 대조군은 인과 GP와 동률인 반면 시퀀스 "
         "백본은 두 모집단 모두에서 4/4로 이긴다. 세그먼트의 과거 전체로 뻗는 도달 범위가 배치 "
         "가능한 가장 강한 기준선을 넘어서게 하는 요소다."),
     (0, "그 비용은 음수다. 샘플별 윈도 조립도 조합적 증강도 없으므로 백본은 윈도 계열의 학습 "
         "시간 중 일부만으로 학습된다.")],
    note_txt="sec:gate. 이후 4개 확정 분할에서 같은 비교는 8/8 양수이고 각 모집단에서 2/4 유의다.",
)

table_slide(
    "결과 ④: skill은 간극 영역에서도 살아 있으며, V_rot는 바로 그곳에서만 이긴다",
    ["Δt", "n (컷 / 포함)", "shot 수", "컷: PCHIP 대비", "포함: PCHIP 대비"],
    [["T_i,  ≤ 15 ms", "134,546 / 135,317", "301", "+0.239 [+0.197, +0.274]", "+0.299 [+0.244, +0.347]"],
     ["T_i,  > 15 ms", "3,422 / 3,334", "265 / 263", "+0.268 [+0.187, +0.337]", "+0.206 [+0.108, +0.290]"],
     ["T_i,  > 45 ms", "460 / 429", "104 / 101", "+0.267 [+0.092, +0.414]", "-0.004 [-0.304, +0.246]"],
     ["V_rot, ≤ 15 ms", "51,689", "197", "+0.233 [+0.020, +0.318]", "+0.233 [+0.020, +0.317]"],
     ["V_rot, > 15 ms", "466 / 456", "130", "+0.418 [+0.104, +0.680]", "+0.432 [+0.128, +0.696]"],
     ["V_rot, > 45 ms", "14", "7", "채점하지 않음", "채점하지 않음"]],
    widths=[2.2, 2.6, 1.6, 3.3, 3.3],
    lead="마지막 관측으로부터의 경과시간으로 계층화하고 4개 test 분할을 통합하였다. Δt가 커질수록 "
         "PCHIP의 과제는 쉬워지고 우리 과제는 어려워진다.",
    tail="Δt > 15 ms에서 V_rot가 PCHIP을 이기는 것이 본 논문에서 유일한 무조건적 V_rot 긍정 "
         "판정이다. T_i는 45 ms를 넘어서면 포함 모집단에서 동률이 되며, 그대로 동률로 보고한다.",
    size=10.5,
    note_txt="표 tab:gap, sec:gap.",
)

content(
    "결과 ⑤: 진짜 결측 지점으로 재가중하면 두 비교가 갈라진다",
    [(0, "참값이 없는 곳에서 오차를 잴 수는 없지만, 관측 행과 결측 행 모두에서 계산 가능한 "
         "공변량(경과시간 구간 × 입력 전용 국소 활동 플래그) 위에서 사후 계층화하고 재가중할 "
         "수는 있다."),
     (0, "먼저 적용 범위. 이 분석의 도메인은 'W=2 안에 진짜 관측이 있음'이며, 진짜 결측인 T_i "
         "행의 54~68 %가 도메인 안이지만 V_rot는 4~6 %뿐이다. 따라서 재가중된 V_rot에서는 "
         "아무것도 끌어내지 않는다."),
     (0, "온라인 시스템이 실제로 돌리는 기준선인 persistence를 상대로는 T_i 우위가 두 모집단 "
         "모두에서 4개 분할 전부 재가중을 견딘다(+0.28~+0.44, 모든 CI가 0을 제외)."),
     (0, "미래를 쓰는 PCHIP을 상대로는 재가중된 점추정이 +0.14~+0.28에 머물지만, 넓은 CI가 컷 "
         "모집단 2개 분할에서 0을 지난다(컷 2/4, 포함 4/4)."),
     (0, "결측 지점은 더 큰 Δt에 놓이는데, 그곳이 바로 양쪽 앵커가 가장 크게 도움이 되고 재가중 "
         "bootstrap이 가장 얇은 곳이다."),
     (0, "한 계층 안에서 결측 행과 관측 행이 교환 가능하다는 가정은 없앤 척하지 않고 명시한다.")],
    note_txt="표 tab:mnar, sec:mnar.",
)

table_slide(
    "결과 ⑥: 캠페인을 가로지르면 윈도 모델의 오프라인 우위는 붕괴하고 시퀀스는 살아남는다",
    ["모집단", "팔", "T_i PCHIP 대비 (초기화 42/1/7/123)", "PASS", "인과 GP 대비"],
    [["컷", "윈도 대조군", "+0.027 / +0.091 / -0.001 / +0.061", "2/4", "0/4"],
     ["컷", "대조군 + shot별 표준화", "+0.103 / +0.107 / +0.094 / +0.107", "4/4", "-"],
     ["컷", "seq_v2", "+0.187 / +0.174 / +0.181 / +0.177", "4/4", "4/4 (+0.11~+0.12)"],
     ["포함", "윈도 대조군", "+0.014 / +0.047 / +0.055 / +0.089", "0/4", "0/4"],
     ["포함", "seq_v2", "+0.173 / +0.202 / +0.198 / +0.184", "4/4", "4/4 (+0.13~+0.16)"]],
    widths=[1.3, 3.2, 4.6, 1.1, 2.6],
    lead="641개 방전을 shot 번호로 정렬해 시간으로 엄격히 잘랐다(train 416 / val 128 / test 97). "
         "바뀐 변수는 분할 규칙 하나뿐이며, 네 실행은 분할이 아니라 초기화만 다르다.",
    tail="이유는 단언이 아니라 측정이다. 고속 진단은 기간 사이에 BES 1.22 σ, ECEI 0.53 σ "
         "이동하는데 타겟 자신은 0.115 σ에 그쳐, 예측 대상보다 5~11배 더 이동한다. 각 방전을 "
         "자기 통계로 표준화하면 누수 없이 그것이 수리된다.",
    size=11.0,
    note_txt="표 tab:campaign, 그림 fig:campaign, sec:campaign. 남는 단서: 하나의 시간 test 블록이며 "
             "컷 실행 4개 중 2개는 30-epoch 상한에서 멈췄다.",
)

figure_slide(
    "같은 그림: 대조군은 무너지고, shot별 표준화가 복구하며, 시퀀스는 흔들리지 않는다",
    ["시퀀스 나우캐스터는 두 모집단 모두에서 4/4 초기화에 대해 PCHIP과 인과 GP를 앞서고 "
         "8/8에서 윈도 대조군을 이긴다.",
     "V_rot도 대조군이 0/4인 자리에서 두 모집단 모두 4/4로 persistence를 이긴다."],
    os.path.join(PAPER_FIG, "fig_campaign.png"),
    note_txt="그림 fig:campaign.",
)

content(
    "결과 ⑦: T_i와 V_rot의 정보 비대칭은 물리에서 예측되었고 절제로 확인되었다",
    [(0, "평가 시점에 한 모달리티 그룹을 0으로 만들고 재학습은 하지 않는 통제된 절제를 윈도 "
         "대조군에 적용하였다."),
     (0, "이력은 두 타겟 모두에 필수불가결하다. CES 이력이 없으면 두 모집단의 두 타겟 모두에서 "
         "PCHIP 대비 -1에서 -4까지 떨어진다. 자기 앵커 없이 100 Hz 진단만으로 예측 가능한 타겟은 "
         "없다."),
     (0, "T_i의 보간 대비 마진은 고속 진단 정보다. 컷 모집단에서 BES/ECEI/MC를 0으로 만들면 "
         "대조군이 보간기 아래로 떨어진다(짝지은 변화 -0.25~-0.43, 4/4 유의)."),
     (1, "포함 모집단에서는 같은 이력 전용 모델이 여전히 PCHIP을 +0.15~+0.23으로 이기며 고속 "
         "채널이 더하는 것은 0.03~0.09뿐이다. 즉 포함 마진에는 스파이크 강건성 성분이 섞여 "
         "있고, 고속 진단의 기여를 분리해 내는 쪽은 컷 모집단이다."),
     (0, "V_rot의 정보는 전적으로 CES 이력이다. 고속 진단을 0으로 만들어도 출력이 8/8 셀에서 "
         "비트 단위로 동일하다. 라우팅이 구조적이기 때문이다."),
     (0, "토로이달 회전은 우리 입력 중 어느 것도 관측하지 않는 운동량 원천이 지배하며, Mirnov는 "
         "100 Hz 샘플링으로 kHz 모드 회전을 앨리어싱으로 잃는다.")],
    note_txt="표 tab:ablation, sec:asym.",
)

figure_slide(
    "절제 결과: 고속 진단을 0으로 만들면 V_rot는 비트 단위로 동일하다",
    ["이는 실패가 아니라 진단 정보량에 관한 발견이며, 모델의 타겟별 라우팅을 정당화한다."],
    os.path.join(PAPER_FIG, "fig_ablation.png"),
    note_txt="그림 fig:ablation.",
)

table_slide(
    "결과 ⑧: 이력은 필수적이며, 그 첫 관측 하나가 사실상 전부를 담는다",
    ["이력 관측 수", "W", "T_i skill", "PASS", "V_rot skill", "PASS"],
    [["0", "4", "-0.026", "0/4", "-0.783", "0/4"],
     ["1", "2", "+0.238", "4/4", "+0.206", "0/4"],
     ["2", "3", "+0.246", "4/4", "+0.203", "1/4"],
     ["3", "4 (구 기본값)", "+0.221", "3/4", "+0.190", "1/4"],
     ["5", "6", "+0.190", "3/4", "+0.205", "1/4"],
     ["7", "8", "+0.216", "4/4", "+0.204", "2/4"]],
    widths=[2.2, 2.4, 2.0, 1.4, 2.2, 1.4],
    lead="24회의 독립 실행(W 다섯 지점 × 4시드 + history-0 × 4시드). 검토되지 않은 기본값 W=4를 "
         "곡선과 명시된 선택 규칙으로 대체하였다.",
    tail="한 지점 안에서의 시드 산포(0.07~0.16)가 곡선 전체보다 넓다. 'plateau에 도달하는 가장 "
         "작은 W'를 적용하면 W=2가 나온다. 더 넓게 갈 방어 가능한 이유는 커버리지 하나뿐이며, "
         "그것은 도달 범위가 세그먼트 전체인 시퀀스 프레이밍의 근거다.",
    size=11.5,
    note_txt="표 tab:window, sec:window. 이 스윕은 컷보다 앞서므로 포함 모집단 기준이다.",
)

table_slide(
    "결과 ⑨: 복잡도 사다리와 폭 스윕이 모델 크기 축을 닫는다",
    ["팔", "파라미터", "컷", "포함"],
    [["Persistence", "0", "-0.264", "-0.288"],
     ["앵커+Δ (명명된 항)", "1,258", "-0.261", "-0.287"],
     ["b3k8 (잠재 8개 유계 보정)", "21,498", "+0.237", "+0.126"],
     ["윈도 대조군", "201,258", "+0.173", "+0.238"],
     ["seq_v2 백본", "357,570", "+0.236", "+0.268"]],
    widths=[4.4, 2.4, 2.6, 2.6],
    lead="held-out TEST T_i skill(4개 분할 평균). b3k8은 예측을 이월값 + 여덟 개의 유계 보정항으로 "
         "정확히 분해하며, 결정 규칙은 TEST 채점 이전에 고정되었다.",
    tail="컷에서는 백본의 T_i skill 전부가 유계 수 여덟 개로 압축된다(짝지은 평균 +0.002, 모든 CI가 "
         "0을 포함). 컷이 없으면 -0.19를 잃는데, 유계 보정으로는 스파이크가 낀 앵커를 복구할 수 없기 "
         "때문이다. 백본을 34k에서 879k로 넓혀도 skill은 ±0.01 미만으로 움직인다.",
    size=11.5,
    note_txt="표 tab:ladder2, 그림 fig:ladder_scaling, sec:ladder. 선형 probe 결과도 같은 절에 있다.",
)

figure_slide(
    "결과 ⑩: 우위는 보간이 가장 약한 고변동 국소 구간에 집중된다",
    ["T_i의 peak 계층 skill은 컷 +0.45~+0.61, 포함 +0.62~+0.72이며 8/8 셀에서 PASS다. "
         "bulk에서는 컷 +0.09~+0.20이다. 매끄러운 bulk에서는 보간이 거의 최적이며, 모델의 가치는 "
         "활동 구간에 있다.",
     "V_rot의 peak 계층 skill은 +0.54~+0.79로 8/8에서 점추정이 양수이고 bulk에서는 약 0이다. "
         "즉 비대칭은 지역적이며, 이를 전역 영가설의 뒤집힘이 아니라 유의성 기반이 얇은 지역적 "
         "강점으로 보고한다."],
    os.path.join(PAPER_FIG, "fig_peak.png"),
    note_txt="그림 fig:peak, sec:peak. 계층 분류는 타겟 행을 제외한 입력 측 대리 지표로만 한다.",
)

table_slide(
    "결과 ⑪: 문맥은 승리를 전형적으로 만들고, 계열은 skill이 아니라 값을 정한다",
    ["문맥", "T_i skill vs 인과 GP [95 % CI]", "방전 단위 승률", "상위 10개 방전 제거 후"],
    [["20 ms", "+0.057 [+0.027, +0.085]", "0.52", "+0.028"],
     ["30 ms", "+0.087 [+0.061, +0.111]", "0.60", "+0.060"],
     ["50 ms", "+0.104 [+0.079, +0.128]", "0.64", "+0.077"],
     ["70 ms", "+0.119 [+0.095, +0.142]", "0.66", "+0.092"],
     ["150 ms", "+0.123 [+0.096, +0.148]", "0.66", "+0.096"],
     ["630 ms", "+0.143 [+0.118, +0.168]", "0.67", "+0.116"]],
    widths=[1.8, 4.4, 2.6, 3.2],
    lead="301개 방전 통합 재채점(순환 계열). 도달 범위는 절단이 아니라 각 문맥에서 다시 학습한 "
         "모델로 잰다. 사전등록 규칙이 반환하는 포화점은 5스텝, 즉 50 ms다.",
    tail="승률이 이 절의 결론이다. 20 ms에서 모델은 평균으로는 이기지만 방전의 절반에서만 이기고, "
         "70 ms부터 3분의 2에서 이기며 평평해진다. 포화점을 '유용해지기 위해 필요한 문맥'으로 "
         "기술해서는 안 된다. 20 ms에서도 persistence를 +0.356으로 앞선다.",
    size=11.5,
    note_txt="표 tab:pooled, sec:context. 4/4 계수는 한 스텝 간격에서 단조가 아니므로 판정은 결손으로 한다.",
)

table_slide(
    "결과 ⑫: 같은 도달 범위에서 계열은 구별되지 않으므로 아키텍처는 비용으로 정한다",
    ["계열", "R에 대한 비용", "R = 15 (ops)", "R = 63 (ops)", "630 ms에서 skill"],
    [["순환 (LSTM, 357,570)", "O(1) - 상태가 나른다", "111", "111", "+0.138"],
     ["확장 인과 합성곱 (3,238)", "O(log R), 층당 +48", "209", "305", "+0.124"],
     ["띠 attention (295,746)", "O(1), 상수 4.3배", "473", "473", "+0.122"],
     ["대각 상태공간 (SSM)", "O(1)", "-", "-", "+0.105 (승률 0.58~0.60)"]],
    widths=[3.6, 3.4, 1.8, 1.8, 3.0],
    lead="도달 범위를 맞추면 110배의 파라미터 범위와 서로 가장 다른 세 연산자가 구별되지 않는 결과를 "
         "낸다. 계열이 T_i를 움직이는 폭은 최대 0.023이고, 도달 범위를 20에서 70 ms로 옮기는 효과는 "
         "+0.060이다.",
    tail="온라인 1스텝 지연은 파라미터 수가 아니라 디스패치되는 연산자 개수를 따라간다(약 2~3 μs/연산자). "
         "이 문제에서 attention은 엄밀히 지배당한다. 10 ms 예산은 어느 팔에도 구속 조건이 아니며, "
         "1 ms 판정은 세션 간 산포 때문에 사전등록 규칙에 따라 보류되었다.",
    size=11.0,
    note_txt="표 tab:family·tab:cost, sec:context. V_rot는 네 계열 모두에서 문맥이 길어질수록 나빠진다.",
)

table_slide(
    "결과 ⑬: 두 개의 스트레스 테스트가 오프라인 주장과 배치 주장을 분리한다",
    ["평가", "PCHIP 대비 (오프라인, 미래 사용)", "인과 기준선 대비"],
    [["무작위 분할, 관측 지점", "4/4 · 4/4, +0.17~+0.32", "인과 GP 대비 4/4 · 4/4, +0.08~+0.17"],
     ["결측 지점으로 재가중", "2/4 · 4/4, +0.14~+0.28", "persistence 대비 4/4 · 4/4, +0.28~+0.44"],
     ["시간 캠페인 분할", "4/4 · 4/4, +0.17~+0.20", "인과 GP 대비 4/4 · 4/4, +0.11~+0.16"]],
    widths=[3.4, 4.6, 5.0],
    lead="seq_v2, T_i. 항목은 PR4 통과 수(컷 · 포함)와 점추정 범위다.",
    tail="실시간 시스템이 실제로 돌릴 수 있는 기준선에 대한 우위는 두 모집단 모두에서 두 테스트를 모두 "
         "견딘다. 이 표의 윈도 모델 판본에서는 오프라인 우위가 두 테스트 모두를 견디지 못했다. 그것을 "
         "바꾼 것이 시퀀스 프레이밍이며, 새 열이 원래부터 거기 있었던 것처럼 제시하지 않는다.",
    size=12.0,
    note_txt="표 tab:stress.",
)

# =============================================================================
# 8. 모델 선택 · 배치 · 개선 여지 · 한계
# =============================================================================
divider("Discussion", "sec:selection / sec:deploy / sec:headroom / sec:limits")

content(
    "모든 모델 결정은 그것이 결정할 수치보다 먼저 적혔다",
    [(0, "윈도 계열의 아키텍처는 통제된 실험의 연속으로 도달하였다. 각 실험은 데이터 계약을 "
         "보존하는 하나의 변경이었고, 깨끗한 비증강 검증 skill로만 채점되었다."),
     (1, "증강된 검증 손실은 보간이 이미 강한 바로 그곳에서 평활화를 보상하므로 결코 쓰지 "
         "않았다."),
     (0, "시퀀스 프레이밍은 네 조건 관문이 먼저 고정되고 그 다음 충족된 뒤에야 채택되었다."),
     (0, "같은 규칙 아래에서 이후 탐색된 단 하나의 아키텍처 후보(관측 마스킹 인과 어텐션 추가)는 "
         "4/4 분할에서 양수였으나 사전 확정 기준 3/4에 대해 1/4에서만 유의했으므로 승격되지 "
         "않았다."),
     (1, "이 후보의 검증 이득은 탐색 분할에서 2/2 유의했으며, 이는 선택 분할 결과의 통상적 "
         "낙관이자 승격 기준을 TEST에 두는 이유다."),
     (0, "사다리 칸과 폭 스윕도 두 갈래 판정을 어떤 TEST 채점 이전에 적어 두었고, 구성상 이 스윕 "
         "위에서 백본을 재선택하는 것은 허용되지 않았다.")],
    note_txt="sec:selection.",
)

content(
    "배치 가능성은 지연시간과 불확실성 두 가지로 판정한다",
    [(0, "순환 나우캐스터는 온라인에서 상태를 유지한 채 실행되며, 새 행마다 배치 1의 순환 스텝 "
         "하나가 든다. 노트북급 CPU에서 중앙값 1.05 ms, p99 1.61 ms로 10 ms 격자의 16 %다."),
     (1, "이 크기에서 GPU는 배치 1에서 아무것도 사주지 않는다(1.21 / 2.31 ms). 실무 지침은 상태 "
         "유지형 나우캐스터를 제어 계산기의 CPU에서 돌리는 것이며, 예산의 80 % 이상이 획득과 "
         "제어에 남는다."),
     (0, "불확실성은 재학습 없이 다룬다. split conformal을 해당 실행 자신의 검증 분할에서 보정하고 "
         "예측기는 아무것도 바꾸지 않으며, 같은 절차를 PCHIP과 persistence에도 적용한다."),
     (0, "모델의 구간이 32개 셀 전부에서 두 기준선의 구간을 이긴다. 폭과 빗나감을 함께 벌하는 "
         "Winkler 점수로 컷 T_i는 1,272 대 1,554(PCHIP) 대 1,727(persistence)이다."),
     (1, "포함 모집단에서 모델의 T_i 구간은 PCHIP보다 오히려 넓은데도 점수가 좋다. 스파이크가 "
         "빗나감 벌점을 부풀리는데 모델이 덜 빗나가기 때문이다."),
     (0, "정직한 실패: 커버리지는 주변적이지 조건부가 아니다. shot 수준의 이동이 교환가능성 "
         "가정을 깨며, 이를 고치려면 현재 shot 수로는 지탱할 수 없는 shot 조건부 보정이 필요하다.")],
    note_txt="sec:deploy. 지연시간은 전원 상태에 따라 실행 간 최대 2배 변동한다.",
)

content(
    "남은 개선 여지는 모델이 아니라 데이터에 있으며, 세 레버는 이름이 지목되어 있다",
    [(0, "용량은 배제되었고 더 긴 윈도도 배제되었으며 도달 범위를 사는 프레이밍은 채택되었다. "
         "남은 것은 데이터이고, 각 레버는 이미 이 데이터셋 안의 증거로 지목된다."),
     (0, "① CES 피팅 품질 메타데이터. 표본별 피팅 χ²이나 신호 수준이 있으면 품질 컷이 값 컷을 "
         "모든 팔에서 대체하고 두 모집단은 하나로 합쳐진다. V_rot 자신의 스파이크도 같은 규칙으로 "
         "처리된다."),
     (0, "② 원시 kHz Mirnov 스트림. 같은 격자에서 lag-1 자기상관은 BES +0.568, ECEI +0.572인 반면 "
         "Mirnov는 -0.009로, 이 격자 위에서 자기 채널은 백색잡음이다. 정보는 플라즈마에 없었던 것이 "
         "아니라 100 Hz 데시메이션이 모델보다 상류에서 파괴한 것이다."),
     (1, "해결책은 모델 변경이 아니라 전처리 변경이다. 원시 시계열에서 계산한 윈도별 RMS·대역 "
         "파워·모드 수·모드 회전 주파수를 V_rot 분기로 라우팅하는 것이 V_rot에 대해 지목할 수 있는 "
         "가장 가치 높은 실험이다."),
     (0, "③ 액추에이터 채널. shot 사이에서 ECE 유래 T_e 대리 지표는 T_i와 r = +0.353(p = 3e-17)로 "
         "상관되지만 V_rot와는 r = +0.024(p = 0.58)이다. 파워는 토크가 아니며, NBI 토크를 더하는 일은 "
         "모델링 과제가 아니라 데이터 획득 과제다."),
     (0, "셋 모두 아카이브된 KSTAR 데이터에서 실행 가능하며, 어느 것도 현재 결과가 천장에 있다는 "
         "진술이 아니다.")],
    note_txt="sec:headroom. 문헌의 양성 대조: 빔 액추에이터를 입력으로 받는 DIII-D 전방전 시뮬레이터.",
)

content(
    "한계는 결과와 같은 자리에 적는다",
    [(0, "통계적 검정력. 재현의 단위는 shot이며 분할당 test shot이 96개(T_i) / 60~66개(V_rot)이므로 "
         "검정력이 모든 유의성 판정을 제약하고, shot별 제곱오차 차이는 두꺼운 꼬리를 갖는다."),
     (0, "MNAR 낙관. skill은 관측 지점에서만 측정되며, 재가중은 결측 T_i 행의 54~68 %와 결측 V_rot "
         "행의 4~6 %에만 닿는다."),
     (0, "오프라인 주장의 상한은 GP 동률이다. '미래를 쓰는 보간을 이긴다'는 '사전등록된 보간들을 "
         "이긴다'를 뜻하며 모든 오프라인 방법으로 확장되지 않는다. 다만 인과 GP로는 확장된다."),
     (0, "값 기준 컷은 대리 지표이고, 캠페인 전이는 4개 분할이 아니라 하나의 시간 블록 위 4개 "
         "초기화에 기댄다. 대조군을 전이시키는 shot별 표준화는 오프라인 형태이며 인과적 러닝 판본은 "
         "여기서 측정되지 않았다."),
     (0, "지표의 비대칭. 보간은 shot 전체의 이웃을 쓰고 백본은 세그먼트의 과거를 본다. 이는 인과 "
         "모델에 의도적으로 불리하지만 직접적인 해석은 어렵게 만든다."),
     (0, "단일 장치·단일 진단 집합이며, 불확실성은 조건부가 아니라 주변적으로 보정된다. 페데스탈 "
         "상단 프레이밍은 데이터 선정에서 물려받았고 사건 위상 분석은 향후 과제다.")],
    note_txt="sec:limits.",
)

# =============================================================================
# 9. 결론
# =============================================================================
divider("Conclusion", "sec:conclusion")

content(
    "결론",
    [(0, "인과적이고 누수 방지된 멀티모달 나우캐스터, 즉 V_rot 분기가 고속 진단을 결코 보지 않는 "
         "전체격자 2분기 순환 모델을 만들고, 사전등록된 shot 군집 두-모집단 프로토콜 아래에서 "
         "의도적으로 불리한 기준에 대해 평가하였다."),
     (0, "이온온도에 대해 모델은 미래를 쓰는 PCHIP을 +0.17~+0.32의 skill로 이기고, 4개 독립 분할과 "
         "두 모집단 전부에서 유의하며, 가장 강한 오프라인 평활기와는 동률이고, 두 타겟 모두에서 모든 "
         "인과 기준선을 이긴다."),
     (0, "배치 주장은 두 스트레스 테스트의 지지를 받는다. 진짜 결측이면서 도메인 안인 시점에서 모든 "
         "인과 방법 대비 우위는 8/8 셀에서 재가중을 견디고, 캠페인 경계를 가로질러서도 시퀀스 "
         "나우캐스터는 두 모집단 모두에서 4/4 초기화로 PCHIP과 인과 GP를 앞선다."),
     (0, "작동하지 않는 곳과 그 이유도 함께 보고한다. 회전은 전역적으로 보간과 동률이고 그 skill은 "
         "과도 구간에 산다. 값 기준 컷은 한쪽 방향 대리 지표이며, 예측 구간은 방전별이 아니라 "
         "주변적으로 보정된다."),
     (0, "이 중 어느 것도 천장 논변이 아니다. 각각은 구체적이고 검정 가능한 변경을 가리킨다: CES 피팅 "
         "품질 메타데이터, 원시 kHz Mirnov 스트림에서 다시 계산한 윈도별 특징, 그리고 회전의 물리적 "
         "원인인 NBI 토크 채널."),
     (0, "모델 크기 축은 닫혔다. 21k 파라미터의 persistence+유계 보정 모델이 컷에서 백본과 대등하고 "
         "26배 폭 스윕은 평평하다. 총합 손실의 한계는 추정기가 아니라 정보다.")],
    note_txt="sec:conclusion.",
)

content(
    "재현성과 가용성",
    [(0, "분할은 디스크에 고정되고 재적재 시 검증된다. 정규화 통계는 학습 파일 전용이며, 시퀀스 계열은 "
         "여기에 shot별 입력 표준화가 더해진다."),
     (0, "모델 선택은 test 데이터를 결코 읽지 않는다. bootstrap은 고정 시드와 10,000회 재표본을 쓰고, "
         "인용된 모든 수치는 단일 수집 스크립트가 동결된 평가 산출물에서 읽어온다."),
     (0, "전체 파이프라인, 사전등록 문서, 통제 실험 러너, 고정된 분할, 그리고 수집기가 공개 저장소에 "
         "있다: github.com/iseungsang01/ml-intern-revision"),
     (1, "백본은 ces_prediction/experiments/seq/model_seq_v2.py, 해석 가능한 사다리 칸은 model_seq_b3.py, "
         "윈도 대조군은 ces_prediction/model_iter009.py이며 테스트 스위트가 내용 해시로 보호한다."),
     (0, "KSTAR 진단 데이터 641개 방전은 코드와 함께 재배포되지 않으며 운영 기관의 데이터 정책을 따른다. "
         "저장소에는 shot 파일이 갖춰진 뒤 분석을 재현하는 데 필요한 모든 것이 들어 있다.")],
    note_txt="sec:conclusion 뒤의 재현성·코드 가용성·데이터 가용성 소절.",
)

# =============================================================================
# 부록 A. 시도한 모델과 문헌 조사
# 표의 내용은 appendix_content.py 한 곳에 있으며 1시간 덱이 같은 것을 읽는다.
# =============================================================================
divider("Appendix", "부록 A. 시도한 모델의 계보와 닫은 이유, 그리고 2026-09-05 문헌 조사.")

table_slide(
    "시도한 모델의 계보와 닫은 이유 (1/3): 윈도 계열",
    AP.TRIED_HEAD, AP.TRIED_WINDOW, widths=[2.7, 1.1, 3.1, 2.8, 2.4],
    lead="아래는 W = 4 시대의 잠정 수치이며 확정 프로토콜의 주장에는 쓰지 않는다.",
    note_txt="THESIS_RESULTS.md 8e / 8b.2 / 8f / 8k / 8u / 8x / 8ad.",
)

table_slide(
    "시도한 모델의 계보와 닫은 이유 (2/3): 시퀀스 계열",
    AP.TRIED_HEAD, AP.TRIED_SEQ, widths=[2.7, 1.1, 3.1, 2.8, 2.4],
    lead="확정 프로토콜(W = 2 · held-free · 두 공동 1차 모집단) 아래에서 짝지어 채점하였다.",
    note_txt="8d / 8t / 8x / 8y / 8z / 8aa / 8ab / 8ai.",
)

for _part, _rows, _tail in ((1, AP.TRIED_MISC[:5], None),
                            (2, AP.TRIED_MISC[5:], AP.TRIED_TAKEAWAY)):
    table_slide(
        "시도한 모델의 계보와 닫은 이유 (3/3, %d쪽): 계열 · 문맥 · 기준선 · 확장 가지" % _part,
        AP.TRIED_HEAD, _rows, widths=[2.7, 1.1, 3.1, 2.8, 2.4],
        tail=_tail,
        note_txt="8ag / 8ai / 8ak / 8af / 8al / 8am 부록 / 8p / 8m / 8ap · 8at, 그리고 B.11.",
    )

_FUSION_W = [1.95, 1.05, 0.85, 1.85, 3.05, 1.55, 3.55]
def _chunk(rows, per=5):
    n = len(rows)
    parts = max(1, -(-n // per))
    size = -(-n // parts)
    return [rows[i:i + size] for i in range(0, n, size)]


_FUSION_PARTS = _chunk(AP.FUSION_ROWS, 5)
for _part, _rows in enumerate(_FUSION_PARTS, start=1):
    table_slide(
        "문헌 조사 (1/4, %d/%d쪽): 핵융합의 진단-대-진단 추정은 여전히 단순한 구조가 주류이다"
        % (_part, len(_FUSION_PARTS)),
        AP.FUSION_HEAD, _rows, widths=_FUSION_W,
        lead="장치와 연도를 따로 두었다.",
        note_txt=AP.sources_note("2026-09-05 조사한 핵융합 12편의 요약이다."),
    )

_GENERAL_W = [1.70, 0.95, 2.60, 3.95, 3.35]
for _part, _rows in ((1, AP.GENERAL_ROWS[:5]), (2, AP.GENERAL_ROWS[5:])):
    table_slide(
        "문헌 조사 (2/4, %d쪽): 일반 시계열 · 센서 예측의 주류와 본 데이터에 대한 판정" % _part,
        AP.GENERAL_HEAD, _rows, widths=_GENERAL_W,
        note_txt="마지막 열은 본 저장소의 통제 실험 판정이며, 새 계열을 도입하기 전에 이 열을 먼저 읽는다.",
    )

table_slide(
    "문헌 조사 (3/4): 구조적으로 동형인 분야에서 반복되는 교훈",
    AP.ISO_HEAD, AP.ISO_ROWS, widths=[2.4, 0.95, 2.35, 3.2, 3.3],
    tail="동형 분야가 반복해 말하는 것은 표현력이 아니라 개체별 보정과 상태추정 프레임이며, "
         "이는 본 연구의 shot별 표준화와 전체격자 인과 프레이밍이 이미 취한 선택이다.",
    note_txt="혼합주기 나우캐스팅 · 저가 센서 보정 · 커프리스 혈압 · 구조 가상 센싱 · "
             "합성 진단 증강의 다섯 계열에서 같은 결론이 반복된다.",
)

table_slide(
    "문헌 조사 (4/4): 문헌이 지목하는 다음 팔은 표현력이 아니라 손실 · 게이팅 · 입력이다",
    AP.PRIORITY_HEAD, AP.PRIORITY_ROWS, widths=[0.7, 4.3, 2.6, 1.8, 2.8],
    lead="각 행은 통제 변수가 하나이며, TEST를 여는 팔은 사전등록 뒤에만 실행한다.",
    note_txt=AP.sources_note("우선순위는 비용 대비 정보량 순이다."),
)

content(
    "권하지 않는 방향과, 열려 있는 V_rot",
    [(0, "권하지 않는 방향은 다음 셋이다.")] +
    [(1, t) for t in AP.NOT_RECOMMENDED] +
    [(0, "회전은 닫힌 결론이 아니라 열린 과제이다.")] +
    [(1, t) for t in AP.VROT_NOTE],
    note_txt="음성 결과는 그것을 뒤집을 측정을 지목할 때만 결론으로 인정한다는 규칙을 따른다.",
)

# =============================================================================
prs.save(OUT)
print("wrote %s  (%d slides)" % (os.path.relpath(OUT, REPO), len(prs.slides._sldIdLst)))
if WARNINGS:
    print("FIT WARNINGS: %d" % len(WARNINGS))
    for w in WARNINGS:
        print("  -", w)
else:
    print("FIT WARNINGS: 0")
