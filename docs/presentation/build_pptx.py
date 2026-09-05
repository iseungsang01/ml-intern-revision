# -*- coding: utf-8 -*-
"""Build the 1-hour KSTAR CES nowcasting thesis presentation (Korean, abstract register).

Output: docs/presentation/KSTAR_CES_발표자료.pptx  (약 60분 학위논문 발표)

2026-08-27 전면 재작성. 두 가지가 바뀌었다.

(1) 문체. 승상님 지시(2026-08-27)에 따라 슬라이드 제목·본문·카드·표의 모든 문장을
    논문 초록 문체(서술형 종결, 객관·비인칭, 배경 → 방법 → 결과 → 결론)로 통일하였다.
    불릿 항목도 완결된 서술문으로 쓴다. 축약어는 첫 등장에서 풀어 쓴다(용어 정리 슬라이드).
(2) 범위. 2026-08-16 §8ab 이후의 기록 — 도달 범위(reach) 사다리와 B.9 4축(§8ac–§8aj),
    밀집 사다리·통합 재채점·승패 방전 분석(§8ak–§8an), μs shot 동결(§8ao),
    양자 가지 종결(§8ap), 논문 프레이밍 §9 — 를 7장(문맥·구조·비용)으로 추가하였다.

확정 프로토콜(THESIS_RESULTS.md §8v–§8ab): W=2 · held-free(학습·평가) · 파일당 500 ·
두 공동 1차 모집단(컷 / 포함) · 인과 GP 기준선. 주 모델은 전체격자 인과 시퀀스 나우캐스터
``seq_v2``(357,570 파라미터), 옛 주 모델(윈도 GRU + 관측마스킹 attention, 201,258)은
W=2 윈도 대조군이다. 6장까지의 수치는 docs/paper/paper_numbers.json(= main_ko.tex)에서,
7장의 수치는 THESIS_RESULTS.md §8ac–§8ap의 표에서 그대로 옮겼다.

레이아웃 헬퍼(slide/box/text/header/footer/bullets/card/add_image_fit/table/divider)와
팔레트는 build_pptx_20min.py · build_pptx_flow.py가 import하므로 시그니처를 바꾸지 않는다.

Figures are read from docs/presentation/figures/ (run make_figures.py first); the
context-family ladder is the paper figure docs/paper/figures/fig_context_family_ladder.png.

Usage (from repo root):
    py docs/presentation/build_pptx.py
    py docs/presentation/preview_pptx.py docs/presentation/KSTAR_CES_발표자료.pptx
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
FIG = os.path.join(HERE, "figures")
if os.path.abspath(HERE) not in sys.path:      # so the build-time fit audit can
    sys.path.insert(0, os.path.abspath(HERE))  # import preview_pptx's metrics

# ---- palette -------------------------------------------------------------
NAVY = RGBColor(0x13, 0x33, 0x5F)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
TEAL = RGBColor(0x1B, 0x9E, 0x8A)
ORANGE = RGBColor(0xE8, 0x74, 0x3B)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x5B, 0x66, 0x70)
LGRAY = RGBColor(0xE8, 0xEC, 0xF1)
MGRAY = RGBColor(0x9A, 0xA5, 0xB1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2B, 0x35)
CARDBG = RGBColor(0xF4, 0xF7, 0xFA)

FONT = "Malgun Gothic"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_pageno = {"n": 0}


# ---- low-level helpers ---------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _set_fill(bg, WHITE)
    bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
        shadow=False, round_=False):
    sp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, italic, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, dict):  # paragraph-level options
            opts = para
            segs = opts["segs"]
            if "align" in opts:
                p.alignment = opts["align"]
            if "space_after" in opts:
                p.space_after = Pt(opts["space_after"])
            if "space_before" in opts:
                p.space_before = Pt(opts["space_before"])
            if "level" in opts:
                p.level = opts["level"]
        else:
            segs = para
        for seg in segs:
            txt, size, color, bold, italic, font = (seg + (None,) * 6)[:6]
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color if color is not None else DARK
            r.font.bold = bool(bold)
            r.font.italic = bool(italic)
            r.font.name = font if font else FONT
    return tb


def _title_size(title, width_in=12.3, sizes=(27, 25, 23, 21, 19.5, 18)):
    """Largest point size at which a one-line slide title fits its box.

    Titles are full declarative sentences (abstract register), so they are
    often longer than a label; the box holds a single line, so the size is
    chosen from real font metrics when preview_pptx is importable, otherwise
    from a character-class estimate.
    """
    try:
        import preview_pptx as PV
        for sz in sizes:
            f = PV.load_font(FONT, True, sz * 110 / 72.0)
            if f.getlength(title) <= (width_in - 0.1) * 110:
                return sz
        return sizes[-1]
    except Exception:
        est = 0.0
        for ch in title:
            est += 1.0 if ord(ch) > 0x2E7F else (0.3 if ch == " " else 0.56)
        for sz in sizes:
            if est * sz <= (width_in - 0.1) * 72:
                return sz
        return sizes[-1]


def header(s, kicker, title, accent=NAVY):
    box(s, Inches(0), Inches(0), EMU_W, Inches(1.28), fill=WHITE)
    box(s, Inches(0.0), Inches(0.0), Inches(0.22), Inches(1.28), fill=accent)
    text(s, Inches(0.55), Inches(0.20), Inches(12.3), Inches(0.34),
         [[(kicker, 13, accent, True, False, None)]])
    text(s, Inches(0.52), Inches(0.50), Inches(12.4), Inches(0.72),
         [[(title, _title_size(title), NAVY, True, False, None)]],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(0.55), Inches(1.18), Inches(12.25), Pt(2), fill=LGRAY)
    footer(s)


def footer(s):
    _pageno["n"] += 1
    text(s, Inches(0.55), Inches(7.06), Inches(8.0), Inches(0.3),
         [[("KSTAR CES Nowcasting — 다중진단 기반 CES 결측 구간 예측", 9, MGRAY, False, False, None)]])
    text(s, Inches(11.6), Inches(7.06), Inches(1.2), Inches(0.3),
         [[(str(_pageno["n"]), 10, MGRAY, False, False, None)]], align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=15, gap=7, color=DARK):
    """items: list of (text, level, color_override, bold)."""
    paras = []
    for it in items:
        txt, lvl, col, bold = (list(it) + [None, None])[:4]
        lvl = lvl or 0
        mark = "●  " if lvl == 0 else ("–  " if lvl == 1 else "·  ")
        sz = size if lvl == 0 else (size - 2 if lvl == 1 else size - 3)
        paras.append({"segs": [(mark + txt, sz,
                                col if col else (NAVY if lvl == 0 else color),
                                bold if bold is not None else (lvl == 0), False, None)],
                      "level": lvl, "space_after": gap})
    return text(s, x, y, w, h, paras, line_spacing=1.08)


def card(s, x, y, w, h, title, lines, accent=BLUE, title_size=14, body_size=12,
         body_color=DARK):
    c = box(s, x, y, w, h, fill=CARDBG, round_=True)
    box(s, x, y, Inches(0.10), h, fill=accent)
    text(s, x + Inches(0.26), y + Inches(0.12), w - Inches(0.4), Inches(0.4),
         [[(title, title_size, accent, True, False, None)]])
    paras = [[(ln, body_size, body_color, False, False, None)] for ln in lines]
    text(s, x + Inches(0.26), y + Inches(0.56), w - Inches(0.42), h - Inches(0.7),
         paras, line_spacing=1.1, space_after=3)
    return c


def add_image_fit(s, path, x, y, w, h):
    """Add image scaled to fit inside (w,h) box, centered."""
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        pic = s.shapes.add_picture(path, x, y, width=w)
        return pic
    ar = iw / ih
    bw, bh = w, h
    if bw / bh > ar:
        nh = bh
        nw = int(bh * ar)
    else:
        nw = bw
        nh = int(bw / ar)
    nx = x + (bw - nw) // 2
    ny = y + (bh - nh) // 2
    return s.shapes.add_picture(path, nx, ny, width=nw, height=nh)


def table(s, x, y, col_w, head, rows, row_h=Inches(0.44), head_h=Inches(0.44),
          head_fill=NAVY, head_color=WHITE, size=13, head_size=12.5,
          zebra=LGRAY, emphasis=None, emphasis_fill=None, label_align_left=True):
    """Simple grid table.

    col_w   : list of column widths (Emu)
    head    : list of header cell strings
    rows    : list of rows; each cell is str or (txt, color, bold, font)
    emphasis: set of row indices drawn with emphasis_fill + bold
    """
    total_w = sum(col_w)
    box(s, x, y, total_w, head_h, fill=head_fill)
    cx = x
    for j, h in enumerate(head):
        text(s, cx + Inches(0.10), y + Inches(0.07), col_w[j] - Inches(0.16),
             head_h - Inches(0.1),
             [[(h, head_size, head_color, True, False, None)]],
             align=PP_ALIGN.LEFT if (j == 0 and label_align_left) else PP_ALIGN.CENTER)
        cx += col_w[j]
    emphasis = emphasis or set()
    yy = y + head_h
    for i, row in enumerate(rows):
        if i in emphasis and emphasis_fill is not None:
            box(s, x, yy, total_w, row_h, fill=emphasis_fill)
        elif zebra is not None and i % 2 == 1:
            box(s, x, yy, total_w, row_h, fill=CARDBG)
        cx = x
        for j, cell in enumerate(row):
            if isinstance(cell, str):
                txt, col, bold, font = cell, DARK, (i in emphasis), None
            else:
                txt, col, bold, font = (tuple(cell) + (None,) * 4)[:4]
                col = col if col is not None else DARK
                bold = bold if bold is not None else (i in emphasis)
            text(s, cx + Inches(0.10), yy + Inches(0.08), col_w[j] - Inches(0.16),
                 row_h - Inches(0.1),
                 [[(txt, size, col, bold, False, font)]],
                 align=PP_ALIGN.LEFT if (j == 0 and label_align_left) else PP_ALIGN.CENTER)
            cx += col_w[j]
        box(s, x, yy + row_h - Pt(0.75), total_w, Pt(0.75), fill=LGRAY)
        yy += row_h
    return yy


# ============================ SLIDES ======================================
PAPERFIG = os.path.join(HERE, "..", "paper", "figures")


# --- 1. Title -------------------------------------------------------------
def s_title():
    s = slide()
    box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
    box(s, 0, Inches(5.6), EMU_W, Inches(1.9), fill=RGBColor(0x0E, 0x26, 0x47))
    box(s, Inches(0.9), Inches(1.7), Inches(2.2), Pt(4), fill=ORANGE)
    text(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(0.5),
         [[("학위논문 발표 · 약 60분", 16, RGBColor(0x9C, 0xC0, 0xE8), True, False, None)]])
    text(s, Inches(0.88), Inches(2.45), Inches(11.7), Inches(2.0),
         [[("KSTAR 다중 진단 기반 인과(causal) 나우캐스팅:", 30, WHITE, True, False, None)],
          [("희소 CES 신호(Tᵢ · V_rot)의 결측 구간 복원", 34, WHITE, True, False, None)]],
         line_spacing=1.12)
    text(s, Inches(0.9), Inches(4.30), Inches(11.5), Inches(1.1),
         [[("본 연구는 빠른 진단(BES · ECEI · Mirnov)과 과거 CES 이력만으로 ", 16, LGRAY, False, False, None),
           ("이온온도 Tᵢ와 토로이달 회전 V_rot", 16, ORANGE, True, False, None),
           ("를 복원하는 인과 모델을 제안하고,", 16, LGRAY, False, False, None)],
          [("미래를 읽는 오프라인 보간과 배치 가능한 최강 기준선(인과 GP)을 상대로 사전등록 프로토콜 아래 검증하였다.",
            16, LGRAY, False, False, None)]],
         line_spacing=1.2)
    text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(1.0),
         [[("이승상  (Seungsang Lee)", 17, WHITE, True, False, None)],
          [("서울대학교 · 원자핵공학  |  확정 프로토콜: W=2 · held-free · 두 모집단 공동 1차 · 백본 seq_v2 (357,570 파라미터) · 2026-08-27",
            13, MGRAY, False, False, None)]],
         line_spacing=1.25)
    return s


# --- 2. Abstract ----------------------------------------------------------
def s_abstract():
    s = slide()
    header(s, "Abstract", "초록")
    paras = [
        ("배경. ", "KSTAR의 전하교환분광(CES)은 페데스탈 물리의 핵심량인 이온온도 Tᵢ와 토로이달 회전 V_rot를 제공하지만, "
                  "광자 적분이 필요하여 10 ms 격자에서 Tᵢ 8.2%, V_rot 23.9%가 결측되고 V_rot 관측값의 54%는 계측기 유지값(held)이다. "
                  "빔방출분광(BES)·전자사이클로트론방출영상(ECEI)·Mirnov 코일은 같은 격자에서 결측 없이 측정된다."),
        ("방법. ", "641개 방전(247,207행)에서 빠른 진단과 과거 CES 이력만을 입력으로 하는 전체격자 인과 시퀀스 모델 seq_v2(357,570 파라미터)를 "
                  "학습하고, 미래 이웃까지 읽는 오프라인 보간(PCHIP)과 배치 가능한 최강 인과 기준선(인과 GP)을 상대로 "
                  "shot 군집 paired bootstrap(B = 10,000)으로 검정하였다. held 값은 학습·평가에서 전부 제거하였고, "
                  "Tᵢ > 3 keV 피팅 실패의 처리 방식(컷 / 포함)을 두 공동 1차 모집단으로 사전등록하였다."),
        ("결과. ", "Tᵢ는 두 모집단 모두에서 4개 독립 분할 전부 PCHIP를 유의하게 능가하였고(skill 컷 +0.17~+0.26, 포함 +0.23~+0.32), "
                  "인과 GP도 8/8 셀에서 이겼으며, 결측 재가중과 캠페인 시간 분할의 두 스트레스를 생존하였다. V_rot는 전역 동률로 나타났고, "
                  "빠른 채널을 0으로 두어도 출력이 동일하여 정보 부재가 구조적으로 확인되었다. 도달 범위 사다리에서는 "
                  "연속 인과 문맥 약 50 ms에서 skill이 포화하고 승리 방전 비율이 0.52에서 0.66으로 상승하였으며, "
                  "순환·확장 합성곱·attention의 세 계열은 같은 문맥에서 0.023 이내로 동률이었다. 온라인 단계 비용은 "
                  "파라미터가 아니라 디스패치 연산자 수에 비례하였다(연산자당 2–3 µs)."),
        ("결론. ", "빠른 진단은 시간 보간이 복원할 수 없는 Tᵢ 정보를 운반하고, 성능 상한은 모델 크기가 아니라 정보에 있으며, "
                  "남은 개선 레버는 CES 피팅 품질 메타데이터·원본 kHz Mirnov·NBI 토크 채널의 데이터 획득이다."),
    ]
    runs = []
    for lead, body in paras:
        runs.append({"segs": [(lead, 13, NAVY, True, False, None), (body, 12.5, DARK, False, False, None)],
                     "space_after": 7})
    text(s, Inches(0.55), Inches(1.42), Inches(12.25), Inches(5.55), runs, line_spacing=1.16)
    return s


# --- 3. Glossary ----------------------------------------------------------
def s_glossary():
    s = slide()
    header(s, "Terms", "용어 정리 — 본문에서 반복되는 표현의 정의")
    cw = [Inches(2.6), Inches(9.65)]
    rows = [
        ["CES / BES / ECEI / Mirnov", "전하교환분광(타깃 Tᵢ·V_rot) / 빔방출분광(밀도요동) / 전자사이클로트론방출영상(전자온도) / 자기요동 코일(dB/dt)이다."],
        ["나우캐스팅(nowcasting)", "미래를 예보하는 것이 아니라, 이미 지난 시점의 빈 값을 그 시점까지의 정보로 채우는 문제이다."],
        ["skill", "1 - MSE_model / MSE_baseline 이다. 0보다 크면 모델이 기준선보다 낫고, 0이면 동률이다."],
        ["PCHIP / 인과 GP / persistence", "미래 이웃까지 읽는 단조 3차 보간 / 과거 이웃 16개만 읽는 가우시안 과정 / 직전 관측값 유지이다."],
        ["held", "직전 관측값과 비트 단위로 동일한 행이며, 독립 측정이 아니므로 학습·평가에서 제거하였다."],
        ["컷 / 포함 모집단", "Tᵢ > 3 keV(피팅 실패)를 결측으로 처리한 모집단 / 처리하지 않은 모집단이며, 둘 다 공동 1차로 보고한다."],
        ["n/4 · PASS", "4개 독립 분할 중 shot 군집 95% 신뢰구간이 0을 제외한 분할의 수이며, 제외하면 PASS이다."],
        ["도달 범위(reach) · 문맥", "모델이 실제로 읽는 연속 과거 스텝 수이며, 10 ms 격자에서 스텝 × 10 ms가 문맥 길이이다."],
        ["디스패치 연산자 수", "온라인 1스텝에서 실행되는 저수준 연산(aten::)의 개수이며, 본 연구에서 지연 시간의 실측 단위이다."],
    ]
    table(s, Inches(0.55), Inches(1.42), cw, ["용어", "정의"], rows,
          row_h=Inches(0.52), head_h=Inches(0.40), size=12, head_size=12.5)
    return s


# --- 4. Agenda ------------------------------------------------------------
def s_agenda():
    s = slide()
    header(s, "Contents", "발표 구성")
    items = [
        ("1", "연구 배경과 문제 정의", "CES 결측의 실태와 두 모집단의 필요성", ORANGE),
        ("2", "접근법", "미래를 읽는 보간과 최강 인과 기준선을 상대로 한 검증", BLUE),
        ("3", "데이터와 파이프라인", "held 전면 제거 · 두 프레이밍 · 누수 삼중 차단", TEAL),
        ("4", "모델", "전체격자 인과 시퀀스 백본 seq_v2와 W=2 윈도 대조군", NAVY),
        ("5", "평가 방법론", "사전등록 · shot 군집 bootstrap · TEST 동결", BLUE),
        ("6", "결과", "Tᵢ 4/4+4/4 · V_rot 동률 · 스트레스 2종 생존", ORANGE),
        ("7", "문맥·구조·비용", "도달 범위 사다리 · 계열 비교 · 연산자 비용 · 승패 방전", TEAL),
        ("8", "결론·한계·향후 연구", "상한은 추정기가 아니라 정보에 있다", GRAY),
    ]
    y = 1.5
    for i, (num, t, sub, col) in enumerate(items):
        col_x = 0.7 if i < 4 else 6.95
        yy = Inches(y + (i % 4) * 1.32)
        box(s, Inches(col_x), yy, Inches(5.7), Inches(1.14), fill=CARDBG, round_=True)
        box(s, Inches(col_x + 0.12), yy + Inches(0.17), Inches(0.8), Inches(0.8),
            fill=col, round_=True)
        text(s, Inches(col_x + 0.12), yy + Inches(0.17), Inches(0.8), Inches(0.8),
             [[(num, 26, WHITE, True, False, None)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(col_x + 1.05), yy + Inches(0.16), Inches(4.5), Inches(0.45),
             [[(t, 15.5, NAVY, True, False, None)]])
        text(s, Inches(col_x + 1.05), yy + Inches(0.62), Inches(4.55), Inches(0.45),
             [[(sub, 11.5, GRAY, False, False, None)]])
    return s


# --- Section divider ------------------------------------------------------
def divider(num, title, sub):
    s = slide()
    box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
    box(s, Inches(0.0), Inches(2.7), Inches(3.6), Inches(2.1), fill=RGBColor(0x0E, 0x26, 0x47))
    text(s, Inches(0.55), Inches(2.55), Inches(3.2), Inches(2.4),
         [[(num, 120, ORANGE, True, False, None)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(4.0), Inches(3.0), Pt(3), Inches(1.6), fill=ORANGE)
    text(s, Inches(4.4), Inches(3.05), Inches(8.4), Inches(1.0),
         [[(title, 36, WHITE, True, False, None)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.42), Inches(4.05), Inches(8.4), Inches(0.8),
         [[(sub, 16, LGRAY, False, False, None)]])
    return s


# --- 1-1. Diagnostics -----------------------------------------------------
def s_diagnostics():
    s = slide()
    header(s, "1. 연구 배경", "진단 장치 구성: CES는 느리고 자주 비며, 빠른 진단은 같은 격자에서 결측이 없다")
    text(s, Inches(0.55), Inches(1.45), Inches(12.3), Inches(0.6),
         [[("KSTAR 페데스탈 상단 플라즈마에서 ", 15, DARK, False, False, None),
           ("CES", 15, ORANGE, True, False, None),
           ("는 핵심 물리량을 제공하지만 느리고 자주 결측된다. ", 15, DARK, False, False, None),
           ("BES·ECEI·Mirnov는 같은 10 ms 격자에서 결측 없이 측정된다.", 15, NAVY, True, False, None)]])
    cards = [
        ("CES  (타깃)", ORANGE,
         ["전하교환분광(Charge Exchange Spectroscopy)이다.",
          "Tᵢ(이온온도)와 V_rot(토로이달 회전)를 제공한다.",
          "광자 적분이 필요하여 느리고 자주 결측된다.",
          "본 연구가 복원하려는 대상이다."]),
        ("BES  (9 ch)", BLUE,
         ["빔방출분광(Beam Emission Spectroscopy)이다.",
          "밀도요동 nₑ의 공간 구조를 측정한다.",
          "10 ms 격자에서 항상 측정된다.",
          "충돌 e–i 결합을 통해 Tᵢ의 단서가 된다."]),
        ("ECEI  (4 ch)", TEAL,
         ["전자사이클로트론방출영상이다.",
          "전자온도 Tₑ의 2차원 영상을 제공한다.",
          "10 ms 격자에서 항상 측정된다.",
          "Tₑ–Tᵢ 결합을 통해 Tᵢ의 단서가 된다."]),
        ("Mirnov coil  (2 ch)", GRAY,
         ["자기요동(MHD 모드) dB/dt를 측정한다.",
          "kHz 신호가 100 Hz로 데시메이트되어 있다.",
          "lag-1 자기상관이 -0.009이다(BES +0.568).",
          "모드 회전 정보는 상류에서 소실되었다."]),
    ]
    x0 = 0.55
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(x0 + i * 3.12), Inches(2.2), Inches(2.95), Inches(3.5),
             t, lines, accent=col, title_size=14.5, body_size=11.5)
    text(s, Inches(0.55), Inches(5.95), Inches(12.3), Inches(0.9),
         [[("핵심 비대칭의 예고: ", 14, NAVY, True, False, None),
           ("빠른 진단은 물리적으로 Tᵢ 정보는 운반하지만 V_rot 정보는 거의 운반하지 않는다고 예측되었고"
            "(NBI 토크 미관측, Mirnov 앨리어싱), 이 가설은 결과에서 그대로 확인되었다.",
            14, DARK, False, False, None)]], line_spacing=1.15)
    return s


# --- 1-2. CES missing problem ---------------------------------------------
def s_problem():
    s = slide()
    header(s, "1. 연구 배경", "문제 정의: CES 결측의 원인과 규모")
    bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(4.6), [
        ("CES는 충분한 신호대잡음비를 위해 광자를 오래 적분해야 한다.", 0),
        ("노출·신호품질 문제로 특정 시점의 측정이 자주 누락된다.", 1),
        ("같은 10 ms 격자에서 Tᵢ 8.2%, V_rot 23.9%가 완전 결측(NaN)이다.", 0),
        ("V_rot는 held(직전값 복사) 41.1%가 더해져 실질 무정보 비율이 65.0%에 이른다.", 1, RED, True),
        ("두 타깃은 서로 독립적으로 결측되므로 타깃별 처리가 필요하다.", 1),
        ("결측은 물리적으로 흥미로운 순간(저신호·ELM·천이)에 집중된다.", 0),
        ("따라서 결측은 무작위가 아니며(MNAR), 관측점 skill은 낙관적 상한이다.", 1, ORANGE, True),
        ("빠른 진단(BES·ECEI·Mirnov)은 같은 격자에서 결측이 없다.", 0),
        ("본 연구는 항상 존재하는 빠른 진단으로 자주 비는 CES를 채운다.", 1, ORANGE, True),
    ], size=14.5)
    box(s, Inches(7.1), Inches(1.55), Inches(5.7), Inches(2.5), fill=CARDBG, round_=True)
    text(s, Inches(7.35), Inches(1.72), Inches(5.3), Inches(0.5),
         [[("데이터 기반 가상 센서의 기대 효과", 16, NAVY, True, False, None)]])
    bullets(s, Inches(7.35), Inches(2.25), Inches(5.25), Inches(1.8), [
        ("결측 시점의 Tᵢ·V_rot를 데이터로 추정한다(gap-filling).", 0, TEAL, True),
        ("강한 역산 가정 없이 축대칭 수준만 가정한다.", 0, TEAL, True),
        ("끊김 없는 CES 가용성은 페데스탈 물리 분석과 실시간 활용을 가능하게 한다.", 0, TEAL, True),
    ], size=12.5, gap=8)
    add_image_fit(s, os.path.join(FIG, "fig_missing.png"),
                  Inches(7.0), Inches(4.2), Inches(6.0), Inches(2.6))
    return s


# --- 1-3. Missingness ledger ----------------------------------------------
def s_missing_table():
    s = slide()
    header(s, "1. 연구 배경", "결측 전수 집계: NaN 결측과 held(직전값 복사)를 합치면 V_rot의 65.0%가 무정보이다", accent=RED)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.62),
         [[("‘V_rot 결측 24%’는 NaN만 센 값이다. ", 14.5, RED, True, False, None),
           ("직전 관측값을 그대로 복사한 held 행을 합치면 실질 무정보 비율은 ", 14.5, DARK, False, False, None),
           ("65.0%", 14.5, RED, True, False, None),
           ("로 나타났다(641 shot · 247,207행 전수 집계).", 14.5, DARK, False, False, None)]],
         line_spacing=1.14)
    cw = [Inches(5.4), Inches(3.4), Inches(3.4)]
    rows = [
        ["전체 10 ms 격자 행 (641 shot)",
         ("247,207", GRAY, False, MONO), ("247,207", GRAY, False, MONO)],
        ["① 완전 결측 — 값이 비어 있음 (NaN)",
         ("20,216  (8.2%)", DARK, False, MONO), ("59,107  (23.9%)", BLUE, True, MONO)],
        ["② held / padding — 직전 관측값과 bit-identical",
         ("1  (0.0%)", DARK, False, MONO), ("101,604  (41.1%)", RED, True, MONO)],
        [("실질 무정보  ① + ②", NAVY, True, None),
         ("20,217  (8.2%)", NAVY, True, MONO), ("160,711  (65.0%)", RED, True, MONO)],
        ["독립 관측 — 실제 정보가 있는 행",
         ("226,990  (91.8%)", GREEN, True, MONO), ("86,496  (35.0%)", DARK, True, MONO)],
        ["관측값(non-NaN) 중 held 비율",
         ("0.0%", GREEN, True, MONO), ("54.0%", RED, True, MONO)],
        ["held이 존재하는 shot 파일 수",
         ("1 / 641", DARK, False, MONO), ("499 / 641", RED, True, MONO)],
        ["연속 held 구간 길이 — 중앙값 / 최대 (행)",
         ("2 / 2", DARK, False, MONO), ("10 / 1,214", RED, True, MONO)],
    ]
    table(s, Inches(0.55), Inches(2.04), cw, ["구분", "CES_TI", "CES_VT"], rows,
          row_h=Inches(0.42), head_h=Inches(0.42), emphasis={3},
          emphasis_fill=RGBColor(0xFD, 0xEC, 0xE8))
    box(s, Inches(0.55), Inches(5.90), Inches(12.25), Inches(1.06),
        fill=CARDBG, round_=True)
    text(s, Inches(0.8), Inches(5.98), Inches(11.75), Inches(0.94),
         [[("의의: ", 12, NAVY, True, False, None),
           ("held 행은 persistence와 보간이 오차 0에 가깝게 맞히는 행이므로, 남겨 두면 기준선의 점수가 부풀려진다.",
            12, DARK, False, False, None)],
          [("확정 프로토콜: ", 12, NAVY, True, False, None),
           ("지도 타깃·이력 입력·정규화 통계·모든 기준선의 보간 앵커에서 held를 동일하게 제거하고 실제 측정만 채점하였다.",
            12, DARK, False, False, None)],
          [("판정 기준: ", 12, NAVY, True, False, None),
           ("연속 블록 안에서 직전 관측값과 부동소수점까지 동일한 행으로 정의하였다. CES_TI는 226,991행 중 1행에 그쳤다.",
            12, DARK, False, False, None)]],
         line_spacing=1.12, space_after=2)
    return s


# --- 1-4. Two populations -------------------------------------------------
def s_two_populations():
    s = slide()
    header(s, "1. 연구 배경", "데이터 품질 감사: Tᵢ 피팅 실패의 처리 방식을 두 공동 1차 모집단으로 사전등록하였다", accent=RED)
    text(s, Inches(0.55), Inches(1.38), Inches(12.3), Inches(0.6),
         [[("관측 Tᵢ의 p99는 2,089 eV, p99.9는 9,601 eV, 최대는 14,984 eV이다. ", 14, DARK, False, False, None),
           ("이 먼 꼬리는 플라즈마가 아니라 실패한 스펙트럼 피팅이다.", 14, RED, True, False, None)]],
         line_spacing=1.14)
    card(s, Inches(0.55), Inches(2.02), Inches(6.0), Inches(2.05),
         "실측 — > 3 keV 행의 구조", [
             "1,197행(0.53%)이 951개 run, 274개 방전에 분포한다.",
             "run의 85%는 단일 행이고 5행 이상 지속은 2%에 그친다.",
             "run 정점은 관측 이웃 평균의 13배(IQR 6–26배)이다.",
             "어떤 방법으로도 예측되지 않으며 보간 앵커를 오염시킨다.",
         ], accent=RED, title_size=14.5, body_size=12.5)
    card(s, Inches(6.8), Inches(2.02), Inches(6.0), Inches(2.05),
         "대응 — 두 처리 모두 비판이 가능하다", [
             "제거하면 어려운 행을 없앴다는 비판을 받는다.",
             "유지하면 스파이크 앵커가 오프라인 기준선을 불리하게 한다.",
             "따라서 두 모집단을 공동 1차로 사전등록하였다.",
             "컷(적재 시 결측 처리, 전 arm 동일) / 포함(컷 없음)이다.",
         ], accent=BLUE, title_size=14.5, body_size=12.5)
    box(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(0.72), fill=NAVY, round_=True)
    text(s, Inches(0.85), Inches(4.46), Inches(11.7), Inches(0.6),
         [[("규칙: ", 14, ORANGE, True, False, None),
           ("무조건부 주장은 두 모집단 모두에서 성립할 때만 한다. 한쪽에서만 성립하면 모집단을 명시하여 보고한다.",
            14, WHITE, False, False, None)]])
    bullets(s, Inches(0.55), Inches(5.35), Inches(12.3), Inches(1.5), [
        ("문턱 민감도는 없었다. 2.5 / 3 / 4 keV로 재학습한 Tᵢ skill은 +0.230 / +0.236 / +0.232이며 PR4는 전부 4/4였다.", 0, GREEN, True),
        ("값 컷은 일방향 프록시이다. 하향 dip 4,965행은 손대지 않고, 2배 이상 상향 이상치의 19%만 제거한다.", 0),
        ("V_rot 스파이크(> 1,000 km/s 119행 / 16 방전, 101행은 한 방전의 한 블록)는 컷 없이 두고 SSE 비중을 병기한다.", 0),
    ], size=13, gap=6)
    return s


# --- 1-5. Research question -----------------------------------------------
def s_idea():
    s = slide()
    header(s, "1. 연구 배경", "연구 질문과 핵심 아이디어")
    box(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.4), fill=NAVY, round_=True)
    text(s, Inches(1.0), Inches(1.72), Inches(11.4), Inches(1.1),
         [[("연구 질문", 13, ORANGE, True, False, None)],
          [("CES가 결측된 10 ms 시점에서, 동시각 빠른 진단(BES·ECEI·Mirnov)과 과거 CES 이력만으로 ",
            16.5, WHITE, False, False, None)],
          [("CES 자체의 시간 보간이 복원할 수 없는 정보를 회복할 수 있는가를 묻는다.",
            16.5, WHITE, True, False, None)]], line_spacing=1.18)
    cards = [
        ("가상 센서 (Virtual Sensor)", BLUE,
         ["빠른 진단으로부터 CES를 데이터 기반으로 추정한다.",
          "역산(inverse mapping) 가정을 두지 않는다.",
          "결측·고장 시점을 온라인으로 메운다."]),
        ("Gap-filling / Nowcasting", TEAL,
         ["미래를 예보하는 forecasting이 아니다.",
          "현재 시점의 빈 값을 채우는 nowcast이다.",
          "초해상(super-resolution)과도 구분된다."]),
        ("검증 방식", ORANGE,
         ["진짜 결측은 참값이 없어 직접 검증되지 않는다.",
          "관측 CES를 가린 뒤 복원 정확도를 측정한다.",
          "두 모집단 각각에서 타깃별로 채점한다."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(0.55 + i * 4.13), Inches(3.35), Inches(3.95), Inches(2.15),
             t, lines, accent=col, title_size=15, body_size=12.5)
    text(s, Inches(0.7), Inches(5.90), Inches(11.9), Inches(0.9),
         [[("결론의 성격: ", 13.5, NAVY, True, False, None),
           ("masking 검증에서 기준선을 이기면 결측 구간에서도 복원이 유효할 것으로 ", 13.5, DARK, False, False, None),
           ("추정", 13.5, ORANGE, True, False, None),
           ("한다. 결측이 무작위라는 보장이 없으므로(MNAR) 결측 지점의 정확도를 단정하지 않고, "
            "결측 분포로 재가중하여 얼마나 살아남는지를 측정하였다.",
            13.5, DARK, False, False, None)]], line_spacing=1.15)
    return s


# --- 2-1. The hard bar ----------------------------------------------------
# --- 1-6. Related work ----------------------------------------------------
def s_related():
    s = slide()
    header(s, "1. 연구 배경", "선행 연구와의 위치: 계보를 잇고 세 축으로 확장한다")
    text(s, Inches(0.55), Inches(1.30), Inches(12.25), Inches(0.30),
         [[("교차 진단 재구성은 활발한 계열이며, 본 연구는 그 프로그램을 이어받아 세 축에서 확장한다.",
            12.5, GRAY, False, False, None)]])
    col_w = [Inches(3.05), Inches(4.55), Inches(4.61)]
    table(s, Inches(0.55), Inches(1.58), col_w,
          ["선행 연구", "무엇을 하는가", "본 연구가 다른 지점"],
          [
              [("Diag2Diag (Nat. Commun. 2025)", NAVY, True, None),
               "동시각 진단들로부터 Thomson Tₑ·nₑ를 요동 진단 속도로 합성한다.",
               "타깃이 전자 물리량이고 CES는 입력이며, 기억이 없다."],
              [("COMPASS 시간 초해상 (PPCF 2026)", NAVY, True, None),
               "고속 복사 진단으로 Thomson 프로파일을 시간 초해상한다.",
               "전자 채널이고 비인과이며 보간을 기준선으로 두지 않는다."],
              [("RTCAKENN (NF 2024, DIII-D)", NAVY, True, None),
               "실시간 운동학 프로파일을 재구성하고 입력 소실을 견딘다.",
               "결측이 입력의 강건성 조건이지 타깃의 희소성이 아니다."],
              [("EAST XCS → Tᵢ·회전 (NF 2024)", NAVY, True, None),
               "같은 두 물리량을 X선 결정 분광기로부터 추론한다.",
               "입력이 또 다른 도플러 분광기이며 순간 매핑이다."],
              [("FusionMAE (HL-3, 2026)", NAVY, True, None),
               "마스킹된 채널을 가상 백업 진단으로 복원한다.",
               "일반 채널 드롭아웃이고 윈도 내 재구성이다."],
              [("KSTAR EPED 재구성 (JKPS 2026)", NAVY, True, None),
               "CES를 입력으로 쓰는 고속 프로파일 재구성이다.",
               "회전은 저자들이 명시한 향후 과제로 남아 있다."],
              [("NN-CES (KSTAR, FED 2025/26)", NAVY, True, None),
               "CES 스펙트럼에서 Tᵢ·회전을 빔 변조 없이 직접 추출한다.",
               "측정이 있는 시각을 읽고, 우리는 없는 시각을 채운다."],
          ], row_h=Inches(0.545), head_h=Inches(0.38), size=10.5, head_size=12)
    axes = [
        ("① 타깃 채널", ["전자 물리량의 재구성에서",
                      "희소 이온 채널(Tᵢ · V_rot)의 복원으로."], TEAL),
        ("② 인과성", ["동시각의 기억 없는 매핑에서",
                    "불규칙한 과거에 조건화된 인과 추정으로."], BLUE),
        ("③ 평가 기준", ["가정된 재구성 가능성에서",
                      "사전등록 기준선에 대한 타깃별 검정으로."], ORANGE),
    ]
    for i, (title, lines, accent) in enumerate(axes):
        card(s, Inches(0.55 + i * 4.13), Inches(5.78), Inches(3.95), Inches(1.24),
             title, lines, accent=accent, title_size=12.5, body_size=11)
    return s


def s_bar():
    s = slide()
    header(s, "2. 접근법", "평가 기준선: 미래를 읽는 오프라인 보간과 배치 가능한 최강 인과 기준선을 함께 두었다")
    text(s, Inches(0.55), Inches(1.42), Inches(12.3), Inches(0.62),
         [[("모델을 ", 15, DARK, False, False, None),
           ("오프라인 CES-only 보간(선형 · PCHIP · 국소 AR · GP)", 15, NAVY, True, False, None),
           ("과 비교하였다. 이 보간들은 타깃 주변의 과거와 미래 CES를 모두 사용한다.",
            15, DARK, False, False, None)]])
    box(s, Inches(0.55), Inches(2.10), Inches(6.0), Inches(2.55), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(2.10), Inches(0.12), Inches(2.55), fill=ORANGE)
    text(s, Inches(0.8), Inches(2.27), Inches(5.6), Inches(0.5),
         [[("본 모델 (인과)", 16, ORANGE, True, False, None)]])
    bullets(s, Inches(0.8), Inches(2.80), Inches(5.5), Inches(1.8), [
        ("타깃 시점까지의 빠른 진단(BES·ECEI·Mirnov)을 읽는다.", 0),
        ("과거 CES 이력을 세그먼트 전체 범위에서 읽는다.", 0),
        ("미래 CES는 전혀 읽지 않는다.", 0, RED, True),
    ], size=13.5, gap=9)
    box(s, Inches(6.8), Inches(2.10), Inches(6.0), Inches(2.55), fill=CARDBG, round_=True)
    box(s, Inches(6.8), Inches(2.10), Inches(0.12), Inches(2.55), fill=BLUE)
    text(s, Inches(7.05), Inches(2.27), Inches(5.6), Inches(0.5),
         [[("보간 기준선 (오프라인)", 16, BLUE, True, False, None)]])
    bullets(s, Inches(7.05), Inches(2.80), Inches(5.5), Inches(1.8), [
        ("타깃 양쪽의 과거와 미래 CES 이웃을 사용한다.", 0),
        ("PCHIP(단조 3차)이 사전등록된 headline 기준선이다.", 0),
        ("세그먼트 경계를 넘지 않으며, 넘어야 하면 persistence로 대체한다.", 0),
    ], size=13.5, gap=9)
    box(s, Inches(0.55), Inches(4.85), Inches(12.25), Inches(1.02), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(4.85), Inches(0.12), Inches(1.02), fill=TEAL)
    text(s, Inches(0.85), Inches(4.94), Inches(11.7), Inches(0.9),
         [[("추가된 팔 — 인과 GP: 배치 가능한 가장 강한 경쟁자", 13.5, TEAL, True, False, None)],
          [("같은 GP를 과거 이웃 16개로 제한하였다(NaN 조건 동일, 모집단 불변). seed 42·컷에서 Tᵢ RMSE는 164.3으로 "
            "persistence 197.2보다 낮다. ‘배치 가능한 모든 인과 방법을 이긴다’는 주장은 이 기준선으로 판정한다.",
            12.5, DARK, False, False, None)]], line_spacing=1.12)
    box(s, Inches(0.55), Inches(6.03), Inches(12.25), Inches(0.85), fill=NAVY, round_=True)
    text(s, Inches(0.85), Inches(6.12), Inches(11.7), Inches(0.72),
         [[("정보 비대칭의 의미", 12.5, ORANGE, True, False, None)],
          [("미래까지 읽는 보간을 인과 모델이 이긴다면, 빠른 진단이 시간 보간으로는 얻을 수 없는 CES 정보를 "
            "운반한다는 강한 증거가 된다.", 14, WHITE, False, False, None)]], line_spacing=1.10, space_after=2)
    return s


# --- 2-2. Validation strategy ---------------------------------------------
def s_validation():
    s = slide()
    header(s, "2. 접근법", "검증 전략과 그 한계")
    cards = [
        ("Masking 복원 검증", BLUE,
         ["관측된 CES 값을 가린 뒤 모델이 복원한다.",
          "물리 단위로 역정규화한 뒤 타깃별로 채점한다.",
          "전 arm이 동일한 (file, row) 집합과 마스크를 쓴다.",
          "모델 입력에서 타깃 시점은 완전히 마스킹된다."]),
        ("Observed-only 측정", TEAL,
         ["관측된 CES 지점에서만 skill을 측정한다.",
          "진짜 결측은 참값이 없어 직접 검증되지 않는다.",
          "관측 지점이 결측 지점보다 쉬울 수 있다.",
          "PR2 폴백률은 Tᵢ 0.3–0.4%, V_rot 40–44%이다."]),
        ("MNAR — 낙관적 상한", ORANGE,
         ["CES 결측은 무작위가 아니다(MNAR).",
          "저 SNR·ELM·천이에서 drop-out이 발생한다.",
          "따라서 observed-only skill은 낙관적 상한이다.",
          "결측 분포로 재가중하여 정량화하였다(결과 ⑤)."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(0.55 + i * 4.13), Inches(1.6), Inches(3.95), Inches(3.0),
             t, lines, accent=col, title_size=15, body_size=12.5)
    box(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.6), fill=CARDBG, round_=True)
    text(s, Inches(0.95), Inches(5.12), Inches(11.4), Inches(1.4),
         [[("Murphy(1988) skill 정의", 14, NAVY, True, False, None)],
          [("skill_vs_pchip = 1 - MSE_model / MSE_pchip", 17, BLUE, True, False, MONO)],
          [("0보다 크면 모델이 보간보다 우수하고 0이면 동률이다. 유의성은 shot 군집 paired bootstrap 95% CI가 0을 제외하는지로 판정하며(PR4), "
            "‘n/4’는 유의한 분할 수이다.", 13, DARK, False, False, None)]], line_spacing=1.2)
    return s


# --- 3-1. Data ------------------------------------------------------------
def s_data():
    s = slide()
    header(s, "3. 데이터 & 파이프라인", "데이터 구성과 세그먼트 구조")
    bullets(s, Inches(0.55), Inches(1.5), Inches(6.6), Inches(4.9), [
        ("641개 KSTAR 방전(shot 30801–32751)을 10 ms 공통 격자로 사용하였다.", 0, NAVY, True),
        ("제공 측 선정 기준은 하드웨어 일관성과 H-mode ELM 억제(RMP) 구간이다.", 1),
        ("총 247,207행이며 파일당 중앙값은 339행이다.", 1),
        ("행당 채널은 BES 9 · ECEI 4 · Mirnov 2 · time · CES_TI · CES_VT이다.", 0),
        ("세그먼트는 0.5 s 이상의 간극에서 분리하였다(이봉 delta 분포의 골).", 0),
        ("(0.1, 0.5) s 구간의 delta는 82개뿐이며 세그먼트 내 스텝의 99.4%가 10 ms이다.", 1),
        ("전형적 파일은 주 세그먼트 1개(중앙값 301행 ≈ 3.0 s, 10–90분위 1.3–7.0 s)이다.", 0),
        ("2개인 파일 28개, 10행 넘는 세그먼트가 없는 파일 20개, 고립 단일행 1,279개가 있다.", 1),
        ("어떤 arm도(보간·모델 입력 모두) 세그먼트 간극을 넘지 않는다.", 0, ORANGE, True),
    ], size=13.5, gap=7)
    box(s, Inches(7.35), Inches(1.5), Inches(5.45), Inches(4.9), fill=NAVY, round_=True)
    text(s, Inches(7.6), Inches(1.66), Inches(5.0), Inches(0.5),
         [[("TEST 규모 (선택 과정에서 읽지 않음)", 15, ORANGE, True, False, None)]])
    rows = [
        ("기준 분할", "seed 42", ""),
        ("Tᵢ (컷)", "32,589 행", "96 방전"),
        ("V_rot (컷)", "10,463 행", "60 방전"),
        ("Tᵢ (포함)", "32,721 행", "96 방전"),
        ("V_rot (포함)", "10,461 행", "60 방전"),
        ("4 분할 Tᵢ", "32.6–35.9k 행", "96 방전"),
        ("4 분할 V_rot", "10.5–14.5k 행", "60–66 방전"),
    ]
    yy = 2.24
    for name, val, shots in rows:
        text(s, Inches(7.6), Inches(yy), Inches(1.85), Inches(0.4),
             [[(name, 13, WHITE, True, False, None)]])
        text(s, Inches(9.45), Inches(yy), Inches(1.9), Inches(0.4),
             [[(val, 12.5, ORANGE, True, False, None)]])
        text(s, Inches(11.4), Inches(yy), Inches(1.35), Inches(0.4),
             [[(shots, 11.5, LGRAY, False, False, None)]])
        yy += 0.48
    text(s, Inches(7.6), Inches(yy + 0.08), Inches(5.0), Inches(0.9),
         [[("모든 수치는 단일 collector가 동결된 run 디렉터리에서 읽으므로 본문·표·그림이 서로 어긋나지 않는다.",
            11.5, LGRAY, False, True, None)]], line_spacing=1.15)
    return s


# --- 3-2. Contract ---------------------------------------------------------
def s_contract():
    s = slide()
    header(s, "3. 데이터 & 파이프라인", "No-Fake-Data 원칙과 데이터 계약")
    cards1 = [
        ("① 가짜 라벨 금지", ORANGE,
         ["학습 행을 만들기 위해 타깃을 대체(impute)하지 않았다.",
          "윈도는 진단 입력이 완전하고 타깃이 1개 이상 관측된 행만 쓴다.",
          "시퀀스는 라벨 없는 행을 맥락으로만 쓴다.",
          "어느 프레이밍도 타깃 행 자신의 값을 읽지 않는다."]),
        ("② 타깃별 masked loss", BLUE,
         ["L = Σ m·(예측 - 실측)² / Σ m 이며 m은 타깃별 관측 마스크이다.",
          "한쪽 타깃만 관측된 행도 그 타깃의 학습에 기여한다.",
          "두 타깃 필수 필터는 라벨 행의 약 28%를 버렸었다.",
          "이 필터의 제거는 순수한 데이터 이득이다."]),
    ]
    for i, (t, col, lines) in enumerate(cards1):
        card(s, Inches(0.55 + i * 6.2), Inches(1.55), Inches(6.0), Inches(2.4),
             t, lines, accent=col, title_size=15, body_size=12.5)
    cards2 = [
        ("③ 누수 삼중 차단", TEAL,
         ["파일(shot) 단위로 분할하여 인접 행의 자기상관 누수를 막았다.",
          "정규화 통계는 학습 파일에서만 추정하였다(희소 타깃은 NaN 인지).",
          "시퀀스 모델은 shot별 입력 표준화를 추가로 적용하였다.",
          "타깃 시점의 값과 관측 flag는 입력에 들어가지 않는다."]),
        ("④ held 전면 제거", NAVY,
         ["관측 V_rot의 54%는 계측기 유지값이며 실제 측정이 아니다.",
          "지도 타깃·이력 입력·정규화 통계에서 제거하였다.",
          "모든 기준선의 보간 앵커에서도 동일하게 제거하였다.",
          "따라서 어떤 arm도 forward-fill로 점수를 얻지 못한다."]),
    ]
    for i, (t, col, lines) in enumerate(cards2):
        card(s, Inches(0.55 + i * 6.2), Inches(4.15), Inches(6.0), Inches(2.4),
             t, lines, accent=col, title_size=15, body_size=12.5)
    return s


# --- 3-3. Split + normalization ------------------------------------------
def s_split():
    s = slide()
    header(s, "3. 데이터 & 파이프라인", "누수 방지 분할과 학습셋 전용 정규화")
    cards = [
        ("File-level split", BLUE,
         ["행 단위가 아니라 CSV(shot) 파일 단위로 분할하였다.",
          "인접 행이 상관되므로 행 단위 분할은 train/val 누수를 낳는다.",
          "고정 split을 디스크에 고정하였다(fixed_*_split.csv).",
          "재로딩 시 데이터가 불일치하면 예외를 발생시킨다."]),
        ("Train-file-only 정규화", TEAL,
         ["BES·ECEI·MC·타깃을 채널별 z-score로 정규화하였다.",
          "통계는 학습 파일에서만 추정하여 val/test 누수를 막았다.",
          "타깃 통계는 NaN 인지 방식으로 관측값만 사용하였다.",
          "시퀀스 모델은 각 방전의 빠른 진단을 그 방전 자신의",
          "통계로 추가 표준화하였다(캠페인 전이의 핵심)."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(0.55 + i * 6.2), Inches(1.55), Inches(6.0), Inches(2.7),
             t, lines, accent=col, title_size=15, body_size=12.5)
    box(s, Inches(0.55), Inches(4.5), Inches(12.25), Inches(2.0), fill=CARDBG, round_=True)
    text(s, Inches(0.8), Inches(4.65), Inches(11.8), Inches(0.5),
         [[("타깃 시점 마스킹 — 누수 차단의 핵심", 15, NAVY, True, False, None)]])
    bullets(s, Inches(0.8), Inches(5.15), Inches(11.8), Inches(1.3), [
        ("윈도 ces_history는 (batch, W=2, 4)이며 이전 정규화 Tᵢ·V_rot와 타깃별 관측 flag 2개로 구성된다.", 0),
        ("Tᵢ·V_rot는 독립적으로 결측되므로(8.2% / 23.9%) 관측을 타깃별로 추적한다.", 1),
        ("타깃 시점은 값과 flag를 모두 0으로 마스킹하여 자기 정답 누수를 차단한다(시퀀스도 동일).", 1, RED, True),
    ], size=13, gap=7)
    return s


# --- 3-4. Two framings ----------------------------------------------------
def s_samples():
    s = slide()
    header(s, "3. 데이터 & 파이프라인", "학습 예제의 두 프레이밍: 윈도(대조군)와 전체격자 시퀀스(주 모델)")
    cards = [
        ("① 연속 세그먼트", BLUE,
         ["time delta 0.5 s 이상을 세그먼트 경계로 정의하였다.",
          "전형 파일은 주 세그먼트 1개(중앙값 301행)이다.",
          "모델 입력과 보간 모두 경계를 넘지 않는다."]),
        ("② 윈도 프레이밍 (대조군, W=2)", TEAL,
         ["타깃 t 앞 W=2행의 bes(2,9)·ecei(2,4)·mc(2,2)를 읽는다.",
          "time_features(2,4)와 ces_history(2,4)를 더한다.",
          "타깃 [Tᵢ, V_rot]와 타깃별 마스크 m ∈ {0,1}²를 가진다.",
          "파일당 샘플 상한은 500으로 사전등록되었다."]),
        ("③ 전체격자 시퀀스 (주 모델)", ORANGE,
         ["세그먼트의 입력 완전 행 전부를 맥락으로 유지한다.",
          "라벨 유무와 무관하며 희소성은 loss가 처리한다.",
          "스텝당 22채널이며 z-score 빠른 채널 15와",
          "log1p(Δt), 타깃별(이월값·신선도·flag) 3×2로 구성된다."]),
        ("④ 시간 특징 4채널 (윈도)", NAVY,
         ["lookback 초, 행간 delta 초, 각각의 log1p이다.",
          "불규칙 관측 패턴을 명시적으로 노출한다.",
          "과거 CES 값의 신뢰도는 10 ms 전인지",
          "200 ms 전인지에 강하게 의존한다."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        r, c = divmod(i, 2)
        card(s, Inches(0.55 + c * 6.2), Inches(1.5 + r * 2.5), Inches(6.0), Inches(2.3),
             t, lines, accent=col, title_size=14.5, body_size=12.5)
    text(s, Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.6),
         [[("두 프레이밍의 결정적 차이는 도달 범위(reach)이다. ", 12.5, NAVY, True, False, None),
           ("윈도는 과거 W-1개 관측을, 시퀀스는 세그먼트 전체를 읽는다. 이 차이는 백본 관문(결과 ③)과 도달 범위 사다리(7장)에서 측정된다.",
            12.5, GRAY, False, False, None)]])
    return s


# --- 3-5. held audit ------------------------------------------------------
def s_stuck():
    s = slide()
    header(s, "3. 데이터 & 파이프라인", "데이터 품질 감사: held / forward-fill된 CES_VT", accent=RED)
    text(s, Inches(0.55), Inches(1.45), Inches(12.3), Inches(0.9),
         [[("감사 결과: ", 15, RED, True, False, None),
           ("관측된 CES_VT 값의 54%가 계측기 유지값이었다", 15, RED, True, False, None),
           ("(같은 연속 블록 안에서 직전 관측과 bit-identical, 최대 1,214행 연속). "
            "V_rot의 고유 측정 주기가 행 주기보다 느려 값이 carry-forward된 것이며, 독립적인 측정이 아니다.",
            14.5, DARK, False, False, None)]], line_spacing=1.16)
    bullets(s, Inches(0.55), Inches(2.55), Inches(6.4), Inches(4.0), [
        ("641개 중 499개 파일이 영향을 받았고, CES_TI는 226,991행 중 1행이었다.", 0),
        ("오탐 통로는 없다. V_rot는 소수점 5자리이며 값 간 최소 간격은 4e-5이다.", 1),
        ("확정 프로토콜은 held를 모든 위치에서 제거한다.", 0, RED, True),
        ("지도 타깃, 이력·이월 입력과 그 관측 flag, 정규화 통계에서 제거한다.", 1),
        ("모든 기준선의 보간 앵커에서도 제거하여 실제 측정만 사용한다.", 1),
        ("held는 평가뿐 아니라 학습도 오염시켰다(짝지은 재학습으로 확인).", 0, NAVY, True),
        ("forward-fill 계단은 이력 복사가 최적이라고 학습시킨다.", 1),
        ("따라서 민감도 한 줄이 아니라 프로토콜로 채택하였다.", 1, NAVY, True),
    ], size=13.5, gap=8)
    box(s, Inches(7.1), Inches(2.55), Inches(5.7), Inches(3.9), fill=CARDBG, round_=True)
    box(s, Inches(7.1), Inches(2.55), Inches(0.12), Inches(3.9), fill=ORANGE)
    text(s, Inches(7.35), Inches(2.70), Inches(5.3), Inches(0.5),
         [[("held 제거의 대가 — PR2 폴백", 14, ORANGE, True, False, None)]])
    bullets(s, Inches(7.35), Inches(3.22), Inches(5.2), Inches(3.1), [
        ("보간은 모델이 채점되는 모든 지점에서 예측해야 한다.", 0),
        ("미래 이웃이 없으면 persistence로 후퇴한다.", 1),
        ("폴백률은 Tᵢ 채점 행의 0.3–0.4%이다.", 0, GREEN, True),
        ("V_rot 채점 행에서는 40–44%이다.", 0, RED, True),
        ("따라서 V_rot의 ‘vs PCHIP’은 5분의 2가 ‘vs persistence’이다.", 1, RED, True),
        ("사전등록은 이 폴백률의 보고를 의무화하였다.", 0),
    ], size=12.5, gap=8)
    return s


# --- 4-1. Architecture: the sequence backbone -----------------------------
def s_arch():
    s = slide()
    header(s, "4. 모델", "주 모델: 전체격자 인과 시퀀스 나우캐스터 seq_v2")
    add_image_fit(s, os.path.join(FIG, "fig_architecture_seq.png"),
                  Inches(0.45), Inches(1.38), Inches(12.45), Inches(5.15))
    text(s, Inches(0.55), Inches(6.62), Inches(12.3), Inches(0.42),
         [[("구조적 라우팅: ", 12.5, NAVY, True, False, None),
           ("Tᵢ 분기(2층 160)는 22채널 전체 상태를 읽고, V_rot 분기(1층 64)는 비-빠른 7채널만 읽는다. 총 파라미터는 357,570이다.",
            12.5, DARK, False, False, None)]], line_spacing=1.12)
    return s


# --- 4-2. Architecture detail --------------------------------------------
def s_arch_detail():
    s = slide()
    header(s, "4. 모델", "핵심 설계 결정과 근거")
    cards = [
        ("라우팅은 head가 아니라 인코더에서 한다", BLUE,
         ["순환 상태를 공유하면 head를 어떻게 배선해도",
          "빠른 진단 정보가 V_rot로 누출된다.",
          "seq_v2는 분기 자체를 분리하였다. 빠른 15채널을",
          "전부 섭동해도 V_rot 출력은 bit-identical이다."]),
        ("도달 범위는 세그먼트 전체이다", ORANGE,
         ["W는 더 이상 하이퍼파라미터가 아니다.",
          "라벨 없는 행도 맥락으로 유지한다(빠른 진단은 조밀하다).",
          "윈도 대조군 대비 pooled Tᵢ는 +0.081이다",
          "(16 run · CI [+0.067, +0.096] · 16/16 양수)."]),
        ("희소성은 loss가 처리한다", TEAL,
         ["세그먼트의 모든 라벨 행에 타깃별 masked MSE를 적용한다.",
          "LayerNorm과 작은 GELU head로 마무리한다.",
          "AdamW 1e-3, batch 16 세그먼트로 학습한다.",
          "val masked MSE로 조기 종료한다(patience 6, 상한 30)."]),
        ("해석 가능한 사다리 칸 b3k8", NAVY,
         ["seq_v2의 프레이밍·라우팅을 유지하고 head만 교체한다.",
          "예측 = 이월값 + Σ w·z + b (z ∈ [-1, 1], K = 8 / 4)이다.",
          "readout을 0으로 초기화하여 정확히 persistence에서 출발한다.",
          "21,498 파라미터로 백본의 6%이다(결과 ⑨)."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        r, c = divmod(i, 2)
        card(s, Inches(0.55 + c * 6.2), Inches(1.5 + r * 2.5), Inches(6.0), Inches(2.3),
             t, lines, accent=col, title_size=14.5, body_size=12)
    text(s, Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.6),
         [[("탐색의 교훈: ", 12, NAVY, True, False, None),
           ("윈도 계열 두 라운드에서 살아남은 메커니즘은 attention pooling 하나였고, 확정 프로토콜에서도 같은 교훈이 "
            "반복되었다(attention 후보 비유의, 폭 스윕 평평). 이것이 모델 절이 짧고 평가 절이 긴 이유이다.",
            12, GRAY, False, False, None)]], line_spacing=1.12)
    return s


# --- 4-3. Paired window control ------------------------------------------
# --- 4-2b. Physical grounding ---------------------------------------------
def s_physics():
    s = slide()
    header(s, "4. 모델", "물리적 근거: 두 갈래 라우팅은 설계 선택이 아니라 두 보존식의 결과이다")
    box(s, Inches(0.55), Inches(1.32), Inches(12.25), Inches(1.16), fill=NAVY, round_=True)
    text(s, Inches(0.85), Inches(1.42), Inches(11.7), Inches(1.00),
         [[("d/dt ( 3/2 · n_i · T_i )  =  Q_ei · (T_e - T_i)  -  ∇·q_i  + ...", 14, WHITE, True, False, None)],
          [("d/dt ( n_i · m_i · <R²> · ω_φ )  =  -∇·Π_φ  +  T_NBI  +  T_NTV  +  T_int  + ...", 14, ORANGE, True, False, None)],
          [("보존식은 값이 아니라 변화율을 구속한다. 따라서 물어야 할 것은 V_rot과의 상관이 아니라 dV_rot/dt와의 상관이다.",
            11, LGRAY, False, False, None)]], line_spacing=1.14, space_after=1)
    table(s, Inches(0.55), Inches(2.62), [Inches(1.70), Inches(2.45), Inches(1.65)],
          ["각운동량 수지의 항", "우리 데이터", "판정"],
          [
              ["LHS  dL/dt", "연속 관측쌍의 차분", ("측정됨", GREEN, True, None)],
              ["T_NBI", "0D 채널 전무", ("부재 §8b.3", RED, True, None)],
              ["T_NTV ~ δB²", "Mirnov 100 Hz 데시메이트", ("음성 §8b.2", RED, True, None)],
              ["∇·Π_turb", "BES 밀도요동", ("측정 → 널", ORANGE, True, None)],
              ["T_intrinsic", "∇T_i 필요, 스칼라뿐", ("도달 불가", GRAY, True, None)],
          ], row_h=Inches(0.44), head_h=Inches(0.40), size=11, head_size=11)
    text(s, Inches(0.60), Inches(5.36), Inches(5.70), Inches(1.55),
         [[("V_rot의 완화 시간은 이 데이터로 잴 수 없다", 12, RED, True, False, None)],
          [("관측의 54%가 계측기 유지값이라 held 제거 시 16 ms,", 10.5, DARK, False, False, None)],
          [("유지 시 300 ms 초과로 19배 이상 벌어진다. Tᵢ는 held가", 10.5, DARK, False, False, None)],
          [("641파일 중 1행뿐이라 159 ms가 깨끗하다. 유지값 병리는", 10.5, DARK, False, False, None)],
          [("값의 54%뿐 아니라 그 동역학을 특징지을 능력까지 파괴하며,", 10.5, DARK, False, False, None)],
          [("NBI 토크·원본 kHz Mirnov에 이은 세 번째 취득 과제이다.", 10.5, DARK, False, False, None)]],
         line_spacing=1.12, space_after=1)
    card(s, Inches(6.60), Inches(2.62), Inches(6.22), Inches(2.24),
         "① BES·ECEI와의 결합 (624 블록, lag 0)", [
             "→ Tᵢ   값 +0.341 / +0.311",
             "         변화율 +0.070 / +0.078",
             "→ V_rot   값 +0.027 / +0.005",
             "         변화율 -0.006 / -0.003",
             "에너지식에서는 신호를 잡으므로 널은 방법의 한계가 아니다.",
         ], accent=TEAL, title_size=13, body_size=11)
    card(s, Inches(6.60), Inches(4.96), Inches(6.22), Inches(1.94),
         "② 시간 척도가 맞는다", [
             "τ_eq = (mᵢ/2mₑ)·τ_e는 Tₑ 0.5–1 keV,",
             "nₑ 2–5×10¹⁹ m⁻³에서 8–59 ms이다.",
             "측정된 문맥 포화 약 50 ms와 같은 자릿수이다(§8al).",
             "Tᵢ·BES·ECEI의 완화 시간은 147–161 ms이다.",
             "주의: τ_eq 파라미터는 인용값이고 자기상관은 추세 미제거이다.",
         ], accent=BLUE, title_size=13, body_size=11)
    return s


def s_arch_window():
    s = slide()
    header(s, "4. 모델", "짝지은 대조군: W=2 윈도 모델(관측 마스킹 attention pooling)")
    add_image_fit(s, os.path.join(FIG, "fig_architecture.png"),
                  Inches(0.45), Inches(1.35), Inches(8.15), Inches(4.55))
    box(s, Inches(8.75), Inches(1.42), Inches(4.05), Inches(4.55), fill=CARDBG, round_=True)
    box(s, Inches(8.75), Inches(1.42), Inches(0.10), Inches(4.55), fill=TEAL)
    text(s, Inches(9.0), Inches(1.55), Inches(3.6), Inches(0.5),
         [[("대조군의 역할", 14.5, TEAL, True, False, None)]])
    bullets(s, Inches(9.0), Inches(2.05), Inches(3.6), Inches(3.8), [
        ("옛 주 모델이며 201,258 파라미터이다.", 0),
        ("모달리티별 time-aware 1-D CNN을 쓴다.", 1),
        ("이력은 양방향 GRU(64)로 부호화한다.", 1),
        ("관측 마스킹 attention pooling을 쓴다.", 0, ORANGE, True),
        ("해당 타깃이 실제 관측된 행에만 질량을 허용한다.", 1),
        ("보간의 귀납 편향을 파라미터 0개로 이식한다.", 1),
        ("나머지는 백본과 동일하게 고정하였다.", 0, NAVY, True),
        ("데이터 계약·held 처리·분할·채점 모집단이 같다.", 1),
        ("따라서 두 모델은 행 단위로 paired 비교된다.", 1, NAVY, True),
    ], size=12, gap=6)
    text(s, Inches(0.55), Inches(6.05), Inches(8.1), Inches(0.9),
         [[("물리 기반 head 라우팅: ", 12, NAVY, True, False, None),
           ("Tᵢ head는 빠른진단+시간+이력을, V_rot head는 이력+시간을 읽는다. 이 구조는 약 40회의 keep/discard 통제 실험의 결과이며, "
            "현재는 (i) 백본 관문의 비교 대상, (ii) 절제 실험의 무대, (iii) 캠페인 붕괴의 재현자로 쓰인다.",
            12, DARK, False, False, None)]], line_spacing=1.14)
    return s


# --- 4-4. Training --------------------------------------------------------
def s_training():
    s = slide()
    header(s, "4. 모델", "학습 절차: 손실 함수와 최적화 설정")
    card(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(3.05),
         "손실 — 타깃별 masked MSE",
         ["L = Σ m·(예측 - 실측)² / Σ m 를 Tᵢ·V_rot 각각에 계산한다.",
          "관측된 타깃만 손실에 기여하고 결측 타깃은 마스킹으로 제외한다.",
          "Tᵢ·V_rot는 독립 결측(8.2% / 23.9%)이므로 한쪽만 있어도 학습한다.",
          "시퀀스는 세그먼트의 모든 라벨 행에 대해 계산한다.",
          "라벨 없는 행은 맥락으로만 기여한다.",
          "출력은 정규화 단위이며 역정규화는 평가에서만 수행한다."],
         accent=BLUE, title_size=14.5, body_size=12.5)
    card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(3.05),
         "최적화 설정",
         ["백본 seq_v2는 AdamW 1e-3, batch 16 세그먼트로 학습하였다.",
          "val masked MSE로 조기 종료하였다(patience 6, 상한 30 epoch).",
          "확정 실행은 14–25 epoch에서 종료되었다.",
          "윈도 대조군은 같은 데이터 계약과 파일당 샘플 상한 500을 쓴다.",
          "분할 seed와 초기화 seed를 분리하였다(분할 4 × 초기화 4).",
          "loss가 비유한이면 즉시 실패하도록 하였다(조용한 NaN 금지)."],
         accent=TEAL, title_size=14.5, body_size=12.5)
    box(s, Inches(0.55), Inches(4.80), Inches(12.25), Inches(1.6), fill=CARDBG, round_=True)
    text(s, Inches(0.8), Inches(4.95), Inches(11.8), Inches(0.5),
         [[("아키텍처의 출처 — 깨끗한 val skill을 게이트로 쓴 keep/discard 통제 루프", 14, NAVY, True, False, None)]])
    bullets(s, Inches(0.8), Inches(5.43), Inches(11.8), Inches(0.9), [
        ("반복마다 구조 변경을 하나만 적용하고 처음부터 재학습한 뒤, 증강 없는 검증셋의 skill_vs_pchip으로 채점하였다.", 0),
        ("증강 val loss는 보간이 이미 강한 곳에서 평활화를 보상하므로 쓰지 않았으며, TEST는 전 과정에서 봉인되었다.", 1, NAVY, True),
    ], size=12.5, gap=6)
    return s


# --- 5-1. Methodology -----------------------------------------------------
def s_methodology():
    s = slide()
    header(s, "5. 평가 방법론", "선택 편향 없는 3-way split과 사전등록")
    cards = [
        ("3-way split · TEST 동결", BLUE,
         ["TEST는 아키텍처 탐색 전에 예약되었고 선택 중에는 열리지 않았다.",
          "모델 선택은 val에서만 이루어져 headline에 winner's curse가 없다.",
          "전 arm이 동일한 (file, row) 집합과 마스크로 채점된다.",
          "짝지은 비교 전에 모집단 키가 bit-identical임을 검증한다.",
          "보간은 타깃 자신의 값을 제외하고 이웃만 읽는다."]),
        ("사전등록 (PR1–PR4 + 확정 프로토콜)", ORANGE,
         ["PR1: headline 기준선은 PCHIP이며 사다리 전체를 함께 보고한다.",
          "PR2: 보간은 모든 채점 지점에서 예측하고 폴백률을 보고한다.",
          "PR3: TEST 하한은 15 방전·3,000 Tᵢ 샘플 이상이다(충족).",
          "PR4: shot 군집 bootstrap 95% CI가 0을 제외하면 PASS이다.",
          "추가: held-free · W=2 · 파일당 500 · 두 모집단 ·",
          "TEST 채점 전 결정 규칙 커밋 · 문턱 민감도이다."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(0.55 + i * 6.2), Inches(1.5), Inches(6.0), Inches(2.5),
             t, lines, accent=col, title_size=15, body_size=12)
    cw = [Inches(3.05), Inches(5.4), Inches(3.8)]
    rows = [
        ["Persistence", "마지막 관측 CES", ("인과", GRAY, True, None)],
        ["AR (국소)", "과거 CES만", ("인과", GRAY, True, None)],
        ["인과 GP", "과거 CES 이웃 16개 (NaN 조건 동일)", ("인과 · 최강 배치 기준선", TEAL, True, None)],
        ["선형 / PCHIP*", "타깃 양쪽의 과거 + 미래 CES 이웃", ("오프라인", BLUE, True, None)],
        ["GP (오프라인)", "Matérn-3/2 + 백색잡음, 이웃 16+16", ("오프라인 · 최강 평활기", BLUE, True, None)],
    ]
    table(s, Inches(0.55), Inches(4.15), cw, ["기준선 사다리", "정보 접근", "분류"], rows,
          row_h=Inches(0.40), head_h=Inches(0.40), size=12.5, head_size=12.5)
    text(s, Inches(0.55), Inches(6.60), Inches(12.3), Inches(0.42),
         [[("* PR1 headline. ", 11.5, GRAY, True, False, None),
           ("보간과 모델 모두 세그먼트 경계를 넘지 않으며, 경계 밖 이웃이 필요하면 보간은 persistence 값을 대신 예측한다(모집단 축소 없음).",
            11.5, GRAY, False, False, None)]])
    return s


# --- 5-2. Bootstrap -------------------------------------------------------
def s_bootstrap():
    s = slide()
    header(s, "5. 평가 방법론", "shot 군집 paired bootstrap: 재현 단위를 shot으로 두었다")
    bullets(s, Inches(0.55), Inches(1.55), Inches(7.0), Inches(3.6), [
        ("한 방전(shot) 내 인접 CES 행은 강하게 상관된다.", 0),
        ("개별 샘플을 독립으로 보면 불확실성이 크게 과소평가된다.", 1, RED, True),
        ("PR4 검정은 샘플별 짝지은 오차 (SE_model - SE_pchip)를", 0),
        ("shot 단위로 묶고 shot 전체를 복원추출한다(B = 10,000, 고정 seed).", 1),
        ("skill 95% CI가 모델에 유리한 방향으로 0을 제외하면 PASS이다.", 0, GREEN, True),
        ("따라서 CI는 within-shot 복제가 아니라 shot 간 일반화를 반영한다.", 0),
        ("모델 대 모델 비교도 같은 행 위에서 같은 paired bootstrap으로 한다.", 0, NAVY, True),
        ("백본 vs 윈도 대조군, 사다리 칸 vs 백본이 그 예이다.", 1),
    ], size=14, gap=10)
    box(s, Inches(7.85), Inches(1.6), Inches(4.95), Inches(4.5), fill=NAVY, round_=True)
    text(s, Inches(8.1), Inches(1.8), Inches(4.5), Inches(0.5),
         [[("재현 단위 = shot", 16, ORANGE, True, False, None)]])
    text(s, Inches(8.1), Inches(2.45), Inches(4.5), Inches(3.5),
         [[("개별 행", 13, LGRAY, True, False, None)],
          [("자기상관으로 인한 가짜 복제이다 (X)", 12.5, MGRAY, False, False, None)],
          [("", 6, WHITE, False, False, None)],
          [("shot 전체", 13, WHITE, True, False, None)],
          [("독립 복제의 단위이다 (O)", 12.5, RGBColor(0x9D, 0xE8, 0xCD), False, False, None)],
          [("", 6, WHITE, False, False, None)],
          [("Tᵢ ≈ 96 방전 · V_rot 60–66 방전", 13, WHITE, True, False, None)],
          [("이 방전 수가 검정력의 상한이며", 12, LGRAY, False, False, None)],
          [("모든 유의성 판정의 구속 조건이다.", 12, LGRAY, False, False, None)]],
         line_spacing=1.18, space_after=3)
    return s


# --- 5-3. Model selection protocol ---------------------------------------
def s_res_protocol():
    s = slide()
    header(s, "5. 평가 방법론", "모델 선택 프로토콜: 규칙을 수치보다 먼저 기록하였다", accent=TEAL)
    text(s, Inches(0.55), Inches(1.4), Inches(12.3), Inches(0.55),
         [[("모든 모델 결정은 검증 데이터 위에서, 또는 ", 14, DARK, False, False, None),
           ("해당 TEST 채점 이전에 문서로 확정된 결정 규칙", 14, TEAL, True, False, None),
           (" 아래에서 이루어졌고, TEST는 결정마다 한 번만 채점되었다.", 14, DARK, False, False, None)]])
    card(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(2.5),
         "백본 관문 — 4조건을 먼저 고정하였다",
         ["① 4개 분할 전부에서 부호가 유지되어야 한다.",
          "② 통합 실행 군집 CI가 0을 제외해야 한다.",
          "③ 예산 균등화(고정 10 epoch)에서도 부호가 유지되어야 한다.",
          "④ V_rot 손실이 없어야 한다.",
          "네 조건이 충족되어 시퀀스 프레이밍을 백본으로 채택하였다(결과 ③)."],
         accent=NAVY, title_size=14.5, body_size=12.5)
    card(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(2.5),
         "유일한 아키텍처 후보 — 미승격",
         ["seq_v2 + 관측 마스킹 인과 attention(0 초기화 사영)이었다.",
          "백본 대비 4/4 분할에서 양수였다:",
          "+0.009 / +0.013 / +0.033 / +0.020.",
          "그러나 유의는 1/4로 사전 기준 3/4 이상에 미달하여 승격하지 않았다.",
          "val에서는 2/2 유의였으며, 이것이 승격 bar를 TEST에 두는 이유이다."],
         accent=ORANGE, title_size=14.5, body_size=12.5)
    box(s, Inches(0.55), Inches(4.75), Inches(12.25), Inches(1.95), fill=CARDBG, round_=True)
    text(s, Inches(0.8), Inches(4.88), Inches(11.8), Inches(0.5),
         [[("나머지 결정 규칙", 14, NAVY, True, False, None)]])
    bullets(s, Inches(0.8), Inches(5.36), Inches(11.8), Inches(1.3), [
        ("윈도 계열은 한 번에 하나만 바꾸고 깨끗한 val skill로 keep/discard하였다. 이력 길이는 24-run 스윕(plateau 최소)으로, held는 감사로 제거하였다.", 0),
        ("사다리 칸·폭 스윕은 두 갈래 판정과 서술적 읽기(상한 / 무릎)를 TEST 채점 전에 문서화하였다.", 0),
        ("스윕 위에서 백본을 재선정하는 것은 구성상 금지된다.", 1, RED, True),
    ], size=12.5, gap=6)
    return s


# --- 6-1. RMSE ladder ------------------------------------------------------
def s_res_ladder():
    s = slide()
    header(s, "6. 결과 ①", "백본은 인과 기준선을 압도하고 오프라인 평활기와 동률이다")
    add_image_fit(s, os.path.join(FIG, "fig_rmse_ladder.png"),
                  Inches(0.55), Inches(1.42), Inches(12.25), Inches(4.35))
    box(s, Inches(0.55), Inches(5.85), Inches(12.25), Inches(1.1), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.95), Inches(11.7), Inches(0.95),
         [[("백본이 두 타깃 모두에서 최저 RMSE를 기록하였다. ", 13, NAVY, True, False, None),
           ("인과 GP보다 Tᵢ 4%, V_rot 18% 낮고, 미래를 읽는 오프라인 GP(153.8)와는 동률(157.8)이다.",
            13, DARK, False, False, None)],
          [("포함 모집단: ", 12.5, NAVY, True, False, None),
           ("seq_v2 363.0 / 23.7, PCHIP 412.4 / 30.2, 인과 GP 394.6 / 28.8, persistence 478.0 / 33.4이다. "
            "스파이크가 Tᵢ RMSE를 두 배 이상 키우지만 순서는 바뀌지 않았다.", 12.5, DARK, False, False, None)]],
         line_spacing=1.14)
    return s


# --- 6-2. Headline forest -------------------------------------------------
def s_res_forest():
    s = slide()
    header(s, "6. 결과 ②", "Headline: Tᵢ는 두 모집단 모두에서 4/4를 기록하였다")
    add_image_fit(s, os.path.join(FIG, "fig_forest.png"),
                  Inches(0.55), Inches(1.38), Inches(12.25), Inches(3.85))
    cards = [
        ("CES_TI — 무조건부 PASS", GREEN,
         ["컷 +0.174 / +0.248 / +0.257 / +0.264(평균 +0.236)로 4/4이다.",
          "포함 +0.225 / +0.238 / +0.292 / +0.316(평균 +0.268)로 4/4이다.",
          "8개 셀 전부에서 인과 GP(+0.08~+0.17)와 persistence(+0.36~+0.46)도 이겼다.",
          "오프라인 GP와는 동률(-0.05~+0.11, 1/8 유의)이며 이를 상한으로 명시한다."]),
        ("CES_VT — 동률 보고", GRAY,
         ["점추정은 8/8 양수이나 PR4는 컷 1/4, 포함 2/4로 잡음 수준이다.",
          "persistence 대비 3/4 양쪽(+0.30~+0.50), 인과 GP 대비 2/4이다.",
          "포함 수치가 높은 것은 스파이크가 보간 앵커를 오염시켜",
          "모든 arm이 PCHIP 대비 좋아 보이기 때문이며 결과 ⑨에서 분해한다."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        card(s, Inches(0.55 + i * 6.2), Inches(5.22), Inches(6.0), Inches(1.72),
             t, lines, accent=col, title_size=13.5, body_size=10.5)
    return s


# --- 6-3. Backbone gate ---------------------------------------------------
def s_res_gate():
    s = slide()
    header(s, "6. 결과 ③", "백본 관문(B.1): 전체격자 프레이밍은 윈도 대조군을 16/16에서 이겼다")
    add_image_fit(s, os.path.join(FIG, "fig_gate_b1.png"),
                  Inches(0.45), Inches(1.40), Inches(8.05), Inches(4.3))
    box(s, Inches(8.65), Inches(1.42), Inches(4.15), Inches(4.35), fill=CARDBG, round_=True)
    box(s, Inches(8.65), Inches(1.42), Inches(0.10), Inches(4.35), fill=NAVY)
    text(s, Inches(8.9), Inches(1.55), Inches(3.7), Inches(0.5),
         [[("16 run paired 결과 (컷)", 14.5, NAVY, True, False, None)]])
    bullets(s, Inches(8.9), Inches(2.05), Inches(3.7), Inches(3.6), [
        ("Tᵢ는 16/16 양수, 13/16 유의였다.", 0, GREEN, True),
        ("분할별 초기화 평균은 다음과 같다.", 0),
        ("+0.129 / +0.059 / +0.078 / +0.058", 1),
        ("pooled 추정치는 +0.081이다.", 0, GREEN, True),
        ("run 군집 CI는 [+0.067, +0.096]이다.", 1),
        ("예산 균등화에서도 4/4 부호가 유지되었다.", 0),
        ("+0.063 / +0.033 / +0.045 / +0.030", 1),
        ("V_rot 유의 열세는 0/16이었다(우세 8/16).", 0, NAVY, True),
    ], size=12, gap=7)
    box(s, Inches(0.55), Inches(5.85), Inches(12.25), Inches(1.1), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.95), Inches(11.7), Inches(0.95),
         [[("확증 4 분할의 같은 비교(seq - 윈도): ", 12.5, NAVY, True, False, None),
           ("컷 +0.130 / +0.058 / +0.062 / +0.044, 포함 +0.053 / +0.024 / +0.047 / +0.029으로 8/8 양수(각 2/4 유의)였다.",
            12.5, DARK, False, False, None)],
          [("의미: ", 12.5, ORANGE, True, False, None),
           ("윈도 대조군은 인과 GP와 동률(1/4)이나 시퀀스 백본은 4/4+4/4이다. 세그먼트 과거 전체로의 도달 범위가 "
            "최강 배치 기준선을 이기게 하며, 윈도 조립과 조합 증강이 없어 학습 비용은 1/10이다.",
            12.5, DARK, False, False, None)]], line_spacing=1.14)
    return s


# --- 6-4. Gap-stratified --------------------------------------------------
def s_res_gap():
    s = slide()
    header(s, "6. 결과 ④", "간극 영역: 비인접 시점에서도 우위가 유지되었다(4 분할 통합)")
    text(s, Inches(0.55), Inches(1.38), Inches(12.3), Inches(0.55),
         [[("분할별로는 넓은 Δt 층의 표본이 수십 개뿐이므로 4개 test 분할을 합치고 방전 단위로 군집화하였다. "
            "PCHIP은 미래 앵커를 읽고 모델은 읽지 않는다.", 13, DARK, False, False, None)]], line_spacing=1.12)
    cw = [Inches(2.35), Inches(2.45), Inches(1.2), Inches(3.15), Inches(3.1)]
    rows = [
        [("Tᵢ  ≤ 15 ms", DARK, True, None), ("134,546 / 135,317", GRAY, False, MONO), ("301", GRAY, False, MONO),
         ("+0.239 [+0.197, +0.274]", GREEN, True, MONO), ("+0.299 [+0.244, +0.347]", GREEN, True, MONO)],
        [("Tᵢ  > 15 ms", DARK, True, None), ("3,422 / 3,334", GRAY, False, MONO), ("265 / 263", GRAY, False, MONO),
         ("+0.268 [+0.187, +0.337]", GREEN, True, MONO), ("+0.206 [+0.108, +0.290]", GREEN, True, MONO)],
        [("Tᵢ  > 45 ms", DARK, True, None), ("460 / 429", GRAY, False, MONO), ("104 / 101", GRAY, False, MONO),
         ("+0.267 [+0.092, +0.414]", GREEN, True, MONO), ("-0.004 [-0.304, +0.246]", RED, True, MONO)],
        [("V_rot  ≤ 15 ms", DARK, True, None), ("51,689", GRAY, False, MONO), ("197", GRAY, False, MONO),
         ("+0.233 [+0.020, +0.318]", GREEN, True, MONO), ("+0.233 [+0.020, +0.317]", GREEN, True, MONO)],
        [("V_rot  > 15 ms", DARK, True, None), ("466 / 456", GRAY, False, MONO), ("130", GRAY, False, MONO),
         ("+0.418 [+0.104, +0.680]", GREEN, True, MONO), ("+0.432 [+0.128, +0.696]", GREEN, True, MONO)],
        [("V_rot  > 45 ms", DARK, True, None), ("14", GRAY, False, MONO), ("7", GRAY, False, MONO),
         ("미채점", MGRAY, False, None), ("미채점", MGRAY, False, None)],
    ]
    table(s, Inches(0.55), Inches(2.02), cw,
          ["Δt 구간", "n (컷 / 포함)", "방전", "컷: vs PCHIP (95% CI)", "포함: vs PCHIP (95% CI)"],
          rows, row_h=Inches(0.46), head_h=Inches(0.44), size=11.5, head_size=12)
    box(s, Inches(0.55), Inches(5.42), Inches(12.25), Inches(1.45), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.54), Inches(11.7), Inches(1.3),
         [[("세 가지 읽기", 13, NAVY, True, False, None)],
          [("① Tᵢ 우위는 인접 이력에 국한되지 않는다. > 15 ms에서 두 모집단 모두 PASS이다(persistence 대비 +0.40 / +0.43).",
            12.5, DARK, False, False, None)],
          [("② > 45 ms는 컷에서 승, 포함에서 동률이다. 429행/101 방전은 스파이크 앵커 몇 행이 층을 지배할 수 있는 규모이다.",
            12.5, DARK, False, False, None)],
          [("③ V_rot는 보간이 가장 어려운 > 15 ms에서 두 모집단 모두 PCHIP를 이겼다. 이것이 유일한 V_rot 무조건부 양성이다.",
            12.5, TEAL, True, False, None)]], line_spacing=1.12, space_after=2)
    return s


# --- 6-5. MNAR reweighting ------------------------------------------------
def s_stress():
    s = slide()
    header(s, "6. 결과 ⑤", "스트레스 ①: 실제 결측점 분포로 재가중(MNAR)하여도 인과 대비 우위는 8/8 생존하였다", accent=ORANGE)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.55),
         [[("층은 Δt(15/25/45 ms) × 입력만의 활동 flag이다. 결측 행의 층 분포로 채점 지점을 재가중하였다(30 미만 층은 기각, "
            "가중 격자에도 컷 적용). 도달 범위는 결측 Tᵢ의 ", 12.5, DARK, False, False, None),
           ("54–68%", 12.5, NAVY, True, False, None),
           (", 결측 V_rot의 ", 12.5, DARK, False, False, None),
           ("4–6%", 12.5, RED, True, False, None),
           ("만 in-domain이므로 재가중 V_rot은 결론을 내지 않는다.", 12.5, DARK, False, False, None)]], line_spacing=1.12)
    cw = [Inches(1.25), Inches(1.0), Inches(1.7), Inches(4.15), Inches(4.15)]
    rows = [
        [("컷", NAVY, True, None), "42", ("+0.174", GRAY, False, MONO),
         ("+0.140 [-0.071, +0.290]", GRAY, False, MONO), ("+0.398 [+0.278, +0.518]", GREEN, True, MONO)],
        [("컷", NAVY, True, None), "1", ("+0.248", GRAY, False, MONO),
         ("+0.164 [+0.062, +0.250]", GREEN, True, MONO), ("+0.366 [+0.299, +0.448]", GREEN, True, MONO)],
        [("컷", NAVY, True, None), "7", ("+0.257", GRAY, False, MONO),
         ("+0.203 [-0.024, +0.346]", GRAY, False, MONO), ("+0.310 [+0.227, +0.398]", GREEN, True, MONO)],
        [("컷", NAVY, True, None), "123", ("+0.264", GRAY, False, MONO),
         ("+0.283 [+0.069, +0.392]", GREEN, True, MONO), ("+0.381 [+0.246, +0.464]", GREEN, True, MONO)],
        [("포함", ORANGE, True, None), "42", ("+0.225", GRAY, False, MONO),
         ("+0.140 [+0.033, +0.250]", GREEN, True, MONO), ("+0.443 [+0.141, +0.623]", GREEN, True, MONO)],
        [("포함", ORANGE, True, None), "1", ("+0.238", GRAY, False, MONO),
         ("+0.217 [+0.086, +0.319]", GREEN, True, MONO), ("+0.383 [+0.273, +0.536]", GREEN, True, MONO)],
        [("포함", ORANGE, True, None), "7", ("+0.292", GRAY, False, MONO),
         ("+0.167 [+0.067, +0.265]", GREEN, True, MONO), ("+0.283 [+0.172, +0.380]", GREEN, True, MONO)],
        [("포함", ORANGE, True, None), "123", ("+0.316", GRAY, False, MONO),
         ("+0.221 [+0.050, +0.320]", GREEN, True, MONO), ("+0.337 [+0.164, +0.455]", GREEN, True, MONO)],
    ]
    table(s, Inches(0.55), Inches(2.02), cw,
          ["모집단", "분할", "비가중 vs PCHIP", "결측 정합: vs PCHIP", "결측 정합: vs persistence"],
          rows, row_h=Inches(0.36), head_h=Inches(0.40), size=11.5, head_size=11.5)
    box(s, Inches(0.55), Inches(5.42), Inches(12.25), Inches(1.45), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.54), Inches(11.7), Inches(1.3),
         [[("Tᵢ 우위는 온라인 시스템의 실제 상대(persistence)에 대해 두 모집단 4/4 생존하였다", 13, GREEN, True, False, None),
           ("(+0.28~+0.44, MNAR 보정 비용 최대 0.12).", 13, DARK, False, False, None)],
          [("PCHIP 대비 점추정은 +0.14~+0.28로 유지되었으나 고정 가중 CI가 컷 2개 분할에서 0을 지났다(컷 2/4 · 포함 4/4).",
            12.5, DARK, False, False, None)],
          [("진술: 실제 결측·in-domain 시점에서 나우캐스터는 모든 인과 CES-only 방법보다 유의하게 낫고, "
            "오프라인 보간보다는 모집단 조건부로 낫다.", 12.5, NAVY, True, False, None)]],
         line_spacing=1.12, space_after=2)
    return s


# --- 6-6. Campaign split --------------------------------------------------
def s_res_campaign():
    s = slide()
    header(s, "6. 결과 ⑥", "스트레스 ②: 캠페인(시간) 분할에서 윈도는 붕괴하고 백본은 4/4+4/4로 생존하였다", accent=ORANGE)
    add_image_fit(s, os.path.join(FIG, "fig_campaign.png"),
                  Inches(0.45), Inches(1.36), Inches(7.85), Inches(3.25))
    box(s, Inches(8.45), Inches(1.36), Inches(4.35), Inches(3.3), fill=CARDBG, round_=True)
    box(s, Inches(8.45), Inches(1.36), Inches(0.10), Inches(3.3), fill=ORANGE)
    text(s, Inches(8.7), Inches(1.47), Inches(3.9), Inches(0.5),
         [[("설계와 수치 (Tᵢ vs PCHIP)", 13.5, ORANGE, True, False, None)]])
    bullets(s, Inches(8.7), Inches(1.94), Inches(3.9), Inches(2.7), [
        ("train 416 (30801–31991)으로 학습하였다.", 0),
        ("val 128 (32002–32310) · test 97 (32312–32751)이다.", 1),
        ("초기화 seed 4개이며 분할 4개가 아니다.", 1, RED, True),
        ("윈도 OFF: 컷 2/4 · 포함 0/4 · 인과 GP 0/4로 붕괴하였다.", 0, RED, True),
        ("윈도 ON(shot별 표준화, 컷)은 4/4로 회복되었다.", 0, TEAL, True),
        ("seq_v2 컷: +0.187/+0.174/+0.181/+0.177이다.", 0, GREEN, True),
        ("포함: +0.173/+0.202/+0.198/+0.184로 4/4+4/4이다.", 1, GREEN, True),
    ], size=11.5, gap=6)
    cw = [Inches(4.35), Inches(4.0), Inches(3.9)]
    rows = [
        [("무작위 분할 · 관측점 (결과 ②)", DARK, True, None),
         ("4/4 · 4/4,  +0.17~+0.32", GREEN, True, None), ("4/4 · 4/4 vs 인과 GP,  +0.08~+0.17", GREEN, True, None)],
        [("결측점 재가중 (결과 ⑤)", DARK, True, None),
         ("2/4 · 4/4,  +0.14~+0.28", ORANGE, True, None), ("4/4 · 4/4 vs persistence,  +0.28~+0.44", GREEN, True, None)],
        [("캠페인 시간 분할 (이 슬라이드)", DARK, True, None),
         ("4/4 · 4/4,  +0.17~+0.20", GREEN, True, None), ("4/4 · 4/4 vs 인과 GP,  +0.11~+0.16", GREEN, True, None)],
    ]
    table(s, Inches(0.55), Inches(4.80), cw,
          ["Tᵢ 평가 (컷 · 포함)", "vs PCHIP (오프라인)", "vs 인과 기준선"], rows,
          row_h=Inches(0.44), head_h=Inches(0.42), size=12, head_size=12)
    text(s, Inches(0.55), Inches(6.52), Inches(12.3), Inches(0.6),
         [[("원인은 측정되었다: ", 12, NAVY, True, False, None),
           ("train→test 드리프트는 BES 1.22σ, ECEI 0.53σ인 반면 타깃은 0.115σ였다. seq-윈도 마진은 8/8 유의였고 "
            "V_rot도 persistence 대비 seq 4/4 양쪽(윈도 0/4)이었다. 남는 주의점은 한 시간 블록 위의 초기화 4개와 컷 run 2/4의 상한 종료이다.",
            12, GRAY, False, False, None)]], line_spacing=1.12)
    return s


# --- 6-7. Asymmetry + ablation --------------------------------------------
def s_res_asym():
    s = slide()
    header(s, "6. 결과 ⑦", "Tᵢ ↔ V_rot 정보 비대칭이 절제 실험으로 확인되었다")
    add_image_fit(s, os.path.join(FIG, "fig_ablation.png"),
                  Inches(0.45), Inches(1.42), Inches(7.4), Inches(5.2))
    box(s, Inches(7.95), Inches(1.5), Inches(4.85), Inches(5.05), fill=CARDBG, round_=True)
    text(s, Inches(8.2), Inches(1.64), Inches(4.4), Inches(0.5),
         [[("무엇을 지우면 무엇이 사라지는가", 14.5, NAVY, True, False, None)]])
    bullets(s, Inches(8.2), Inches(2.15), Inches(4.4), Inches(4.3), [
        ("이력은 두 타깃 모두에 필수적이다.", 0, NAVY, True),
        ("no_history: Tᵢ -2.11 / -1.16, V_rot -2.89 / -3.51이다.", 1),
        ("Tᵢ: 컷 모집단의 마진은 빠른 진단 정보이다.", 0, ORANGE, True),
        ("no_fast는 -0.125(paired 4/4 유의 감소)이다.", 1),
        ("물리 채널은 충돌 e–i 결합(ECEI Tₑ · BES nₑ)이다.", 1),
        ("포함 모집단에는 스파이크 강건성 성분이 섞인다.", 0, RED, True),
        ("이력 전용도 PCHIP를 +0.15~+0.23 이겼다.", 1),
        ("빠른 채널이 더하는 몫은 0.03–0.09이다.", 1),
        ("V_rot: 정보는 전부 CES 이력에서 온다.", 0, BLUE, True),
        ("빠른 채널을 0으로 두면 출력이 bit-identical이다(8/8).", 1),
        ("NBI 토크 미관측과 Mirnov 100 Hz 앨리어싱이 원인이다.", 1),
        ("V_rot의 비-승리는 실패가 아니라 발견이다.", 0, NAVY, True),
    ], size=11.5, gap=5)
    return s


# --- 6-8. Window sweep ----------------------------------------------------
def s_window_sweep():
    s = slide()
    header(s, "6. 결과 ⑧", "이력은 관측 하나로 충분하다: 윈도 스윕이 W=2의 근거이다", accent=TEAL)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.56),
         [[("W ∈ {2,3,4,6,8} × seed 4개와 history-0으로 ", 13, DARK, False, False, None),
           ("독립 run 24회", 13, NAVY, True, False, None),
           ("를 수행하였다. held-free 학습·평가, 파일당 500, 컷 없음(이 스윕의 동결 W=2 run이 포함 모집단 대조군)이며, "
            "각 run은 자기 자신의 held-out TEST skill_vs_pchip으로 채점하였다.", 13, DARK, False, False, None)]],
         line_spacing=1.12)
    add_image_fit(s, os.path.join(FIG, "fig_window_sweep.png"),
                  Inches(0.45), Inches(1.98), Inches(7.75), Inches(4.15))
    box(s, Inches(8.35), Inches(1.98), Inches(4.45), Inches(4.15), fill=CARDBG, round_=True)
    text(s, Inches(8.6), Inches(2.12), Inches(4.0), Inches(0.5),
         [[("곡선이 보여주는 세 가지", 15, NAVY, True, False, None)]])
    bullets(s, Inches(8.6), Inches(2.66), Inches(4.0), Inches(3.4), [
        ("history가 없으면 무너진다.", 0, RED, True),
        ("Tᵢ -0.026 (0/4), V_rot -0.783이다.", 1),
        ("과거 관측 1개가 skill의 전부를 만든다.", 0, GREEN, True),
        ("W=2에서 Tᵢ +0.238 (4/4), V_rot +0.206이다.", 1),
        ("이후 곡선은 평평하다(Tᵢ 0.190–0.246).", 1),
        ("점 내부 seed 산포 0.07–0.16이 곡선 전체보다 크다.", 1),
        ("따라서 plateau 최소 W는 2이다.", 0, TEAL, True),
        ("넓은 W의 유일한 논거는 커버리지인데,", 1),
        ("그것은 시퀀스 프레이밍의 논거이다.", 1),
    ], size=12, gap=7)
    text(s, Inches(0.55), Inches(6.22), Inches(12.3), Inches(0.62),
         [[("결론  ", 12.5, TEAL, True, False, None),
           ("W=3 +0.246(4/4), W=4 +0.221(3/4), W=6 +0.190(3/4), W=8 +0.216(4/4)으로 곡선은 평평하다. "
            "이력을 늘려 얻는 것은 skill이 아니라 > 15 ms 채점 커버리지(456→1,958)이며, 시퀀스 프레이밍이 W 없이 해결한다. "
            "7장은 ‘관측 수’와 별개의 자원인 ‘연속 문맥 길이’를 측정한다.", 12, DARK, False, False, None)]], line_spacing=1.12)
    return s


# --- 6-9. Scaling ---------------------------------------------------------
def s_res_scaling():
    s = slide()
    header(s, "6. 결과 ⑨", "복잡도 사다리와 크기 축: 상한은 추정기가 아니라 정보에 있다")
    add_image_fit(s, os.path.join(FIG, "fig_ladder_scaling.png"),
                  Inches(0.45), Inches(1.36), Inches(7.55), Inches(3.35))
    box(s, Inches(8.15), Inches(1.36), Inches(4.65), Inches(3.35), fill=CARDBG, round_=True)
    box(s, Inches(8.15), Inches(1.36), Inches(0.10), Inches(3.35), fill=NAVY)
    text(s, Inches(8.4), Inches(1.47), Inches(4.2), Inches(0.5),
         [[("Tᵢ skill (컷 / 포함)", 13.5, NAVY, True, False, None)]])
    rowsx = [
        ("Persistence", "0", "-0.264 / -0.288", GRAY),
        ("Anchor+Δ", "1,258", "-0.261 / -0.287", GRAY),
        ("b3k8", "21,498", "+0.237 / +0.126", GREEN),
        ("윈도 대조군", "201,258", "+0.173 / +0.238", DARK),
        ("seq_v2 백본", "357,570", "+0.236 / +0.268", NAVY),
    ]
    yy = 2.02
    for name, par, val, col in rowsx:
        text(s, Inches(8.4), Inches(yy), Inches(1.65), Inches(0.36),
             [[(name, 12, col, True, False, None)]])
        text(s, Inches(10.0), Inches(yy), Inches(1.0), Inches(0.36),
             [[(par, 11.5, GRAY, False, False, MONO)]], align=PP_ALIGN.RIGHT)
        text(s, Inches(11.05), Inches(yy), Inches(1.65), Inches(0.36),
             [[(val, 11.5, col, True, False, MONO)]], align=PP_ALIGN.RIGHT)
        yy += 0.42
    text(s, Inches(8.4), Inches(yy + 0.06), Inches(4.2), Inches(0.8),
         [[("b3 - anchor는 컷 +0.35~+0.42, 포함 +0.29~+0.34로 양쪽 4/4 유의이다.",
            11, DARK, False, False, None)]], line_spacing=1.12)
    card(s, Inches(0.55), Inches(4.85), Inches(6.0), Inches(2.05),
         "컷: 백본의 Tᵢ skill 전부가 유계 수 8개와 persistence로 압축된다", [
             "b3 - seq_v2 평균 +0.002(CI 전부 0 포함), PR4 4/4, 인과 GP 4/4이다.",
             "probe 결과 latent은 직전 Tᵢ(R² 0.47–0.75)와 ECEI Tₑ(0.31–0.48)를",
             "분산 부호화하며 보정은 예측 분산의 25–39%이다.",
             "포함에서는 -0.194(4/4 유의)이다. 유계 보정으로는 스파이크 앵커를",
             "살리지 못하며, ≈1% 행이 모든 arm SSE의 70–83%를 차지한다.",
         ], accent=TEAL, title_size=13, body_size=11.5)
    card(s, Inches(6.8), Inches(4.85), Inches(6.0), Inches(2.05),
         "크기 축은 닫혔다", [
             "Tᵢ 인코더 폭 24 / 40 / 80 / 160 / 260은",
             "34k / 49k / 114k / 358k / 879k 파라미터에 해당한다.",
             "skill은 +0.230 / +0.236 / +0.235 / +0.236 / +0.230(컷)이다.",
             "160 대비 ±0.008이며 최대 폭의 유의 우세는 1/4, V_rot는 불변이다.",
             "따라서 남은 분산은 모델 크기가 아니라 분할 분산이다.",
         ], accent=ORANGE, title_size=13, body_size=11.5)
    return s


# --- 6-10. Reproducibility ceiling ---------------------------------------
def s_noise_floor():
    s = slide()
    header(s, "6. 결과 ⑩", "재현성 상한: 모델 오차는 타깃 자신의 10 ms 재현성과 같은 크기이다")
    text(s, Inches(0.55), Inches(1.34), Inches(12.25), Inches(0.34),
         [[("타깃은 광자 적분 스펙트럼 피팅이므로 스스로 잡음을 지닌다. 차분 기반 추정기로 그 재현성을 재면 상한이 상대 비교가 아닌 절대 단위로 표현된다.",
            12.5, GRAY, False, False, None)]])
    col_w = [Inches(3.30), Inches(1.75), Inches(1.85)]
    table(s, Inches(0.55), Inches(1.80), col_w,
          ["재현성 추정기 (전부 상한)", "Tᵢ (eV)", "V_rot (km/s)"],
          [
              ["1차 차분 (Rice)", "151.8", "28.8"],
              ["2차 차분 (GSJS)", "139.9", "22.8"],
              ["3차 차분", "133.8", "19.5"],
              ["4차 차분", ("129.9", NAVY, True, MONO), ("16.4", NAVY, True, MONO)],
              ["4차 · 상위 5% 절사", "68.7", "6.8"],
              ["1차 · MAD (본류)", ("46.4", TEAL, True, MONO), ("4.3", TEAL, True, MONO)],
              [("백본 seq_v2의 RMSE", NAVY, True, None),
               ("157.8", ORANGE, True, MONO), ("23.6", ORANGE, True, MONO)],
          ], row_h=Inches(0.47), head_h=Inches(0.44), size=12, head_size=12,
          emphasis={6}, emphasis_fill=LGRAY)
    text(s, Inches(0.55), Inches(5.55), Inches(6.9), Inches(0.85),
         [[("컷 모집단, 641파일 전수, held 제외·스파이크 컷 후 연속 관측 구간에서 측정하였다. 차수가 오를수록 신호 편향이 벗겨지며, 4차에서도 아직 하강 중이므로 129.9 eV는 수렴값이 아니라 상한이다.",
            11, GRAY, False, False, None)]], line_spacing=1.12)
    card(s, Inches(7.75), Inches(1.78), Inches(5.05), Inches(1.66),
         "① 절대 기준이 생겼다", [
             "백본은 Tᵢ를 157.8 eV로 복원하고, 타깃 자신의",
             "10 ms 재현성은 46~130 eV이다. 같은 자릿수이며",
             "skill 점수가 만들 수 없는 물리 단위의 문장이다.",
         ], accent=NAVY, title_size=13.5, body_size=11.5)
    card(s, Inches(7.75), Inches(3.53), Inches(5.05), Inches(1.66),
         "② 상한은 꼬리에 있다", [
             "4차 차분 제곱질량의 46.6%(V_rot 65.6%)가 상위",
             "1%에 있다. 본류에서 모델은 타깃 산포의 2.3~3.4배로",
             "MSE의 7~9%만 환원 불가능하며, 여지가 남아 있다.",
         ], accent=TEAL, title_size=13.5, body_size=11.5)
    card(s, Inches(7.75), Inches(5.28), Inches(5.05), Inches(1.66),
         "③ 크기 축이 평평했던 이유", [
             "폭 26배·계열 3종의 평평함은 정보 고갈이 아니라",
             "지표가 꼬리 통계인 결과이다. 지목된 다음 측정은",
             "강건 공동 지표의 재채점이며 사전등록 후 실행한다.",
         ], accent=ORANGE, title_size=13.5, body_size=11.5)
    return s


# --- 6-11. Peak -----------------------------------------------------------
def s_res_peak():
    s = slide()
    header(s, "6. 결과 ⑪", "우위는 고변동(peak) 구간에 집중된다")
    add_image_fit(s, os.path.join(FIG, "fig_peak.png"),
                  Inches(0.5), Inches(1.45), Inches(7.4), Inches(5.0))
    box(s, Inches(8.0), Inches(1.55), Inches(4.8), Inches(4.85), fill=CARDBG, round_=True)
    text(s, Inches(8.25), Inches(1.70), Inches(4.35), Inches(0.5),
         [[("가치가 발생하는 위치", 15, NAVY, True, False, None)]])
    bullets(s, Inches(8.25), Inches(2.25), Inches(4.35), Inches(4.1), [
        ("peak는 입력 기반 고-국소활동 이웃으로 정의한다.", 0),
        ("타깃 행 자체는 제외한 보수적 영역 프록시이다.", 1),
        ("Tᵢ peak는 컷 +0.45~+0.61, 포함 +0.62~+0.72이다.", 0, TEAL, True),
        ("8/8 셀 PASS로 무조건부이다.", 1, GREEN, True),
        ("Tᵢ 본류는 +0.09~+0.20(컷 4/4), +0.06~+0.19(포함 2/4)이다.", 1),
        ("V_rot peak는 +0.54~+0.79(8/8 양수, PASS 각 2/4)이다.", 0, BLUE, True),
        ("persistence 대비는 +0.75~+0.86으로 8/8 PASS이다.", 1),
        ("본류는 ≈0(-0.07~+0.15, 0/8)이다.", 1),
        ("보간은 매끄러운 본류에서 이미 최적에 가깝다.", 0, NAVY, True),
        ("따라서 비대칭은 전역이 아니라 지역적이다.", 1, NAVY, True),
    ], size=12, gap=7)
    return s


# --- 6-11. Transient ------------------------------------------------------
def s_res_transient():
    s = slide()
    header(s, "6. 결과 ⑫", "급변 구간의 사례: held-out TEST shot #31815")
    add_image_fit(s, os.path.join(FIG, "fig_transient_seq_31815.png"),
                  Inches(0.45), Inches(1.40), Inches(7.15), Inches(5.4))
    box(s, Inches(7.75), Inches(1.55), Inches(5.05), Inches(4.85), fill=CARDBG, round_=True)
    text(s, Inches(8.0), Inches(1.70), Inches(4.6), Inches(0.5),
         [[("한 shot에서 관찰되는 현상", 15, NAVY, True, False, None)]])
    bullets(s, Inches(8.0), Inches(2.25), Inches(4.6), Inches(4.1), [
        ("빠른 진단이 급변을 먼저 관측한다.", 0, NAVY, True),
        ("BES 급락(빨간 점선)이 CES crash와 정렬된다.", 1),
        ("PCHIP는 스파이크마다 overshoot한다.", 0, GRAY, True),
        ("과거와 미래를 모두 읽는 오프라인 보간임에도 그러하다.", 1),
        ("모델은 세그먼트 과거와 빠른 진단만 쓰는 인과 모델이다.", 0, GREEN, True),
        ("Tᵢ RMSE는 199.2로 PCHIP 262.3보다 낮고 skill은 +0.42이다.", 0, ORANGE, True),
        ("n = 395 실측점(genuine, 컷 후)이다.", 1),
        ("V_rot는 17.3 vs 19.5로 skill +0.21이다(n = 149).", 0, BLUE, True),
        ("V_rot 관측이 끊기면 이력만으로 완만히 감쇠한다.", 1),
        ("우위는 gap·peak에 집중되며 결과 ④⑩과 일관된다.", 0, GRAY, True),
    ], size=12, gap=6)
    return s


# --- 6-12. Deployment -----------------------------------------------------
def s_deploy():
    s = slide()
    header(s, "6. 결과 ⑬", "배치 가능성: 지연과 불확실성을 측정하였다")
    card(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(2.65),
         "지연 — 상태 유지 1-step은 CPU 10 ms 예산에 여유 있게 든다", [
             "온라인에서는 은닉 상태를 격자를 따라 이월하여 새 행마다 1-step만 계산한다.",
             "seq_v2 CPU 중앙값 1.05 ms / p99 1.61 ms로 예산의 16%이다(유휴 세션).",
             "같은 세션에서 인과 GP는 p99 2.34 ms, 윈도 대조군 W=2는 4.46 ms(44.6%)였다.",
             "순서는 세션마다 불변이다: seq_v2 < 인과 GP < 윈도 W=2 < 윈도 W=4.",
             "절댓값은 세션 간 2–21.8배 변동하므로 절댓값이 아니라 순서와",
             "연산자 개수를 주장한다(7장).",
         ], accent=TEAL, title_size=14, body_size=11.5)
    card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.65),
         "불확실성 — split conformal (재학습 없음, α = 0.10)", [
             "해당 run 자신의 val에서 보정하며 예측기는 건드리지 않는다.",
             "같은 절차를 PCHIP·persistence에도 적용한 공정 비교이다.",
             "모델 구간이 32/32 셀에서 두 기준선을 Winkler 점수로 이겼다.",
             "Tᵢ 1,272 vs 1,554 vs 1,727(컷), 2,290 / 2,851 / 3,120(포함)이다.",
             "V_rot는 150 vs 164 vs 179이다.",
             "Mondrian 보정은 Tᵢ 팔을 4–5% 조이며 판정은 불변이다.",
         ], accent=BLUE, title_size=14, body_size=11.5)
    card(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(2.2),
         "실무 지침과 인정하는 한계", [
             "상태 유지형 나우캐스터를 제어 계산기의 CPU에서 실행하면 10 ms 예산의 80% 이상이 획득과 제어에 남는다.",
             "한계 ①: coverage는 marginal이며 조건부가 아니다. Tᵢ 0.87–0.92(목표 0.90), V_rot 0.91–0.94이고 shot별로는 넓게 흩어진다.",
             "한계 ②: 지연 절댓값은 기계·전원 상태에 종속된다. 5세션 프로토콜에서 p99 산포가 21.84배까지 관측되어 1 ms 판정은 보류되었다.",
             "포함 모집단에서는 모델 Tᵢ 구간이 PCHIP보다 넓은데도(반폭 224–255 vs 211–241 eV) 점수가 더 좋다. 스파이크가 miss 페널티를 키우고 모델이 덜 놓친다.",
         ], accent=ORANGE, title_size=14, body_size=12)
    return s


# =========================== 7. B.9 (context / family / cost) ==============

# --- 7-1. Reach: truncation vs trained-at-reach ---------------------------
def s_reach_design():
    s = slide()
    header(s, "7. 문맥·구조·비용 ①", "도달 범위(reach)의 측정: 절단 사다리의 결손은 대부분 cold start였다", accent=TEAL)
    text(s, Inches(0.55), Inches(1.38), Inches(12.3), Inches(0.82),
         [[("질문: ", 12.5, TEAL, True, False, None),
           ("합성곱 대신 순환이 정당한가, 그리고 윈도 계열이 충분해 보이는데 무한 순환 상태가 필요한가. 두 질문은 "
            "‘학습된 모델이 실제로 몇 스텝의 연속 과거를 쓰는가’로 환원된다. 동결된 B.1 백본 4개의 순환 상태를 ctx 스텝 전에 "
            "리셋하여 재학습 없이 재채점하였고(§8ac), 이어서 각 도달 범위에서 학습한 모델과 비교하였다(§8ae, §8af).",
            12, DARK, False, False, None)]], line_spacing=1.12)
    cw = [Inches(1.9), Inches(2.4), Inches(1.15), Inches(2.6), Inches(1.15), Inches(2.0)]
    rows = [
        [("2 (20 ms)", DARK, True, None), ("-0.510", RED, True, MONO), ("4/4", RED, False, MONO),
         ("-0.065", ORANGE, True, MONO), ("4/4", ORANGE, False, MONO), ("87%", NAVY, True, MONO)],
        [("7 (70 ms)", DARK, True, None), ("-0.064", RED, False, MONO), ("4/4", RED, False, MONO),
         ("+0.002", GREEN, True, MONO), ("0/4", GREEN, False, MONO), ("103%", NAVY, False, MONO)],
        [("15 (150 ms)", DARK, True, None), ("-0.018", GRAY, False, MONO), ("2/4", GRAY, False, MONO),
         ("+0.005", GREEN, False, MONO), ("1/4", GRAY, False, MONO), ("—", MGRAY, False, MONO)],
        [("31 (310 ms)", DARK, True, None), ("-0.009", GRAY, False, MONO), ("1/4", GRAY, False, MONO),
         ("+0.023", GREEN, False, MONO), ("0/4", GREEN, False, MONO), ("—", MGRAY, False, MONO)],
        [("63 (630 ms)", DARK, True, None), ("+0.000", GRAY, False, MONO), ("0/4", GREEN, False, MONO),
         ("+0.029", GREEN, False, MONO), ("0/4", GREEN, False, MONO), ("—", MGRAY, False, MONO)],
    ]
    table(s, Inches(0.55), Inches(2.25), cw,
          ["도달 범위(스텝)", "절단(상태 리셋) vs 전체", "유의", "그 범위에서 학습 vs 전체", "유의", "warm-up 비중"],
          rows, row_h=Inches(0.40), head_h=Inches(0.44), size=12, head_size=11.5)
    box(s, Inches(0.55), Inches(4.85), Inches(12.25), Inches(2.08), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(4.93), Inches(11.7), Inches(1.95),
         [[("읽기 (CES_TI, 전체 블록 백본 대비 paired, 4 분할 평균)", 12.5, NAVY, True, False, None)],
          [("① 절단 사다리는 63스텝이 되어야 자기 자신과 같아지지만, 그 범위에서 학습한 모델은 7스텝이면 충분하였다. "
            "20 ms 결손 -0.510의 87%는 학습되지 않은 상태 재구축(cold start)이었다.", 11.5, DARK, False, False, None)],
          [("② 백본이 윈도 대조군보다 앞서는 +0.081(§8x)은 도달 범위 -0.065와 구조 ≈ -0.016으로 분해된다. "
            "백본이 사는 것의 4/5는 얼마나 멀리 보는가이고 1/5는 무엇인가이다.", 11.5, DARK, False, False, None)],
          [("③ 무한 문맥은 필요하지 않으며 더 낫지도 않다. 310·630 ms는 전체 블록보다 명목상 위(+0.023, +0.029, 0/4)이다. "
            "V_rot는 20 ms에서만 -0.013(2/4)의 작은 결손을 보였다.", 11.5, DARK, False, False, None)]],
         line_spacing=1.1, space_after=3)
    return s


# --- 7-2. Dense ladder + pooled re-scoring --------------------------------
def s_reach_ladder():
    s = slide()
    header(s, "7. 문맥·구조·비용 ②", "밀집 사다리와 통합 재채점: 문맥 약 50 ms에서 포화하며, 문맥이 사는 것은 전형성이다", accent=TEAL)
    add_image_fit(s, os.path.join(PAPERFIG, "fig_context_family_ladder.png"),
                  Inches(0.45), Inches(1.38), Inches(6.55), Inches(5.55))
    cw = [Inches(1.15), Inches(2.55), Inches(1.05), Inches(1.0)]
    rows = [
        [("20 ms", DARK, True, None), ("+0.057 [+0.027, +0.085]", GREEN, False, MONO), ("0.52", ORANGE, True, MONO), ("+0.028", GRAY, False, MONO)],
        [("30 ms", DARK, True, None), ("+0.087 [+0.061, +0.111]", GREEN, False, MONO), ("0.60", DARK, False, MONO), ("+0.060", GRAY, False, MONO)],
        [("50 ms", DARK, True, None), ("+0.104 [+0.079, +0.128]", GREEN, False, MONO), ("0.64", DARK, False, MONO), ("+0.077", GRAY, False, MONO)],
        [("70 ms", DARK, True, None), ("+0.119 [+0.095, +0.142]", GREEN, False, MONO), ("0.66", GREEN, True, MONO), ("+0.092", GRAY, False, MONO)],
        [("150 ms", DARK, True, None), ("+0.123 [+0.096, +0.148]", GREEN, False, MONO), ("0.66", GREEN, True, MONO), ("+0.096", GRAY, False, MONO)],
        [("630 ms", DARK, True, None), ("+0.143 [+0.118, +0.168]", GREEN, False, MONO), ("0.67", GREEN, True, MONO), ("+0.116", GRAY, False, MONO)],
    ]
    text(s, Inches(7.15), Inches(1.40), Inches(5.65), Inches(0.42),
         [[("CES_TI vs 인과 GP, 301 방전 통합 (순환 계열, §8am)", 12.5, NAVY, True, False, None)]])
    table(s, Inches(7.15), Inches(1.82), cw, ["문맥", "skill [95% CI]", "승률", "-top10"], rows,
          row_h=Inches(0.36), head_h=Inches(0.38), size=11, head_size=11)
    box(s, Inches(7.15), Inches(4.50), Inches(5.65), Inches(2.45), fill=CARDBG, round_=True)
    text(s, Inches(7.4), Inches(4.58), Inches(5.2), Inches(2.35),
         [[("결과", 12.5, NAVY, True, False, None)],
          [("① 사전등록 규칙(전체 대비 결손 < 0.02, 유의 결손 ≤ 1/4)은 4·5·6·10스텝을 채운 뒤 50 ms를 반환하였다. "
            "70 ms는 4·5·6이 학습되지 않아 생긴 값이었다.", 11.5, DARK, False, False, None)],
          [("② 4/4 계수는 한 스텝 간격에서 4/4→3/4→4/4로 요동하므로 문턱 위치 추정기에서 퇴출하고, 통합 추세 검정으로 대체하였다. "
            "문맥 10배당 skill +0.050 [+0.036, +0.064]로 상승한다.", 11.5, DARK, False, False, None)],
          [("③ 모델은 20 ms에서도 인과 GP를 이긴다. 문맥이 사는 것은 평균이 아니라 전형성이며, 승리 방전 비율이 0.52에서 0.66으로 "
            "오른 뒤 평평해진다.", 11.5, DARK, False, False, None)]],
         line_spacing=1.12, space_after=3)
    return s


# --- 7-3. Family comparison -----------------------------------------------
def s_family():
    s = slide()
    header(s, "7. 문맥·구조·비용 ③", "계열 비교: 같은 문맥에서 순환·확장 합성곱·attention은 0.023 이내로 동률이었다", accent=TEAL)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.55),
         [[("세 계열을 seq_v2의 라우팅을 유지한 채 각자의 수용 범위에서 학습·채점하고, 같은 도달 범위의 LSTM 칸과 행 단위로 paired 비교하였다"
            "(§8ag, §8ak, §8al). 계열이 유일한 변수가 되도록 문맥을 고정하였다.", 12.5, DARK, False, False, None)]], line_spacing=1.12)
    cw = [Inches(1.3), Inches(1.0), Inches(1.2), Inches(1.35), Inches(0.9), Inches(1.25), Inches(1.35), Inches(0.9), Inches(1.0)]
    rows = [
        [("tcn3", DARK, True, MONO), "30 ms", ("71,442", GRAY, False, MONO), ("-0.004", DARK, False, MONO), "0/0", ("동률", GREEN, True, None), ("+0.001", DARK, False, MONO), "0/1", ("동률", GREEN, False, None)],
        [("tcn7", DARK, True, MONO), "70 ms", ("128,034", GRAY, False, MONO), ("-0.004", DARK, False, MONO), "0/1", ("동률", GREEN, True, None), ("+0.003", DARK, False, MONO), "0/0", ("동률", GREEN, False, None)],
        [("tcn15", DARK, True, MONO), "150 ms", ("184,626", GRAY, False, MONO), ("+0.014", DARK, False, MONO), "2/0", ("미결", GRAY, False, None), ("+0.000", DARK, False, MONO), "0/0", ("동률", GREEN, False, None)],
        [("tcn63", DARK, True, MONO), "630 ms", ("297,810", GRAY, False, MONO), ("-0.016", DARK, False, MONO), "0/2", ("미결", GRAY, False, None), ("-0.036", DARK, False, MONO), "0/1", ("미결", GRAY, False, None)],
        [("xfmr7", DARK, True, MONO), "70 ms", ("295,746", GRAY, False, MONO), ("-0.023", RED, True, MONO), ("0/3", RED, False, None), ("차이", RED, True, None), ("+0.007", DARK, False, MONO), "0/0", ("동률", GREEN, False, None)],
        [("xfmr15", DARK, True, MONO), "150 ms", ("295,746", GRAY, False, MONO), ("+0.002", DARK, False, MONO), "0/0", ("동률", GREEN, True, None), ("+0.022", DARK, False, MONO), "0/0", ("미결", GRAY, False, None)],
        [("xfmr63", DARK, True, MONO), "630 ms", ("295,746", GRAY, False, MONO), ("-0.019", DARK, False, MONO), "0/1", ("동률", GREEN, True, None), ("-0.024", DARK, False, MONO), "0/1", ("미결", GRAY, False, None)],
    ]
    table(s, Inches(0.55), Inches(1.95), cw,
          ["arm", "문맥", "파라미터", "Tᵢ paired", "승/패", "판정", "V_rot paired", "승/패", "판정"],
          rows, row_h=Inches(0.36), head_h=Inches(0.40), size=11, head_size=11)
    text(s, Inches(0.55), Inches(4.95), Inches(12.3), Inches(0.3),
         [[("같은 도달 범위의 LSTM 칸 대비 paired skill(4 분할 평균). 승/패는 CI가 0을 제외한 분할 수이며, 판정은 사전등록 §3.2 규칙(|Δ| < 0.02, 유의 ≤ 1/4 = 동률)이다.",
            10.5, GRAY, False, False, None)]])
    card(s, Inches(0.55), Inches(5.3), Inches(6.0), Inches(1.62), "판정", [
        "순환과 확장 합성곱은 같은 칸(70 ms)에서 전환하며 공유 칸 전부에서 동률이다.",
        "attention만 70 ms에서 -0.023(3/4 유의)으로 뒤지며 150 ms부터 동률이다.",
        "계열 효과의 최댓값 0.023은 문맥 효과 +0.060(20→70 ms)의 1/2.6이다.",
    ], accent=TEAL, title_size=13, body_size=11)
    card(s, Inches(6.8), Inches(5.3), Inches(6.0), Inches(1.62), "10k 파라미터 아래에서는 동률이 깨진다 (§8ai)", [
        "크기를 맞춘 순환 arm 대비 합성곱 arm은 +0.027~+0.040(3–4/4 유의)이다.",
        "tcn2k(1,808 파라미터)는 인과 GP를 4/4로 이기며 백본과 +0.001 차이이다.",
        "계열은 상한이 아니라 상한에 이르는 데 필요한 파라미터 수를 정한다.",
    ], accent=NAVY, title_size=13, body_size=11)
    return s


# --- 7-4. Cost = dispatched operator count --------------------------------
def s_cost():
    s = slide()
    header(s, "7. 문맥·구조·비용 ④", "비용 모델: 온라인 1스텝의 지연은 파라미터가 아니라 디스패치 연산자 수에 비례한다", accent=TEAL)
    cw = [Inches(1.85), Inches(1.65), Inches(1.1), Inches(1.4), Inches(1.65), Inches(0.95)]
    rows = [
        [("seq_v2", DARK, True, MONO), "stock nn.LSTM", ("357,570", GRAY, False, MONO), ("118", DARK, False, MONO), ("793", DARK, False, MONO), ("6.72", RED, True, MONO)],
        [("seq_v2_lean", DARK, True, MONO), "명시적 연산", ("357,570", GRAY, False, MONO), ("161", DARK, False, MONO), ("378", DARK, False, MONO), ("2.35", GREEN, False, MONO)],
        [("seq_v2_tight", DARK, True, MONO), "연산자 융합", ("357,570", GRAY, False, MONO), ("111", GREEN, True, MONO), ("352", GREEN, True, MONO), ("3.17", GREEN, False, MONO)],
        [("v2m2k_tight", DARK, True, MONO), "연산자 융합", ("2,362", GRAY, False, MONO), ("86", GREEN, True, MONO), ("241", DARK, False, MONO), ("2.80", GREEN, False, MONO)],
        [("tcn15_lean", DARK, True, MONO), "명시적 연산", ("184,626", GRAY, False, MONO), ("209", DARK, False, MONO), ("488", DARK, False, MONO), ("2.33", GREEN, False, MONO)],
        [("tcn63_lean", DARK, True, MONO), "명시적 연산", ("297,810", GRAY, False, MONO), ("305", ORANGE, False, MONO), ("799", DARK, False, MONO), ("2.62", GREEN, False, MONO)],
        [("tcn2k_tight", DARK, True, MONO), "연산자 융합", ("1,808", GRAY, False, MONO), ("106", GREEN, True, MONO), ("—", MGRAY, False, MONO), ("—", MGRAY, False, MONO)],
        [("xfmr63_tight", DARK, True, MONO), "KV 캐시 융합", ("295,746", GRAY, False, MONO), ("473", RED, True, MONO), ("1,504", RED, False, MONO), ("3.18", GREEN, False, MONO)],
    ]
    table(s, Inches(0.55), Inches(1.42), cw,
          ["arm", "구현", "파라미터", "ops / 스텝", "최소 중앙값 µs", "µs/op"],
          rows, row_h=Inches(0.36), head_h=Inches(0.40), size=11, head_size=11)
    box(s, Inches(9.3), Inches(1.42), Inches(3.5), Inches(3.65), fill=NAVY, round_=True)
    text(s, Inches(9.5), Inches(1.55), Inches(3.15), Inches(3.45),
         [[("계열별 비용 법칙 (§8aj)", 13, ORANGE, True, False, None)],
          [("순환: 도달 범위에 O(1)이다. 상태가 과거를 실어 나르므로 111 ops가 모든 칸에서 같다.", 11.5, WHITE, False, False, None)],
          [("확장 합성곱: O(log R)이다. 층당 정확히 +48 ops이며 RF = 2^(L+1) - 1이다.", 11.5, WHITE, False, False, None)],
          [("attention: O(1)이나 상수가 크다. KV 캐시로 도달 범위는 무료이지만 473 ops로 융합 순환의 4.3배이다.", 11.5, WHITE, False, False, None)],
          [("변환율은 2.1–3.2 µs/op로 거의 상수이다(파라미터 151배 범위, 세 계열 전부).", 11.5, LGRAY, False, False, None)]],
         line_spacing=1.14, space_after=4)
    box(s, Inches(0.55), Inches(5.2), Inches(12.25), Inches(1.72), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.28), Inches(11.7), Inches(1.6),
         [[("판정", 12.5, NAVY, True, False, None)],
          [("① 비용 모델은 t ≈ N_ops × 2–3 µs이다. 파라미터 151배는 지연 1.6배만 사며, v2m2k의 스텝당 산술은 약 2.4k MAC으로 마이크로초 수준이다. "
            "백본의 구현 사다리는 0.793 → 0.688 → 0.378 → 0.352 ms(2.25배)이며 융합은 거의 소진되었다.", 11.5, DARK, False, False, None)],
          [("② 10 ms 예산은 어느 arm에도 구속 조건이 아니다(§8ah, 전 arm 2배 이상 여유). 1 ms 판정은 5세션 p99 산포 21.84배로 인해 "
            "사전등록 규칙이 보류하였으며, 조용한 기계에서의 재측정이 지목된 측정이다.", 11.5, DARK, False, False, None)],
          [("③ 따라서 아키텍처는 skill이 아니라 비용으로 선택한다. tcn3k(3,238)와 v2m4k(3,898)는 인과 GP를 4/4로 이기는 최저 비용 운용점이다.",
            11.5, DARK, False, False, None)]],
         line_spacing=1.12, space_after=3)
    return s


# --- 7-5. Which discharges are won -----------------------------------------
def s_wins():
    s = slide()
    header(s, "7. 문맥·구조·비용 ⑤", "승패 방전 분석: 모델은 타깃이 움직이는 방전에서 이기며, V_rot의 우위는 shot-general이 아니다", accent=TEAL)
    cw = [Inches(2.7), Inches(1.0), Inches(1.5), Inches(1.5)]
    rows = [
        [("Tᵢ vs 인과 GP", DARK, True, None), ("96", GRAY, False, MONO), ("0.695", GREEN, True, MONO), ("9–16%", GREEN, False, MONO)],
        [("Tᵢ vs PCHIP", DARK, True, None), ("96", GRAY, False, MONO), ("0.557", DARK, False, MONO), ("—", MGRAY, False, MONO)],
        [("V_rot vs 인과 GP", DARK, True, None), ("~62", GRAY, False, MONO), ("0.481", RED, True, MONO), ("49–164%", RED, True, MONO)],
        [("V_rot vs PCHIP", DARK, True, None), ("~62", GRAY, False, MONO), ("0.329", RED, False, MONO), ("—", MGRAY, False, MONO)],
    ]
    text(s, Inches(0.55), Inches(1.40), Inches(6.7), Inches(0.4),
         [[("방전 단위 승률 (동결 B.1 산출물, §8al §4)", 12.5, NAVY, True, False, None)]])
    table(s, Inches(0.55), Inches(1.80), cw, ["비교", "방전", "승리 비율", "top-1 방전 기여"], rows,
          row_h=Inches(0.38), head_h=Inches(0.40), size=11.5, head_size=11)
    text(s, Inches(0.55), Inches(3.85), Inches(6.7), Inches(0.9),
         [[("V_rot는 약 62개 방전 중 상위 5개를 제거하면 4개 분할 전부에서 통합 우위가 0 이하로 떨어진다(split 42: +0.331 → -0.006). "
            "Tᵢ는 같은 처리에서 +0.045~+0.113을 유지한다. 따라서 V_rot는 검정력 부족이 아니라 소수 방전에 집중된 우위이다.",
            11.5, DARK, False, False, None)]], line_spacing=1.12)
    cw2 = [Inches(1.25), Inches(1.0), Inches(1.05), Inches(1.05), Inches(1.05)]
    rows2 = [
        [("조용함", DARK, True, None), ("42%", RED, True, MONO), ("-0.051", RED, False, MONO), ("34%", RED, True, MONO), ("-0.091", RED, False, MONO)],
        [("중간", DARK, True, None), ("83%", GREEN, False, MONO), ("—", MGRAY, False, MONO), ("48%", DARK, False, MONO), ("—", MGRAY, False, MONO)],
        [("변동 큼", DARK, True, None), ("85%", GREEN, True, MONO), ("+0.156", GREEN, False, MONO), ("55%", ORANGE, True, MONO), ("+0.031", ORANGE, False, MONO)],
    ]
    text(s, Inches(7.4), Inches(1.40), Inches(5.4), Inches(0.4),
         [[("타깃 산포 3분위별 승률 (v2r63, 통합, §8an)", 12.5, NAVY, True, False, None)]])
    table(s, Inches(7.4), Inches(1.80), cw2, ["3분위", "Tᵢ 승률", "Tᵢ 중앙", "V_rot 승률", "V_rot 중앙"], rows2,
          row_h=Inches(0.38), head_h=Inches(0.40), size=11, head_size=10.5)
    text(s, Inches(7.4), Inches(3.45), Inches(5.4), Inches(1.3),
         [[("11개 공변량 중 ‘방전 내 타깃 산포’만이 승패를 예측하였다(순열 검정, Bonferroni: Tᵢ ρ = +0.401, V_rot ρ = +0.281). "
            "캠페인 위치·간극 길이·peak 비율·held 비율·독립 관측 수·Tₑ 수준은 전부 null이었다.",
            11.5, DARK, False, False, None)]], line_spacing=1.12)
    box(s, Inches(0.55), Inches(4.95), Inches(12.25), Inches(1.95), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.03), Inches(11.7), Inches(1.85),
         [[("메커니즘", 12.5, NAVY, True, False, None)],
          [("① 평평한 방전에서는 인과 GP가 이미 최적에 가까우므로 모델이 분산을 더해 근소하게 진다. 대부분의 방전이 조용하므로 "
            "V_rot의 통합 승률 0.46–0.48은 이길 것이 없는 방전의 비율이다.", 11.5, DARK, False, False, None)],
          [("② 변동이 큰 3분위에서도 Tᵢ는 85%, V_rot는 55%만 이긴다. 이 잔차는 회전의 구동 변수(NBI 토크)가 데이터에 없다는 §2.3d의 결론과 "
            "일치한다(Tₑ~V_rot r = +0.024, p = 0.58가 세 번째 독립 증거). 모델은 입력이 없는 운동을 추적할 수 없다.", 11.5, DARK, False, False, None)],
          [("③ 문장은 ‘V_rot는 검정력 부족’이 아니라 ‘V_rot의 우위는 회전이 실제로 변하는 방전에 집중되며, 거기서도 작다’로 바뀐다. "
            "이를 뒤집을 측정은 B.6이며, μs 재획득이 모드 회전 주파수를 제공하면 변동 3분위의 승률이 먼저 올라야 한다.", 11.5, DARK, False, False, None)]],
         line_spacing=1.12, space_after=3)
    return s


# --- 7-6. Extension programme: μs shot set + quantum arm -------------------
def s_extensions():
    s = slide()
    header(s, "7. 문맥·구조·비용 ⑥", "확장 프로그램의 현황: μs 재획득 shot 집합은 동결되었고, 양자 가지는 음성으로 종결되었다", accent=TEAL)
    card(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(4.0),
         "B.6 — 마이크로초 재획득 shot 집합 (§8ao, 2026-08-21 동결)", [
             "641/641 방전을 OpenAlex 전문 검색으로 스캔하여 문헌 등장 방전 11개를 확정하였다.",
             "5자리 번호의 위양성 4부류(타 장치 shot, DOI 조각, 비플라스마 논문,",
             "AIP 논문번호 03xxxx)를 판별 체계로 기록하였다. 하루 15건이 위양성이었다.",
             "역할 배정: test 4(31921·31873·31114·31902), pool 6(31097·31359·31747·",
             "32027·32092·32097), companion 2(31923·31357)이다.",
             "test = 4는 V_rot 유효 군집 3개(296/311/412행)를 확보하며 측정 검정력은",
             "V_rot 0.750 / Tᵢ 0.368이다.",
             "#32092는 NF 2026 QH-mode 대표 방전으로 confirmed 되었고(EHO n = 1,",
             "~4/~8 kHz), B.6의 다섯 번째 양성 대조가 된다.",
             "μs 입력은 10 ms 예측을 개선하는가를 묻는 것이며, 타깃 격자는 바뀌지 않는다.",
         ], accent=NAVY, title_size=13, body_size=11)
    card(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(4.0),
         "양자 가지 — 하드웨어 검증 후 종결 (§8ap, 2026-07-26 → 08-24)", [
             "변분 양자 회로(VQC)를 파라미터 수를 맞춘 고전 MLP와",
             "PCA 축소 입력에서 비교하고, IonQ Forte에서 55개 회로를 실행하였다($1,418.44).",
             "정확도: 무잡음 시뮬레이터에서도 persistence에 진다(471.5 vs 449.6 eV).",
             "지연: 예측당 22.9 s로 10 ms 예산의 약 1,800배이다.",
             "하드웨어: 탈분극 계수 λ = 0.661 ± 0.040(16.7σ)으로 진폭의 34%만 반환된다.",
             "확장성: 입력 차원과 파라미터를 늘리면 고전 대조군도 나빠져 여유가 없다.",
             "QML 4계열 중 3계열(변분·고정 특징맵·저장소)은 닫혔다.",
             "양자 커널은 고전 대조군과 동률(+0.3279 vs RBF +0.3307)이며,",
             "학습형 양자 인코더는 미결로 표기하였다.",
             "어떤 수치도 백본의 skill과 나란히 인용되지 않으며, 코드는 experiments/quantum/에 격리된다.",
         ], accent=GRAY, title_size=13, body_size=11)
    box(s, Inches(0.55), Inches(5.65), Inches(12.25), Inches(1.25), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(5.73), Inches(11.7), Inches(1.15),
         [[("위치 설정 (§9.5): ", 12.5, NAVY, True, False, None),
           ("논문은 10 ms 격자에서 완결되며, μs 프로그램은 의존성이 아니라 확장이다. 현재 격자는 99.46% 균일하므로 "
            "‘불규칙 샘플링을 다룬다’는 주장은 이 데이터로 뒷받침되지 않으며, 다중 속도 획득이 그것을 실재하게 한다. "
            "attention은 도달 범위가 스텝이 아니라 경과 시간으로 정의되는 유일한 계열이므로, 여기서 패한 계열이 격자가 비균일해질 때 "
            "이길 것이라는 반증 가능한 예측이 성립한다.", 12, DARK, False, False, None)]],
         line_spacing=1.14)
    return s


# =========================== 8. Conclusion =================================

# --- 8-1. Conclusion ------------------------------------------------------
def s_conclusion():
    s = slide()
    header(s, "8. 결론", "결론 (5가지)")
    items = [
        ("1", "Tᵢ는 미래를 읽는 보간을 두 모집단 모두에서 유의하게 능가하였다", BLUE,
         "skill_vs_pchip 컷 +0.17~+0.26, 포함 +0.23~+0.32로 4개 독립 분할 전부 shot 군집 95% CI가 0을 제외하였다(4/4+4/4). 인과 GP는 8/8 셀에서 이겼고, 최강 오프라인 평활기(GP)와는 동률이며 이를 상한으로 보고한다."),
        ("2", "배치 주장은 두 스트레스를 견뎠다", GREEN,
         "실제 결측 in-domain 시점에서 인과 방법 대비 8/8 생존, 캠페인 경계 너머 PCHIP·인과 GP 대비 4/4+4/4이다. 대체된 윈도 대조군은 오프라인 우위를 잃었고(2/4·0/4) 그 원인은 드리프트로 측정되었다."),
        ("3", "약 50 ms의 연속 인과 문맥이 우위를 전형적으로 만들며, 어떻게 넘느냐는 비용을 정한다", TEAL,
         "20 ms에서도 인과 GP를 이기지만 승리 방전 비율은 0.52이며 50–70 ms에서 0.66으로 평평해진다. 세 계열은 같은 문맥에서 0.023 이내로 동률이므로 아키텍처는 비용(연산자 수: 순환 O(1), 합성곱 O(log R), attention 4.3배 상수)으로 선택한다."),
        ("4", "V_rot는 전역 동률이며, 우위는 회전이 실제로 변하는 방전에 집중된다", GRAY,
         "PCHIP 대비 1/4·2/4이나 > 15 ms 간극과 peak 층에서는 이겼다. 빠른 채널을 0으로 두어도 출력이 동일하고, 승률은 조용한 방전 34%·변동 방전 55%로 구동 변수(NBI 토크)의 부재를 가리킨다. 검정력 문제가 아니다."),
        ("5", "총합 손실의 상한은 추정기가 아니라 정보에 있으며, 남은 레버는 데이터이다", ORANGE,
         "21,498 파라미터 b3k8이 컷에서 백본과 동급이고 폭 34k→879k는 평평하다. 다만 타깃 자신의 10 ms 재현성이 46~130 eV이므로 이 진술은 총합 MSE에 한정된다. 남은 레버는 CES 피팅 품질·원본 kHz Mirnov·NBI 토크이다."),
    ]
    yy = 1.45
    for num, t, col, body in items:
        box(s, Inches(0.6), Inches(yy), Inches(12.2), Inches(1.02), fill=CARDBG, round_=True)
        box(s, Inches(0.72), Inches(yy + 0.14), Inches(0.74), Inches(0.74), fill=col, round_=True)
        text(s, Inches(0.72), Inches(yy + 0.14), Inches(0.74), Inches(0.74),
             [[(num, 24, WHITE, True, False, None)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.65), Inches(yy + 0.07), Inches(11.0), Inches(0.38),
             [[(t, 14, NAVY, True, False, None)]])
        text(s, Inches(1.65), Inches(yy + 0.44), Inches(11.0), Inches(0.56),
             [[(body, 10.5, DARK, False, False, None)]], line_spacing=1.06)
        yy += 1.1
    return s


# --- 8-2. Mirnov ----------------------------------------------------------
def s_mirnov():
    s = slide()
    header(s, "8. 추가 검증", "Mirnov 재가공으로 V_rot를 개선할 수 있는가: 정보는 전처리 단계에서 소실되었다")
    pic = s.shapes.add_picture(os.path.join(FIG, "fig_mirnov.png"),
                               Inches(0.75), Inches(1.45), Inches(4.74), Inches(3.6))
    pic.crop_right = 0.5   # frame aspect matches the (a) panel of the source figure
    box(s, Inches(5.95), Inches(1.45), Inches(6.85), Inches(3.6), fill=CARDBG, round_=True)
    box(s, Inches(5.95), Inches(1.45), Inches(0.10), Inches(3.6), fill=RED)
    text(s, Inches(6.2), Inches(1.58), Inches(6.4), Inches(0.5),
         [[("진단은 실측이다", 14.5, RED, True, False, None)]])
    bullets(s, Inches(6.2), Inches(2.08), Inches(6.4), Inches(2.9), [
        ("같은 10 ms 격자의 연속 블록 내 lag-1 자기상관을 측정하였다.", 0),
        ("BES +0.568, ECEI +0.572에 대해 Mirnov는 -0.009였다.", 1, RED, True),
        ("블록의 82%가 |r| < 0.1이므로 이 격자 위에서는 백색잡음이다.", 1),
        ("이는 kHz dB/dt를 안티앨리어싱 없이 100 Hz로 데시메이트한 서명이다.", 0),
        ("연속 표본의 상대 위상이 무작위가 된다.", 1),
        ("진단 집합에서 회전의 유일한 대리(모드 회전 주파수)가", 0, NAVY, True),
        ("모델보다 상류에서 버려졌다.", 1, NAVY, True),
    ], size=12, gap=7)
    card(s, Inches(0.55), Inches(5.20), Inches(3.90), Inches(1.62), "① 시도한 것", [
        "적분·PCHIP 적분·|MC|·이동 RMS 특징을 학습하였다.",
        "4-seed paired 비교에서 개선은 없었다.",
        "잃은 정보는 하류에서 복구되지 않는다.",
    ], accent=ORANGE, title_size=13.5, body_size=10.5)
    card(s, Inches(4.72), Inches(5.20), Inches(3.90), Inches(1.62), "② 해야 할 것", [
        "모델 변경이 아니라 전처리 변경이 필요하다.",
        "원본 kHz 스트림에서 윈도 RMS, 대역 파워,",
        "모드 수, 모드 회전 주파수를 계산해야 한다.",
    ], accent=TEAL, title_size=13.5, body_size=11)
    card(s, Inches(8.89), Inches(5.20), Inches(3.90), Inches(1.62), "③ 방법", [
        "V_rot 분기로 라우팅한다(사전등록된",
        "pilot-then-expand 규칙 아래).",
        "B.6의 μs shot 집합이 이 검정의 무대이다.",
    ], accent=NAVY, title_size=13.5, body_size=11)
    return s


# --- 8-3. Te ~ NBI --------------------------------------------------------
def s_te_nbi():
    s = slide()
    header(s, "8. 추가 검증",
           "Tₑ가 NBI 가열을 대리하여 V_rot 정보를 담는가: 가설을 검정하고 기각하였다", accent=RED)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.60),
         [[("이론적으로 가능한 경로이다. ", 14, RED, True, False, None),
           ("NBI가 들어가면 전자가 가열되므로 Tₑ가 NBI 주입량을 간접 반영하고, "
            "그렇다면 ECEI가 회전 정보를 간접적으로 운반할 수 있다. 이를 데이터로 직접 검정하였다.",
            14, DARK, False, False, None)]], line_spacing=1.14)
    box(s, Inches(0.55), Inches(2.00), Inches(5.05), Inches(3.62), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(2.00), Inches(0.10), Inches(3.62), fill=ORANGE)
    text(s, Inches(0.80), Inches(2.13), Inches(4.6), Inches(0.4),
         [[("가설의 인과 사슬", 14.5, ORANGE, True, False, None)]])
    bullets(s, Inches(0.80), Inches(2.62), Inches(4.6), Inches(2.9), [
        ("NBI 주입은 전자를 가열하여 Tₑ를 올린다.", 0),
        ("따라서 Tₑ는 NBI power의 간접 대리 변수이다.", 0),
        ("power가 크면 토크도 크다고 가정하면", 0),
        ("ECEI가 V_rot 정보를 운반해야 한다.", 1, NAVY, True),
        ("데이터에 NBI 컬럼이 없으므로", 0),
        ("Tₑ 대리변수로 경로의 존재 여부를 검정하였다.", 1),
    ], size=12.5, gap=9)
    text(s, Inches(5.85), Inches(2.03), Inches(7.0), Inches(0.4),
         [[("shot 간 상관 (538 shot, ECEI 채널 평균 = Tₑ 대리)",
            13, NAVY, True, False, None)]])
    cwn = [Inches(3.30), Inches(1.85), Inches(1.80)]
    nrows = [
        [("Tₑ ~ CES_TI", DARK, True, None), ("+0.353", GREEN, True, MONO), ("3e-17", GREEN, True, MONO)],
        [("Tₑ ~ CES_VT", DARK, True, None), ("+0.024", RED, True, MONO), ("0.58", RED, True, MONO)],
        [("Tₑ ~ |CES_VT|", DARK, False, None), ("+0.001", RED, False, MONO), ("0.98", RED, False, MONO)],
        [("Tₑ 변동성 ~ |CES_VT|", DARK, False, None), ("-0.026", GRAY, False, MONO), ("0.55", GRAY, False, MONO)],
        [("BES 변동성 ~ |CES_VT|", DARK, False, None), ("-0.059", GRAY, False, MONO), ("0.17", GRAY, False, MONO)],
    ]
    table(s, Inches(5.85), Inches(2.48), cwn, ["관계", "Pearson r", "p"], nrows,
          row_h=Inches(0.38), head_h=Inches(0.36), size=12, head_size=11.5,
          emphasis_fill=None)
    text(s, Inches(5.85), Inches(4.86), Inches(6.95), Inches(0.80),
         [[("shot 내부에서도 같다: ", 12, NAVY, True, False, None),
           ("Tₑ~CES_TI는 블록 평균 r = +0.246으로 부호가 일관되게 양수인 반면, "
            "Tₑ~CES_VT는 +0.006으로 부호조차 무작위이다(|r| > 0.3 블록: 42.7% vs 14.8%).",
            12, DARK, False, False, None)]], line_spacing=1.14)
    box(s, Inches(0.55), Inches(5.72), Inches(12.25), Inches(1.16), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(5.72), Inches(0.12), Inches(1.16), fill=RED)
    text(s, Inches(0.85), Inches(5.81), Inches(11.75), Inches(1.02),
         [[("결론 — 경로의 전반부는 참이고 후반부에서 끊긴다.", 12, RED, True, False, None)],
          [("Tₑ가 가열 수준의 대리로 작동하는 것은 사실이다(Tᵢ와 r = +0.35). 그러나 같은 Tₑ가 "
            "V_rot과는 무관하다(r = +0.02, p = 0.58). 끊기는 지점은 ",
            12, DARK, False, False, None),
           ("power ≠ torque", 12, RED, True, False, None),
           ("이다. 토크는 빔 에너지·접선 반경·주입 기하에 의존하여 power와 분리된다. NBI 토크 확보는 모델링이 아니라 "
            "데이터 획득 과제이며, 문헌에 양성 대조(DIII-D 전방전 시뮬레이터)가 있다. 7장 ⑤의 승패 분석이 세 번째 독립 증거이다.",
            12, DARK, False, False, None)]],
         line_spacing=1.12)
    return s


# --- 8-4. Limitations + future --------------------------------------------
def s_limits():
    s = slide()
    header(s, "8. 한계 & 향후 연구", "한계와 다음 단계")
    box(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(4.95), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(1.5), Inches(0.12), Inches(4.95), fill=RED)
    text(s, Inches(0.8), Inches(1.64), Inches(5.5), Inches(0.5),
         [[("한계 — 논문이 먼저 인정하는 것", 15.5, RED, True, False, None)]])
    bullets(s, Inches(0.8), Inches(2.18), Inches(5.55), Inches(4.2), [
        ("검정력: test 방전 96(Tᵢ) / 60–66(V_rot)이 모든 유의성의 구속조건이다.", 0),
        ("포함 모집단에서는 ≈1% 행이 SSE의 70–83%를 차지한다.", 1),
        ("MNAR 낙관: 재가중은 인과 비교엔 양쪽, 오프라인 비교엔 모집단 조건부이다.", 0),
        ("도달 범위도 Tᵢ 54–68%, V_rot 4–6%에 그친다.", 1),
        ("오프라인 주장의 상한은 GP 동률(1/8 유의)이다.", 0),
        ("값 컷은 일방향 프록시이며 V_rot 스파이크는 컷되지 않은 채 남는다.", 0),
        ("캠페인은 한 시간 블록(초기화 4개)이고 컷 run 2/4가 상한 종료였다.", 0),
        ("통합 재채점은 방법의 기대 skill을 추정하며 단일 체크포인트의 배치 주장은 아니다.", 0),
        ("승패 공변량 분석은 탐색적이며 사전등록 규칙이 없다.", 1),
        ("지연 절댓값은 기계 종속이며 1 ms 판정은 보류되었다.", 0),
        ("단일 장치 · conformal은 marginal · 페데스탈 상단 프레이밍에 한정된다.", 0),
    ], size=11.5, gap=6)
    box(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.95), fill=CARDBG, round_=True)
    box(s, Inches(6.8), Inches(1.5), Inches(0.12), Inches(4.95), fill=TEAL)
    text(s, Inches(7.05), Inches(1.64), Inches(5.5), Inches(0.5),
         [[("향후 — 남은 레버는 전부 데이터이다", 15.5, TEAL, True, False, None)]])
    bullets(s, Inches(7.05), Inches(2.18), Inches(5.55), Inches(4.2), [
        ("음성 결과는 그것을 뒤집을 측정을 지목할 때만 결론으로 인정한다.", 0, NAVY, True),
        ("① CES 피팅 품질 메타데이터(χ², 신호 수준)를 확보한다.", 0, ORANGE, True),
        ("값 컷을 품질 컷으로 대체하면 두 모집단이 하나로 합쳐진다.", 1),
        ("② 원본 kHz Mirnov 특징을 계산한다(V_rot 최상위 레버).", 0, ORANGE, True),
        ("B.6 shot 집합(test 4 · pool 6 · companion 2)이 동결되어 있다.", 1),
        ("예측: 변동 3분위의 V_rot 승률이 먼저 올라야 메커니즘이 맞다.", 1),
        ("③ NBI 토크 채널을 확보한다(회전의 원인 변수).", 0, ORANGE, True),
        ("모델링이 아니라 데이터 획득 과제이며 양성 대조가 존재한다.", 1),
        ("④ 조용한 기계에서 5세션 지연 프로토콜을 재실행한다(1 ms 판정).", 0, GRAY, True),
        ("용량·긴 문맥·계열 교체는 이미 배제되었다. 크기 축은 닫혔다.", 0, GRAY, True),
    ], size=11.5, gap=6)
    return s


# --- 8-5. Closing ---------------------------------------------------------
def s_closing():
    s = slide()
    box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
    box(s, Inches(0.9), Inches(0.95), Inches(2.2), Pt(4), fill=ORANGE)
    text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.6),
         [[("요약", 26, WHITE, True, False, None)]])
    points = [
        ("항상 존재하는 빠른 진단으로 자주 비는 CES를 채우는, 엄격히 인과적인 가상 센서를 구축하였다.", ORANGE),
        ("Tᵢ는 미래를 읽는 PCHIP를 두 모집단 4/4+4/4로 능가하였고(+0.17~+0.32), 인과 GP는 8/8에서 이겼다.", GREEN),
        ("결측 재가중 인과 대비 8/8과 캠페인 분할 4/4+4/4의 두 스트레스를 생존하였다(윈도 대조군은 붕괴).", BLUE),
        ("연속 문맥 약 50 ms에서 skill이 포화하며 승리 방전 비율이 0.52→0.66으로 올랐고, 세 계열은 동률이었다.", TEAL),
        ("V_rot는 전역 동률이며 우위는 회전이 변하는 방전에 집중되었다. 구동 변수의 부재가 원인이다.", ORANGE),
        ("비용은 연산자 수(2–3 µs/op)로 정해지며, 상한은 정보에 있다. 다음 단계는 데이터 레버 3종이다.", BLUE),
    ]
    yy = 2.1
    for t, col in points:
        box(s, Inches(0.95), Inches(yy + 0.05), Inches(0.28), Inches(0.28), fill=col)
        text(s, Inches(1.45), Inches(yy - 0.04), Inches(11.0), Inches(0.55),
             [[(t, 14.5, WHITE, False, False, None)]], line_spacing=1.1)
        yy += 0.62
    box(s, Inches(0.9), Inches(6.15), Inches(11.5), Pt(2), fill=RGBColor(0x2A, 0x47, 0x6E))
    text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7),
         [[("감사합니다.  ", 20, WHITE, True, False, None),
           ("Q & A", 20, ORANGE, True, False, None)]])
    return s


# ============================ APPENDIX ====================================
# 부록 A - 시도한 모델과 닫은 이유, 그리고 2026-09-05 문헌 조사.
# 표의 내용은 appendix_content.py 한 곳에 있으며 졸논정리 덱이 같은 것을 읽는다.
import appendix_content as AP  # noqa: E402


def _note(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip("\n")
    return s


try:
    from PIL import Image as _PILImage, ImageDraw as _PILDraw
    _MEASURE = _PILDraw.Draw(_PILImage.new("RGB", (8, 8)))
except Exception:      # pragma: no cover - QC only
    _MEASURE = None


def _wrapped_lines(txt, avail_in, size_pt, bold=False):
    """Line count after wrapping, with preview_pptx's tokenizer and font metrics."""
    fallback = max(1, int(len(txt) * size_pt / (avail_in * 72.0)) + 1)
    if _MEASURE is None:
        return fallback
    try:
        import preview_pptx as PV
    except Exception:  # pragma: no cover
        return fallback
    font = PV.load_font(FONT, bold, size_pt)     # px == pt -> lengths in points
    avail_px = max(avail_in, 0.1) * 72.0
    n, cur = 1, 0.0
    for tok in PV._TOKEN.findall(txt):
        tw = _MEASURE.textlength(tok, font=font)
        if tok.isspace():
            if cur + tw > avail_px:
                n, cur = n + 1, 0.0
                continue
            cur += tw
            continue
        if cur + tw > avail_px and cur > 0:
            n, cur = n + 1, 0.0
        cur += tw
    return n


def table_fit(s, x, y, w, avail_h, head, rows, weights, size=11.0, min_size=7.5,
              head_fill=NAVY, head_color=WHITE, first_col_color=NAVY):
    """Variable-row-height table that shrinks its font until it fits ``avail_h``.

    ``table()`` above uses one fixed row height, which is right for short cells;
    the appendix tables carry full sentences, so every row is measured and the
    point size is reduced until the block fits. Same tokenizer and metrics as
    preview_pptx, so a clean build implies a clean preview.
    """
    total_wt = float(sum(weights))
    col_w_in = [w / 914400.0 * wt / total_wt for wt in weights]
    pad_in = 0.16
    # preview_pptx: line height = pt * 1.24 * line_spacing; text() adds 2 pt side
    # margins. Cells are drawn with space_after 0 so a row is exactly its lines.
    LS = 1.06
    SLACK = 0.17        # 0.05 top offset + 0.08 box inset + a hair
    SAFE = 0.94         # measure narrower than the cell so a row is never short

    def heights(sz):
        line_h = sz * 1.24 * LS / 72.0
        out = []
        for r, row in enumerate([head] + rows):
            n = 1
            for j, cell in enumerate(row):
                txt = cell if isinstance(cell, str) else cell[0]
                avail = (col_w_in[j] - pad_in - 2 * (2.0 / 72.0)) * SAFE
                n = max(n, _wrapped_lines(txt, avail, sz, r == 0 or j == 0))
            out.append(max(0.30, n * line_h + SLACK))
        return out

    while size > min_size and sum(heights(size)) > avail_h / 914400.0:
        size -= 0.5
    hs = heights(size)

    col_w = [Inches(c) for c in col_w_in]
    yy = y
    box(s, x, yy, w, Inches(hs[0]), fill=head_fill)
    cx = x
    for j, h in enumerate(head):
        text(s, cx + Inches(0.08), yy + Inches(0.05), col_w[j] - Inches(pad_in),
             Inches(hs[0]) - Inches(0.08),
             [[(h, size, head_color, True, False, None)]],
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
             space_after=0, line_spacing=1.06)
        cx += col_w[j]
    yy += Inches(hs[0])
    for i, row in enumerate(rows):
        rh = Inches(hs[i + 1])
        if i % 2 == 1:
            box(s, x, yy, w, rh, fill=CARDBG)
        cx = x
        for j, cell in enumerate(row):
            if isinstance(cell, str):
                txt = cell
                col = first_col_color if j == 0 else DARK
                bold = (j == 0)
            else:
                txt, col, bold = (tuple(cell) + (None,) * 3)[:3]
                col = col if col is not None else DARK
            text(s, cx + Inches(0.08), yy + Inches(0.05), col_w[j] - Inches(pad_in),
                 rh - Inches(0.08),
                 [[(txt, size, col, bold, False, None)]], align=PP_ALIGN.LEFT,
                 space_after=0, line_spacing=1.06)
            cx += col_w[j]
        box(s, x, yy + rh - Pt(0.75), w, Pt(0.75), fill=LGRAY)
        yy += rh
    return yy


_TRIED_W = [2.55, 1.10, 3.15, 2.80, 2.65]
_APX_X, _APX_W = Inches(0.55), Inches(12.25)


def s_apx_tried_window():
    s = slide()
    header(s, "부록 A-1", "시도한 모델의 계보와 닫은 이유 (1/3): 윈도 계열", accent=GRAY)
    text(s, Inches(0.55), Inches(1.32), Inches(12.25), Inches(0.28),
         [[("아래는 W = 4 시대의 잠정 수치이며 확정 프로토콜의 주장에는 쓰지 않는다(§8v).",
            11.5, GRAY, False, False, None)]])
    table_fit(s, _APX_X, Inches(1.68), _APX_W, Inches(5.10),
              AP.TRIED_HEAD, AP.TRIED_WINDOW, _TRIED_W)
    return _note(s, "출처는 THESIS_RESULTS.md 8e / 8b.2 / 8f / 8k / 8u / 8x / 8ad이다. "
                    "연속시간 인코더는 코드가 2026-08-09에 제거되었고 판정만 남아 있다.")


def s_apx_tried_seq():
    s = slide()
    header(s, "부록 A-1", "시도한 모델의 계보와 닫은 이유 (2/3): 시퀀스 계열", accent=GRAY)
    text(s, Inches(0.55), Inches(1.32), Inches(12.25), Inches(0.28),
         [[("확정 프로토콜(W = 2 · held-free · 두 공동 1차 모집단) 아래에서 짝지어 채점하였다.",
            11.5, GRAY, False, False, None)]])
    table_fit(s, _APX_X, Inches(1.68), _APX_W, Inches(5.10),
              AP.TRIED_HEAD, AP.TRIED_SEQ, _TRIED_W)
    return _note(s, "출처는 8d / 8t / 8x / 8y / 8z / 8aa / 8ab / 8ai이다. "
                    "v3는 val에서 2/2 유의였으나 TEST에서 1/4에 그쳐 규칙대로 미승격하였다.")


def _s_apx_tried_misc(part, rows, takeaway):
    s = slide()
    header(s, "부록 A-1",
           "시도한 모델의 계보와 닫은 이유 (3/3, %d쪽): 계열 · 문맥 · 기준선 · 확장 가지" % part,
           accent=GRAY)
    table_fit(s, _APX_X, Inches(1.44), _APX_W,
              Inches(4.95 if takeaway else 5.45), AP.TRIED_HEAD, rows, _TRIED_W)
    if takeaway:
        text(s, Inches(0.55), Inches(6.55), Inches(12.25), Inches(0.62),
             [[(AP.TRIED_TAKEAWAY, 12, NAVY, True, False, None)]], line_spacing=1.10,
             space_after=0)
    return _note(s, "출처는 8ag / 8ai / 8ak / 8af / 8al / 8am 부록 / 8p / 8m / 8ap · 8at이며, "
                    "마지막 행은 진행 중인 B.11(PREREGISTRATION_B11.md)이다.")


def s_apx_tried_misc_a():
    return _s_apx_tried_misc(1, AP.TRIED_MISC[:5], False)


def s_apx_tried_misc_b():
    return _s_apx_tried_misc(2, AP.TRIED_MISC[5:], True)


_FUSION_W = [1.95, 1.05, 0.85, 1.85, 3.05, 1.55, 3.55]


def _s_apx_lit_fusion(part, rows):
    s = slide()
    header(s, "부록 A-2",
           "문헌 조사 (1/4, %d쪽): 핵융합의 진단-대-진단 추정은 여전히 단순한 구조가 주류이다" % part,
           accent=TEAL)
    text(s, Inches(0.55), Inches(1.32), Inches(12.25), Inches(0.28),
         [[("장치와 연도를 따로 두었다. 2026-09-05 조사한 12편 중 %d쪽이다." % part,
            11.5, GRAY, False, False, None)]])
    table_fit(s, _APX_X, Inches(1.68), _APX_W, Inches(5.10),
              AP.FUSION_HEAD, rows, _FUSION_W)
    return _note(s, AP.sources_note("2026-09-05 조사한 핵융합 12편의 요약이다."))


def chunk(rows, per=5):
    """Balanced slices of at most `per` rows, so no slide carries a lone row."""
    n = len(rows)
    parts = max(1, -(-n // per))
    size = -(-n // parts)
    return [rows[i:i + size] for i in range(0, n, size)]


_FUSION_PARTS = chunk(AP.FUSION_ROWS, 5)


_GENERAL_W = [1.70, 0.95, 2.60, 3.95, 3.35]


def _s_apx_lit_general(part, rows):
    s = slide()
    header(s, "부록 A-2",
           "문헌 조사 (2/4, %d쪽): 일반 시계열 · 센서 예측의 주류와 본 데이터에 대한 판정" % part,
           accent=TEAL)
    table_fit(s, _APX_X, Inches(1.42), _APX_W, Inches(5.35),
              AP.GENERAL_HEAD, rows, _GENERAL_W)
    return _note(s, "마지막 열은 본 저장소의 통제 실험 판정이며, 새 계열을 도입하기 전에 "
                    "이 열을 먼저 읽는다.")


def s_apx_lit_general_a():
    return _s_apx_lit_general(1, AP.GENERAL_ROWS[:5])


def s_apx_lit_general_b():
    return _s_apx_lit_general(2, AP.GENERAL_ROWS[5:])


def s_apx_lit_iso():
    s = slide()
    header(s, "부록 A-2", "문헌 조사 (3/4): 구조적으로 동형인 분야에서 반복되는 교훈",
           accent=TEAL)
    table_fit(s, _APX_X, Inches(1.42), _APX_W, Inches(4.20),
              AP.ISO_HEAD, AP.ISO_ROWS, [2.35, 0.95, 2.35, 3.25, 3.30])
    card(s, Inches(0.55), Inches(5.80), Inches(12.25), Inches(1.25),
         "공통 교훈",
         ["동형 분야가 반복해 말하는 것은 표현력이 아니라 개체별 보정과 상태추정 프레임이며, 이는 "
          "본 연구의 shot별 표준화(§8s)와 전체격자 인과 프레이밍(§8t)이 이미 취한 선택이다."],
         accent=TEAL, title_size=12.5, body_size=11.5)
    return _note(s, "혼합주기 나우캐스팅 · 저가 센서 보정 · 커프리스 혈압 · 구조 가상 센싱 · "
                    "합성 진단 증강의 다섯 계열에서 같은 결론이 반복된다.")


def s_apx_next():
    s = slide()
    header(s, "부록 A-2",
           "문헌 조사 (4/4): 문헌이 지목하는 다음 팔은 표현력이 아니라 손실 · 게이팅 · 입력이다",
           accent=ORANGE)
    table_fit(s, _APX_X, Inches(1.40), _APX_W, Inches(3.40),
              AP.PRIORITY_HEAD, AP.PRIORITY_ROWS, [0.60, 4.40, 2.55, 1.85, 2.85])
    card(s, Inches(0.55), Inches(4.98), Inches(6.00), Inches(2.10),
         "권하지 않는 방향", AP.NOT_RECOMMENDED, accent=RED,
         title_size=12.5, body_size=10.5)
    card(s, Inches(6.80), Inches(4.98), Inches(6.00), Inches(2.10),
         "V_rot는 열린 과제이다", AP.VROT_NOTE, accent=ORANGE,
         title_size=12.5, body_size=10.5)
    return _note(s, AP.sources_note(
        "각 행은 통제 변수가 하나이며, TEST를 여는 팔은 사전등록 뒤에만 실행한다(§8j)."))


# ======================= build ===========================================
def _fit_report():
    """Re-run preview_pptx's layout math at build time — FIT WARNING must be 0.

    Uses the *same* tokenizer, font metrics and overflow rule as
    ``preview_pptx.py`` (dpi 110), so a clean build implies a clean preview.
    """
    try:
        from PIL import Image, ImageDraw
        import preview_pptx as PV
    except Exception as exc:  # pragma: no cover - QC only
        return ["fit audit unavailable (%s)" % exc]
    scale, emu = 110, 914400
    W = int(prs.slide_width / emu * scale)
    H = int(prs.slide_height / emu * scale)
    warns = []
    for idx, sl in enumerate(prs.slides, start=1):
        draw = ImageDraw.Draw(Image.new("RGB", (W, H), (255, 255, 255)))
        for sh in sl.shapes:
            if sh.left is None:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                if PV.draw_text_frame(draw, sh, scale):
                    warns.append("slide %d: text overflows its box @ (%.2f, %.2f) -> %r"
                                 % (idx, sh.left / emu, sh.top / emu,
                                    sh.text_frame.text[:44]))
            x, y = sh.left / emu * scale, sh.top / emu * scale
            w, h = sh.width / emu * scale, sh.height / emu * scale
            if x < -1 or y < -1 or x + w > W + 1 or y + h > H + 1:
                warns.append("slide %d: shape outside slide @ (%.2f, %.2f)"
                             % (idx, sh.left / emu, sh.top / emu))
    return warns




def build():
    try:  # slide text is Korean; the console may be cp949
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    s_title()
    s_abstract()
    s_glossary()
    s_agenda()
    divider("1", "연구 배경과 문제 정의", "CES 결측의 실태와 두 모집단의 필요성")
    s_diagnostics()
    s_problem()
    s_missing_table()
    s_two_populations()
    s_idea()
    s_related()
    divider("2", "접근법", "미래를 읽는 보간과 배치 가능한 최강 기준선(인과 GP)을 상대로 검증한다")
    s_bar()
    s_validation()
    divider("3", "데이터와 파이프라인", "held 전면 제거 · 두 프레이밍 · 누수 삼중 차단")
    s_data()
    s_contract()
    s_split()
    s_samples()
    s_stuck()
    divider("4", "모델", "전체격자 인과 시퀀스 백본 seq_v2와 W=2 윈도 대조군")
    s_arch()
    s_arch_detail()
    s_physics()
    s_arch_window()
    s_training()
    divider("5", "평가 방법론", "사전등록 · shot 군집 bootstrap · TEST 동결")
    s_methodology()
    s_bootstrap()
    s_res_protocol()
    divider("6", "결과", "Tᵢ 4/4+4/4 · V_rot 동률 · 스트레스 2종 생존")
    s_res_ladder()
    s_res_forest()
    s_res_gate()
    s_res_gap()
    s_stress()
    s_res_campaign()
    s_res_asym()
    s_window_sweep()
    s_res_scaling()
    s_noise_floor()
    s_res_peak()
    s_res_transient()
    s_deploy()
    divider("7", "문맥·구조·비용", "도달 범위 사다리 · 계열 비교 · 연산자 비용 · 승패 방전 (B.9, 2026-08-17 ~ 08-24)")
    s_reach_design()
    s_reach_ladder()
    s_family()
    s_cost()
    s_wins()
    s_extensions()
    divider("8", "결론 · 한계 · 향후 연구", "무엇을 주장하고 무엇을 인정하는가")
    s_conclusion()
    s_mirnov()
    s_te_nbi()
    s_limits()
    s_closing()
    divider("A", "부록: 시도한 모델과 문헌 조사",
            "무엇을 시도하여 무엇을 닫았는가, 그리고 2026-09-05 문헌 조사가 지목하는 다음 팔")
    s_apx_tried_window()
    s_apx_tried_seq()
    s_apx_tried_misc_a()
    s_apx_tried_misc_b()
    for _i, _rows in enumerate(_FUSION_PARTS, start=1):
        _s_apx_lit_fusion(_i, _rows)
    s_apx_lit_general_a()
    s_apx_lit_general_b()
    s_apx_lit_iso()
    s_apx_next()

    warns = _fit_report()
    out = os.path.join(HERE, "KSTAR_CES_발표자료.pptx")
    prs.save(out)
    print("SAVED:", out, "| slides:", len(prs.slides._sldIdLst))
    for w in warns:
        print("  FIT WARNING:", w)
    print("FIT WARNING count:", len(warns))


if __name__ == "__main__":
    build()
