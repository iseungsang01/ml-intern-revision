# -*- coding: utf-8 -*-
"""Build the 20-minute KSTAR CES nowcasting talk (Korean, abstract register).

Output: docs/presentation/KSTAR_CES_발표자료_20분.pptx  (24 slides)

2026-08-27 전면 재작성. (1) 승상님 지시에 따라 슬라이드 본문과 발표자 노트를 모두 논문 초록
문체(서술형 종결, 객관·비인칭)로 통일하였다. (2) 2026-08-16 이후의 기록(B.9 도달 범위 사다리 ·
계열 비교 · 연산자 비용 · 승패 방전 분석, §8ac–§8an)을 압축한 슬라이드 한 장(t_context)을
결과 뒤에 추가하였고, 결론·한계 슬라이드를 §9 프레이밍에 맞추었다.

확정 프로토콜(THESIS_RESULTS.md §8v–§8ab): W=2 · held-free(학습·평가) · 파일당 500 ·
두 공동 1차 모집단(컷 / 포함) · 인과 GP 기준선. 주 모델은 전체격자 인과 시퀀스 나우캐스터
``seq_v2``(357,570 파라미터), 옛 주 모델(윈도 GRU + 관측마스킹 attention, 201,258)은
W=2 윈도 대조군이다. 수치는 docs/paper/paper_numbers.json(= main_ko.tex)과
THESIS_RESULTS.md §8ac–§8an에서 그대로 옮겼다.

이 덱은 ``build_pptx.py``(1시간 학위논문 발표)의 팔레트·레이아웃 헬퍼를 그대로 쓰고,
짧은 발표에 그대로 들어가는 슬라이드는 재사용하며(kicker()로 장 번호만 다시 매김),
나머지는 압축·병합본으로 새로 쓴다:

    1h deck                                          -> 20min deck
    s_diagnostics + s_problem + s_idea               -> t_background
    s_data + s_contract + s_split + s_stuck          -> t_pipeline
    s_bootstrap + s_validation + s_res_protocol      -> t_eval2
    s_stress + s_res_campaign                        -> t_stress   (결과 ④)
    s_reach_* + s_family + s_cost + s_wins           -> t_context  (결과 ⑨, 신설)
    s_limits + 개선 여지 + 결정 기록                    -> t_limits
    s_abstract + s_agenda                            -> t_overview (message-first)
    (제외) s_glossary · s_arch_detail · s_res_gap · s_window_sweep · s_deploy ·
           s_extensions · s_mirnov · s_te_nbi · 모든 divider

모든 슬라이드에 러닝 클록이 붙은 발표자 노트가 있다(총 19:20 + Q&A).
그림은 docs/presentation/figures/ 에서 읽는다 (make_figures.py 먼저 실행).

Usage (from repo root):
    py docs/presentation/build_pptx_20min.py
    py docs/presentation/preview_pptx.py "docs/presentation/KSTAR_CES_발표자료_20분.pptx" --out docs/presentation/.preview20
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
if HERE not in sys.path:
    sys.path.insert(0, HERE)

import build_pptx as B  # noqa: E402  (path bootstrap must run first)
from build_pptx import (  # noqa: E402
    prs, slide, box, text, header, bullets, card, add_image_fit, table,
    NAVY, BLUE, TEAL, ORANGE, GREEN, RED, GRAY, LGRAY, MGRAY, WHITE, DARK, CARDBG,
    MONO, EMU_W, EMU_H, FIG, PAPERFIG,
)
from pptx.util import Inches, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402


def note(s, txt):
    """Attach speaker notes to a slide."""
    s.notes_slide.notes_text_frame.text = txt.strip("\n")
    return s


def kicker(s, txt):
    """Rewrite the header kicker (small accent line above the slide title).

    Slides reused from the 1-hour deck carry that deck's section numbering
    ("6. 결과 ①", "8. 결론", …); this renumbers them for the 20-minute running order.
    """
    for sh in s.shapes:
        if not sh.has_text_frame or sh.top is None:
            continue
        if abs(sh.top - Inches(0.20)) < Inches(0.03) and abs(sh.left - Inches(0.55)) < Inches(0.03):
            runs = sh.text_frame.paragraphs[0].runs
            if runs:
                runs[0].text = txt
                for extra in runs[1:]:
                    extra.text = ""
            return s
    raise RuntimeError(f"kicker textbox not found (wanted to set {txt!r})")


# ---- post-edit helpers for slides reused from the 1-hour deck -------------
# Cross-references inside reused slides point at the 1-hour deck's result
# numbering; rather than fork those slide builders we rebuild them and patch the
# offending paragraph here, so build_pptx.py keeps producing the 1-hour deck
# unchanged.

def _text_shapes(s):
    return [sh for sh in s.shapes if sh.has_text_frame]


def drop_shape(s, needle):
    """Delete the whole textbox whose text contains `needle`."""
    for sh in _text_shapes(s):
        if needle in sh.text_frame.text:
            sh._element.getparent().remove(sh._element)
            return s
    raise RuntimeError(f"drop_shape: no shape containing {needle!r}")


def drop_para(s, needle):
    """Delete the single paragraph containing `needle`."""
    for sh in _text_shapes(s):
        for p in sh.text_frame.paragraphs:
            if needle in p.text:
                p._p.getparent().remove(p._p)
                return s
    raise RuntimeError(f"drop_para: no paragraph containing {needle!r}")


def set_run(s, needle, new_text):
    """Rewrite only the run containing `needle`, leaving the paragraph's other runs
    (and their colours/weights) alone — used to shorten a line that wraps."""
    for sh in _text_shapes(s):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if needle in r.text:
                    r.text = new_text
                    return s
    raise RuntimeError(f"set_run: no run containing {needle!r}")


def set_para(s, needle, new_text):
    """Rewrite the paragraph containing `needle`, keeping its first run's format."""
    for sh in _text_shapes(s):
        for p in sh.text_frame.paragraphs:
            if needle in p.text:
                runs = p.runs
                if not runs:
                    continue
                runs[0].text = new_text
                for extra in runs[1:]:
                    extra.text = ""
                return s
    raise RuntimeError(f"set_para: no paragraph containing {needle!r}")


# ========================= NEW / COMPRESSED SLIDES ========================

# --- 1. Title -------------------------------------------------------------
def t_title():
    s = slide()
    box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
    box(s, 0, Inches(5.6), EMU_W, Inches(1.9), fill=RGBColor(0x0E, 0x26, 0x47))
    box(s, Inches(0.9), Inches(1.7), Inches(2.2), Pt(4), fill=ORANGE)
    text(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(0.5),
         [[("원자핵공학과 대학원 내부 발표 · 20분", 16, RGBColor(0x9C, 0xC0, 0xE8), True, False, None)]])
    text(s, Inches(0.88), Inches(2.45), Inches(11.7), Inches(1.9),
         [[("KSTAR 다중 진단 기반 인과(causal) 나우캐스팅:", 30, WHITE, True, False, None)],
          [("희소 CES 신호(Tᵢ · V_rot)의 결측 구간 복원", 34, WHITE, True, False, None)]],
         line_spacing=1.12)
    text(s, Inches(0.9), Inches(4.30), Inches(11.5), Inches(1.25),
         [[("본 연구는 빠른 진단(BES · ECEI · Mirnov)과 과거 CES 이력만으로 ", 16, LGRAY, False, False, None),
           ("이온온도 Tᵢ와 토로이달 회전 V_rot", 16, ORANGE, True, False, None),
           ("를 복원하는 인과 모델을 제안하고,", 16, LGRAY, False, False, None)],
          [("미래를 읽는 오프라인 보간과 배치 가능한 최강 기준선(인과 GP)을 상대로 사전등록 프로토콜 아래 검증하였다.",
            16, LGRAY, False, False, None)]],
         line_spacing=1.2)
    text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(1.0),
         [[("이승상  (Seungsang Lee)", 17, WHITE, True, False, None)],
          [("서울대학교 · 원자핵공학  |  확정 프로토콜: W=2 · held-free · 두 모집단 공동 1차 · 백본 seq_v2 (357,570 파라미터) · 2026-08-27",
            13, MGRAY, False, False, None)]],
         line_spacing=1.25)
    return note(s, """
⏱ 00:00–00:30  (30초)

본 발표는 KSTAR의 빠른 진단으로 CES 결측 구간을 채우는 연구를 20분간 보고한다.

핵심 한 문장을 먼저 제시한다: 느리고 자주 비는 CES를 항상 측정되는 빠른 진단으로 채울 수 있는지,
그리고 미래까지 읽는 오프라인 보간을 인과 모델이 이길 수 있는지를 통계적으로 검증하였다.

용어: forecasting(예보)이 아니라 nowcasting(현재 시점의 결측 채우기)이다.
재실험(B.1–B.5)과 B.9 이후의 확정 프로토콜 수치이며, 옛 W=4 발표 수치는 전부 폐기되었다.
""")


# --- 2. Message-first overview -------------------------------------------
def t_overview():
    s = slide()
    header(s, "Overview", "결론을 먼저 제시한다 — 세 문장")
    msgs = [
        ("1", "빠른 진단은 Tᵢ 정보를 실제로 운반한다", GREEN,
         ["과거만 읽는 인과 모델이 미래까지 읽는 오프라인 보간(PCHIP)을 Tᵢ에서 유의하게 능가하였다.",
          "컷 +0.17~+0.26, 포함 +0.23~+0.32로 두 모집단 × 4개 독립 분할 전부 PASS(4/4 + 4/4)였다.",
          "배치 가능한 최강 인과 기준선(인과 GP)도 8개 셀 전부에서 이겼고, 두 스트레스를 생존하였다."]),
        ("2", "V_rot는 전역 동률이며, 그 원인은 물리적으로 설명된다", ORANGE,
         ["PR4 통과는 컷 1/4, 포함 2/4로 잡음 수준이므로 회전의 승리를 주장하지 않는다.",
          "빠른 채널을 전부 0으로 두어도 출력이 bit-identical이므로 회전 정보는 입력에 없다.",
          "우위는 회전이 실제로 변하는 방전에 집중되며(승률 34% → 55%), 구동 변수 NBI 토크가 미관측이다."]),
        ("3", "문맥이 우위를 전형적으로 만들고, 비용이 아키텍처를 정하며, 상한은 정보에 있다", BLUE,
         ["연속 인과 문맥 약 50 ms에서 skill이 포화하고 승리 방전 비율이 0.52에서 0.66으로 올랐다.",
          "순환·확장 합성곱·attention은 같은 문맥에서 0.023 이내로 동률이므로 비용(연산자 수)으로 선택한다.",
          "21,498 파라미터 모델이 백본과 동급(컷)이고 폭 26배 스윕이 평평하므로, 남은 레버는 데이터이다."]),
    ]
    yy = 1.5
    for num, t, col, body in msgs:
        box(s, Inches(0.6), Inches(yy), Inches(12.2), Inches(1.66), fill=CARDBG, round_=True)
        box(s, Inches(0.72), Inches(yy + 0.33), Inches(0.9), Inches(0.9), fill=col, round_=True)
        text(s, Inches(0.72), Inches(yy + 0.33), Inches(0.9), Inches(0.9),
             [[(num, 28, WHITE, True, False, None)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.82), Inches(yy + 0.14), Inches(10.8), Inches(0.45),
             [[(t, 16.5, NAVY, True, False, None)]])
        text(s, Inches(1.82), Inches(yy + 0.60), Inches(10.8), Inches(1.0),
             [[(line, 11.5, DARK, False, False, None)] for line in body], line_spacing=1.12)
        yy += 1.78
    return note(s, """
⏱ 00:30–01:15  (45초)

이 슬라이드가 발표 전체의 요약이며, 결론을 먼저 제시한 뒤 근거로 들어간다.

첫째, 빠른 진단은 이온온도 정보를 실제로 운반한다. 과거만 읽는 모델이 미래까지 읽는 오프라인
보간을 Tᵢ에서 유의하게 능가하였고, 데이터 처리 방식을 둘로 나눈 검증(스파이크 컷/포함)에서
양쪽 4개 분할 전부 통과하였다.
둘째, 회전 속도는 전역적으로 동률이다. 이는 실패가 아니라 물리적으로 예측된 비대칭이며,
승패 방전 분석이 그 메커니즘(구동 변수의 부재)을 지목한다.
셋째, 약 50 ms의 연속 문맥이 우위를 전형적으로 만들고, 아키텍처는 비용으로 정해지며,
성능 상한은 모델이 아니라 데이터에 있다.

시간이 부족하면 이 슬라이드를 30초로 줄이고 결론에서 회수한다.
""")


# --- 3. Background --------------------------------------------------------
def t_background():
    s = slide()
    header(s, "1. 배경 & 문제", "CES 결측의 원인과 규모, 그리고 이를 채울 입력")
    text(s, Inches(0.55), Inches(1.38), Inches(12.3), Inches(0.60),
         [[("CES는 페데스탈 물리의 핵심량 Tᵢ · V_rot를 제공하지만 ", 14, DARK, False, False, None),
           ("광자 적분이 필요하여 느리고 자주 결측된다", 14, ORANGE, True, False, None),
           (". 빠른 진단은 같은 10 ms 격자에서 결측 없이 측정된다.", 14, DARK, False, False, None)]])
    card(s, Inches(0.55), Inches(1.98), Inches(6.0), Inches(3.35),
         "진단 구성 — 같은 10 ms 격자", [
             "CES(타깃)는 Tᵢ · V_rot를 제공하며 느리고 자주 결측된다.",
             "BES 9 ch는 밀도요동 nₑ의 공간 구조를 제공하여 Tᵢ의 단서가 된다.",
             "ECEI 4 ch는 전자온도 Tₑ의 2D 영상을 제공하여 Tᵢ의 단서가 된다.",
             "물리 경로는 충돌 e–i 결합이다.",
             "Mirnov 2 ch는 kHz dB/dt를 100 Hz로 데시메이트한 신호이며,",
             "lag-1 자기상관이 -0.009(BES +0.568)로 이 격자에서 백색잡음이다.",
             "빠른 진단은 격자에서 100% 측정된다.",
         ], accent=BLUE, title_size=15, body_size=12.5)
    box(s, Inches(6.75), Inches(1.98), Inches(4.55), Inches(3.35), fill=NAVY, round_=True)
    text(s, Inches(7.0), Inches(2.10), Inches(4.1), Inches(0.4),
         [[("결측 실태 (641 방전 · 247,207행)", 14.5, ORANGE, True, False, None)]])
    for i, (lab, pct, col) in enumerate([("CES_TI 완전결측(NaN)", "8.2 %", RGBColor(0x9D, 0xE8, 0xCD)),
                                         ("CES_VT 완전결측(NaN)", "23.9 %", ORANGE)]):
        xx = 7.0 + i * 2.05
        text(s, Inches(xx), Inches(2.58), Inches(2.0), Inches(0.34),
             [[(lab, 11, LGRAY, False, False, None)]])
        text(s, Inches(xx), Inches(2.86), Inches(2.0), Inches(0.56),
             [[(pct, 26, col, True, False, None)]])
    box(s, Inches(7.0), Inches(3.44), Inches(4.05), Inches(0.78),
        fill=RGBColor(0x8E, 0x2B, 0x22), round_=True)
    text(s, Inches(7.18), Inches(3.52), Inches(3.75), Inches(0.66),
         [[("held(직전값 복사) CES_VT 41.1 %가 더해진다", 12, RGBColor(0xFF, 0xD5, 0xCE), True, False, None)],
          [("V_rot 실질 무정보 비율은 65.0 %이다", 13.5, WHITE, True, False, None)]],
         line_spacing=1.12, space_after=0)
    text(s, Inches(7.0), Inches(4.32), Inches(4.05), Inches(0.95),
         [[("· 두 타깃은 ", 11, WHITE, False, False, None),
           ("독립적으로", 11, ORANGE, True, False, None),
           (" 결측된다(타깃별 처리).", 11, WHITE, False, False, None)],
          [("· 결측은 저 SNR · ELM · 천이에 집중된다.", 11, WHITE, False, False, None)],
          [("  MNAR이므로 관측점 skill은 낙관적 상한이다.", 11, LGRAY, False, False, None)]],
         line_spacing=1.12, space_after=1)
    add_image_fit(s, os.path.join(FIG, "fig_raw_csv_missing.png"),
                  Inches(11.45), Inches(1.99), Inches(1.75), Inches(3.3))
    box(s, Inches(0.55), Inches(5.45), Inches(12.25), Inches(1.42), fill=NAVY, round_=True)
    text(s, Inches(0.85), Inches(5.56), Inches(11.7), Inches(1.25),
         [[("연구 질문", 13, ORANGE, True, False, None)],
          [("CES가 결측된 10 ms 시점에서, 동시각 빠른 진단과 과거 CES 이력만으로 "
            "CES 자체의 시간 보간이 복원할 수 없는 정보를 회복할 수 있는가를 묻는다.", 15, WHITE, True, False, None)],
          [("핵심 비대칭의 예고: 빠른 진단은 Tᵢ 정보는 운반하지만 V_rot 정보는 거의 운반하지 않는다고 예측되었고, 결과에서 확인되었다.",
            12, LGRAY, False, False, None)]], line_spacing=1.14, space_after=3)
    return note(s, """
⏱ 01:15–02:10  (55초)

CES는 이온온도와 토로이달 회전이라는 페데스탈 물리의 두 핵심량을 제공한다. 그러나 충분한
신호대잡음비를 얻기 위해 광자를 오래 적분해야 하므로 느리고 자주 결측된다. 같은 10 ms 격자에서
Tᵢ는 8.2%, 회전은 23.9%가 값 자체가 비어 있다. 회전은 여기에 더해, 값이 채워진 행 중 54%가
직전 값을 복사한 held 값이며 독립 측정이 아니다. NaN과 held를 합치면 회전은 전체 행의 65%가
실질적으로 정보가 없다. Tᵢ의 held는 226,991행 중 1행으로 사실상 0%이다. 두 타깃은 서로
독립적으로 결측된다.

반면 BES, ECEI, Mirnov 같은 빠른 진단은 같은 격자에서 100% 측정된다. 따라서 연구 질문은
항상 존재하는 입력으로 자주 비는 타깃을 채울 수 있는가, 정확히는 CES 자신의 시간 보간이
복원할 수 없는 정보를 회복할 수 있는가이다.

물리적 연결(질문 대비): ECEI(Tₑ)와 BES(nₑ)는 충돌 e–i 결합을 통해 Tᵢ와 연결된다. 회전은
주로 NBI 토크가 구동하며 그 토크가 데이터에 없다. MNAR은 결과 ④에서 정량화된다.
""")


# --- 7. Data pipeline + contract (merged) --------------------------------
def t_pipeline():
    s = slide()
    header(s, "3. 데이터 & 모델", "데이터 계약: No-Fake-Data · held 전면 제거 · 누수 삼중 차단")
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.72),
         [[("641 방전(shot 30801–32751), 10 ms 격자, 247,207행이다. 세그먼트는 0.5 s 이상의 간극에서 분리되고 "
            "전형적 파일은 주 세그먼트 1개(중앙값 301행 ≈ 3.0 s)이며, 모델 입력도 보간도 이 경계를 넘지 않는다. ",
            12.5, DARK, False, False, None),
           ("TEST(seed 42, 컷)는 Tᵢ 32,589행 / 96 방전, V_rot 10,463행 / 60 방전이며 선택이 끝날 때까지 봉인되었다.",
            12.5, NAVY, True, False, None)]], line_spacing=1.12)
    cards = [
        ("① 가짜 라벨 금지 (No Fake Data)", ORANGE,
         ["학습 행을 만들기 위해 타깃을 대체(impute)하지 않았다.",
          "윈도는 진단 입력이 완전하고 타깃이 1개 이상 관측된 행만 쓴다.",
          "시퀀스는 라벨 없는 행을 맥락으로만 쓴다.",
          "어느 프레이밍도 타깃 행 자신의 값을 읽지 않는다."]),
        ("② 타깃별 masked loss", BLUE,
         ["L = Σ m·(예측 - 실측)² / Σ m 이며 m은 타깃별 관측 마스크이다.",
          "한쪽 타깃만 관측된 행도 그 타깃의 학습에 기여한다.",
          "두 타깃 필수 필터는 라벨 행의 약 28%를 버렸었다.",
          "이 필터의 제거는 순수한 데이터 이득이다."]),
        ("③ 누수 삼중 차단", TEAL,
         ["파일(shot) 단위로 분할하여 인접 행의 자기상관 누수를 막았다.",
          "정규화 통계는 학습 파일에서만 추정하였다(희소 타깃은 NaN 인지).",
          "시퀀스 모델은 shot별 입력 표준화를 추가로 적용하였다.",
          "타깃 시점의 값과 관측 flag는 입력에 들어가지 않는다."]),
        ("④ held 전면 제거", NAVY,
         ["관측 V_rot의 54%는 계측기 유지값이다(499/641 파일).",
          "지도 타깃·이력 입력·정규화 통계·보간 앵커에서 동일하게 제거하였다.",
          "대가는 PR2 폴백률 Tᵢ 0.3–0.4%, V_rot 40–44%이다.",
          "따라서 어떤 arm도 forward-fill로 점수를 얻지 못한다."]),
    ]
    for i, (t, col, lines) in enumerate(cards):
        r, c = divmod(i, 2)
        card(s, Inches(0.55 + c * 6.2), Inches(2.20 + r * 2.35), Inches(6.0), Inches(2.2),
             t, lines, accent=col, title_size=14.5, body_size=12)
    return note(s, """
⏱ 04:30–05:15  (45초)

데이터는 641개 방전, 10 ms 격자의 247,207행이다. 파일 하나는 대개 3초 길이의 측정 세그먼트
하나이며, 모델도 보간도 그 경계를 넘지 않아 정보 조건이 양쪽에 동일하게 맞추어졌다.

데이터 계약은 네 가지이다. 첫째, 가짜 라벨을 만들지 않았다. 결측을 보간으로 메워 학습에 쓰면
‘보간을 이기는가’라는 질문 자체가 무의미해진다. 둘째, 두 타깃이 독립적으로 결측되므로 손실을
타깃별로 마스킹하였다. 이를 하지 않으면 한쪽만 관측된 행(라벨의 약 28%)이 버려진다. 셋째,
누수를 세 곳에서 막았다: shot 파일 단위 분할, 학습 파일 전용 정규화, 타깃 시점의 값과 관측
플래그의 완전 마스킹이다. 넷째, held 값을 지도 타깃·이력 입력·정규화 통계·모든 기준선의
보간 앵커에서 전부 제거하였다. 그 대가로 회전 채점 행의 40~44%에서 보간이 persistence로
후퇴하며, 이 폴백률은 사전등록에 따라 보고된다.

질문 대비: TEST는 아키텍처 탐색 전에 예약되었고 선택 중에는 열리지 않았다.
""")


# --- 12. Evaluation methodology, part 2 ----------------------------------
def t_eval2():
    s = slide()
    header(s, "4. 평가 방법론", "shot 군집 paired bootstrap · 두 모집단 규칙 · 모델 선택")
    bullets(s, Inches(0.55), Inches(1.5), Inches(6.5), Inches(3.0), [
        ("한 방전 안의 인접 CES 행은 강하게 상관된다.", 0),
        ("개별 샘플을 독립으로 보면 불확실성이 크게 과소평가된다.", 1, RED, True),
        ("PR4 검정은 샘플별 짝지은 오차 (SE_model - SE_pchip)를", 0),
        ("shot 단위로 묶고 shot 전체를 복원추출한다(B = 10,000).", 1),
        ("95% CI가 0을 제외하면 PASS이다.", 0, GREEN, True),
        ("이 CI는 ‘새로운 방전에서도 이기는가’에 답한다.", 1),
        ("유효 표본은 방전 수(Tᵢ ≈96 · V_rot 60–66)이며 검정력의 상한이다.", 0),
        ("모델 대 모델 비교도 같은 행 위에서 같은 paired bootstrap으로 한다.", 0, NAVY, True),
    ], size=13, gap=7)
    box(s, Inches(7.35), Inches(1.5), Inches(5.45), Inches(3.0), fill=NAVY, round_=True)
    text(s, Inches(7.6), Inches(1.62), Inches(5.0), Inches(1.35),
         [[("Murphy (1988) skill score", 13.5, ORANGE, True, False, None)],
          [("skill = 1 - MSE_model / MSE_baseline", 15, WHITE, True, False, MONO)],
          [("0보다 크면 모델이 우수하고 0이면 동률이다. 오차는 물리 단위로 역정규화하여 타깃별로 계산한다.",
            12, LGRAY, False, False, None)]], line_spacing=1.2)
    box(s, Inches(7.6), Inches(3.02), Inches(4.95), Pt(2), fill=RGBColor(0x2A, 0x47, 0x6E))
    text(s, Inches(7.6), Inches(3.16), Inches(5.0), Inches(1.25),
         [[("두 공동 1차 모집단", 13.5, ORANGE, True, False, None)],
          [("컷은 Tᵢ > 3 keV(피팅 실패)를 결측 처리한 모집단, 포함은 컷 없음이다. 전 arm에 동일 적용한다.",
            12, WHITE, False, False, None)],
          [("무조건부 주장은 두 모집단 모두에서 성립할 때만 한다.", 12.5, WHITE, True, False, None)]],
         line_spacing=1.2)
    box(s, Inches(0.55), Inches(4.70), Inches(12.25), Inches(2.0), fill=CARDBG, round_=True)
    text(s, Inches(0.85), Inches(4.82), Inches(11.7), Inches(0.45),
         [[("모델 선택 프로토콜 — 규칙을 수치보다 먼저 기록하였다", 14.5, NAVY, True, False, None)]])
    bullets(s, Inches(0.85), Inches(5.30), Inches(11.7), Inches(1.35), [
        ("백본 관문은 4조건(4 분할 부호 유지 · 통합 CI 0 제외 · 예산 균등화 · V_rot 손실 없음)을 먼저 고정하고 그다음 충족 여부를 판정하였다.", 0),
        ("유일한 아키텍처 후보(seq_v2 + 관측마스킹 인과 attention)는 4/4 양수(+0.009/+0.013/+0.033/+0.020)였으나 유의 1/4로 승격되지 않았다.", 0),
        ("val에서는 2/2 유의였으며, 이것이 승격 bar를 TEST에 두는 이유이다. 스윕 위에서 백본을 재선정하는 것은 구성상 금지된다.", 1, RED, True),
        ("사다리 칸·폭 스윕의 판정 규칙도 TEST 채점 전에 문서화되었고, TEST는 결정마다 한 번만 채점되었다.", 0),
    ], size=12.5, gap=6)
    return note(s, """
⏱ 08:05–08:50  (45초)

신뢰구간은 행이 아니라 방전 단위로 계산하였다. 한 방전 안의 10 ms 간격 측정들은 거의 복사본이므로,
이를 수만 개의 독립 증거로 세면 확신이 과장된다. 따라서 짝지은 오차를 shot으로 묶고 shot을 통째로
10,000회 재추출하였다. 모델에 불리한 계산이지만 ‘새로운 방전에서도 이기는가’에 답하는 방식이며,
검정력의 상한은 샘플 수가 아니라 방전 수(Tᵢ 96, 회전 60~66)이다.

두 모집단: Tᵢ 관측값의 0.53%가 3 keV를 넘으며 이는 플라즈마가 아니라 실패한 스펙트럼 피팅이다.
제거하면 어려운 행을 없앴다는 비판을, 유지하면 스파이크가 보간 앵커를 오염시킨다는 비판을 받으므로
둘 다 공동 1차로 사전등록하였고, 무조건부 주장은 양쪽에서 성립할 때만 한다.

모델 선택: 백본을 바꾸는 관문의 네 조건을 먼저 고정한 뒤 충족시켰고, 이후 유일한 아키텍처 후보는
4개 분할 전부 양수였음에도 유의가 1개뿐이어서 승격하지 않았다. val에서는 2/2 유의였다는 점이
승격 기준을 TEST에 두는 이유이다.

시간이 부족하면 아래 밴드(모델 선택)를 한 줄로 요약한다.
""")


# --- 17. Result 4: two stress tests (MNAR + campaign) --------------------
def t_stress():
    s = slide()
    header(s, "5. 결과 ④", "스트레스 2종: 실제 결측점 재가중(MNAR)과 캠페인(시간) 분할을 모두 생존하였다",
           accent=ORANGE)
    text(s, Inches(0.55), Inches(1.36), Inches(12.3), Inches(0.5),
         [[("관측점 skill은 낙관적 상한이고 무작위 분할은 시간 이동을 검사하지 않는다. "
            "배치 주장을 가르는 두 스트레스를 사전에 정해 두고 통과 여부를 측정하였다.",
            13, DARK, False, False, None)]], line_spacing=1.12)
    card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(2.5),
         "① 실제 결측점으로 재가중 (MNAR)", [
             "층은 Δt(15/25/45 ms) × 입력만의 활동 flag이며,",
             "결측 행의 층 분포로 채점 지점을 재가중하였다.",
             "도달: 결측 Tᵢ의 54–68%가 in-domain이고 V_rot은 4–6%이다.",
             "따라서 재가중 V_rot은 결론을 내지 않는다(결측 질량의 1/20).",
             "Tᵢ vs persistence: 컷 4/4 · 포함 4/4 (+0.28~+0.44)이다.",
             "Tᵢ vs PCHIP: 컷 2/4 · 포함 4/4 (점추정 +0.14~+0.28)이다.",
         ], accent=ORANGE, title_size=14.5, body_size=12)
    card(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(2.5),
         "② 캠페인(시간) 분할 — shot 번호로 자른다", [
             "train 416 (30801–31991) / val 128 (32002–32310) /",
             "test 97 (32312–32751)이며 초기화 seed 4개이다.",
             "윈도 대조군은 컷 2/4 · 포함 0/4 · 인과 GP 0/4로 붕괴하였다.",
             "seq_v2 컷은 +0.187/+0.174/+0.181/+0.177로 4/4이다.",
             "seq_v2 포함은 +0.173/+0.202/+0.198/+0.184로 4/4이다.",
             "원인은 측정되었다: 드리프트 BES 1.22σ·ECEI 0.53σ vs 타깃 0.115σ.",
         ], accent=BLUE, title_size=14.5, body_size=12)
    cw = [Inches(4.35), Inches(4.0), Inches(3.9)]
    rows = [
        [("무작위 분할 · 관측점 (결과 ②)", DARK, True, None),
         ("4/4 · 4/4,  +0.17~+0.32", GREEN, True, None),
         ("4/4 · 4/4 vs 인과 GP,  +0.08~+0.17", GREEN, True, None)],
        [("결측점 재가중 (MNAR)", DARK, True, None),
         ("2/4 · 4/4,  +0.14~+0.28", ORANGE, True, None),
         ("4/4 · 4/4 vs persistence,  +0.28~+0.44", GREEN, True, None)],
        [("캠페인 시간 분할", DARK, True, None),
         ("4/4 · 4/4,  +0.17~+0.20", GREEN, True, None),
         ("4/4 · 4/4 vs 인과 GP,  +0.11~+0.16", GREEN, True, None)],
    ]
    table(s, Inches(0.55), Inches(4.78), cw,
          ["Tᵢ 평가 (컷 · 포함)", "vs PCHIP (오프라인)", "vs 인과 기준선"], rows,
          row_h=Inches(0.44), head_h=Inches(0.42), size=12, head_size=12)
    text(s, Inches(0.55), Inches(6.54), Inches(12.3), Inches(0.40),
         [[("진술: ", 12, NAVY, True, False, None),
           ("실제 결측 시점에서 나우캐스터는 모든 인과 CES-only 방법보다 유의하게 낫고, "
            "오프라인 보간보다는 모집단 조건부로 낫다.", 12, DARK, False, False, None)]],
         line_spacing=1.12)
    return note(s, """
⏱ 11:15–12:20  (65초)  ★ 배치 주장을 가르는 슬라이드

관측된 지점에서만 채점하면 낙관적이다. 결측은 어려운 순간에 집중되기 때문이다. 따라서 결측 행의
층 분포로 채점 지점을 재가중하였다(왼쪽). 온라인 시스템이 실제로 경쟁하는 persistence 대비로는
두 모집단 4개 분할 전부 생존하였다. 오프라인 보간 대비로는 점추정은 유지되지만 컷 모집단에서
2개 분할의 신뢰구간이 0을 지나므로 모집단 조건부로 기술한다. 재가중의 도달 범위는 Tᵢ가 결측의
54~68%인 반면 회전은 4~6%뿐이므로 재가중 회전은 결론을 내지 않는다.

오른쪽은 시간 분할이다. 방전 번호로 잘라 과거로 학습하고 미래를 채점하였다. 옛 윈도 모델은
오프라인 우위를 완전히 잃었으나(컷 2/4, 포함 0/4) 시퀀스 백본은 4/4와 4/4로 생존하였다.
원인은 측정되었다: 빠른 진단의 드리프트가 1.22σ인 반면 타깃은 0.115σ이다. 학습 파일 전용
정규화가 무작위 분할에서는 옳지만 캠페인 이동에서 깨지며, 백본은 정의상 shot별 표준화를 하므로
이 함정을 피한다.

주의(먼저 말할 것): 캠페인은 한 시간 블록 위의 초기화 4개이며 분할 4개가 아니다.
""")


# --- 21. Result 9: context / family / cost / wins (B.9, compressed) --------
def t_context():
    s = slide()
    header(s, "5. 결과 ⑨", "문맥·구조·비용: 약 50 ms의 문맥이 우위를 전형적으로 만들고, 계열은 비용으로 정한다", accent=TEAL)
    add_image_fit(s, os.path.join(PAPERFIG, "fig_context_family_ladder.png"),
                  Inches(0.45), Inches(1.38), Inches(6.2), Inches(5.3))
    cards = [
        ("① 문맥 — 포화 50 ms, 문맥이 사는 것은 전형성", TEAL,
         ["각 도달 범위에서 학습한 seq_v2의 전체 대비 결손은 50 ms부터 0.02 아래이다.",
          "301 방전 통합에서 20 ms도 인과 GP를 이기나(+0.057) 승률은 0.52이며,",
          "70 ms에서 0.66으로 평평해진다. 문맥 10배당 skill은 +0.050이다.",
          "절단 사다리 결손의 87%는 cold start였다."]),
        ("② 계열 — 같은 문맥에서 동률, 비용은 연산자 수", NAVY,
         ["순환·확장 합성곱·attention은 같은 문맥의 LSTM과 0.023 이내로 동률이다.",
          "지연은 t ≈ N_ops × 2–3 µs이다. 순환 111 ops(O(1)), 합성곱 +48 ops/층,",
          "attention 473 ops이다. 10 ms 예산은 어느 arm에도 구속 조건이 아니다.",
          "1,808 파라미터 tcn2k도 인과 GP를 4/4로 이긴다."]),
        ("③ 승패 방전 — 타깃이 움직이는 방전에서 이긴다", ORANGE,
         ["방전 단위 승률은 Tᵢ 0.695, V_rot 0.481이며, V_rot는 상위 5개 방전을 빼면",
          "4 분할 전부 0 이하이다. 11개 공변량 중 ‘방전 내 타깃 산포’만 승패를 예측하였다",
          "(Tᵢ ρ +0.401, V_rot +0.281). 변동 큰 3분위에서도 Tᵢ 85%, V_rot 55%이며,",
          "잔차는 구동 변수(NBI 토크)의 부재를 가리킨다."]),
    ]
    yy = 1.38
    for t, col, lines in cards:
        card(s, Inches(6.8), Inches(yy), Inches(6.0), Inches(1.78), t, lines,
             accent=col, title_size=12.5, body_size=10.5)
        yy += 1.86
    return note(s, """
⏱ 15:15–16:45  (90초)  ★ 2026-08-16 이후 추가된 결과

백본이 세그먼트 전체를 읽는 것이 정당한가라는 질문에서 출발하여 세 축을 측정하였다.

문맥. seq_v2를 2·3·4·5·6·7·10·15·31·63 스텝의 도달 범위에서 각각 학습하고 채점하였다.
사전등록 규칙(전체 대비 결손 0.02 미만, 유의 결손 1/4 이하)은 50 ms를 반환한다. 301개 방전을
통합하면 모델은 20 ms에서도 인과 GP를 이기지만 승리 방전 비율은 0.52이며, 70 ms에서 0.66으로
평평해진다. 즉 문맥이 사는 것은 평균이 아니라 전형성이다. 동결 백본의 상태를 절단한 첫 사다리의
결손 중 87%는 학습되지 않은 cold start였다.

계열. 같은 문맥에서 순환·확장 합성곱·attention은 0.023 이내로 동률이므로 skill로는 구분되지
않는다. 지연은 파라미터가 아니라 디스패치 연산자 수에 비례하며(연산자당 2~3 µs), 순환은 도달
범위에 O(1), 합성곱은 O(log R), attention은 상수가 4.3배이다. 10 ms 예산은 구속 조건이 아니다.

승패 방전. Tᵢ는 방전의 70%에서 이기고 V_rot는 48%에서 이긴다. 승패를 예측하는 유일한
변수는 방전 내 타깃의 산포이며, 조용한 방전에서는 GP가 이미 최적이다. 변동이 큰 방전에서도
V_rot는 55%만 이기므로, 잔차는 회전의 구동 변수가 데이터에 없다는 결론과 일치한다.

질문 대비: 1 ms 판정은 5세션 p99 산포 21.84배로 인해 사전등록 규칙이 보류하였다.
""")


# --- 23. Limitations + future --------------------------------------------
def t_limits():
    s = slide()
    header(s, "7. 한계 & 향후", "무엇을 인정하고, 다음에 무엇을 측정하는가")
    box(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(3.9), fill=CARDBG, round_=True)
    box(s, Inches(0.55), Inches(1.5), Inches(0.12), Inches(3.9), fill=RED)
    text(s, Inches(0.8), Inches(1.62), Inches(5.5), Inches(0.45),
         [[("한계 — 논문이 먼저 인정하는 것", 15, RED, True, False, None)]])
    bullets(s, Inches(0.8), Inches(2.12), Inches(5.55), Inches(3.1), [
        ("검정력: test 방전 96(Tᵢ) / 60–66(V_rot)이 유의성의 구속조건이다.", 0),
        ("포함 모집단에서는 ≈1% 행이 SSE의 70–83%를 차지한다.", 1),
        ("MNAR 낙관: 재가중 도달은 Tᵢ 54–68%, V_rot 4–6%에 그친다.", 0),
        ("오프라인 주장의 상한은 GP 동률(1/8 유의)이다.", 0),
        ("값 컷은 일방향 프록시이며 V_rot 스파이크는 남는다.", 0),
        ("캠페인은 한 시간 블록 위 초기화 4개이며 컷 run 2/4가 상한 종료였다.", 0),
        ("통합 재채점은 방법의 기대 skill이며 단일 체크포인트의 주장이 아니다.", 0),
        ("지연 절댓값은 기계 종속이며 1 ms 판정은 보류되었다.", 0),
        ("단일 장치 · conformal은 marginal · 승패 공변량 분석은 탐색적이다.", 1),
    ], size=11.5, gap=5)
    box(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.9), fill=CARDBG, round_=True)
    box(s, Inches(6.8), Inches(1.5), Inches(0.12), Inches(3.9), fill=TEAL)
    text(s, Inches(7.05), Inches(1.62), Inches(5.5), Inches(0.45),
         [[("향후 — 남은 레버는 전부 데이터이다", 15, TEAL, True, False, None)]])
    bullets(s, Inches(7.05), Inches(2.12), Inches(5.55), Inches(3.1), [
        ("음성 결과는 그것을 뒤집을 측정을 지목할 때만 결론으로 인정한다.", 0, NAVY, True),
        ("① CES 피팅 품질 메타데이터(χ² · 신호 수준)를 확보한다.", 0, ORANGE, True),
        ("값 컷을 품질 컷으로 대체하면 두 모집단이 하나로 합쳐진다.", 1),
        ("② 원본 kHz Mirnov 특징을 계산한다(V_rot 최상위 레버).", 0, ORANGE, True),
        ("B.6 shot 집합(test 4 · pool 6 · companion 2)이 동결되어 있다.", 1),
        ("예측: 변동 3분위의 V_rot 승률이 먼저 올라야 메커니즘이 맞다.", 1),
        ("③ NBI 토크 채널을 확보한다(회전의 원인 변수).", 0, ORANGE, True),
        ("Tₑ~Tᵢ r = +0.353 vs Tₑ~V_rot r = +0.024이므로 power ≠ torque이다.", 1),
        ("크기·문맥·계열 축은 닫혔다: 21k = 백본(컷), 폭 26배 평평, 계열 동률.", 0, GRAY, True),
    ], size=12, gap=6)
    box(s, Inches(0.55), Inches(5.56), Inches(12.25), Inches(0.94), fill=NAVY, round_=True)
    text(s, Inches(0.85), Inches(5.65), Inches(11.7), Inches(0.82),
         [[("결정 기록 (2026-08-16 · 08-21 · 08-24)", 12.5, ORANGE, True, False, None)],
          [("① 두 모집단 공동 1차 유지 ② V_rot 프로토콜 불변(재학습 없음, anchored 비교엔 SSE 비중 병기) "
            "③ B.6 μs shot 집합 동결(test 4 / pool 6 / companion 2) ④ 양자 가지는 하드웨어 검증 후 음성으로 종결.", 12.5, WHITE, False, False, None)]],
         line_spacing=1.14, space_after=2)
    return note(s, """
⏱ 17:55–18:50  (55초)

한계를 먼저 명시한다. 가장 큰 제약은 검정력이며, test 방전이 Tᵢ 96개, 회전 60~66개로 모든
유의성 판정을 구속한다. 채점은 관측 지점에서 이루어지므로 낙관적 상한이고, 재가중이 닿는 범위도
Tᵢ 54~68%, 회전 4~6%뿐이다. 오프라인 주장의 상한도 명시한다: 가장 강한 오프라인 평활기인
GP와는 동률이다. 통합 재채점은 방법의 기대 skill을 추정하며 단일 체크포인트의 배치 주장이 아니다.

향후 과제는 전부 데이터 쪽이다. 원칙은 ‘음성 결과는 그것을 뒤집을 측정을 지목할 때만 결론이
된다’이며 셋을 지목하였다. CES 피팅 품질 메타데이터를 받으면 값 컷을 품질 컷으로 바꿔 두 모집단을
하나로 합칠 수 있다. 원본 kHz Mirnov에서 모드 회전 주파수를 뽑는 것이 회전의 최상위 레버이며,
그 검정의 무대인 B.6 shot 집합은 동결되어 있고 예측(변동 방전에서 먼저 개선)도 기록되어 있다.
NBI 토크 채널은 모델링이 아니라 데이터 획득 과제이다. 크기·문맥·계열 축은 이미 닫혔다.

질문 대비: Mirnov 파생 특징 재가공(적분·|MC|·이동 RMS)은 이미 시도되어 개선이 없었다.
이미 잃은 정보는 하류에서 복구되지 않으므로 전처리(원본 kHz)로 올라가야 한다.
""")


# ============================== build =====================================
def build():
    try:  # slide text is Korean; the console may be cp949
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

    # ---- 도입 (1–2) -----------------------------------------------------
    t_title()
    t_overview()

    # ---- 1. 배경 & 문제 (3–5) -------------------------------------------
    t_background()
    note(kicker(B.s_missing_table(), "1. 배경 & 문제"), """
⏱ 02:10–02:45  (35초)

결측을 비율이 아니라 전수 집계로 제시한다(641 shot, 247,207행 전부).

위 두 줄이 핵심이다. 값이 비어 있는 NaN은 Tᵢ 8.2%, 회전 23.9%이다. 그 아래, 직전 관측값과
비트 단위로 동일한 held 행이 회전에서 41.1%이며 최대 1,214행이 연속으로 같은 값이다. 회전 진단의
실제 측정 주기가 행 주기보다 느려 값이 carry-forward된 것이며 독립 측정이 아니다. 둘을 합치면
회전은 65%가 실질 무정보이고 실제 정보가 있는 행은 35%이다. Tᵢ는 held가 226,991행 중 1행이다.

held 행은 persistence나 보간이 오차 0에 가깝게 맞히는 행이므로, 확정 프로토콜에서는 학습 타깃,
이력 입력, 정규화 통계, 모든 기준선의 보간 앵커에서 전부 제거하고 실제 측정만 채점하였다.

이 슬라이드는 데이터 감사의 신뢰를 확보하는 지점이다. 발견하여 고쳤다는 사실을 명확히 전달한다.
""")
    note(kicker(B.s_two_populations(), "1. 배경 & 문제"), """
⏱ 02:45–03:25  (40초)

두 번째 감사이다. 관측 Tᵢ의 p99가 2,089 eV인 반면 최댓값은 14,984 eV이다. 이 먼 꼬리는
플라즈마가 아니라 실패한 스펙트럼 피팅이다. 3 keV를 넘는 행은 1,197행(0.53%)이며 85%가 단일
행이고 정점이 이웃 평균의 13배이다. 어떤 방법으로도 예측되지 않으며 보간의 앵커를 오염시킨다.

제거하면 어려운 행을 없앴다는 비판을 받고, 유지하면 스파이크 앵커가 오프라인 기준선에 핸디캡을
준다. 어느 쪽도 안전하지 않으므로 두 대응을 모두 공동 1차 모집단(컷과 포함)으로 사전등록하였다.
규칙은 하나이다: 무조건부 주장은 두 모집단 모두에서 성립할 때만 한다. 문턱 민감도도 검사하였으며
2.5·3·4 keV에서 결과가 같았다.

질문 대비: 값 컷은 일방향 프록시이다. 하향 dip 4,965행은 손대지 않았고 상향 이상치의 19%만
제거된다. 회전 스파이크는 컷하지 않고 SSE 비중을 병기한다.
""")

    # ---- 2. 접근법 (6) ---------------------------------------------------
    note(B.s_bar(), """
⏱ 03:25–04:30  (65초)  ★ 이 발표에서 가장 중요한 슬라이드

평가 기준선의 설정이 이 연구의 핵심이다. 흔한 방식은 persistence(직전 값 유지)와의 비교이나
이는 너무 쉬운 상대이므로, 의도적으로 훨씬 어려운 상대인 CES-only 오프라인 보간(선형·PCHIP·
국소 AR·GP)을 선택하였다.

정보 비대칭에 주목한다. 오른쪽의 보간은 맞히려는 시점의 과거와 미래 CES를 모두 읽는다. 반면
본 모델은 그 시점까지의 빠른 진단과 세그먼트 과거 이력만 읽고 미래는 전혀 읽지 않는다. 엄격히
인과적이다. 따라서 미래를 읽는 보간을 과거만 읽는 모델이 이긴다면, 빠른 진단이 시간 보간으로는
얻을 수 없는 CES 정보를 실제로 운반한다는 강한 증거가 된다.

이번 개정에서 팔을 하나 더 넣었다. 인과 GP는 같은 GP를 과거 이웃 16개로 제한한 것이며, 실제로
배치할 수 있는 방법 중 가장 강한 경쟁자이다. ‘배치 가능한 모든 인과 방법을 이긴다’는 주장은
persistence가 아니라 이 인과 GP로 판정한다.

청중이 ‘왜 굳이 불리하게 설계하는가’를 물을 수 있으므로 마지막 두 문장을 강조한다.
""")

    # ---- 3. 데이터 & 모델 (7–10) ----------------------------------------
    t_pipeline()
    note(kicker(B.s_samples(), "3. 데이터 & 모델"), """
⏱ 05:15–05:55  (40초)

학습 예제를 만드는 방식은 두 가지이며, 이 대조가 결과 ③의 핵심이다.

윈도 프레이밍은 맞히려는 시점 앞 두 행을 잘라 텐서로 만든다. 옛 주 모델의 방식이며 현재는
대조군이다. 전체격자 시퀀스는 세그먼트 안에서 입력이 온전한 행을 라벨 유무와 무관하게 전부
맥락으로 유지하고, 희소성은 loss 마스킹으로 처리한다.

결정적 차이는 도달 범위이다. 윈도는 과거 관측 W-1개만 읽고 시퀀스는 세그먼트 전체를 읽는다.
시퀀스에서는 W가 더 이상 하이퍼파라미터가 아니다. 이 도달 범위가 실제로 얼마나 필요한지는
결과 ⑨에서 측정된다.

질문 대비: 윈도의 시간 특징 4채널(lookback·간격·각 log1p)은 불규칙한 관측 간격을 명시적으로
노출하기 위한 것이다. 과거 값의 신뢰도는 10 ms 전인지 200 ms 전인지에 의존한다.
""")
    note(kicker(B.s_arch(), "3. 데이터 & 모델"), """
⏱ 05:55–06:50  (55초)

주 모델은 22채널 격자 시퀀스 위의 독립적인 인과 LSTM 두 개이다.

Tᵢ 분기는 2층 160으로 전체 상태(빠른 진단, 두 타깃의 이월값, 신선도, 시간 간격)를 읽는다.
V_rot 분기는 1층 64이며 빠른 진단이 아닌 7채널만 읽는다. 라우팅을 head가 아니라 인코더에서
하였다. 순환 상태를 공유하면 head를 어떻게 배선해도 빠른 진단 정보가 회전 쪽으로 누출된다.
분기 자체를 분리하였으므로 빠른 채널 15개를 전부 섭동해도 회전 출력이 비트 단위로 동일하며,
이것이 뒤의 비대칭 결론을 구조적으로 보증한다.

희소성은 loss가 처리한다(세그먼트의 라벨 있는 모든 행에 대한 타깃별 masked MSE). 전체
357,570 파라미터이며 AdamW와 조기 종료로 14~25 에폭에서 학습이 끝난다.

질문 대비: 라벨 없는 행을 남기는 이유는 빠른 진단이 그 행에서도 조밀하게 관측되기 때문이며,
버리면 인과 문맥이 끊긴다.
""")
    note(kicker(B.s_arch_window(), "3. 데이터 & 모델"), """
⏱ 06:50–07:20  (30초)

짝지은 대조군은 옛 주 모델이다. 진단별 시간 인지 1D CNN, 양방향 GRU 이력 인코더, 그 위의 관측
마스킹 attention pooling으로 구성된 201,258 파라미터 모델이다. attention은 해당 타깃이 실제
관측된 행에만 질량을 허용하며, 보간의 귀납 편향을 파라미터 0개로 이식한 것이다. 약 40회의
keep/discard 통제 실험의 산물이다.

현재 역할은 셋이다: 백본 관문의 비교 대상, 절제 실험의 무대, 캠페인 붕괴의 재현자이다. 데이터
계약·held 처리·분할·채점 모집단이 백본과 동일하므로 행 단위로 짝지어 비교된다.

시간이 부족하면 ‘옛 주 모델이 현재는 대조군’이라는 사실만 전달한다.
""")

    # ---- 4. 평가 방법론 (11–12) ------------------------------------------
    note(kicker(B.s_methodology(), "4. 평가 방법론"), """
⏱ 07:20–08:05  (45초)  ★ 통계 질문이 나오는 구간

결과에 앞서 수치의 신뢰 근거를 제시한다.

왼쪽은 TEST 동결이다. test 방전은 아키텍처 탐색 전에 분리되었고 선택이 끝날 때까지 열리지
않았다. 모델 선택은 전부 val에서 이루어졌으므로 headline 수치에 winner's curse가 없다.

오른쪽은 사전등록이다. 비교 상대는 PCHIP으로 확정하고, 보간이 예측하지 못하는 지점은
persistence로 채점하되 폴백률을 보고하며, test 최소 규모를 두고, 신뢰구간이 0을 제외할 때만
‘이겼다’고 판정한다. 이번 개정에서 held-free, W=2, 파일당 500, 두 모집단, TEST 채점 전
결정 규칙 커밋이 추가되었다.

아래는 기준선 사다리이다. 위 세 개가 인과, 아래 두 개가 미래를 읽는 오프라인이며 인과 GP가
배치 가능한 최강 상대이다.

예상 질문 ‘왜 PCHIP인가’: 단조성을 보존하여 overshoot이 적은 보수적이고 강한 기준선이며,
타깃 양쪽의 과거와 미래를 모두 읽는다. 즉 모델에 불리하게 설계된 기준선이다.
""")
    t_eval2()

    # ---- 5. 결과 ①~⑧ (13–20) -------------------------------------------
    note(kicker(B.s_res_ladder(), "5. 결과 ①"), """
⏱ 08:50–09:25  (35초)

가장 단순한 결과인 RMSE 사다리이며 낮을수록 좋다.

두 타깃 모두 백본이 최저 오차이다. persistence나 국소 AR 같은 인과 기준선은 큰 마진으로
이기고, 배치 가능한 최강 상대인 인과 GP보다도 Tᵢ 4%, 회전 18% 낮다. 미래까지 읽는 오프라인
GP와는 153.8 대 157.8로 동률이며, 이것이 오프라인 주장의 상한이고 함께 보고한다.

포함 모집단에서는 스파이크 때문에 Tᵢ RMSE가 두 배 이상 커지지만 순서는 유지된다.

다음 슬라이드로의 연결: 점추정이 앞선다는 것만으로 이겼다고 하지 않으며 신뢰구간을 본다.
""")
    s_forest = kicker(B.s_res_forest(), "5. 결과 ②")
    # cross-reference fix: the scaling slide is 결과 ⑨ in the 1-hour deck, ⑥ here.
    set_para(s_forest, "결과 ⑨에서 분해",
             "모든 arm이 PCHIP 대비 좋아 보이기 때문이며 결과 ⑥에서 분해한다.")
    note(s_forest, """
⏱ 09:25–10:30  (65초)  ★ headline 결과

forest plot이며 가로축이 PCHIP 대비 skill이다. 점이 추정치, 가로 막대가 shot 군집 95%
신뢰구간이며, 막대가 0선을 넘지 않으면 통계적으로 유의하게 이긴 것이다.

Tᵢ는 4개의 독립 분할 전부에서 신뢰구간이 0을 제외한다. 컷 모집단 +0.17~+0.26, 포함 모집단
+0.23~+0.32로 4/4와 4/4이다. 데이터 처리 방식에 무관하게 성립하므로 무조건부 주장이다.
8개 셀 전부에서 인과 GP와 persistence도 이긴다.

회전은 다르다. 점추정은 8개 셀 전부 양수이나 유의는 컷 1개, 포함 2개뿐이다. 4개 중 1~2개는
잡음이 만들 수 있는 수준이므로 ‘이겼다’고 쓰지 않고 동률로 보고한다.

포함 모집단의 수치가 더 높아 보이는 것은 모델이 더 잘해서가 아니라 스파이크가 보간 앵커를
오염시켜 모든 arm이 좋아 보이기 때문이며, 그 성분은 결과 ⑥에서 분해된다.
""")
    note(kicker(B.s_res_gate(), "5. 결과 ③"), """
⏱ 10:30–11:15  (45초)

주 모델을 바꾼 근거인 백본 관문이다.

조건 네 개를 먼저 고정하였다: 4개 분할 전부 부호 유지, 통합 신뢰구간이 0 제외, 예산 균등화에서도
부호 유지, 회전에 손해 없음. 분할 4개 × 초기화 4개, 총 16회 학습하여 각 run을 자기 분할의 윈도
대조군과 짝지었다. Tᵢ는 16/16 양수, 13개 유의이며 통합 +0.081, 신뢰구간 [+0.067, +0.096]이다.
학습 예산을 고정해도 4개 분할 부호가 유지되고 회전의 유의한 열세는 없다.

의미: 윈도 대조군은 인과 GP와 동률(1/4)인 반면 시퀀스 백본은 4/4와 4/4이다. 세그먼트 과거
전체로의 도달 범위가 최강 배치 기준선을 이기게 하였고, 학습 비용은 윈도 조립과 조합 증강이 없어
1/10이다. 그 도달 범위 중 실제로 필요한 양은 결과 ⑨에서 측정된다.
""")
    t_stress()
    note(kicker(B.s_res_asym(), "5. 결과 ⑤"), """
⏱ 12:20–13:15  (55초)  ★ 과학적 발견

비대칭의 원인을 절제로 확인하였다. 평가 시점에 modality를 지우는 실험이다.

공통점: 이력을 지우면 두 타깃 모두 무너진다(-1에서 -4까지).

차이는 빠른 진단이다. Tᵢ는 컷 모집단에서 빠른 채널을 0으로 두면 보간 아래로 떨어진다(-0.125).
즉 보간을 이기는 마진 자체가 빠른 진단이 운반한 정보이다. 물리적으로 ECEI의 Tₑ와 BES의 nₑ가
충돌 e–i 결합으로 Tᵢ와 연결된다.

회전은 정반대이다. 빠른 채널을 전부 0으로 두어도 8개 셀 전부에서 출력이 비트 단위로 동일하다.
회전 정보는 100% CES 이력에서 온다. 이유는 둘이다: 회전을 구동하는 NBI 토크가 데이터에 없고,
Mirnov는 kHz dB/dt를 100 Hz로 순간샘플하여 위상이 무작위이다(lag-1 자기상관 BES +0.568,
Mirnov -0.009).

포함 모집단에서는 이력만 쓰는 모델도 PCHIP을 +0.15~+0.23 이기므로 그 마진에 스파이크 강건성
성분이 섞여 있으며, 빠른 진단의 기여를 분리하려면 컷 모집단을 보아야 한다. 두 모집단을 함께
보고하는 이유이다.

질문 대비: ‘Tₑ가 NBI 가열을 대리하니 회전 정보도 담기지 않는가’는 검정되었다. Tₑ~Tᵢ는
r = +0.353(p = 3e-17)로 경로가 실재하지만 Tₑ~V_rot는 r = +0.024(p = 0.58)이다. power와
torque가 다르기 때문이다.
""")
    note(kicker(B.s_res_scaling(), "5. 결과 ⑥"), """
⏱ 13:15–14:00  (45초)

모델을 더 키우면 되는가라는 질문을 두 실험으로 닫았다.

왼쪽은 복잡도 사다리이다. b3k8은 persistence 예측에 유계 latent 8개의 선형 보정을 더한
21,498 파라미터 모델인데, 컷 모집단에서 357k 백본과 동급이다(짝지은 차이 평균 +0.002, 신뢰구간
전부 0 포함). 즉 컷 조건에서 백본의 Tᵢ skill 전부가 숫자 8개로 압축된다. probe 결과 latent은
직전 Tᵢ와 ECEI의 Tₑ를 담는다. 포함 모집단에서는 -0.194로 백본에 지는데, 유계 보정으로는
스파이크 앵커를 살리지 못하기 때문이다(1% 미만의 행이 모든 arm SSE의 70~83%).

오른쪽은 크기 축이다. Tᵢ 인코더 폭을 24에서 260까지(34k에서 879k 파라미터, 26배) 키워도
skill이 +0.230에서 +0.236 사이에서 평평하다. 남은 분산은 모델 크기가 아니라 분할 분산이다.
""")
    note(kicker(B.s_res_peak(), "5. 결과 ⑦"), """
⏱ 14:00–14:40  (40초)

모델이 이득을 얻는 위치를 peak 분석으로 확인하였다.

peak는 예측값이 아니라 입력 기준으로 정의한다: 타깃 행을 제외한 이웃에서 국소 활동도가 높은
구간이며, 정답을 읽지 않는 보수적 정의이다.

Tᵢ는 peak 구간에서 컷 +0.45~+0.61, 포함 +0.62~+0.72로 8개 셀 전부 PASS이며 본류에서는
+0.09~+0.20이다. 회전은 전역 동률이지만 peak 구간에서 +0.54~+0.79이고 persistence 대비
8/8 PASS이며 본류는 0에 가깝다.

해석: 매끄러운 본류에서는 보간이 이미 거의 최적이라 이길 여지가 없고, 급변 구간에서는 보간이
구조적으로 실패하여 빠른 진단이 가치를 갖는다. 회전의 열세도 전역적이 아니라 지역적이다.
""")
    s_tr = kicker(B.s_res_transient(), "5. 결과 ⑧")
    # cross-reference fix: the 1-hour deck points at its 결과 ④(간극)·⑩(peak).
    set_para(s_tr, "결과 ④⑩", "●  우위는 gap·peak에 집중되며 결과 ⑦과 일관된다.")
    note(s_tr, """
⏱ 14:40–15:15  (35초)

held-out TEST shot #31815의 사례이다.

위가 이온온도이다. 검은 점이 실측, 파란 선이 모델, 회색이 PCHIP이다. 빨간 점선은 BES가
급락한 시점이며 CES crash와 정렬되어 있다. 빠른 진단이 급변을 먼저 관측한다. PCHIP은 미래까지
읽고도 스파이크마다 overshoot한다. 이 shot에서 Tᵢ RMSE는 199.2 대 262.3으로 skill +0.42이며
실측점 395개 기준이다. 회전도 이 shot에서는 +0.21로 이기며, 관측이 끊기면 이력만으로 완만히
감쇠한다.

이런 shot이 전부는 아니다. 우위는 고르게 퍼져 있지 않고 간극과 급변 구간에 집중되며, 앞의 peak
분석과 같은 그림이다. 어떤 방전에서 이기는가는 결과 ⑨에서 정량화된다.

시간이 부족하면 이 슬라이드는 생략해도 논지가 성립한다.
""")

    # ---- 5. 결과 ⑨ — B.9 (21) -------------------------------------------
    t_context()

    # ---- 6~7. 결론 · 한계 · 요약 (22–24) --------------------------------
    note(kicker(B.s_conclusion(), "6. 결론"), """
⏱ 16:45–17:55  (70초)

결론은 다섯 가지이다.

하나. Tᵢ는 미래를 읽는 보간을 두 모집단 모두에서 유의하게 능가하였다. 4개 독립 분할 전부이며,
배치 가능한 최강 인과 방법인 인과 GP는 8개 셀 전부에서 이겼다. 상한도 함께 보고한다: 최강
오프라인 평활기와는 동률이다.

둘. 배치 주장은 두 스트레스를 모두 견뎠다. 실제 결측 지점으로 재가중해도 인과 방법 대비 8/8,
캠페인 경계를 넘어서도 4/4와 4/4이다. 대체된 윈도 대조군은 오프라인 우위를 잃었고 그 원인은
드리프트로 측정되었다.

셋. 약 50 ms의 연속 인과 문맥이 우위를 전형적으로 만들며, 세 계열은 같은 문맥에서 동률이므로
아키텍처는 비용(연산자 수)으로 선택한다.

넷. 회전은 전역 동률이며, 우위는 회전이 실제로 변하는 방전에 집중된다. 이는 검정력 문제가 아니라
구동 변수의 부재이다.

다섯. 상한은 추정기가 아니라 정보에 있다. 21,498 파라미터가 컷에서 백본과 같고 폭은 평평하며
1,808 파라미터 합성곱도 인과 GP를 4/4로 이긴다. 따라서 다음 단계는 더 큰 모델이 아니라 데이터이다.
""")
    t_limits()
    note(B.s_closing(), """
⏱ 18:50–19:20  (30초) + Q&A

요약: 항상 존재하는 빠른 진단으로 자주 비는 CES를 채우는 엄격히 인과적인 가상 센서를 구축하였고,
미래까지 읽는 보간이라는 의도적으로 어려운 기준선으로 검증하였다. 이온온도는 두 모집단 모두에서
유의하게 이겼고 두 스트레스를 견뎠으며, 약 50 ms의 문맥이 우위를 전형적으로 만든다. 회전은 전역
동률이며 그 이유가 물리와 승패 방전 분석으로 설명된다.

── 예상 질문 대비 ──────────────────────────────
Q. 왜 forecasting이 아니라 nowcasting인가?
 → 미래를 예보하는 것이 아니라 이미 지난 시점의 빈 값을 채우는 문제이며, 오프라인 물리 분석용이다.

Q. 진짜 결측 구간에서도 이만큼 복원되는가?
 → 단정하지 않는다. MNAR이므로 관측점 성능은 낙관적 상한이며, 결측 분포로 재가중하면 인과 대비는
    8/8 생존, 오프라인 대비는 모집단 조건부(컷 2/4·포함 4/4)이다.

Q. 보간이 미래를 읽는데 이기는 것이 가능한가?
 → 가능하며 그것이 논지의 핵심이다. 빠른 진단이 시간 보간으로 얻을 수 없는 정보를 운반하며,
    보간이 구조적으로 실패하는 급변 구간에서 이긴다.

Q. 왜 LSTM인가? 합성곱이나 attention이 더 낫지 않은가?
 → 같은 문맥에서 세 계열은 0.023 이내로 동률이다(B.9). 선택은 비용으로 하며, 순환은 도달 범위에
    O(1), 합성곱은 O(log R), attention은 상수가 4.3배이다. 10k 파라미터 아래에서는 합성곱이 낫다.

Q. 얼마나 긴 문맥이 필요한가?
 → 약 50 ms에서 포화한다. 20 ms에서도 인과 GP를 이기지만 승리 방전 비율이 0.52이며 70 ms에서
    0.66으로 평평해진다. 무한 문맥은 필요하지 않다.

Q. V_rot을 살리려면?
 → NBI 토크 채널 확보가 1순위, 원본 kHz Mirnov(모드 회전 주파수)가 그다음이다. B.6 shot 집합이
    동결되어 있고, 변동 방전에서 먼저 개선되어야 한다는 예측이 기록되어 있다.

Q. 실시간에 쓸 수 있는가?
 → 상태 유지 1-step이 CPU 유휴 세션에서 중앙값 1.05 ms / p99 1.61 ms이며, 10 ms 예산은 어느
    arm에도 구속 조건이 아니다. 절댓값은 기계 종속이므로 연산자 수와 순서를 주장한다.

Q. seed에 따라 흔들리지 않는가?
 → Tᵢ는 4개 독립 분할 × 두 모집단 전부 PASS로 강건하다. V_rot은 방전 단위 승률이 0.48로 소수
    방전에 집중되며, 이것이 ‘현 규모에서 V_rot 단일 수치는 shot-general이 아니다’라는 근거이다.
""")

    warns = B._fit_report()
    out = os.path.join(HERE, "KSTAR_CES_발표자료_20분.pptx")
    prs.save(out)
    print("SAVED:", out, "| slides:", len(prs.slides._sldIdLst))
    for w in warns:
        print("  FIT WARNING:", w)
    print("FIT WARNING count:", len(warns))


if __name__ == "__main__":
    build()
